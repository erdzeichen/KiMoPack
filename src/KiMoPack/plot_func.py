# -*- coding: utf-8 -*-
version = "7.12.7"
Copyright = '@Jens Uhlig'
if 1: #Hide imports	
	import os
	import sys
	import pandas
	import numpy as np
	import numbers
	import matplotlib
	import matplotlib.pyplot as plt
	import matplotlib.cm as cm
	import matplotlib.image as mpimg
	from matplotlib.gridspec import GridSpec
	from matplotlib.ticker import FuncFormatter
	from matplotlib.colors import BoundaryNorm,Normalize,SymLogNorm
	from matplotlib.ticker import LinearLocator	
	import math
	import re
	import scipy
	import scipy.constants
	from scipy.signal import savgol_filter
	from scipy.special import erf
	from scipy.stats import binned_statistic
	import scipy.stats	
	import pathlib
	from pathlib import Path
	from tkinter import filedialog
	import tkinter
	import time as tm #sorry i use time in my code
	import lmfit
	
	try:
		import h5py
	except:
		print('the h5py module was not imported that allows to save the projects files on windows you can install that with pip install h5py')
	try:
		import keyboard
	except:
		print('the keyboard module was not imported. on Windows this allows to stop the fit by pressing q/n you can install it with "pip install keyboard" ')
	try:	
		from pptx import Presentation
		from pptx.util import Inches
	except:
		print('We need python-pptx to create a powerpoint file. Not essential. Either use pip or for anaconda: conda install -c conda-forge python-pptx')
	try:
		import urllib3
		import shutil
	except:
		print('We need the packages urllib3 and shutil to download files from the web') 
	plt.ion()
	pandas.options.mode.chained_assignment = None  # I use this a lot and think I can ignore it
	FWHM = 2.35482
	shading = 'auto'  # gouraud
	standard_map = cm.jet
	halfsize=False
	start_time=tm.time()
print('Plot_func version %s\nwas imported from path:\n %s' % (version, os.path.dirname(os.path.realpath(__file__))))
print('The current working folder is:\n %s' % os.getcwd())

#use this to trigger a real error for DeprecationWarnings
#np.warnings.filterwarnings('error', category=np.VisibleDeprecationWarning)   
               
def download_notebooks(libraries_only=False):
	'''function loads the workflow notebooks into the active folder
	if libraries_only is set to True, only the function library and the import library are loaded'''
	http = urllib3.PoolManager()
	list_of_tools=['Function_library_overview.pdf',
					'function_library.py',
					'import_library.py',
					'TA_Advanced_Fit.ipynb',
					'TA_comparative_plotting_and_data_extraction.ipynb',
					'TA_Raw_plotting.ipynb',
					'TA_Raw_plotting_and_Simple_Fit.ipynb',
					'TA_single_scan_handling.ipynb',
					'Streak_camera_analysis.ipynb',
					'XES_Raw_plotting_and_Simple_Fit.ipynb']
	print('Now downloading the workflow tools')
	for i,f in enumerate(list_of_tools):
		url = "https://raw.githubusercontent.com/erdzeichen/KiMoPack/main/Workflow_tools/%s"%f
		print('Downloading Workflow Tools/%s'%f)
		with open(check_folder(path = 'Workflow_tools', current_path = os.getcwd(), filename = f), 'wb') as out:
			r = http.request('GET', url, preload_content=False)
			shutil.copyfileobj(r, out)
		if libraries_only and i==2:
			break
def download_all(single_tutorial=None):
	''' function loads workflow notebooks and example files and tutorials'''
	http = urllib3.PoolManager()
	if single_tutorial is None:
		download_notebooks()
		print('Now downloading the workflow tools and tutorials')
	else:
		download_notebooks(libraries_only=True)
		print('Libraries downloaded')
	list_of_example_data=['sample_1_chirp.dat',
							'Sample_2_chirp.dat',
							'sample_1.hdf5',
							'sample_2.hdf5',
							'Sample_1.SIA',
							'Sample_2.SIA',
							'XES_diff.SIA',
							'XES_on.SIA',
							'FeCM02-266nm-4mw-QB390-t6-G63-w450-s150-556ms-E100.dat',
							'FeCM02-266nm-4mw-QB390-t6-G63-w450-s150-556ms-E100_chirp.dat']
	print('Now downloading the example files')
	if (single_tutorial is None) or (single_tutorial == 'workflow'): #we do not use this to download data for Colab 
		for f in list_of_example_data:
			url = "https://raw.githubusercontent.com/erdzeichen/KiMoPack/main/Workflow_tools/Data/%s"%f
			print('Downloading Workflow Tools/Data/%s'%f)
			with open(check_folder(path = 'Workflow_tools'+os.sep+'Data', current_path = os.getcwd(), filename = f), 'wb') as out:
				r = http.request('GET', url, preload_content=False)
				shutil.copyfileobj(r, out)
	
	list_of_tutorials=['function_library.py',
						'Function_library_overview.pdf',
						'import_library.py',
						'KiMoPack_tutorial_0_Introduction.ipynb',
						'KiMoPack_tutorial_1_Fitting.ipynb',
						'KiMoPack_tutorial_2_Fitting.ipynb',
						'KiMoPack_tutorial_3_CompareFit.ipynb',
						'KiMoPack_tutorial_4_ScanHandling.ipynb',
						'KiMoPack_tutorial_5_MultiModal.ipynb']
	if single_tutorial is None: #we do not use this to download data for Colab 
		print('Now downloading tutorials')
		for f in list_of_tutorials:
			url = "https://raw.githubusercontent.com/erdzeichen/KiMoPack/main/Tutorial_Notebooks/%s"%f
			print('Downloading tutorial %s'%f)
			with open(check_folder(path = 'Tutorial_Notebooks', current_path = os.getcwd(), filename = f), 'wb') as out:
				r = http.request('GET', url, preload_content=False)
				shutil.copyfileobj(r, out)
	tutorial_data={'Compare':['TA_Ru-dppz_400nm_DCM_paral.hdf5','TA_Ru-dppz_400nm_H2O_paral.hdf5','UVvis_SEC_Rudppz_ACN.dat'],
				   'Master':['TA_Ru-dppz_400nm_ACN_paral.hdf5'],
				   'Fitting-1':['TA_Ru-dppz_400nm_ACN.SIA','TA_Ru-dppz_400nm_ACN_chirp.dat','TA_Ru-dppz_400nm_DCM.SIA','TA_Ru-dppz_400nm_DCM_chirp.dat','TA_Ru-dppz_400nm_H2O.SIA','TA_Ru-dppz_400nm_H2O_chirp.dat'],
					'Fitting-2':['TA_Ru-dppz_400nm_ACN.SIA','TA_Ru-dppz_400nm_ACN_chirp.dat'],
					'Introduction':['catalysis1.SIA','catalysis2.SIA','con_1.SIA','con_1_solved.hdf5','con_2.SIA','con_2_chirp.dat','con_3.SIA','con_4.SIA','con_5.SIA','con_6.SIA','con_6_chirp.dat','full_consecutive_fit.hdf5','full_consecutive_fit_with_GS.hdf5','sample_1_chirp.dat'],
					'Scan':['ACN_001.SIA','ACN_002.SIA','ACN_003.SIA','ACN_004.SIA','ACN_005.SIA','ACN_006.SIA','ACN_007.SIA','ACN_008.SIA','ACN_009.SIA','TA_Ru-dppz_400nm_ACN_mean.SIA','TA_Ru-dppz_400nm_ACN_mean_chirp.dat'],
					'MultiModal':['combined_optical_spectrum.SIA','XES_on.SIA']}
	url = "https://raw.githubusercontent.com/erdzeichen/KiMoPack/main/Tutorial_Notebooks/Data"
	for key in tutorial_data.keys():
		if single_tutorial is not None:   #this is a shortcut to download data fror Colab use
			if not key==single_tutorial:
				continue
		for f in tutorial_data[key]:
			if 'Master' in key:
				url = "https://raw.githubusercontent.com/erdzeichen/KiMoPack/main/Tutorial_Notebooks/Data/Compare/Master/%s"%f
				with open(check_folder(path = os.sep.join(['Tutorial_Notebooks','Data','Compare','Master']), current_path = os.getcwd(), filename = f), 'wb') as out:
					r = http.request('GET', url, preload_content=False)
					shutil.copyfileobj(r, out)
			else:
				url = "https://raw.githubusercontent.com/erdzeichen/KiMoPack/main/Tutorial_Notebooks/Data/%s/%s"%(key,f)
				with open(check_folder(path = os.sep.join(['Tutorial_Notebooks','Data',key]), current_path = os.getcwd(), filename = f), 'wb') as out:
					r = http.request('GET', url, preload_content=False)
					shutil.copyfileobj(r, out)
	tutorial_images=['Cor_Chirp.gif','Fig1_parallel_model.png','Fig2_consecutive_model.png','Fig3_complex_model.png','Intro_tutorial.png','Model_selection.jpg']
	if single_tutorial is None:
		for f in tutorial_images:
			url = "https://raw.githubusercontent.com/erdzeichen/KiMoPack/main/Tutorial_Notebooks/img/%s"%f
			with open(check_folder(path = os.sep.join(['Tutorial_Notebooks','img']), current_path = os.getcwd(), filename = f), 'wb') as out:
				r = http.request('GET', url, preload_content=False)
				shutil.copyfileobj(r, out)


def changefonts(weight='bold', font='standard', SMALL_SIZE=11, MEDIUM_SIZE=13, LARGE_SIZE=18):
	'''
	Small function that sets the matplotlib font sizes and fonts, written as conveniens to not need to remember all the 
	codes and what is names what. Calling the function will change the matplotlib *rc* settings
	
	Parameters
	------------
	
	weight : str, optional
		'bold' or  'normal'
	
	font : str, optional
		this is a meta switch that changes the family. known are:
		'standard'='DejaVu Sans'\n
		'arial'='Arial'\n
		'helvetica'= 'Helvetica'\n
		'garamond'='Garamond'\n
		'verdana'='Verdana'\n
		'bookman'='Bookman'\n
		'times'='Times New Roman'

	SMALL_SIZE : int, optional
		(DEFAULT = 11)\n 
		all written text, legend title and face size

	MEDIUM_SIZE : int, optional
		(DEFAULT = 13)\n 
		tick size and tick numbers

	LARGE_SIZE : int, optional
		(DEFAULT = 18)\n 
		axis titles, figure titles, axis labels

	'''
	global halfsize
	
	if halfsize:
		SMALL_SIZE=6 
		MEDIUM_SIZE=7 
		LARGE_SIZE=9
	
	font_dict = {
				'standard': {'weight': weight, 'size': SMALL_SIZE, 'family': 'DejaVu Sans'},
				'arial': {'weight': weight, 'size': SMALL_SIZE, 'family': 'Arial'},
				'helvetica': {'weight': weight, 'size': SMALL_SIZE, 'family': 'Helvetica'},
				'garamond': {'weight': weight, 'size': SMALL_SIZE, 'family': 'Garamond'},
				'verdana': {'weight': weight, 'size': SMALL_SIZE, 'family': 'Verdana'},
				'bookman': {'weight': weight, 'size': SMALL_SIZE, 'family': 'Bookman'},
				'times': {'weight': weight, 'size': SMALL_SIZE, 'family': 'Times New Roman'},
	}
	plt.rc('font', **font_dict[font])
	plt.rc('axes', titlesize=LARGE_SIZE, labelweight=weight)  # fontsize of the axes title
	plt.rc('axes', labelsize=LARGE_SIZE, labelweight=weight)  # fontsize of the x and y labels
	plt.rc('axes', linewidth=1)	 # linewidth of all axes
	plt.rc('axes', facecolor=(1, 1, 1, 0))
	plt.rc('xtick', labelsize=MEDIUM_SIZE)	# fontsize of the tick labels
	plt.rc('ytick', labelsize=MEDIUM_SIZE)	# fontsize of the tick labels
	plt.rc('legend', fontsize=SMALL_SIZE)  # legend fontsize
	plt.rc('legend', title_fontsize=SMALL_SIZE)
	plt.rc('legend', facecolor=(1, 1, 1, 0))
	plt.rc('legend', edgecolor=(1, 1, 1, 0))
	plt.rc('legend', framealpha=0)
	plt.rc('figure', titlesize=LARGE_SIZE)	# fontsize of the figure title
	plt.rc('figure', facecolor=(1, 1, 1, 0))  # fontsize of the figure title
changefonts() #we need to apply the font settings


def clean_double_string(filename, path=None):
	'''Stupid function that reads and changes!!! the file. It searchers for double lines and double dots and replaces them with single'''
	import re
	if path is None: path = os.path.dirname(os.path.realpath(__file__))
	with open(Path(os.sep.join([path, filename])), 'r+') as f:
		text = f.read()
		text = re.sub('--', '-', text)
		text = re.sub(r'\.+', '.', text)
		f.seek(0)
		f.write(text)
		f.truncate()
		
def mouse_move(event):
    x, y = event.xdata, event.ydata
    print(x, y)

def flatten(mainlist):
    return [entry for sublist in mainlist for entry in sublist]

def nearest_neighbor_method3(X, q):
	'''returns nearest neighbour value to q''' 
	X = X.T
	return np.argmin(np.sum((X - q) ** 2, axis=1))


def log_and(x, y, *args):
	"""Returns the logical and of all 2+ arguments."""
	result = np.logical_and(x, y)
	for a in args:
		result = np.logical_and(result, a)
	return result
	
	
def s2_vs_smin2(Spectral_points = 512, Time_points = 130, number_of_species = 3, fitted_kinetic_pars = 7, target_quality = 0.95):
	'''dfn is numerator and number of fitted parameters, dfd is denominator and number of degrees of freedom,
	F-test is deciding if a set of parameters gives a statistical significant difference. T-test is if a single parameter gives	 statistical difference. 
	Null hypothesis, all parameter are zero, if significant, the coefficients improve the fit
	the f-statistics compares the number of 
	"fitted parameter"=number of species*number of spectral points + number of kinetic parameter
	"free points"=number of species*number of spectral points*number of time points - fitted parameter
	within the target quality, meaning, what fraction do my variances need to have, so that I'm 100% * target_quality sure that they are different from zero'''
	data_points = Spectral_points*Time_points
	fitted_parameter = Spectral_points*number_of_species+fitted_kinetic_pars
	Free_points = data_points-fitted_parameter
	f_stat = scipy.stats.f.ppf(q = target_quality, dfn = fitted_parameter, dfd = Free_points)
	#print('fitted points:%g\n Free points:%g\n f-stats: %g'%(fitted_parameter,Free_points,f_stat))
	return 1+(fitted_parameter*f_stat/Free_points)	  


def GUI_open(project_list = None, path = None, filename_part = None, fileending = 'hdf5', sep = "\t", decimal = '.', 
			index_is_energy = False, transpose = False, sort_indexes = False, divide_times_by = None, 
			shift_times_by = None, external_time = None, external_wave = None, use_same_name = True, data_type = None, 
			units = None, baseunit = None, conversion_function = None):
								
	'''	This Function 
		1. opens a gui and allows the selection of multiple saved projects, which are returned as a list
		2. if given a list of project names opens them
		3. if given the word 'all', opens all files in a given folder 
		The general behavior is selected by the first parameter (project_list)
		
		it is designed to open combined files that contain both the wavelength and the time. (e.g. SIA files as recorded by Pascher instruments software) or hdf5 projects saved by this software
		There are however a lot of additional options to open other ascii type files and adapt their format internally
		Important, as default the parameter "fileending" selects hdf5 files only, which are used as project files (see :meth:`plot_func.TA.Save_project`) 
		for opening of other files the fileending parameter needs to be changed. 

		Parameters
		----------
		
		project_list : list (of str) or 'all', optional
			Give a list of filenames that will be opened and returned as a list of objects
			if the project list is 'all' then all files in the folder specified in path. The parameter "filename_part" 
			and "fileending" can be used to specify this selection 
		
		path : str or path object (optional)
			if path is a string without the operation system dependent separator, it is treated as a relative path, 
			e.g. data will look from the working directory in the sub director data. Otherwise this has to be a 
			full path in either strong or path object form.
		
		filename_part : str, optional
			This parameter is only used for the option 'all', the (Default) None means do nothing. if a string is given then only  
			files that start with this string will be read. 
		
		fileending : str, optional
			this string is used to select the filetype that is suppose to open. For the GUI, only these files will be shown,
			with the option 'all' this selects the files that will be read in the folder, 'hdf5' (Default) 
		
		sep : str (optional)
			is the separator between different numbers, typical are tap '\t' (Default) ,one or 
			multiple white spaces '\s+' or comma ','.
		
		decimal : str (optional) 
			sets the ascii symbol that is used for the decimal sign. In most countries this is '.'(Default) 
			but it can be ',' in countries like Sweden or Germany
		
		index_is_energy : bool (optional)
			switches if the wavelength is given in nm (Default) or in eV (if True), currently everything 
			is handled as wavelength in nm internally
		
		transpose : bool (optional)
			if this switch is False (Default) the wavelength are the columns and the rows the times.
		
		data_type: str (optional)
			data_type is the string that represents the intensity measurements. Usually this contains if absolute 
			of differential data. This is used for the color intensity in the 2d plots and the y-axis for the 1d plots
			
		units: str (optional)
			this is used to identify the units on the energy axis and to label the slices, recognized is 'nm', 'eV' and 'keV' 
			but if another unit like 'cm^-1' is used it will state energy in 'cm^-1'. Pleas observe that if you use the index_is_energy
			switch the program tries to convert this energy into wavelength. 
		
		baseunit: str (optional)
			this is used to identify the units on the developing/time axis. This is name that is attached to the index of the dataframe. 
			setting this during import is equivalent to ta.baseunit
		
		sort_indexes : bool (optional)
			For False (Default) I assume that the times and energies are already in a rising order. 
			with this switch, both are sorted again. 
		
		divide_times_by : None or float (optional) 
			here a number can be given that scales the time by an arbitary factor. This is actually dividing
			the times by this value. Alternatively there is the variable self.baseunit. The latter only affects 
			what is written on the axis, while this value is actually used to scale the times. None (Default) 
			ignores this
		
		shift_times_by : None, float (optional)
			This a value by which the time axis is shifted during import. This is a useful option of e.g. 
			the recording software does not compensate for t0 and the data is always shifted. 
			None (Default) ignores this setting
		
		external_time : None or str (optional)
			Here a filename extension (string) can be given that contains the time vector. 
			The file is assumed to be at the same path as the data and to contain a single 
			type of separated data without header. 
			If use_same_name = True (default)
			It assumes that this is the ending for the file. The filename itself is taken from the filename. 
			e.g. if samp1.txt is the filename and external_time='.tid' the program searches 
			samp1.tid for the times. The transpose setting is applied and sets where the times are 
			to be inserted (row or column indexes)
			If use_same_name = False this should be the file containing the vector for the time (in the same format as the main file)
			
		external_wave : None or str (optional) 
			Here a filename extension (string) can be given that contains the wavelength vector. 
			If use_same_name = True (default)
			The file is assumed to be at the same path as the data and to contain a single type 
			of separated data without header. This is the ending for the file. The filename itself 
			is taken from the filename. e.g. if samp1.txt is the filename and external_wave='.wav' 
			then the program searches samp1.wav for the wavelength. The transpose setting is applied 
			and sets where the wavelength are to be inserted (columns or row indexes)
			If use_same_name = False
			this should be a full filename that contains the vector
			
		use_same_name : bool, optional
			this switches if the external filename included the loaded filename or is a separate file True(default)
		
		
		conversion_function: function(optional)
			function that receives should have the shape:
			return pandas Dataframe with time/frames  in rows and wavelength/energy in columns,
			The function is tested to accept (in that order) a 
			my_function(filename, external_time,external_wave), 
			my_function(filename, external_time), 
			my_function(filename,external_wave), 
			my_function(filename) and 
			return: the dataframe ds with the time_axis as rows and spectral axis as columns 
			if the ds.index.name ia not empty the "time axis" is in to that name the spectral axis is in ds.columns.name
			the return is investigated if it is one, two, or three things. 
			if two are returned then the second must be the name of what the intensity axis is. This value will then be set to data_type
			if three are returned the third is the baseunit (for the time axis) this allows to use the automatic naming in ps or nanosecond
			If the values units, data_type or baseunit are (manually) set in the import function the corresponding entries in
			datafram will be overwritten
			shift_times_by and divide_times_by will be applied if not None (useful to adjust for offset before chirp correction)
		
		Returns
		--------------
		
		List of opened TA objects
		
		Examples
		--------------
		
		>>> import plot_func as pf
		>>> project_list=pf.GUI_open() #start the GUI to open project Files
		>>> project_list=pf.GUI_open(fileending='SIA') #start the GUI to open SIA Files
		
		Opening a list of files using the file names 
		
		>>> project_list=pf.GUI_open(project_list = ['file1.SIA', 'file2.SIA'])
		
		Opening all files in the folder "all_data" (relative to where the notebook is with the ending "hdf5" 
		
		>>> project_list=pf.GUI_open('all',path="all_data")
		
		Opening a list of files with external time vector (of the same name) so it looks for a data
		file "file1.txt" and a file with the time information "file1.tid"
		
		>>> project_list=pf.GUI_open(project_list = ['file1.txt', 'file2.txt'], external_time = 'tid')
		'''
	if project_list is None:
		root_window = tkinter.Tk()
		root_window.withdraw()
		root_window.attributes('-topmost',True)
		root_window.after(1000, lambda: root_window.focus_force())
		path_list = filedialog.askopenfilename(initialdir=os.getcwd(),multiple=True,filetypes=[('TA project files','*.%s'%fileending)])									 
		if project_list is None:
			project_list=[]
	elif project_list=='all':
		scan_path=check_folder(path = path, current_path = os.getcwd())
		if filename_part is not None:#we specified a specific name and want only the files with this name in it
			path_list = sorted([os.path.join(scan_path, name) for name in os.listdir(scan_path) if
									name.endswith(fileending) and filename_part in name])
		else:#we have not specified a specific name and want all files in the folder
			path_list = sorted([currentFile for currentFile in scan_path.glob("*.%s"%fileending)])
	else:
		if len(project_list)<1:	
			raise ValueError('The use_gui switch is ment to bypass the gui, but you still need at least some files as a list')
		else:
			if isinstance(project_list, str):project_list=[project_list]
			if not hasattr(project_list, '__iter__'):project_list=[project_list]
			path_list = []
			for filename in project_list:
				ta=check_folder(path=path, filename=filename, current_path=os.getcwd())
				path_list.append(ta)
	return_list = []
	for entrance in path_list:
		try:
			listen=os.path.split(entrance)
			path=os.path.normpath(listen[0])
			filename=listen[1]
			ta = TA(filename = filename, path = path, sep = sep, decimal = decimal, 
					index_is_energy = index_is_energy, transpose = transpose, sort_indexes = sort_indexes, 
					divide_times_by = divide_times_by, shift_times_by = shift_times_by, external_time = external_time, 
					external_wave = external_wave, use_same_name = use_same_name, data_type = data_type, units = units, 
					baseunit = baseunit, conversion_function = conversion_function)					
			return_list.append(ta)
		except:
			print('Problem with entrance:\n %s'%entrance)
	return return_list


def check_folder(path = None, current_path = None, filename = None):
	'''Helper function using robust path determination.\n 
	In any case if a valif file name is given it is attached to the total path\n
	The path can be string or windows/linux path or pure path or byte type paths.\n
	paths that do not exists (including parents) are created\n
	1. if path is given absolute, it is returned\n_colors
	2. if path is a string (relative) the current_path + path is returned.\n
	3. if current_path is not absolute or None, the current working directory is assumed as path.\n
	4. IF all is None, the current working directory is returned
	
	Parameters
	-----------
	
	path : str, purePath, absolute or relative, optional
		the final part of the path used
	
	current_path : None, str, purePath, absolute, optional
		path that sits before the "path variable, is filled with current working directory if left None
	
	filename: None, str, optional
		attached after path and returned if not None
		
	'''
	
	if isinstance(path,bytes):
		path = '%s'%path
	if path is not None:
		path = pathlib.Path(path)

	if isinstance(current_path, bytes):
		current_path = '%s'%current_path
	if current_path is not None:
		current_path=pathlib.Path(current_path)
		
	if isinstance(filename, bytes):
		filename='%s'%filename
	if filename is not None:
		filename = pathlib.Path(filename)
	if path is None:
		if current_path is None:
			directory = Path.cwd()
		elif current_path.is_absolute():
			directory=current_path
		else:
			print('attention, current_path was given but not absolute, replaced by cwd')
			directory = Path.cwd()
	elif path.is_absolute():
		directory = path
	else:
		if current_path is None:
			directory = Path.cwd().joinpath(path)
		elif current_path.is_absolute():
			directory = current_path.joinpath(path)
		else:
			print('attention, current_path was given but not absolute, replaced by cwd')
			directory = Path.cwd().joinpath(path)
	directory.mkdir( parents=True, exist_ok=True)
	if filename is None:
		return directory
	else:
		return directory.joinpath(filename)



def rebin(ori_df,new_x):
	'''interpolation of values to new index'''
	if isinstance(ori_df,pandas.DataFrame):
		dum={'dummy':new_x}
		new_df=pandas.DataFrame(dum,index=new_x)
		for col in ori_df.columns:
			new_df[col]=np.interp(new_x,ori_df.index.values.astype('float'),ori_df[col].values)
		new_df=new_df.drop(['dummy'],axis=1)
		return new_df
	elif isinstance(ori_df,pandas.Series):
		new_df=np.interp(new_x,ori_df.index.values.astype('float'),ori_df.values)
		return pandas.Series(new_df,index=new_x)


def savitzky_golay(y, window_size, order, deriv=0, rate=1):
	'''Ported from a previous function'''
	return savgol_filter(x=y, window_length=window_size, polyorder=order, deriv=deriv, delta=rate)


def Frame_golay(df, window=5, order=2,transpose=False):
	'''Convenience method that returns the Golay smoothed data for each column (DataFrame) or the series
	
	Parameters
	-----------
	
	df : pandas.DataFrame,pandas.Series
		the DataFrame that has to be interpolated
	
	window_size : int,optional
		5(Default) an integer that indicates how many units are to be interpolated
		
	order : int, optional
		2 (Default) an integer that indicates what orderpolynoninal is to be used to interpolate the points. 
		order=1 effectively turns this into a floating average
		
	transpose : bool,optional 
		in which orientation is the interpolation to be done. Default is in within the column (usually timepoints)
	
	Returns
	---------
	
	pandas.DataFrame or pandas.Series
		DataFrame or Series with the interpolation applied
	
	'''
	#df=df.fillna(0)

	if transpose:
		df=df.T
	window= min(len(df.index.values), window)  # Ensure it's an odd number
	order= min(len(df.index.values), order)
	if window % 2 == 0:
		window -= 1
	if isinstance(df,pandas.DataFrame):
		for col in df.columns:
			try:
				df.loc[:,col]=savitzky_golay(df.loc[:,col].values, window, order)
			except Exception as e:
				print(col)
				print('was not smoothed')
				print(e)
		if transpose:
			df=df.T
		return df
	elif isinstance(df,pandas.Series):
		return pandas.Series(savitzky_golay(df.values, window, order),index=df.index)
	else:
		raise TypeError('must be series or DataFrame')


def find_nearest(arr,value,con_str=False):
	'''returns the value in the array closest to value'''
	return arr[find_nearest_index(arr,value,con_str=False)]	


def find_nearest_index(arr,value,con_str=False):
	'''returns the index in the array closest to value (the first one'''
	if con_str:
		temp_array=np.array(arr,dtype='float')
		idx = (np.abs(temp_array-value)).argmin()
	else:
		idx = (np.abs(arr-value)).argmin()
	return idx


def rise(x,sigma=0.1,begin=0):
	'''	 my own implementation of the instrument response function. 
		 Based upon an error function from 0 to 1. 
		 Sigma is the width (after which it has 50%) 
		 and begin is 10% of height'''
	return (erf((x-begin-sigma)*np.sqrt(2)/(sigma))+1)/2


def gauss(t,sigma=0.1,mu=0):
	'''Gauss function'''
	y=np.exp(-0.5*((t-mu)**2)/sigma**2)
	y/=sigma*np.sqrt(2*np.pi)
	return y


def norm(df):
	'''Min max norming of a dataframe'''
	return df.apply(lambda x: (x - np.min(x)) / (np.max(x) - np.min(x)))


def shift(df,name = None,shift = None):
	'''Shifts a dataframe along the columns, interpolate and then resample'''
	if name is None:name = df.columns
	if isinstance(name,type('hello')):name = [name]
	for nam in name:
		ori_dat = df[nam].values
		ori_en = np.array(df.index,dtype = 'float')
		if ori_en[0]>ori_en[1]:#oh we have inverse order
			dat = np.interp(ori_en[::-1],ori_en[::-1]+shift,ori_dat[::-1])
			dat = dat[::-1]
		else:
			dat = np.interp(ori_en,ori_en+shift,ori_dat)
		df[nam] = dat
	return df
	

def colm(k,cmap = standard_map):
	'''If a colour map is given, this is used.'''
	if isinstance(cmap,type(cm.jet)) or isinstance(cmap,type(cm.viridis)):
		if hasattr(k,'__iter__'):
			if min(k) >0:#we got a color offset
				mini = min(k)/(min(k)+1)
			else:
				mini = 0
			out = [ cmap(x) for x in np.linspace(mini, 1, len(k)+1) ]
			out = out[:-1]
			return out
		else:# get me 10 colors
			out = [cmap(x) for x in np.linspace(0, 1, 10)]
			ret = out[k]
			return ret
	else:  #we assume it is a iterable thingy
		if not hasattr(k,'__iter__'):k = [k]
		if isinstance(cmap,pandas.DataFrame):
			out = [list(cmap.iloc[ent,:].values) for ent in k]
		elif isinstance(cmap,np.ndarray):
			out = [cmap[int(ent),:] for ent in k]
		elif isinstance(cmap,list):
			out = [cmap[int(ent)] for ent in k]
		else:
			print('didn\'t find the right ')
		return out


def Summarize_scans(list_of_scans = None, path_to_scans = 'Scans', list_to_dump = 'range', window1 = None, window2 = None, 
					save_name = 'combined.SIA', fileending = 'SIA', filename_part = 'Scan', return_removed_list = False, 
					sep = "\t", decimal = '.', index_is_energy = False, transpose = False, sort_indexes = False, 
					divide_times_by = None, shift_times_by = None, external_time = None, external_wave = None, use_same_name = True,
					return_ds_only=False, data_type = None, units = None, baseunit = None, conversion_function = None, fitcoeff = None,
					base_TA_object = None, value_filter = None, zscore_filter_level = None, zscore_in_window = True,
					dump_times = True, replace_values = None, drop_scans = False):
	'''
	Average single scans. Uses single scans of the data set and plots them as average after different conditions. Usually one defines one or two windows in which the intensity is integrated. This integrated number is then displayed for each scan in the list. There are different tools to select certain scans that are excluded from the summary. These are defined in the list_to_dump. This list can take either be a list with the number, or a string with the words 'single' or 'range' (see below) 
	

	Parameters
	-----------
	
	list_of_scans : None, 'gui' or list
		'gui' (choose scans via gui)\n
		None (Default) load scan files from the specified folder (path_to_scans) with the specified file-ending
		(file_ending), if filename_part is a string than only files with this string in the name are taken\n
		list of names (strings) loads this list of files
		list of integers (that will be directly attached to the filename_part) to form the file name
	
	path_to_scans : None, str or path object, optional
		specify relative or absolute path to the scan-files (Default:'Scans')
	
	file_ending : str, optional
		specify the file extension of the single scan files. The Gui will only show this fileending
		(Default: '.SIA')

	filename_part : str
		specify a part of the string included in all scan-files (Default: 'Scan')
		
	window1: None or list of 4 floats, optional
		window in time and wavelength over which each scan is averaged.\n 
		window must have the shape [start time, end time, start wavelength, end wavelength]
		(Default: None)
	
	window2: list of 4 floats, optional
		window in time and wavelength over which each scan is averaged.\n 
		window must have the shape [start time, end time, start wavelength, end wavelength]
		(Default: None) IF not given then only one window will be displayed
		
	list_to_dump : list, 'single' or 'range', or None, optional
		takes a list of scans to be excluded from the average, this list can be indexes (order) 
		in which the scans come, or a list of names. if this is given as a list the option "range" 
		is offered, which allows to add additional selection to the cut.\n 
		**'single'** allows you (in a GUI) to click on single points in plotted window1 or two that 
		is to be removed, useful for spike removal and makes only sense in conjunction with at least 
		a defined window1, if none is defined window1 = [0.5,10,300,1200] will be set automatically. 
		A right click removes the last selection a middle click applies it. An empty middle click 
		(without selecting anything) finishes the gui\n
		**'range'** allows you (in a GUI) to click and define regions.\n
		first left click is the left side of the window, second left click the ride side of the window. 
		Third left click the left side of the second window,... A right click removes the last set point.
		a middle click finishes and applies the selection\n 
		An **empty middle click** (without selecting anything) finishes the gui\n
		useful for spike removal and definition of exclusion region (e.g. where the sample died) 
		This makes only sense in conjunction with at least a defined window1 , 
		if none is defined window1 = [0.5,10,300,1200] will be set automatically
		if None then it is not filtered, but simply returned										 

	data_type: str (optional)
		data_type is the string that represents the intensity measurements. Usually this contains if absolute 
		of differential data. This is used for the color intensity in the 2d plots and the y-axis for the 1d plots
		
	units: str (optional)
		this is used to identify the units on the energy axis and to label the slices, recognized is 'nm', 'eV' and 'keV' 
		but if another unit like 'cm^-1' is used it will state energy in 'cm^-1'. Pleas observe that if you use the index_is_energy
		switch the program tries to convert this energy into wavelength. 

	baseunit: str (optional)
		this is used to identify the units on the developing/time axis. This is name that is attached to the index of the dataframe. 
		setting this during import is equivalent to ta.baseunit

	save_name : str, optional
		specify name for saving the combined scans (Default) 'combined.SIA')
	
	return_removed_list : bool, optional
		(Default) False, returns the list of removed scans instead of the averaged data set. (this list could then be given as "list_to_dump" to get the averaged datafile too. If a file name is given for saved file (which is Default) then the file is saved anyways.
	
	sep : str (optional)
		is the separator between different numbers, typical are tap (Backslash t) (Default) ,one or 
		multiple white spaces 'backslash s+' or comma ','.
	
	decimal : str (optional) 
		sets the ascii symbol that is used for the decimal sign. In most countries this is '.'(Default) 
		but it can be ',' in countries like Sweden or Germany
	
	index_is_energy : bool (optional)
		switches if the wavelength is given in nm (Default) or in eV (if True), currently everything 
		is handled as wavelength in nm internally
	
	transpose : bool (optional)
		if this switch is False (Default) the wavelength are the columns and the rows the times.
	
	sort_indexes : bool (optional)
		For False (Default) I assume that the times and energies are already in a rising order. 
		with this switch, both are sorted again. 
	
	divide_times_by : None or float (optional) 
		here a number can be given that scales the time by an arbitary factor. This is actually dividing
		the times by this value. Alternatively there is the variable self.baseunit. The latter only affects 
		what is written on the axis, while this value is actually used to scale the times. None (Default) 
		ignores this
	
	shift_times_by : None, float (optional)
		This a value by which the time axis is shifted during import. This is a useful option of e.g. 
		the recording software does not compensate for t0 and the data is always shifted. 
		None (Default) ignores this setting
	
	external_time : None or str (optional)
		Here a filename extension (string) can be given that contains the time vector. 
		The file is assumed to be at the same path as the data and to contain a single 
		type of separated data without header. 
		If use_same_name = True (default)
		It assumes that this is the ending for the file. The filename itself is taken from the filename. 
		e.g. if samp1.txt is the filename and external_time='.tid' the program searches 
		samp1.tid for the times. The transpose setting is applied and sets where the times are 
		to be inserted (row or column indexes)
		If use_same_name = False this should be the file containing the vector for the time (in the same format as the main file)
		
	external_wave : None or str (optional) 
		Here a filename extension (string) can be given that contains the wavelength vector. 
		If use_same_name = True (default)
		The file is assumed to be at the same path as the data and to contain a single type 
		of separated data without header. This is the ending for the file. The filename itself 
		is taken from the filename. e.g. if samp1.txt is the filename and external_wave='.wav' 
		then the program searches samp1.wav for the wavelength. The transpose setting is applied 
		and sets where the wavelength are to be inserted (columns or row indexes)
		If use_same_name = False
		this should be a full filename that contains the vector
			
	use_same_name : bool, optional
		this switches if the external filename included the loaded filename or is a separate file True(default)
		
	conversion_function: function(optional)
		function that receives should have the shape:
		return pandas Dataframe with time/frames  in rows and wavelength/energy in columns,
		The function is tested to accept (in that order) a 
		my_function(filename, external_time,external_wave), 
		my_function(filename, external_time), 
		my_function(filename,external_wave), 
		my_function(filename) and 
		return: the dataframe ds with the time_axis as rows and spectral axis as columns 
		if the ds.index.name ia not empty the "time axis" is in to that name the spectral axis is in ds.columns.name
		the return is investigated if it is one, two, or three things. 
		if two are returned then the second must be the name of what the intensity axis is. This value will then be set to data_type
		if three are returned the third is the baseunit (for the time axis) this allows to use the automatic naming in ps or nanosecond
		If the values units, data_type or baseunit are (manually) set in the import function the corresponding entries in
		datafram will be overwritten
		shift_times_by and divide_times_by will be applied if not None (useful to adjust for offset before chirp correction)
	
	return_ds_only: boolean,(optional)
		if False (Dafault) returns a TA object, otherwise just a DataFrame
	
	fitcoeff: list, optional
		these should be the shirp parameteres that are to be applied to all sub scans in the list. 
	
	base_TA_object: TA object, optional
		instead of the fit_coefficients a Ta object can be provided that is then used as a template, meaning that the scattercuts and bordercuts will be applied before the filtering.
		
	value_filter : None, float or iterable with two entries, optional
		if float, everything above that value or below -abs(value_filter) will be filtered replaced with replace_values
		if iterable, then first is lower treshold, second is upper treshold
		
	zscore_filter_level : float, optional
		if this value is set then the manual selection will be replaced with an automatic filter, the following options, dump_times = True, 
		replace_values = None, drop_scans = False decide what is done to the values that are filtered
		typical value would be e.g. 3
		
	zscore_in_window : bool, 
		decides if the filter is applied in the windows or over the whole matrix (using statistics on the values)
		
	dump_times : bool,optional 
		Standard True means that if the zscore filter filters a file the bad time is droped for the average
	
	replace_values : None, float, optional
		if dump times is False the values will be replaced with this value. = None, drop_scans = False
		
	drop_scans : bool,optional
		Default: = False. This is the harshest type to filter and means that the whole scan is dropped
	
	Returns
	---------
	
	TA object if return_ds_only is False(Default) averaged dataset (ds) of the selected scans or
	(if return_removed_list = True) the list of removed scans. 
	
	Examples
	----------
	
	Use use a range to select the rejected scans, look on the scans by integrating the window 0.5ps to 1ps and 450nm to 470nm
	
	>>> import plot_func as pf #import the module
	>>> window1=[0.5,1,450,470] #define the window
	>>> #use a 'GUI' to select the files
	>>> pf.Summarize_scans(list_of_scans='gui',window1=window1) 
	>>> #use all scans in the subfolder scans that have the word 'Scan' in them and use the ending 'SIA'
	>>> pf.Summarize_scans(path_to_scans = 'Scans', filepart_name = 'Scan', window1=window1)
	>>> #This does the same as these are standard
	>>> pf.Summarize_scans(window1=window1)

	'''
	if (base_TA_object is not None) and (conversion_function is None):
		if units is None:units=base_TA_object.ds.columns.name
		if baseunit is None:baseunit=base_TA_object.ds.index.name
	debug = True
	if list_of_scans is None:
		scan_path=check_folder(path = path_to_scans, current_path = os.getcwd())
		if filename_part is not None:#we specified a specific name and want only the files with this name in it
			list_of_scans = sorted([os.path.join(scan_path, name) for name in os.listdir(scan_path) if
									name.endswith(fileending) and filename_part in name])
		else:#we have not specified a specific name and want all files in the folder
			list_of_scans = sorted([currentFile for currentFile in scan_path.glob("*.%s"%fileending)])
	elif list_of_scans == 'gui':
		root_window = tkinter.Tk()
		root_window.withdraw()
		root_window.attributes('-topmost',True)
		root_window.after(1000, lambda: root_window.focus_force())
		path_list = filedialog.askopenfilename(initialdir = os.getcwd(),multiple = True,filetypes = [('Raw scan files',"*.%s"%fileending)])
		list_of_scans = path_list
	elif not hasattr(list_of_scans,'__iter__'):
		raise ValueError('We need something to iterate for the list')
	
	if not isinstance(list_of_scans[0],TA):#we do not have opened file but most likely a list of names
		try:
			list_of_projects = []
			for entrance in list_of_scans:
				listen = os.path.split(entrance)
				path = os.path.normpath(listen[0])
				filename = listen[1]
				new_ds=TA(filename = filename,path = path, sep = sep, decimal = decimal, 
											index_is_energy = index_is_energy, transpose = transpose, 
											sort_indexes = sort_indexes, divide_times_by = divide_times_by, 
											shift_times_by = shift_times_by, external_time = external_time, 
											external_wave = external_wave, use_same_name = use_same_name, 
											data_type = data_type, units = units, baseunit = baseunit, 
											conversion_function = conversion_function).ds
				if base_TA_object is None:
					if fitcoeff is not None:
						new_ds=Fix_Chirp(ds=new_ds,fitcoeff=fitcoeff)
					list_of_projects.append(new_ds.values)
				else:
					if fitcoeff is not None:
						new_ds=Fix_Chirp(ds=new_ds,fitcoeff=fitcoeff)
					try:
						new_ds=sub_ds(new_ds, ignore_time_region = base_TA_object.ignore_time_region, wave_nm_bin = base_TA_object.wave_nm_bin, baseunit = base_TA_object.baseunit, 
										scattercut = base_TA_object.scattercut, bordercut = base_TA_object.bordercut, timelimits = base_TA_object.timelimits, time_bin = base_TA_object.time_bin, 
										equal_energy_bin = base_TA_object.equal_energy_bin)
						if (base_TA_object.wave_nm_bin is not None) or (base_TA_object.equal_energy_bin is not None):
							print('in the original TA objec the data was rebinned, which is now also done for the single scans. To avoid that use "ta.wave_nm_bin = None" and / or "ta.equal_energy_bin = None" before handing it to base_TA_object')
					except:
						print('applying the base_TA_object slices failed')
					list_of_projects.append(new_ds.values)
			if base_TA_object is None:
				ds = TA(filename = filename,path = path,  sep = sep, decimal = decimal, 
						index_is_energy = index_is_energy, transpose = transpose, sort_indexes = sort_indexes, 
						divide_times_by = divide_times_by, shift_times_by = shift_times_by, 
						external_time = external_time, external_wave = external_wave, 
						use_same_name = use_same_name, data_type = data_type, units = units, 
						baseunit = baseunit, conversion_function = conversion_function).ds

			else:
				ds=base_TA_object.ds
				ds = sub_ds(ds, ignore_time_region = base_TA_object.ignore_time_region, wave_nm_bin = base_TA_object.wave_nm_bin, baseunit = base_TA_object.baseunit, 
							scattercut = base_TA_object.scattercut, bordercut = base_TA_object.bordercut, timelimits = base_TA_object.timelimits, time_bin = base_TA_object.time_bin, 
							equal_energy_bin = base_TA_object.equal_energy_bin)
			######################
			try:
				list_of_projects = np.transpose(np.array(list_of_projects),(1, 2, 0))
			except:
				print('the stacking of the scans failed, are you sure that all are have the same shape')
			#######################
		except:
			raise ValueError('Sorry did not understand the project_list entry, use GUI_open to create one')
	else:
		try:
			list_of_projects = []
			list_of_scans_names = []
			for entrance in list_of_scans:
				list_of_projects.append(entrance.ds.values)
				list_of_scans_names.append(entrance.filename)
			if base_TA_object is None:
				ds = list_of_scans[0]
			else:
				ds=base_TA_object.ds
			list_of_scans = list_of_scans_names
			
			##########################
			try:
				list_of_projects = np.transpose(np.array(list_of_projects),(1, 2, 0))
			except:
				print('the stacking of the scans failed, are you sure that all are have the same shape')
			#########################
		except:
			raise ValueError('Sorry did not understand the project_list entry, use GUI_open to create one')
			
	if window1 is None:
		window1 = [ds.index.values.min(),ds.index.values.max(),ds.columns.values.min(),ds.columns.values.max()]
	
	#### automatic filtering#####
	if (zscore_filter_level is not None) or (value_filter is not None):
		if replace_values is not None:
			cut_bad_times=False
		if replace_values is None:
			replace_values = np.nan
		dataset=list_of_projects
		if value_filter is not None:
			if hasattr(value_filter,'__iter__'):
				lowervalue=value_filter[0]
				uppervalue=value_filter[1]
			else:
				uppervalue = np.abs(value_filter)
				lowervalue = -np.abs(value_filter)
			outside_range=np.invert(log_and(dataset>lowervalue,dataset<uppervalue))
			if dump_times:#this is default
				outside_range=np.tile(outside_range.all(axis=1,keepdims=True),(1,dataset.shape[1],1))
			elif drop_scans:
				outside_range=np.tile(outside_range.any(axis=1,keepdims=True),(1,dataset.shape[1],1))
				outside_range=np.tile(outside_range.any(axis=0,keepdims=True),(dataset.shape[0],1,1))
			dataset[outside_range]=replace_values
		if zscore_filter_level is not None:
			if zscore_in_window:
				window1_index = [find_nearest_index(ds.index.values,window1[0]),find_nearest_index(ds.index.values,window1[1]),find_nearest_index(ds.columns.values,window1[2]),find_nearest_index(ds.columns.values,window1[3])]
				vector=np.nanmean(np.nanmean(dataset[window1_index[0]:window1_index[1],window1_index[2]:window1_index[3],:],axis=0),axis=1)
				good=log_and(vector>(np.nanmean(vector) - zscore_filter_level*np.nanstd(vector)),vector<(np.nanmean(vector) + zscore_filter_level*np.nanstd(vector)))
				if window2 is not None:
					window2_index = [find_nearest_index(ds.index.values,window2[0]),find_nearest_index(ds.index.values,window2[1]),find_nearest_index(ds.columns.values,window2[2]),find_nearest_index(ds.columns.values,window2[3])]
					vector=np.nanmean(np.nanmean(dataset[window2_index[0]:window2_index[1],window2_index[2]:window2_index[3],:],axis=0),axis=1)
					good2=log_and(vector>(np.nanmean(vector) - zscore_filter_level*np.nanstd(vector)),vector<(np.nanmean(vector) + zscore_filter_level*np.nanstd(vector)))
					good=log_and(good,good2)
				
				print(dataset.shape)
				print(good.shape)
				dataset[:,:,np.invert(good)]=replace_values
			else:
				mean=np.nanmean(dataset,axis=2)
				var=np.nanstd(dataset,axis=2)
				lower=(mean - zscore_filter_level*var).T
				upper=(mean + zscore_filter_level*var).T
				lower=np.array([lower for i in range(dataset.shape[2])]).T
				upper=np.array([upper for i in range(dataset.shape[2])]).T
				outside_range=np.invert(log_and(dataset>lower,dataset<upper))
				if drop_scans:
					outside_range=np.tile(outside_range.any(axis=1,keepdims=True),(1,dataset.shape[1],1))
					outside_range=np.tile(outside_range.any(axis=0,keepdims=True),(dataset.shape[0],1,1))
				elif dump_times:
					outside_range=np.tile(outside_range.any(axis=1,keepdims=True),(1,dataset.shape[1],1))
				dataset[outside_range]=replace_values
		list_of_projects=dataset
		
	#############manual filtering################
	else:	
		if baseunit is None:baseunit=ds.index.name
		if units is None:units=ds.columns.name
		if list_to_dump is not None:
			if list_to_dump == 'single':
				print('we will use a gui to select single scans to extract')
			elif list_to_dump == 'range':
				print('we will use a gui to select the first and last scan to remove')
			else:
				if not hasattr(list_to_dump,'__iter__'):#we have only a single number/name in there
					list_to_dump = [list_to_dump]
				filenames_to_dump = []
				for entry in list_to_dump:
					try:
						filenames_to_dump.append(list_of_scans[entry].filename) #list_of_scans is a list of TA objects that have filename and if entry is an index of this list this goes well
					except:
						filenames_to_dump.append(entry)# we assume it is already a filename
				list_to_dump = []
				for filename in filenames_to_dump:
					list_to_dump.append(list_of_scans.index(filename))
			for i in range(30):#we make a maximum of 30 rounds
				window1_index = [find_nearest_index(ds.index.values,window1[0]),find_nearest_index(ds.index.values,window1[1]),find_nearest_index(ds.columns.values,window1[2]),find_nearest_index(ds.columns.values,window1[3])]
				series1 = pandas.Series(list_of_projects[window1_index[0]:window1_index[1],window1_index[2]:window1_index[3],:].mean(axis = (0,1)))
				series1.name = '%.3g:%.3g %s at %.1f:%.1f %s'%(window1[0],window1[1],baseunit,window1[2],window1[3],units)
				if window2 is not None:
					window2_index = [find_nearest_index(ds.index.values,window2[0]),find_nearest_index(ds.index.values,window2[1]),find_nearest_index(ds.columns.values,window2[2]),find_nearest_index(ds.columns.values,window2[3])]
					series2 = pandas.Series(list_of_projects[window2_index[0]:window2_index[1],window2_index[2]:window2_index[3],:].mean(axis = (0,1)))
					series2.name = '%.3g:%.3g %s at %.1f:%.1f %s'%(window2[0],window2[1],baseunit,window2[2],window2[3],units)
					fig,(ax,ax2) = plt.subplots(2,1,sharex = True,figsize = (16,12))
					series1.plot(ax = ax,color = colm(1),use_index = False)
					series2.plot(ax = ax2,color = colm(3),use_index = False)
					if len(series1) >15:
						gol=Frame_golay(series1,window=11,order=1)
						gol.plot(ax=ax,use_index=False,color=colm(2))
						ax.fill_between(x=range(len(series1)), y1=gol-series1.var(), y2=gol+series1.var(),color=colm(2),alpha=0.3)	
						gol=Frame_golay(series2,window=11,order=1)
						gol.plot(ax=ax2,use_index=False,color=colm(4))
						ax2.fill_between(x=range(len(series1)), y1=gol-2*series1.var(), y2=gol+2*series1.var(),color=colm(4),alpha=0.3)	
				else:
					fig,ax=plt.subplots(1,1,sharex=True,figsize=(16,12))
					series1.plot(ax=ax,color=colm(1),use_index=False)
					if len(series1) >15:
						gol=Frame_golay(series1,window=11,order=1)
						gol.plot(ax=ax,use_index=False,color=colm(2))
						#try:
						ax.fill_between(x=range(len(series1)), y1=gol-2*np.nanvar(series1.values), y2=gol+2*np.nanvar(series1.values),color=colm(2),alpha=0.3)
						#except:
						#	pass
				if list_to_dump == 'single':
					ax.set_title('click on the scans that should be dropped\n left click to chose, right click to delete last point, middle click finishes selection\n an empty middle click ends the process')
					polypts=np.asarray(plt.ginput(n=int(len(series1)/2),timeout=300, show_clicks=True,mouse_add=1, mouse_pop=3, mouse_stop=2))
					if len(polypts)<1:break
					to_remove=[int(a) for a in np.array(polypts)[:,0]]
					remove=pandas.Series(np.arange(len(series1)))+1	
					remove[to_remove]=0
					to_remove=list(remove[remove<1].index.values)
					to_keep=list(remove[remove>1].index.values)
				elif (list_to_dump == 'range') or (i>0):
					ax.set_title('click on the first and last scan to be removed, repeat as long as necessary\n an empty middle click ends the process')
					polypts=np.asarray(plt.ginput(n=int(len(series1)/2),timeout=300, show_clicks=True,mouse_add=1, mouse_pop=3, mouse_stop=2))
					if len(polypts)<1:break
					polypts=np.array(polypts)[:,0]
					remove=pandas.Series(np.arange(len(series1)))+1
					for i in range(int(len(polypts)/2)):
						remove.loc[polypts[2*i]:polypts[2*i+1]]=0
					to_remove=list(remove[remove<1].index.values)
					to_keep=list(remove[remove>1].index.values)
				elif i == 0:
					to_keep=list(range(len(series1)))
					to_keep.remove(list_to_dump)
				else:
					raise ValueError('Something is weired')
				list_of_projects=list_of_projects[:,:,to_keep]
				plt.close('all')
	plt.close('all')
	if 0:
		try:
			df=pandas.DataFrame(np.any(np.isnan(dataset),axis=1),index=ds.index)
			plot2d(df,levels = 2,use_colorbar = False,intensity_range=[0,1],title='rejected are red')
		except Exception as e:
			print('plotting of filtered went wrong\n Error message was: \n')
			print(e)
	try:
		ds=pandas.DataFrame(np.nanmean(list_of_projects,axis=2),index=ds.index,columns=ds.columns)
	except:
		print('nanmean failed, assume single wavelength')
		ds=pandas.DataFrame(np.mean(list_of_projects),index=ds.index,columns=ds.columns)																	  
	if save_name is not None:
		path = str(check_folder(path=path,filename=save_name))
		ds.to_csv(path,sep='\t')
		ta=TA(path)
	else:
		path = str(check_folder(path=path,filename='temp_combined.SIA'))
		ds.to_csv(path,sep='\t')
		ta=TA(path)
		try:
			os.remove(path)
		except:
			print('could not remove temp_combined.SIA')
	if return_ds_only:
		return ds
	elif return_removed_list:
		return filenames_to_dump
	else:
		if base_TA_object is not None:
			ta=base_TA_object.Copy()
			ta.ds_ori=ds
			ta.ds=ds
		return ta
		

def sub_ds(ds, times = None, time_width_percent = 0, ignore_time_region = None, drop_ignore=False, wave_nm_bin = None, 
			baseunit = None, scattercut = None, drop_scatter=False, bordercut = None, timelimits = None, wavelength_bin = None, 
			wavelength = None, time_bin = None, equal_energy_bin = None, from_fit = False):
	'''This is the main function that creates all the slices of the data matrix
	
	Parameters
	---------------
	
	ds : DataFrame
		This dataframe contains the data to be plotted. It is copied and sliced into the
		regions defined. The dataframe expects the time to be in Index and the wavelength/energy 
		to be in the columns. The spectra is plotted with a second (energy) axis
	
	times : float or list/vector (of floats), optional
		For each entry in rel_time a spectrum is plotted. If time_width_percent=0 (Default) the 
		nearest measured timepoint is chosen. For other values see 'time_width_percent'
	
	time_width_percent : float
		"rel_time" and "time_width_percent" work together for creating spectral plots at 
		specific timepoints. For each entry in rel_time a spectrum is plotted. 
		If however e.g. time_width_percent=10 the region between the timepoint closest 
		to the  1.1 x timepoint and 0.9 x timepoint is averaged and shown 
		(and the legend adjusted accordingly). This is particularly useful for the densly
		sampled region close to t=0. Typically for a logarithmic recorded kinetics, the 
		timepoints at later times will be further appart than 10 percent of the value, 
		but this allows to elegantly combine values around time=0 for better statistics. 
		This averaging is only applied for the plotting function and not for the fits.

	ignore_time_region : None or list (of two floats or of lists), optional
		cut set a time range with a low and high limit from the fits. (Default) None nothing happens
		The region will be removed during the fitting process (and will be missing in the fit-result
		plots) 
		Usage single region: [lower region limit,upper region limit],
		use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]
		
	drop_ignore : Bool, True or False, optional
		If set to True the values in ignore_time_region are removed from the dataset instead of set to zero
		
	wave_nm_bin : None or float, optional
		rebins the original data into even intervals. If set to None the original data will be used. 
		If set to a width (e.g. 2nm), the wavelength axis will be divided into steps of this size
		and the mean of all measurements in the interval is taken. The re-binning stops as soon as
		the measured stepsize is wider than given here, then the original bins are used. 
		This function is particularly useful for spectrometer with non-linear dispersion, 
		like a prism in the infrared.
		
	equal_energy_bin : None or float(optional)
		if this is set the wave_nm_bin is ignored and the data is rebinned into equal energy bins.
		
	baseunit : str
		baseunit is a neat way to change the unit on the time axis of the plots. (Default) 'ps', but they 
		can be frames or something similarly. This is changing only the label of the axis. 
		During the import there is the option to divide the numbers by a factor. 
		I have also used frames or fs as units. Important is that all time units will be labeled with
		this unit.
				
	scattercut : None or iterable (of floats or other iterable, always pairs!), optional
		intented to "cut" one or multiple scatter regions. (if (Default) None nothing
		happens) If it is set the spectral region between the limits is set to zero. 
		Usage single region: [lower region limit,upper region limit], 
		use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]
		
	drop_scatter : Bool, True or False, optional
		If set to True the values in scattercut are removed from the dataset instead of set to zero
		
	bordercut : None or iterable (with two floats), optional
		cut spectra at the low and high wavelength limit. (Default) None 
		uses the limits of measurement 
 
	timelimits : None or list (of 2 floats), optional
		cut times at the low and high time limit. (Default) None uses the limits of measurement
		Important: If either the background or the chirp is to be fit this must include the 
		time before zero! Useful: It is useful to work on different regions, starting with
		the longest (then use the ta.Backgound function prior to fit) and expand from there
	
	wavelength : float or list (of floats), optional 
		'wavelength' and 'wavelength_bin' work together for the creation 
		of kinetic plots. When plotting kinetic spectra one line will be plotted for each entrance
		in the list/vector rel_wave. During object generation the vector np.arange(300,1000,100) 
		is set as standard. Another typical using style would be to define a list of interesting 
		wavelength at which a kinetic development is to be plotted. At each selected wavelength 
		the data between wavelength+ta.wavelength_bin and wavelength-ta.wavelength_bin is averaged 
		for each timepoint returned 
	
	wavelength_bin : float, optional
		the width used in kinetics, see below (Default) None
		
	time_bin : None or int, optional
		is dividing the points on the time-axis in even bins and averages the found values in between. 
		This is a hard approach that also affects the fits. I do recommend to use this carefully, 
		it is most useful for modulated data. A better choice for transient absorption that only 
		affects the kinetics is 'time_width_percent'
	'''
	time_label=ds.index.name
	energy_label=ds.columns.name
	if (wavelength is not None) and (times is not None):raise ValueError('can not get wavelength and times back')
	if (bordercut is not None) and not from_fit:
		ds.columns=ds.columns.astype('float')
		ds=ds.loc[:,bordercut[0]:bordercut[1]]
	if (equal_energy_bin is not None) and (wavelength is None):# we work with optical data but want to bin in equal energy
		x=ds.columns.values.astype('float')
		y=ds.index.values.astype('float')
		energy_label='Energy in eV'
		x=scipy.constants.h*scipy.constants.c/(x*1e-9*scipy.constants.electron_volt)
		if from_fit:#they are already binned
			ds.columns=x
			ds.sort_index(axis=1,ascending=False)
		elif (x[1:]-x[:-1]>equal_energy_bin).all():
			raise ValueError("equal_energy_bin bins are to small for the data")
		else:
			rebin_max=np.argmin((x[1:]-x[:-1])<equal_energy_bin)#find the position where the difference is larger than the wave_nm_bin
			if rebin_max==0:rebin_max=len(x)# we get 0 when all teh values are ok
			if rebin_max<len(x):
				if (x[1:]-x[:-1]>equal_energy_bin).all():raise ValueError("equal_energy_bin bins are to small for the data")
				bins=np.arange(x.min(),x[rebin_max],equal_energy_bin)
				bin_means,bin_edges = binned_statistic(x[:rebin_max], ds.values[:,:rebin_max], statistic='mean',bins=bins)[:2]  
				bins=(bin_edges[1:]+bin_edges[:-1])/2.
				ds=pandas.concat((pandas.DataFrame(bin_means,index=y,columns=bins),ds.iloc[:,rebin_max:]), axis=1, join='outer')   
			else:
				bins=np.arange(x.min(),x.max()+equal_energy_bin,equal_energy_bin)
				bin_means,bins = binned_statistic(x, ds.values, statistic='mean',bins=bins)[:2]
				bins=(bins[1:]+bins[:-1])/2.
				ds=pandas.DataFrame(bin_means,index=y,columns=bins)
	elif (wave_nm_bin is not None) and (wavelength is None):# bin in wavelength
		x=ds.columns.values.astype('float')
		y=ds.index.values.astype('float')
		if (x[1:]-x[:-1]>wave_nm_bin).all():raise ValueError("wavelength_nm_bins bins are to small for the data")
		rebin_max=np.argmin((x[1:]-x[:-1])<wave_nm_bin)#find the position where the difference is larger than the wave_nm_bin
		if rebin_max==0:rebin_max=len(x)# we get 0 when all teh values are ok
		if rebin_max<len(x):
			if (x[1:]-x[:-1]>wave_nm_bin).all():raise ValueError("wavelength_nm_bins bins are to small for the data")
			bins=np.arange(x.min(),x[rebin_max],wave_nm_bin)
			bin_means,bin_edges = binned_statistic(x[:rebin_max], ds.values[:,:rebin_max], statistic='mean',bins=bins)[:2]  
			bins=(bin_edges[1:]+bin_edges[:-1])/2.
			ds=pandas.concat((pandas.DataFrame(bin_means,index=y,columns=bins),ds.iloc[:,rebin_max:]), axis=1, join='outer')
					   
		else:
			bins=np.arange(x.min(),x.max()+wave_nm_bin,wave_nm_bin)
			bin_means,bins = binned_statistic(x, ds.values, statistic='mean',bins=bins)[:2]
			bins=(bins[1:]+bins[:-1])/2.
			ds=pandas.DataFrame(bin_means,index=y,columns=bins)

	if time_bin is not None:
		time=ds.index.values.astype('float')
		y=ds.columns.values.astype('float')
		time_bin=int(time_bin)
		time_bins=time[::time_bin]
		bin_means,bins = binned_statistic(time, ds.values.T, statistic='mean',bins=time_bins)[:2]
		bins=(bins[1:]+bins[:-1])/2.
		ds=pandas.DataFrame(bin_means,index=y,columns=bins)
		ds=ds.T

	if timelimits is not None:
		ds.index=ds.index.astype('float')
		ds=ds.loc[timelimits[0]:timelimits[1],:]	
	if ignore_time_region is not None:
		ds=ds.fillna(value=0)
		ds.index=ds.index.astype('float')
		if isinstance(ignore_time_region[0], numbers.Number):
			if drop_ignore:
				ds.loc[ignore_time_region[0]:ignore_time_region[1],:]=np.nan
			else:
				ds.loc[ignore_time_region[0]:ignore_time_region[1],:]=0
		else:
			try:
				for entries in ignore_time_region:
					if drop_ignore:
						ds.loc[entries[0]:entries[1],:]=np.nan
					else:
						ds.loc[entries[0]:entries[1],:]=0
			except:
				pass
		ds=ds.dropna(axis=0)
		
	if scattercut is not None:
		ds=ds.fillna(value=0)
		x=ds.columns.values.astype('float')
		if isinstance(scattercut[0], numbers.Number):
			if (equal_energy_bin is not None):
				scattercut=[scipy.constants.h*scipy.constants.c/(a*1e-9*scipy.constants.electron_volt) for a in scattercut]
				scattercut=scattercut[::-1]
			lower=find_nearest_index(x,scattercut[0])
			upper=find_nearest_index(x,scattercut[1])
			if drop_scatter:
				ds.iloc[:,lower:upper]=np.nan
			else:
				ds.iloc[:,lower:upper]=0
		else:
			try:
				for entries in scattercut:
					if equal_energy_bin is not None:
						scattercut=[scipy.constants.h*scipy.constants.c/(a*1e-9*scipy.constants.electron_volt) for a in scattercut]
						scattercut=scattercut[::-1]
					lower=find_nearest_index(x,entries[0])
					upper=find_nearest_index(x,entries[1])
					if drop_scatter:
						ds.iloc[:,lower:upper]=np.nan
					else:
						ds.iloc[:,lower:upper]=0
			except:
				pass
		ds=ds.dropna(axis=1)
				
	#until here we always have the same matrix
	ds.index.name=time_label
	ds.columns.name=energy_label
	
	if wavelength is not None:#ok we want to have singular wavelength
		if not hasattr(wavelength,'__iter__'):wavelength=np.array([wavelength])
		if len(wavelength)>1:wavelength.sort()
		for i,wave in enumerate(wavelength):
			upper=wave+wavelength_bin/2
			lower=wave-wavelength_bin/2
			if equal_energy_bin is not None and from_fit:
				upper=scipy.constants.h*scipy.constants.c/(lower*1e-9*scipy.constants.electron_volt)
				lower=scipy.constants.h*scipy.constants.c/(upper*1e-9*scipy.constants.electron_volt)
				wave=scipy.constants.h*scipy.constants.c/(wave*1e-9*scipy.constants.electron_volt)
			if i == 0:
				out=ds.loc[:,lower:upper].mean(axis='columns').to_frame()
				out.columns = [wave]
			else:
				if wave in out.columns:continue
				out[wave] = ds.loc[:,lower:upper].mean(axis='columns')
		out.columns=out.columns.astype('float')
		out.columns.name=energy_label
		out.index.name=time_label
		ds=out
	if times is not None:  #ok we want to have single times
		if not hasattr(times, '__iter__'):times=np.array([times])
		if baseunit is None:baseunit = 'ps'
		time_scale=ds.index.values
		if time_width_percent>0:
			for i,time in enumerate(times):			
				if time<0:
					limits = [find_nearest_index(time_scale,time+time*time_width_percent/100.),
							  find_nearest_index(time_scale,time-time*time_width_percent/100.)]		 
				else:
					limits = [find_nearest_index(time_scale,time-time*time_width_percent/100.),
							  find_nearest_index(time_scale,time+time*time_width_percent/100.)]						 
				time_lower = time_scale[limits[0]]
				time_upper = time_scale[limits[1]]
				time_mean = (time_lower+time_upper)/2
				if i == 0:
					out=ds.iloc[limits[0]:limits[1],:].mean(axis='rows').to_frame()
					out.columns = ['%.3g %s (%.3g - %.3g %s)'%(time_mean,baseunit,time_lower,time_upper,baseunit)]
				else:
					out['%.3g %s (%3g - %.3g %s)'%(time_mean,baseunit,time_lower,time_upper,baseunit)]=ds.iloc[limits[0]:limits[1],:].mean(axis='rows').to_frame()			 
		else:
			for i,time in enumerate(times):
				index=find_nearest_index(time_scale,time)
				if i == 0:
					out=ds.iloc[index,:].to_frame()
					out.columns=['%.3g %s'%(time_scale[index],baseunit)]
				else:
					out['%.3g %s'%(time_scale[index],baseunit)]=ds.iloc[index,:]
		out.columns.name=time_label
		out.index.name=energy_label
		ds=out
		#ds.index.name='Wavelength in nm'
	ds.fillna(value=0,inplace=True)#lets fill nan values with zero to catch problems
	if equal_energy_bin is not None:
		ds.sort_index(axis=1,inplace=True,ascending=False)
	return ds


def plot2d(ds, ax = None, title = None, intensity_range = None, baseunit = 'ps', timelimits = None,
			scattercut = None, bordercut = None, wave_nm_bin = None, ignore_time_region = None,
			time_bin = None, log_scale = False, plot_type = 'symlog', lintresh = 1, 
			wavelength_bin = None, levels = 256, use_colorbar = True, cmap = None, 
			data_type = 'differential Absorption in $\mathregular{\Delta OD}$', equal_energy_bin = None, from_fit = False):
	'''function for plotting matrix of TA data.
	
	Parameters
	---------------
	
	ds : DataFrame
		This dataframe contains the data to be plotted. It is copied and sliced into the
		regions defined. The dataframe expects the time to be in Index and the wavelength/energy 
		to be in the columns. The spectra is plotted with a second (energy) axis
	
	ax : None, matplotlib axis object optional
		If None (Default) a new plot is is created and a new axis, otherwise ax needs to be Matplotlib Axis

	data_type : str
		this is the datatype and effectively the unit put on the intensity axis 
		(Default)'differential Absorption in $\mathregular{\Delta OD}$
	
	title : None or str
		title to be used on top of each plot
		The (Default) None triggers  self.filename to be used. Setting a specific title as string will
		be used in all plots. To remove the title all together set an empty string with this command title="" 
		
	intensity_range : None, float or list [of two floats]
		intensity_range is a general switch that governs what intensity range the plots show. 
		For the 1d plots this is the y-axis for the 2d-plots this is the colour scale. 
		This parameter recognizes three settings. If set to "None" (Default) this uses the minimum and
		maximum of the data. A single value like in the example below and the intended use is the symmetric
		scale while a list with two entries an assymmetric scale e.g. 
		intensity_range=3e-3 is converted into intensity_range=[-3e-3,3e-3]
		
	baseunit : str
		baseunit is a neat way to change the unit on the time axis of the plots. (Default) 'ps', but they 
		can be frames or something similarly. This is changing only the label of the axis. 
		During the import there is the option to divide the numbers by a factor. 
		I have also used frames or fs as units. Important is that all time units will be labeled with
		this unit.
		
	timelimits : None or list (of 2 floats), optional
		cut times at the low and high time limit. (Default) None uses the limits of measurement
		Important: If either the background or the chirp is to be fit this must include the 
		time before zero! Useful: It is useful to work on different regions, starting with
		the longest (then use the ta.Backgound function prior to fit) and expand from there
		
	scattercut : None or iterable (of floats or other iterable, always pairs!), optional
		intented to "cut" one or multiple scatter regions. (if (Default) None nothing
		happens) If it is set the spectral region between the limits is set to zero. 
		Usage single region: [lower region limit,upper region limit], 
		use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]
		
	bordercut : None or iterable (with two floats), optional
		cut spectra at the low and high wavelength limit. (Default) None 
		uses the limits of measurement 
		
	wave_nm_bin : None or float, optional
		rebins the original data into even intervals. If set to None the original data will be used. 
		If set to a width (e.g. 2nm), the wavelength axis will be divided into steps of this size
		and the mean of all measurements in the interval is taken. The re-binning stops as soon as
		the measured stepsize is wider than given here, then the original bins are used. 
		This function is particularly useful for spectrometer with non-linear dispersion, 
		like a prism in the infrared.
		
	equal_energy_bin : None or float(optional)
		if this is set the wave_nm_bin is ignored and the data is rebinned into equal energy bins (based upon that the data is in nm. 
		If dual axis  is on then the lower axis is energy and the upper is wavelength
	
	ignore_time_region : None or list (of two floats or of lists), optional
		cut set a time range with a low and high limit from the fits. (Default) None nothing happens
		The region will be removed during the fitting process (and will be missing in the fit-result
		plots) 
		Usage single region: [lower region limit,upper region limit],
		use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]
		
	time_bin : None or int, optional
		is dividing the points on the time-axis in even bins and averages the found values in between. 
		This is a hard approach that also affects the fits. I do recommend to use this carefully, 
		it is most useful for modulated data. A better choice for transient absorption that only 
		affects the kinetics is 'time_width_percent'
		 
	log_scale : bool, optional
		If True (Default), The 2D plots (Matrix) is plotted with a pseudo logarithmic intensity scale. 
		This usually does not give good results unless the intensity scale is symmetric 
		
	plot_type : None or str
		is a general setting that can influences what time axis will be used for the plots. 
		"symlog" (linear around zero and logarithmic otherwise) "lin" and "log" are valid options.
		
	lintresh : float
		The pseudo logratihmic range "symlog" is used for most time axis. Symlog plots a range around
		time zero linear and beyond this linear treshold 'lintresh' on a logarithmic scale. (Default) 0.3 
		
	wavelength_bin : float, optional
		the width used in kinetics, see below (Default) 10nm
	
	levels : int, optional
		how many different colours to use in the description. less makes for more contrast but less 
		intensity details (Default) 256
		
	use_colorbar : bool, optional
		if True (Default) a colour bar is added to the 2d plot for intensity explanation, switch
		mostely used for creating multiple plots
	
	cmap : None or matplotlib color map, optional
		is a powerfull variable that chooses the colour map applied for all plots. If set to 
		None (Default) then the self.cmap is used.
		As standard I use the color map "jet" from matplotlib. There are a variety of colormaps 
		available that are very usefull. Beside "jet", "viridis" is a good choice as it is well 
		visible under red-green blindness. Other useful maps are "prism" for high fluctuations 
		or diverging color maps like "seismic". 
		See https://matplotlib.org/3.1.0/tutorials/colors/colormaps.html for a comprehensive 
		selection. In the code the colormaps are imported so if plot_func is imported as pf then 
		self.cmap=pf.cm.viridis sets viridis as the map to use. Internally the colors are chosen 
		with the "colm" function. The 2d plots require a continuous color map so if something 
		else is give 2d plots are shown automatically with "jet". For all of the 1d plots however 
		I first select a number of colors before each plot. If cmap is a continous map then these
		are sampled evenly over the colourmap. Manual iterables of colours 
		cmap=[(1,0,0),(0,1,0),(0,0,1),...] are also accepted, as are vectors or dataframes that 
		contain as rows the colors. There must be of course sufficient colors present for 
		the numbers of lines that will be plotted. So I recommend to provide at least 10 colours 
		(e.g.~your university colors). colours are always given as a, list or tuple with RGA or RGBA
		(with the last A beeing the Alpha=transparency. All numbers are between 0 and 1. 
		If a list/vector/DataFrame is given for the colours they will be used in the order provided.
		
	from_fit : bool optional
		it needed this swtich to avoid re-slicing of data in spectal axis for equal energy bins
	'''
	global halfsize
	if cmap is None:
		cmap=standard_map
	elif not np.array([isinstance(cmap,type(cm.viridis)),isinstance(cmap,type(cm.jet)),isinstance(cmap,type(cm.Blues)),isinstance(cmap,type(cm.coolwarm)),isinstance(cmap,type(cm.terrain))]).any():#we must have a 
		cmap=standard_map

	if ax is None:
		ax_ori=False
		if halfsize:
			fig,ax=plt.subplots(figsize=(5,3),dpi=100)
		else:
			fig,ax=plt.subplots(figsize=(10,6),dpi=100)
	else:
		ax_ori=True
		fig=ax.get_images()
	if timelimits is None:
		timelimits=(ds.index.min(),ds.index.max())
	ds = sub_ds(ds, scattercut = scattercut, bordercut = bordercut, timelimits = timelimits, wave_nm_bin = wave_nm_bin, 
				wavelength_bin = wavelength_bin, time_bin = time_bin, ignore_time_region = ignore_time_region, 
				drop_scatter = False, drop_ignore = False, equal_energy_bin = equal_energy_bin, from_fit = from_fit)		
	
	if intensity_range is None:
		try:
			maxim=max([abs(ds.values.min()),abs(ds.values.max())])
			intensity_range=[-maxim,maxim]
		except:
			intensity_range=[-1e-2,1e-2]
	else:
		if not hasattr(intensity_range,'__iter__'):#lets have an lazy option
			intensity_range=[-intensity_range,intensity_range]
		else:
					   
			if log_scale:print('I highly recommend to make a symmetric intensity distribution for logarithmic scale, the colorbar might look strange otherwise')
	if log_scale: 
		if 0:# old manual symlog
			bounds0 = list(-1*np.logspace(np.log10(-intensity_range[0]), np.log10(-intensity_range[0]/(levels/2)), levels))
			bounds1 = np.logspace(np.log10(intensity_range[1]/(levels/2)),np.log10(intensity_range[1]),	 levels)
			bounds0.append(0)
			for a in bounds1:
				bounds0.append(a)
			norm = BoundaryNorm(boundaries=bounds0, ncolors=len(bounds0))
			mid_color=colm(k=range(levels),cmap=cmap)
		else:
			norm=SymLogNorm(abs(max(intensity_range)-min(intensity_range))/100, linscale=1.0, vmin=min(intensity_range), vmax=max(intensity_range), clip=True, base=10)
			mid_color=cmap(0.5)
	else:
		if 0:
			nbins=levels
			levels = LinearLocator(numticks=levels).tick_values(vmin=min(intensity_range), vmax=max(intensity_range))
			norm = BoundaryNorm(levels,clip=True,ncolors=cmap.N)
			mid_color_index=find_nearest_index(0,levels)
			mid_color=colm(k=range(nbins),cmap=cmap)
			mid_color=mid_color[mid_color_index]
		else:
			norm = Normalize(vmin=min(intensity_range),vmax=max(intensity_range))
			mid_color=cmap(0.5)
	#print(ds.head())
	x = ds.columns.values.astype('float')
	y = ds.index.values.astype('float')
	X, Y = np.meshgrid(x, y)
	img=ax.pcolormesh(X,Y,ds.values,norm=norm,cmap=cmap,shading=shading)
	#img=ax.pcolormesh(X,Y,ds.values,cmap=cmap,shading=shading)
	if ignore_time_region is None:
		pass
	elif isinstance(ignore_time_region[0], numbers.Number):
		ds.index=ds.index.astype(float)
		try:
			upper=ds.loc[ignore_time_region[1]:,:].index.values.min()
			lower=ds.loc[:ignore_time_region[0],:].index.values.max()
			if equal_energy_bin is not None:
				rect = plt.Rectangle((x.max(),lower), width=abs(ax.get_xlim()[0]-ax.get_xlim()[1]), height=abs(upper-lower),facecolor=mid_color,alpha=1)#mid_color)
			else:
				rect = plt.Rectangle((x.min(),lower), width=abs(ax.get_xlim()[1]-ax.get_xlim()[0]), height=abs(upper-lower),facecolor=mid_color,alpha=1)#mid_color)
			ax.add_patch(rect)
		except:
			pass
	else:
		ignore_time_region_loc=flatten(ignore_time_region)
		for k in range(int(len(ignore_time_region_loc)/2+1)):
			try:
				upper=ds.loc[ignore_time_region[k+1]:,:].index.values.min()
				lower=ds.loc[:ignore_time_region[k],:].index.values.max()
				if equal_energy_bin is not None:
					rect = plt.Rectangle((x.max(),lower), width=abs(ax.get_xlim()[0]-ax.get_xlim()[1]), height=abs(upper-lower),facecolor=mid_color,alpha=1)
				else:
					rect = plt.Rectangle((x.min(),lower), width=abs(ax.get_xlim()[1]-ax.get_xlim()[0]), height=abs(upper-lower),facecolor=mid_color,alpha=1)
				ax.add_patch(rect)
			except:
				pass

	if scattercut is None:
		pass
	elif isinstance(scattercut[0], numbers.Number):
		try:
			if equal_energy_bin is not None:
				scattercut=[scipy.constants.h*scipy.constants.c/(a*1e-9*scipy.constants.electron_volt) for a in scattercut]
				scattercut=scattercut[::-1]
			upper=ds.loc[:,scattercut[1]:].columns.values.min()
			lower=ds.loc[:,:scattercut[0]].columns.values.max()
			width=abs(upper-lower)
			rect = plt.Rectangle((lower,y.min()), height=abs(ax.get_ylim()[1]-ax.get_ylim()[0]), width=width, facecolor=mid_color,alpha=1)#mid_color)
			ax.add_patch(rect)
		except:
			pass
	else:
		scattercut=flatten(scattercut)
		if equal_energy_bin is not None:
			scattercut=[scipy.constants.h*scipy.constants.c/(a*1e-9*scipy.constants.electron_volt) for a in scattercut]
			scattercut=scattercut[::-1]
		for k in range(int(len(scattercut)/2+1)):
			try:
				upper=ds.loc[:,scattercut[k][1]:].columns.values.min()
				if upper==0:raise
				lower=ds.loc[:,:scattercut[k][0]].columns.values.max()
				rect = plt.Rectangle((lower.min()), height=abs(ax.get_ylim()[1]-ax.get_ylim()[0]), width=abs(upper-lower),facecolor=mid_color,alpha=1)#mid_color)
				ax.add_patch(rect)	
			except:
				pass	
	if use_colorbar:
		mid=(intensity_range[1]+intensity_range[0])/2
		if log_scale:
			values=[intensity_range[0],mid-abs(intensity_range[0]-mid)/10,mid-abs(intensity_range[0]-mid)/100,mid,mid+abs(intensity_range[1]-mid)/100,mid+abs(intensity_range[1]-mid)/10,intensity_range[1]]
		else:
			values=[intensity_range[0],intensity_range[0]+abs(intensity_range[0]-mid)/2,mid,intensity_range[1]-abs(intensity_range[1]-mid)/2,intensity_range[1]]
																				  
		labels=['%.2g'%(a) for a in values]
		labels[0]='<' + labels[0]
		labels[-1]='>'+labels[-1]
		cbar=plt.colorbar(img, ax=ax,ticks=values,pad=0.01)
		cbar.ax.set_yticklabels(labels)
		a=ax.yaxis.label
		fontsize=a.get_fontsize()
		fontsize-=4
		if data_type is not None:#we use this as a switch to enable a flexible avoidance of the label setting.
			if log_scale:
				if ax_ori:cbar.set_label(data_type + '\nLog-scale', rotation=270,labelpad=20,fontsize=fontsize)
				else:cbar.set_label(data_type + '\nLog-scale', rotation=270,labelpad=20,fontsize=fontsize)		
			else:
				if ax_ori:cbar.set_label(data_type, rotation=270,labelpad=20,fontsize=fontsize)
				else:cbar.set_label(data_type, rotation=270,labelpad=20,fontsize=fontsize)

	
	if "symlog" in plot_type:

		ax.plot(ax.get_xlim(),[lintresh,lintresh],'black',lw=0.5,alpha=0.3)
		ax.plot(ax.get_xlim(),[-1.0*lintresh,-1.0*lintresh],'black',lw=0.5,alpha=0.3)
		ax.plot(ax.get_xlim(),[0,0],'black',lw=0.5,alpha=0.6)
		if 1:
			ax.set_yscale('symlog', linthresh=lintresh)
			locmaj = matplotlib.ticker.LogLocator(base=10.0, subs=(0.1,1.0,10.,1e2,1e3,1e4))
			ax.yaxis.set_major_locator(locmaj)
			locmin = matplotlib.ticker.LogLocator(base=10.0, subs=np.arange(0.1,1,0.1)) 
			ax.yaxis.set_minor_locator(locmin)
			ticks=list(ax.get_yticks())	
			ticks.append(lintresh)
			[ticks.append(a) for a in [-0.3,-1,-2,-5,-10]]
			ticks.sort()
			if timelimits[1]>100:
				ticks=np.array(ticks)
				ticks=np.concatenate((ticks.clip(min=0.1),np.zeros(1),ticks.clip(max=-0.1,min=timelimits[0])),axis=0)  
			ax.set_yticks(ticks)
		else:
			print('here2')
			ax.set_yscale('symlog', linthresh=lintresh,subsy=range(2,9),linscaley=lintresh)
		ax.set_ylim(y.min(),y.max())
	elif "log" in plot_type:
		lower_time=max(1e-6,timelimits[0])
		ax.set_ylim(lower_time,y.max())
		ax.set_yscale('log')
	else:
		ax.set_yscale('linear')
		ax.set_ylim(timelimits)
	if bordercut is not None:
		try:
			if equal_energy_bin is not None:
				bordercut=[scipy.constants.h*scipy.constants.c/(a*1e-9*scipy.constants.electron_volt) for a in bordercut]
			ax.set_xlim(bordercut[0],bordercut[1])
		except:
			print('bordercut failed')
			pass
	if equal_energy_bin is not None and False:
		temp=np.array(ax.get_xlim())
		ax.set_xlim(temp.max(),temp.min())
	ax.set_xlabel(ds.columns.name)
	ax.set_ylabel(ds.index.name)
	if title:
		ax.set_title(title)	
	if ax_ori:return ax
	return fig


def plot2d_fit(re, error_matrix_amplification=5, use_images=True, patches=False, title = None, 
				intensity_range = None, baseunit = 'ps', timelimits = None,
				scattercut = None, bordercut = None, wave_nm_bin = None, ignore_time_region = None,
				time_bin = None, log_scale = False, scale_type = 'symlog', lintresh = 1, 
				wavelength_bin = None, levels = 256, plot_with_colorbar = True, cmap = None, 
				data_type = 'differential Absorption in $\mathregular{\Delta OD}$', equal_energy_bin = None):													   
	'''Plots the fit output as a single plot with meas,fitted and difference. 
	The differnece used err_matrix_amplification as a factor. patches moves the labels from the
	title into white patches in the top of the figure
	
	Parameters
	---------------
	
	re : dict
		Dictionary that contains the fit results and  specific the dataframes A, AC and AE 
	
	error_matrix_amplification : int, optional
		the error matrix AE is multiplied by this factor for the plot.
	
	use_images : bool:
		(Default)True converts the matrix into images, to reduce the filesize. 
	
	data_type : str
		this is the datatype and effectively the unit put on the intensity axis 
		(Default)'differential Absorption in $\mathregular{\Delta OD}$
		
	patches : bool, optional
		If False (Default) the names "measured" "fitted" "difference" will be placed above the images.
		If True, then they will be included into the image (denser)
	
	title : None or str
		title to be used on top of each plot
		The (Default) None triggers  self.filename to be used. Setting a specific title as string will
		be used in all plots. To remove the title all together set an empty string with this command title="" 
		
	intensity_range : None, float or list [of two floats]
		intensity_range is a general switch that governs what intensity range the plots show. 
		For the 1d plots this is the y-axis for the 2d-plots this is the colour scale. 
		This parameter recognizes three settings. If set to "None" (Default) this uses the minimum and
		maximum of the data. A single value like in the example below and the intended use is the symmetric
		scale while a list with two entries an assymmetric scale e.g. 
		intensity_range=3e-3 is converted into intensity_range=[-3e-3,3e-3]
		
	baseunit : str
		baseunit is a neat way to change the unit on the time axis of the plots. (Default) 'ps', but they 
		can be frames or something similarly. This is changing only the label of the axis. 
		During the import there is the option to divide the numbers by a factor. 
		I have also used frames or fs as units. Important is that all time units will be labeled with
		this unit.
		
	timelimits : None or list (of 2 floats), optional
		cut times at the low and high time limit. (Default) None uses the limits of measurement
		Important: If either the background or the chirp is to be fit this must include the 
		time before zero! Useful: It is useful to work on different regions, starting with
		the longest (then use the ta.Backgound function prior to fit) and expand from there
		
	scattercut : None or iterable (of floats or other iterable, always pairs!), optional
		intented to "cut" one or multiple scatter regions. (if (Default) None nothing
		happens) If it is set the spectral region between the limits is set to zero. 
		Usage single region: [lower region limit,upper region limit], 
		use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]
		
	bordercut : None or iterable (with two floats), optional
		cut spectra at the low and high wavelength limit. (Default) None 
		uses the limits of measurement 
		
	wave_nm_bin : None or float, optional
		rebins the original data into even intervals. If set to None the original data will be used. 
		If set to a width (e.g. 2nm), the wavelength axis will be divided into steps of this size
		and the mean of all measurements in the interval is taken. The re-binning stops as soon as
		the measured stepsize is wider than given here, then the original bins are used. 
		This function is particularly useful for spectrometer with non-linear dispersion, 
		like a prism in the infrared.
	
	equal_energy_bin : None or float(optional)
		if this is set the wave_nm_bin is ignored and the data is rebinned into equal energy bins (based upon that the data is in nm. 
		If dual axis  is on then the lower axis is energy and the upper is wavelength
	
	ignore_time_region : None or list (of two floats or of lists), optional
		cut set a time range with a low and high limit from the fits. (Default) None nothing happens
		The region will be removed during the fitting process (and will be missing in the fit-result
		plots) 
		Usage single region: [lower region limit,upper region limit],
		use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]
		
	time_bin : None or int, optional
		is dividing the points on the time-axis in even bins and averages the found values in between. 
		This is a hard approach that also affects the fits. I do recommend to use this carefully, 
		it is most useful for modulated data. A better choice for transient absorption that only 
		affects the kinetics is 'time_width_percent'
		 
	log_scale : bool, optional
		If True (Default), The 2D plots (Matrix) is plotted with a pseudo logarithmic intensity scale. 
		This usually does not give good results unless the intensity scale is symmetric 
		
	Scale_type : None or str
		is a general setting that can influences what time axis will be used for the plots. 
		"symlog" (linear around zero and logarithmic otherwise) "lin" and "log" are valid options.
		
	lintresh : float
		The pseudo logratihmic range "symlog" is used for most time axis. Symlog plots a range around
		time zero linear and beyond this linear treshold 'lintresh' on a logarithmic scale. (Default) 0.3 
		
	wavelength_bin : float, optional
		the width used in kinetics, see below (Default) 10nm
	
	levels : int, optional
		how many different colours to use in the description. less makes for more contrast but less 
		intensity details (Default) 256
		
	plot_with_colorbar : bool, optional
		if True (Default) a colour bar is added to the 2d plot for intensity explanation, switch
		mostely used for creating multiple plots
	
	cmap : None or matplotlib color map, optional
		is a powerfull variable that chooses the colour map applied for all plots. If set to 
		None (Default) then the self.cmap is used.
		As standard I use the color map "jet" from matplotlib. There are a variety of colormaps 
		available that are very usefull. Beside "jet", "viridis" is a good choice as it is well 
		visible under red-green blindness. Other useful maps are "prism" for high fluctuations 
		or diverging color maps like "seismic". 
		See https://matplotlib.org/3.1.0/tutorials/colors/colormaps.html for a comprehensive 
		selection. In the code the colormaps are imported so if plot_func is imported as pf then 
		self.cmap=pf.cm.viridis sets viridis as the map to use. Internally the colors are chosen 
		with the "colm" function. The 2d plots require a continuous color map so if something 
		else is give 2d plots are shown automatically with "jet". For all of the 1d plots however 
		I first select a number of colors before each plot. If cmap is a continous map then these
		are sampled evenly over the colourmap. Manual iterables of colours 
		cmap=[(1,0,0),(0,1,0),(0,0,1),...] are also accepted, as are vectors or dataframes that 
		contain as rows the colors. There must be of course sufficient colors present for 
		the numbers of lines that will be plotted. So I recommend to provide at least 10 colours 
		(e.g.~your university colors). colours are always given as a, list or tuple with RGA or RGBA
		(with the last A beeing the Alpha=transparency. All numbers are between 0 and 1. 
		If a list/vector/DataFrame is given for the colours they will be used in the order provided.
	
	'''
	global halfsize
	if intensity_range is None:intensity_range=5e-3
	if halfsize:
		fig,ax=plt.subplots(3,figsize=(4.5,5.5))
	else:
		fig,ax=plt.subplots(3,figsize=(9,11))
	if patches:			
		plot2d(re['A'], cmap = cmap, log_scale = log_scale, intensity_range = intensity_range, ax = ax[0], 
				baseunit = baseunit, use_colorbar = plot_with_colorbar, levels = levels, plot_type = scale_type, 
				ignore_time_region = ignore_time_region,  lintresh = lintresh, bordercut = bordercut, 
				scattercut = scattercut, timelimits = timelimits, data_type = data_type, equal_energy_bin = equal_energy_bin, from_fit = True)
		plot2d(re['AC'], cmap = cmap, log_scale = log_scale, intensity_range = intensity_range, ax = ax[1], 
				baseunit = baseunit, use_colorbar = plot_with_colorbar, levels = levels, plot_type = scale_type, 
				ignore_time_region = ignore_time_region,  lintresh = lintresh, bordercut = bordercut, 
				scattercut = scattercut, timelimits = timelimits, data_type = data_type, equal_energy_bin = equal_energy_bin, from_fit = True)
		plot2d(re['AE'], cmap = cmap, log_scale = log_scale, intensity_range = np.array(intensity_range)/error_matrix_amplification, ax = ax[2], 
				baseunit = baseunit, use_colorbar = plot_with_colorbar, levels = levels, plot_type = scale_type, 
				ignore_time_region = ignore_time_region,  lintresh = lintresh, bordercut = bordercut, scattercut = scattercut, 
				timelimits = timelimits, data_type = data_type, equal_energy_bin = equal_energy_bin, from_fit = True)
		for i in range(3):
			ax[i].set_title(label='')
			stringen=['measured','calculated','difference']
			x_width=(ax[i].get_xlim()[1]-ax[i].get_xlim()[0])/4
			if 'lin' in scale_type:
				y_width=(ax[i].get_ylim()[1])/8
			else:
				y_width=(ax[i].get_ylim()[1])/1.5
			rect = plt.Rectangle((ax[i].get_xlim()[1]-x_width, ax[i].get_ylim()[1]-y_width), x_width, y_width,facecolor="white", alpha=0.5)
			ax[i].add_patch(rect)
			ax[i].text(ax[i].get_xlim()[1]-x_width+x_width*0.1,ax[i].get_ylim()[1]-y_width+y_width*0.1,stringen[i],fontsize=16)								 
			fig.subplots_adjust(left=0.15, bottom=0.067, right=0.97, top=0.985, wspace=0.0, hspace=0.258)
	else:
		plot2d(re['A'], cmap = cmap, title = 'Measured', log_scale = log_scale, intensity_range = intensity_range, 
				ax = ax[0], baseunit = baseunit, use_colorbar = plot_with_colorbar, levels = levels, plot_type = scale_type, 
				ignore_time_region = ignore_time_region,  lintresh = lintresh, bordercut = bordercut, scattercut = scattercut, 
				timelimits = timelimits, data_type = data_type, equal_energy_bin = equal_energy_bin, from_fit = True)																					
		plot2d(re['AC'], cmap = cmap, title = 'Calculated', log_scale = log_scale, intensity_range = intensity_range, 
				ax = ax[1], baseunit = baseunit, use_colorbar = plot_with_colorbar, levels = levels, plot_type = scale_type, 
				ignore_time_region = ignore_time_region,  lintresh = lintresh, bordercut = bordercut, scattercut = scattercut, 
				timelimits = timelimits , data_type = data_type, equal_energy_bin = equal_energy_bin, from_fit = True)
		plot2d(re['AE'], cmap = cmap, title = 'Difference', log_scale = log_scale, intensity_range = np.array(intensity_range)/error_matrix_amplification, 
				ax = ax[2], baseunit = baseunit, use_colorbar = plot_with_colorbar, levels = levels, plot_type = scale_type, 
				ignore_time_region = ignore_time_region,  lintresh = lintresh, bordercut = bordercut, scattercut = scattercut, 
				timelimits = timelimits, data_type = data_type, equal_energy_bin = equal_energy_bin, from_fit = True)
		#fig.subplots_adjust(left=0.15, bottom=0.067, right=0.97, top=0.97, wspace=0.0, hspace=0.398)
	fig.tight_layout()
	return fig
	
	
def plot_fit_output( re, ds, cmap = standard_map, plotting = range(7), title = None, path = None, filename = None, f = 'standard', 
					intensity_range = 1e-2, baseunit = 'ps', timelimits = None, scattercut = None, bordercut = None, 
					error_matrix_amplification = 20, wave_nm_bin = 5, rel_wave = None, width = 10, rel_time = [1, 5, 10], 
					time_width_percent = 10, ignore_time_region = None, save_figures_to_folder = True, return_figures_handles=False,
					log_fit = False, mod = None, subplot = False, color_offset = 0, log_scale = True, savetype = 'png', 
					evaluation_style = False, lintresh = 1, scale_type = 'symlog', patches = False, print_click_position = False, 
					data_type = 'differential Absorption in $\mathregular{\Delta OD}$', plot_second_as_energy = True, units = 'nm',
					equal_energy_bin = None):
	'''Purly manual function that plots all the fit output figures. Quite cumbersome, 
	but offers a lot of manual options. The figures can be called separately 
	or with a list of plots. e.g. range(6) call plots 0-5 Manual plotting of certain type:
	This is a wrapper function that triggers the plotting of all the fitted plots.
	The parameter in this plot call are to control the general look and features of the plot.
	Which plots are printed is defined by the command (plotting)
	The plots are generated from the fitted Matrixes and as such only will work after a fit was actually
	completed (and the "re" dictionary attached to the object.)
	In all plots the RAW data is plotted as dots and the fit with lines 
	
	*Contents of the plots*

	0. DAC contains the assigned spectra for each component of the fit. For
	   a modelling with independent exponential decays this corresponds to
	   the "Decay Associated Spectra" (DAS). For all other models this
	   contains the "Species Associated Spectra" (SAS). According to the
	   model the separate spectra are labeled by time (process) or name, if
	   a name is associated in the fitting model. The spectra are shown in
	   the extracted strength in the right pane and normalized in the left.
	   Extracted strength means that the measured spectral strength is the
	   intensity (concentration matrix) times this spectral strength. As the
	   concentration maxima for all DAS are 1 this corresponds to the
	   spectral strength for the DAS. (please see the documentation for the
	   fitting algorithm for further details)

	1. summed intensity. All wavelength of the spectral axis are summed for
	   data and fit. The data is plotted in a number of ways vs linear and
	   logarithmic axis. This plot is not ment for publication but very
	   useful to evaluate the quality of a fit.

	2. plot kinetics for selected wavelength (see corresponding RAW plot)

	3. plot spectra at selected times (see corresponding RAW plot)

	4. plots matrix (measured, modelled and error Matrix). The parameter are
	   the same as used for the corresponding RAW plot with the addition of
	   "error_matrix_amplification" which is a scaling factor multiplied
	   onto the error matrix. I recommend to play with different "cmap",
	   "log_scale" and "intensity_scale" to create a pleasing plot

	5. concentrations. In the progress of the modelling/fitting a matrix is
	   generated that contains the relative concentrations of the species
	   modelled. This plot is showing the temporal development of these
	   species. Further details on how this matrix is generated can be found
	   in the documentation of the fitting function. The modeled spectra are
	   the convolution of these vectors (giving the time-development) and
	   the DAS/SAS (giving the spectral development).
	
	Parameters
	---------------
	
	ds : DataFrame
		This dataframe contains the data to be plotted. It is copied and sliced into the
		regions defined. The dataframe expects the time to be in Index and the wavelength/energy 
		to be in the columns. The spectra is plotted with a second (energy) axis
	
	re : dict
		Dictionary that contains the fit results and  specific the dataframes A, AC and AE 

	data_type : str
		this is the datatype and effectively the unit put on the intensity axis 
		(Default)'differential Absorption in $\mathregular{\Delta OD}$
	error_matrix_amplification : int, optional
		the error matrix AE is multiplied by this factor for the plot.
	
	plotting : int or iterable (of integers), optional
		This parameter determines which figures are plotted 
		the figures can be called separately with plotting = 1
		or with a list of plots (Default) e.g.~plotting=range(6) calls plots 0,1,2,3,4,5
		The plots have the following numbers:\
		0 - DAS or SAS\
		1 - summed intensity\
		2 - Kinetics\
		3 - Spectra\
		4 - Matrixes\
		5 - Concentrations (the c-object)\
		6 - Residuals\ 
		The plotting takes all parameter from the "ta" object unless otherwise specified
	
	path : None, str or path object, optional
		This defines where the files are saved if the safe_figures_to_folder parameter is True, 
		quite useful if a lot of data sets are to be printed fast. 
		If a path is given, this is used. If a string like the (Default) "result_figures" is given, 
		then a subfolder of this name will be used (an generated if necessary) 
		relative to self.path. Use and empty string to use the self.path
		If set to None, the location of the plot_func will be used and
		a subfolder with title "result_figures" be generated here
	
	savetype : str or iterable (of str), optional 
		matplotlib allows the saving of figures in various formats. (Default) "png", 
		typical and recommendable options are "svg" and "pdf".
	
	return_figures_handles
		
	evaluation_style : bool, optional
		True (Default = False) adds a lot of extra information in the plot
	
	title : None or str, optional
	   "title=None" is in general the filename that was loaded. Setting a
	   specific title will be used in all plots. To remove the title all
	   together set an empty string with title=""
	
	scale_type : str, optional
	   refers to the time-axis and takes, ’symlog’ (Default)(linear around zero and logarithmic otherwise)
	   and ‘lin’ for linear and  ‘log’ for logarithmic, switching all the time axis to this type
	
	patches : bool, optional
		If False (Default) the names "measured" "fitted" "difference" will be placed above the images.
		If True, then they will be included into the image (denser)

	filename : str, optional
		offers to replace the base-name used for all plots (to e.g.~specify what sample was used). 
		if (Default) None is used, the self.filename is used as a base name. The filename plays only a 
		role during saving, as does the path and savetype	


	save_figures_to_folder : bool, optional
		(Default) is True, if True the Figures are automatically saved
		
	return_figures_handles : bool, optional
		(Default) is False, if True the Figure handles are returned as a dictionary.

	log_scale : bool, optional
		If True (Default), The 2D plots (Matrix) is plotted with a pseudo logarithmic intensity scale. 
		This usually does not give good results unless the intensity scale is symmetric 


	subplot : bool, optional
		If False (Default) axis labels and such are set. If True, we plot into the same axis and 
		do not set labels
		

	color_offset : int, optional
		At the (Default) 0 the colours are chose from the beginning, for a larger value Color_offset 
		colors are skipped. Usually only used if multiple plots are created, and the data/or fit is
		only shown for some of them.
		
	lintresh : float
		The pseudo logratihmic range "symlog" is used for most time axis. Symlog plots a range around
		time zero linear and beyond this linear treshold 'lintresh' on a logarithmic scale. (Default) 1 


	rel_time : float or list/vector (of floats), optional
		For each entry in rel_time a spectrum is plotted. If time_width_percent=0 (Default) the 
		nearest measured timepoint is chosen. For other values see 'time_width_percent'

	time_width_percent : float
		"rel_time" and "time_width_percent" work together for creating spectral plots at 
		specific timepoints. For each entry in rel_time a spectrum is plotted. 
		If however e.g. time_width_percent=10 the region between the timepoint closest 
		to the  1.1 x timepoint and 0.9 x timepoint is averaged and shown 
		(and the legend adjusted accordingly). This is particularly useful for the densly
		sampled region close to t=0. Typically for a logarithmic recorded kinetics, the 
		timepoints at later times will be further appart than 10 percent of the value, 
		but this allows to elegantly combine values around time=0 for better statistics. 
		This averaging is only applied for the plotting function and not for the fits.
		
	ignore_time_region : None or list (of two floats or of lists), optional
		cut set a time range with a low and high limit from the fits. (Default) None nothing happens
		The region will be removed during the fitting process (and will be missing in the fit-result
		plots) 
		Usage single region: [lower region limit,upper region limit],
		use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]

	width : float, optional
		the width used in kinetics, see below (Default) 10nm

	rel_wave : float or list (of floats), optional 
		'rel_wave' and 'width' (in the object called 'wavelength_bin' work together for the creation 
		of kinetic plots. When plotting kinetic spectra one line will be plotted for each entrance
		in the list/vector rel_wave. During object generation the vector np.arange(300,1000,100) 
		is set as standard. Another typical using style would be to define a list of interesting 
		wavelength at which a kinetic development is to be plotted. At each selected wavelength 
		the data between wavelength+ta.wavelength_bin and wavelength-ta.wavelength_bin is averaged 
		for each timepoint returned 

	timelimits : None or list (of 2 floats), optional
		cut times at the low and high time limit. (Default) None uses the limits of measurement
		Important: If either the background or the chirp is to be fit this must include the 
		time before zero! Useful: It is useful to work on different regions, starting with
		the longest (then use the ta.Backgound function prior to fit) and expand from there
		
	scattercut : None or iterable (of floats or other iterable, always pairs!), optional
		intented to "cut" one or multiple scatter regions. (if (Default) None nothing
		happens) If it is set the spectral region between the limits is set to zero. 
		Usage single region: [lower region limit,upper region limit], 
		use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]
		
	bordercut : None or iterable (with two floats), optional
		cut spectra at the low and high wavelength limit. (Default) None 
		uses the limits of measurement 
		
	wave_nm_bin : None or float, optional
		rebins the original data into even intervals. If set to None the original data will be used. 
		If set to a width (e.g. 2nm), the wavelength axis will be divided into steps of this size
		and the mean of all measurements in the interval is taken. The re-binning stops as soon as
		the measured stepsize is wider than given here, then the original bins are used. 
		This function is particularly useful for spectrometer with non-linear dispersion, 
		like a prism in the infrared.
		
	equal_energy_bin : None or float(optional)
		if this is set the wave_nm_bin is ignored and the data is rebinned into equal energy bins (based upon that the data is in nm. 
		If dual axis  is on then the lower axis is energy and the upper is wavelength

	intensity_range : None, float or list [of two floats]
		intensity_range is a general switch that governs what intensity range the plots show. 
		For the 1d plots this is the y-axis for the 2d-plots this is the colour scale. 
		This parameter recognizes three settings. If set to "None" (Default) this uses the minimum and
		maximum of the data. A single value like in the example below and the intended use is the symmetric
		scale while a list with two entries an assymmetric scale e.g. 
		intensity_range=3e-3 is converted into intensity_range=[-3e-3,3e-3]

	baseunit : str
		baseunit is a neat way to change the unit on the time axis of the plots. (Default) 'ps', but they 
		can be frames or something similarly. This is changing only the label of the axis. 
		During the import there is the option to divide the numbers by a factor. 
		I have also used frames or fs as units. Important is that all time units will be labeled with
		this unit.	
	
	f : str
		f is a replacement title that is set instead of the title. mainly used to have some options 
		(Default) is 'standard'
			
	log_fit : bool, optional
		(default)= False Used for legend generation, tells if the fit was in log or lin space 
	
	mod : str, optional
		Used for legend generation, tells what model was used for fitting 

	cmap : None or matplotlib color map, optional
		is a powerfull variable that chooses the colour map applied for all plots. If set to 
		None (Default) then the self.cmap is used.
		As standard I use the color map "jet" from matplotlib. There are a variety of colormaps 
		available that are very usefull. Beside "jet", "viridis" is a good choice as it is well 
		visible under red-green blindness. Other useful maps are "prism" for high fluctuations 
		or diverging color maps like "seismic". 
		See https://matplotlib.org/3.1.0/tutorials/colors/colormaps.html for a comprehensive 
		selection. In the code the colormaps are imported so if plot_func is imported as pf then 
		self.cmap=pf.cm.viridis sets viridis as the map to use. Internally the colors are chosen 
		with the "colm" function. The 2d plots require a continuous color map so if something 
		else is give 2d plots are shown automatically with "jet". For all of the 1d plots however 
		I first select a number of colors before each plot. If cmap is a continous map then these
		are sampled evenly over the colourmap. Manual iterables of colours 
		cmap=[(1,0,0),(0,1,0),(0,0,1),...] are also accepted, as are vectors or dataframes that 
		contain as rows the colors. There must be of course sufficient colors present for 
		the numbers of lines that will be plotted. So I recommend to provide at least 10 colours 
		(e.g.~your university colors). colours are always given as a, list or tuple with RGA or RGBA
		(with the last A beeing the Alpha=transparency. All numbers are between 0 and 1. 
		If a list/vector/DataFrame is given for the colours they will be used in the order provided.		
	
	print_click_position : bool, optional
		if True then the click position is printed for the spectral plots
	
	Examples
	------------
			
	>>> ta.plot_fit_output(ta.re,ta.ds)
		
	'''
	global halfsize 
	if baseunit != 'ps':
		if baseunit == 'ns':baseunit = 'Time in ns'
		re['A'].index.name=baseunit
		re['AC'].index.name=baseunit
		re['AE'].index.name=baseunit
		ds.index.name=baseunit
		re['c'].index.name=baseunit
	if width is None:width=wave_nm_bin
	stringen=[]
	timedf=re['fit_results_times']
	if mod is not None:
			if not isinstance(mod,type('hello')):mod='ext. func.'
	if evaluation_style:
		if mod is not None:
			stringen.append('Fit with Model: %s'%mod)
		timedf.rename(index={'resolution': "res"},inplace=True)
		timedf.rename(columns={'init_value': "init"},inplace=True)
		try:	  
			stringen.append(timedf.to_string(columns = ['value','init','vary','min','max','expr'],
						float_format = '{:.3g}'.format, justify = 'center'))
		except:
			print('something strange happened, most likely one value went "inf" or is set unexpectedly to None')
	else:
		if mod is not None:
			if mod in ['paral','exponential']:stringen.append('Fit with ind.\nexpon. decays:')	
			else:stringen.append('Fit with time parameters:')
		try:
			timedf.drop(index=['resolution','t0'],inplace=True)
		except:
			pass
		stringen.append(timedf.to_string(columns=['value'],float_format='{:.3g}'.format,header=False))
	stringen='\n'.join(stringen)
	times=timedf[timedf.is_rate].loc[:,'value'].values
	time_string=timedf[timedf.is_rate].to_string(columns=['value'],float_format='{:.3g}'.format,header=False)
	if not hasattr(plotting,'__iter__'):plotting=[plotting]
	if 0 in plotting:#DAC
		#-------plot DAC------------
		#for i,col in enumerate(re['DAC']):
			#re['DAC'].iloc[:,i]=re['DAC'].iloc[:,i].values*re['c'].max().iloc[i]
		if halfsize:
			fig1,(ax1a,ax1b,ax1c)=plt.subplots(1,3,figsize=(6,2.5),dpi=100)
		else:
			fig1,(ax1a,ax1b,ax1c)=plt.subplots(1,3,figsize=(12,5),dpi=100)
		n_colors=len(re['DAC'].columns)
		DAC=re['DAC']
		DAC_copy=DAC.copy()
		normed=(DAC/DAC.abs().max())
		for i,col in enumerate(DAC_copy):
			DAC_copy.iloc[:,i]=DAC_copy.iloc[:,i].values*re['c'].abs().max().iloc[i]
		if scattercut is None:	
			DAC.plot(ax=ax1b,color=colm(range(n_colors),cmap=cmap))
			normed.plot(ax=ax1a,color=colm(range(n_colors),cmap=cmap))
			DAC_copy.plot(ax=ax1c,color=colm(range(n_colors),cmap=cmap))
		elif isinstance(scattercut[0], numbers.Number):
			DAC.loc[:scattercut[0],:].plot(ax=ax1b,color=colm(range(n_colors),cmap=cmap))
			DAC.loc[scattercut[1]:,:].plot(ax=ax1b,color=colm(range(n_colors),cmap=cmap), label='_nolegend_')
			normed.loc[:scattercut[0],:].plot(ax=ax1a,color=colm(range(n_colors),cmap=cmap))
			normed.loc[scattercut[1]:,:].plot(ax=ax1a,color=colm(range(n_colors),cmap=cmap), label='_nolegend_')			
			DAC_copy.loc[:scattercut[0],:].plot(ax=ax1c,color=colm(range(n_colors),cmap=cmap))
			DAC_copy.loc[scattercut[1]:,:].plot(ax=ax1c,color=colm(range(n_colors),cmap=cmap), label='_nolegend_')	
		else:
			try:
				scattercut=flatten(scattercut)
				for i in range(len(scattercut)/2+1):
					if i==0:
						DAC.loc[:scattercut[0],:].plot(ax=ax1b,color=colm(range(n_colors),cmap=cmap))
						normed.loc[:scattercut[0],:].plot(ax=ax1a,color=colm(range(n_colors),cmap=cmap))
						DAC_copy.loc[:scattercut[0],:].plot(ax=ax1c,color=colm(range(n_colors),cmap=cmap))
					elif i<(len(scattercut)/2):
						DAC.loc[scattercut[2*i-1]:scattercut[2*i],:].plot(ax=ax1b,color=colm(range(n_colors),cmap=cmap), label='_nolegend_')
						normed.loc[scattercut[2*i-1]:scattercut[2*i],:].plot(ax=ax1a,color=colm(range(n_colors),cmap=cmap), label='_nolegend_')			
						DAC_copy.loc[scattercut[2*i-1]:scattercut[2*i],:].plot(ax=ax1c,color=colm(range(n_colors),cmap=cmap), label='_nolegend_')	
					else:
						DAC.loc[scattercut[-1]:,:].plot(ax=ax1b,color=colm(range(n_colors),cmap=cmap), label='_nolegend_')
						normed.loc[scattercut[-1]:,:].plot(ax=ax1a,color=colm(range(n_colors),cmap=cmap), label='_nolegend_')			
						DAC_copy.loc[scattercut[-1]:,:].plot(ax=ax1c,color=colm(range(n_colors),cmap=cmap), label='_nolegend_')	
			except:
				DAC.plot(ax=ax1b,color=colm(range(n_colors),cmap=cmap))
				normed.plot(ax=ax1a,color=colm(range(n_colors),cmap=cmap))
				DAC_copy.plot(ax=ax1c,color=colm(range(n_colors),cmap=cmap))	
		
		if mod in ['paral','exponential']:
			try:
				names=['decay %i: %.3g %s'%(i,a,baseunit) for i,a in enumerate(times)]
			except:
				print('something strange happened, most likely one value went "inf" or is set unexpectedly to None')
				names=['decay %i: %s %s'%(i,a,baseunit) for i,a in enumerate(times)]
			if 'background' in list(re['DAC'].columns):names.append('background')
			if 'Non Decaying' in list(re['DAC'].columns):names.append('Non Decaying')
			ax1a.legend(names,title='Model: {}'.format(mod))
			ax1b.legend(names,title='Model: {}'.format(mod))
			ax1c.legend(names,title='Model: {}'.format(mod))
		elif mod in ['exp']:
			names=['species %i'%i for i,a in enumerate(re['DAC'].columns.values)]
			if 'background' in list(re['DAC'].columns):
				if 'Non Decaying' in list(re['DAC'].columns):
					names[-1]='background'
					names[-2]='Non Decaying'
				else:
					names[-1]='background'
			else:
				if 'Non Decaying' in list(re['DAC'].columns):
					names[-1]='Non Decaying'
			ax1a.legend(names,title='Model: {}'.format(mod))
			ax1b.legend(names,title='Model: {}'.format(mod))
			ax1c.legend(names,title='Model: {}'.format(mod))
		else:
			names=['%s'%a for a in re['DAC'].columns.values]
			ax1a.legend(names,title='Model: {}'.format(mod))
			ax1b.legend(names,title='Model: {}'.format(mod))
			ax1c.legend(names,title='Model: {}'.format(mod))
		#ax1c.legend(time_string,title='Model: {}'.format(mod))
		if title is None:
			ax1a.set_title(f)
		else:
			if len(title)>0:
				ax1a.set_title(title)
		if title is None:
			ax1b.set_title(f)
		else:
			if len(title)>1:ax1b.set_title(title)
		if title is None:
			ax1c.set_title(f)
		else:
			if len(title)>1:ax1c.set_title(title)
		ax1a.plot(ax1a.get_xlim(),[0,0],'black',zorder=10)
		ax1b.plot(ax1b.get_xlim(),[0,0],'black',zorder=10)
		ax1c.plot(ax1b.get_xlim(),[0,0],'black',zorder=10)
		ax1a.set_xlabel(ds.columns.name)
		ax1b.set_xlabel(ds.columns.name)
		ax1c.set_xlabel(ds.columns.name)
		ax1a.set_ylabel('intensity norm.')
		ax1b.set_ylabel('intensity in arb. units')
		ax1c.set_ylabel('intensity*max(c) in arb. units')
		fig1.tight_layout()
	if 1 in plotting:  #-------plot sum_sum------------
		if halfsize:
			fig2 = plt.figure(figsize = (9, 2.5), dpi = 100)
		else:
			fig2 = plt.figure(figsize = (18, 5), dpi = 100)
		ax2a=[plt.subplot2grid((3, 3), (0, i)) for i in range(3)]
		ax2=[plt.subplot2grid((3, 3), (1, i), rowspan=2) for i in range(3)]
		dat = [pandas.DataFrame(re['A'], index = re['A'].index, columns = re['A'].columns).abs().sum(axis = 1)]
		dat.append(pandas.DataFrame(re['AC'], index = re['AC'].index, columns = re['AC'].columns).abs().sum(axis = 1))
		dat.append(pandas.DataFrame(re['AE'], index = re['AE'].index, columns = re['AE'].columns).abs().sum(axis = 1))
		dat_names=['measured','calculated','error']
		dat_styles=['*','-','-']
		dat_cols=colm(range(3), cmap = cmap)
		
		limits = (dat[0].min(), dat[0].max())
		xlimits = (dat[0].index.min(), dat[0].index.max())
		if ignore_time_region is None:
			for i in range(3):
				for j in range(3):
					if i==2:
						_ = dat[i].plot(ax = ax2a[j], label = dat_names[i], style = dat_styles[i], color = dat_cols[i])
					else:
						_ = dat[i].plot(ax = ax2[j], label = dat_names[i], style = dat_styles[i], color = dat_cols[i])
				
		elif isinstance(ignore_time_region[0], numbers.Number):
			x=dat[0].index.values.astype('float')
			lower=find_nearest_index(x,ignore_time_region[0])
			upper=find_nearest_index(x,ignore_time_region[1])
			for i in range(3):
				for j in range(3):
					if i==2:
						_ = dat[i].iloc[:lower].plot(ax = ax2a[j], label = dat_names[i], style = dat_styles[i], color = dat_cols[i])
						_ = dat[i].iloc[upper:].plot(ax = ax2a[j], label = '_nolegend_', style = dat_styles[i], color = dat_cols[i])						
					else:
						_ = dat[i].iloc[:lower].plot(ax = ax2[j], label = dat_names[i], style = dat_styles[i], color = dat_cols[i])
						_ = dat[i].iloc[upper:].plot(ax = ax2[j], label = '_nolegend_', style = dat_styles[i], color = dat_cols[i])
		else:
			try:
				ignore_time_region_loc=flatten(ignore_time_region)
				for k in range(len(ignore_time_region_loc)/2+1):
					if k==0:
						for i in range(3):
							for j in range(3):
								if i==2:
									_ = dat[i].loc[:ignore_time_region_loc[k]].plot(ax = ax2a[j], label = dat_names[i], style = dat_styles[i], color = colm(i, cmap = cmap))
								else:
									_ = dat[i].loc[:ignore_time_region_loc[k]].plot(ax = ax2[j], label = dat_names[i], style = dat_styles[i], color = colm(i, cmap = cmap))			
					elif k<(len(ignore_time_region)/2):
						for i in range(3):
							for j in range(3):
								if i==2:
									_ = dat[i].loc[ignore_time_region_loc[2*k-1]:ignore_time_region_loc[2*k]].plot(ax = ax2a[j], label = '_nolegend_', style = dat_styles[i], color = colm(i, cmap = cmap))
								else:
									_ = dat[i].loc[ignore_time_region_loc[2*k-1]:ignore_time_region_loc[2*k]].plot(ax = ax2[j], label = '_nolegend_', style = dat_styles[i], color = colm(i, cmap = cmap))	
					else:
						for i in range(3):
							for j in range(3):
								if i==2:
									_ = dat[i].loc[ignore_time_region_loc[-1]:].plot(ax = ax2a[j], label = '_nolegend_', style = dat_styles[i], color = colm(i, cmap = cmap))									
								else:
									_ = dat[i].loc[ignore_time_region_loc[-1]:].plot(ax = ax2[j], label = '_nolegend_', style = dat_styles[i], color = colm(i, cmap = cmap))	
			except:
				for i in range(3):
					for j in range(3):
						if i==2:
							_ = dat[i].plot(ax = ax2a[j], label = dat_names[i], style = dat_styles[i], color = colm(i, cmap = cmap))
						else:
							_ = dat[i].plot(ax = ax2[j], label = dat_names[i], style = dat_styles[i], color = colm(i, cmap = cmap))
		ax2[0].set_xlim(xlimits)
		ax2[0].set_xscale('symlog', linscale=0.1)
		ax2[0].autoscale(axis='y', tight=True)
		ax2a[0].set_xlim(xlimits)
		ax2a[0].set_xscale('symlog', linscale=0.1)
		ax2a[0].autoscale(axis='y', tight=True)
		
		ax2[1].set_xlim(xlimits)
		ax2[1].set_xscale('linear')
		ax2[1].set_ylim(np.nanmax([limits[0],limits[1]/10000]),limits[1])
		ax2[1].set_yscale('log')
		ax2a[1].set_xlim(xlimits)
		ax2a[1].set_xscale('linear')	

		ax2[2].set_xscale('log')
		ax2[2].set_ylim(np.nanmax([limits[0],limits[1]/10000]),limits[1])
		ax2[2].set_yscale('log')
		ax2[2].set_xlim(max(0.1, xlimits[0]), xlimits[1])
		ax2a[2].set_xscale('log')
		ax2a[2].set_xlim(max(0.1, xlimits[0]), xlimits[1])
		
		#draw a black line at zero
		ax2a[0].plot(ax2[0].get_xlim(), [0, 0], 'black', zorder=10, label = '_nolegend_')
		ax2a[1].plot(ax2[0].get_xlim(), [0, 0], 'black', zorder=10, label = '_nolegend_')
		ax2a[2].plot(ax2[0].get_xlim(), [0, 0], 'black', zorder=10, label = '_nolegend_')
		#plot empty to get the labels right
		ax2[0].plot([], [], ' ', label=stringen)
		ax2[0].legend(title='Model: {}'.format(mod),frameon=False)
		if title is None:
			ax2[1].legend(labels=[],title=f,frameon=False)
		else:
			if not len(title)==0:
				ax2[1].legend(labels=[],title=title,frameon=False)
		for t in times:
			if isinstance(t,float):
				ax2[0].plot([t,t],ax2[0].get_ylim(),lw=0.5,zorder=10,alpha=0.5,color='black')
				ax2[1].plot([t,t],ax2[1].get_ylim(),lw=0.5,zorder=10,alpha=0.5,color='black')
				ax2[2].plot([t,t],ax2[2].get_ylim(),lw=0.5,zorder=10,alpha=0.5,color='black')
		x_label=ds.index.name
		ax2[0].set_xlabel(x_label)
		ax2[1].set_xlabel(x_label)
		ax2[2].set_xlabel(x_label)
		ax2a[0].set_xlabel(x_label)
		ax2a[1].set_xlabel(x_label)
		ax2a[2].set_xlabel(x_label)
		fig2.tight_layout()
	if 2 in plotting:#---plot single wavelength----------
		if halfsize:
			fig3,ax3 = plt.subplots(figsize = (7.5,3),dpi = 100)
		else:
			fig3,ax3 = plt.subplots(figsize = (15,6),dpi = 100)
		#fig3,ax3 = plt.subplots(figsize = (8,4),dpi = 100)
		_=plot1d( ds = re['AC'], cmap = cmap, ax = ax3, width = width, wavelength = rel_wave, 
					  lines_are = 'fitted', plot_type = scale_type, baseunit = baseunit, lintresh = lintresh, 
					  timelimits = timelimits, text_in_legend = time_string, title = title, 
					  ignore_time_region = ignore_time_region,  data_type = data_type, units = units, from_fit = True)
		_=plot1d( ds = re['A'], cmap = cmap,ax = ax3, subplot = True, width = width, 
					  wavelength = rel_wave,lines_are = 'data', plot_type = scale_type, 
					  baseunit = baseunit , lintresh = lintresh , timelimits = timelimits, 
					  ignore_time_region = ignore_time_region,  data_type = data_type, units = units, from_fit = True)
		ax3.autoscale(axis = 'y',tight = True)
		for t in times:
			if isinstance(t, float):
				ax3.plot([t, t], ax3.get_ylim(), lw = 0.5, zorder = 10, alpha = 0.5, color = 'black')
		fig3.tight_layout()
	if 3 in plotting:		#---plot at set time----------
		if halfsize:
			fig4, ax4 = plt.subplots(figsize = (7.5, 3), dpi = 100)	
		else:
			fig4, ax4 = plt.subplots(figsize = (15, 6), dpi = 100)	
		_=plot_time(ds=re['A'],cmap=cmap,ax=ax4,subplot=False, rel_time=rel_time, title=title, baseunit=baseunit, 
					time_width_percent=time_width_percent, lines_are='data', scattercut=scattercut, 
					bordercut = bordercut, intensity_range = intensity_range,  data_type = data_type, 
					plot_second_as_energy = plot_second_as_energy, units = units, equal_energy_bin = equal_energy_bin, from_fit = True )
		_=plot_time(ds=re['AC'],cmap=cmap,ax=ax4, subplot=False, rel_time=rel_time, title=title, 
						baseunit=baseunit, time_width_percent=time_width_percent, lines_are='fitted',
						color_offset=color_offset, scattercut=scattercut,  data_type = data_type, 
						plot_second_as_energy = plot_second_as_energy, units = units, equal_energy_bin = equal_energy_bin, from_fit = True )
		try:
			if equal_energy_bin is not None:
				bordercut=[scipy.constants.h*scipy.constants.c/(a*1e-9*scipy.constants.electron_volt) for a in bordercut]
			ax4.set_xlim(bordercut)
		except:
			pass
		fig4.tight_layout()
	if 4 in plotting:#---matrix with fit and error, as figures----------
		fig5 = plot2d_fit(re, cmap = cmap, intensity_range = intensity_range, baseunit = baseunit, 
							error_matrix_amplification = error_matrix_amplification, wave_nm_bin = None, equal_energy_bin = equal_energy_bin,
							use_images = True, log_scale = log_scale, scale_type = scale_type, patches = patches, 
							lintresh = lintresh, bordercut = bordercut, ignore_time_region = ignore_time_region, 
							scattercut = scattercut, timelimits = timelimits, levels = 200,  data_type = data_type)	
		plt.ion()
		plt.show()
	if 5 in plotting:
		if halfsize:
			fig6=plt.figure(figsize=(6,3))
		else:
			fig6=plt.figure(figsize=(12,6))
		G = GridSpec(4, 4)
		ax6=[]
		ax6.append(fig6.add_subplot(G[1:,0]))
		ax6.append(fig6.add_subplot(G[1:,1]))
		ax6.append(fig6.add_subplot(G[1:,2]))
		ax6.append(fig6.add_subplot(G[1:,3]))
		ax6.append(fig6.add_subplot(G[0,2:]))
		ax6.append(fig6.add_subplot(G[0,0:2]))
		n_colors=len(re['c'].columns)
		for i in range(4):
			ax6[i]=re['c'].plot(ax=ax6[i],color=colm(range(n_colors),cmap=cmap),legend=False)
		if re['c'].index.name == 'time':
			for i in range(4):
				ax6[i].set_xlabel('time in %s'%baseunit)
		ax6[1].set_yscale('log')
		ax6[1].set_ylim(1e-5,1.1)
		ax6[2].set_yscale('log')
		ax6[2].set_ylim(1e-5,1.1)
		ax6[2].set_xscale('log')
		ax6[2].set_xlim(0.05,ax6[2].get_xlim()[1])
		ax6[3].set_xscale('log')
		ax6[3].set_xlim(0.05,ax6[3].get_xlim()[1])
		handles, labels = ax6[3].get_legend_handles_labels()
		ax6[4].axis('off')
		ax6[5].axis('off')
		if title is None:
			title=f
		else:
			if not len(title)==0:
				title=title
			else:
				title=''
		if len(handles)<5:
			ncol=2
		elif len(handles)<7:
			ncol=3
		else:
			ncol=4
		leg=ax6[4].legend(handles,labels,loc=3, ncol=ncol,edgecolor=(0,0,0,1),framealpha=1,frameon=True,title=title)   
		for i in range(6):
			ax6[i].set_title('')
		ax6[5].text(0,0,'This factor represents the temporal evolution\n of the components in the fit.\nThis time dependent factor multiplied with the \nspectral intensity from the SAS/DAS is re[\"AC\"]',fontsize=float(plt.rcParams['legend.fontsize'])-1)
		fig6.tight_layout()
	if 6 in plotting:#---residuals----------
		if halfsize:
			fig7,ax7 = plt.subplots(figsize = (6,3),dpi = 100)
		else:
			fig7,ax7 = plt.subplots(figsize = (12,6),dpi = 100)
		#fig3,ax3 = plt.subplots(figsize = (8,4),dpi = 100)
		_=plot1d( ds = re['AE'], cmap = cmap, ax = ax7, width = width, wavelength = rel_wave, 
					  lines_are = 'smoothed', plot_type = scale_type, baseunit = baseunit, lintresh = lintresh, 
					  timelimits = timelimits, text_in_legend = time_string, title = 'residuals', 
					  ignore_time_region = ignore_time_region,  data_type = data_type, units = units, from_fit = True)
		_=plot1d( ds = re['AE'], cmap = cmap,ax = ax7, subplot = True, width = width, 
					  wavelength = rel_wave,lines_are = 'data', plot_type = scale_type, 
					  baseunit = baseunit , lintresh = lintresh , timelimits = timelimits, 
					  ignore_time_region = ignore_time_region,  data_type = data_type, units = units, from_fit = True)
		ax7.autoscale(axis = 'y',tight = True)
		for t in times:
			if isinstance(t, float):
				ax7.plot([t, t], ax7.get_ylim(), lw = 0.5, zorder = 10, alpha = 0.5, color = 'black')
		fig7.tight_layout()							
	plt.ion()
	plt.show()
	if print_click_position:
		plt.connect('button_press_event', mouse_move)
	if save_figures_to_folder:
		if path is None:path=os.path.dirname(os.path.realpath(__file__))
		figure_path=check_folder(path=path)
		if filename is None:
			filename='test.fig'
		fi=filename.split('.')[0]
		name_extension=['DAC','SUM','SEL','SPEC','FIG_MAT','concentrations','RESIDUAL']
		for a in range(7):
			try:
				eval('fig%i'%(a+1)).savefig(check_folder(path=figure_path,filename='%s_%s.%s'%(fi,name_extension[a],savetype)),bbox_inches='tight')
			except Exception as e:
				print(e)
	if return_figures_handles:
		name_extension=['DAC','SUM','SEL','SPEC','FIG_MAT','concentrations','RESIDUAL']
		returning_dict={}
		for a in range(7):
			try:
				returning_dict[name_extension[a]]=eval('fig%i'%(a+1))
			except Exception as e:
				print(e)
		return returning_dict

def plot_raw(ds = None, plotting = range(4), title = None, intensity_range = 1e-2, baseunit = 'ps',
			timelimits = None, scattercut = None, bordercut = None, wave_nm_bin = None, width = 10,
			rel_wave = np.arange(400, 900, 100), rel_time = [1, 5, 10], time_width_percent = 10,
			ignore_time_region = None, time_bin = None, cmap = None, color_offset = 0, 
			log_scale = True, plot_type = 'symlog', lintresh = 0.3, times = None, 
			save_figures_to_folder = False, savetype = 'png', path = None, filename = None, 
			print_click_position = False, data_type = 'differential Absorption in $\mathregular{\Delta OD}$',
			plot_second_as_energy = True, units = 'nm', return_plots = False, equal_energy_bin = None,
			return_figures_handles=False):
	'''This is the extended plot function, for convenient object based plotting see TA.Plot_RAW 
	This function plotts of various RAW (non fitted) plots. Based on the DataFrame ds a number of 
	cuts are created using the shaping parameters explained below.
	In all plots the RAW data is plotted as dots and interpolated with lines 
	(using Savitzky-Golay window=5, order=3 interpolation). 

	Parameters
	---------------
	
	ds : DataFrame
		This dataframe contains the data to be plotted. It is copied and sliced into the
		regions defined. The dataframe expects the time to be in Index and the wavelength/energy 
		to be in the columns. The spectra is plotted with a second (energy) axis
		
	plotting : int or iterable (of integers), optional
		This parameter determines which figures are plotted 
		the figures can be called separately with plotting = 1
		or with a list of plots (Default) e.g.plotting=range(4) calls plots 0,1,2,3
		The plots have the following numbers: 
		
			0. Matrix
			1. Kinetics
			2. Spectra
			3. SVD 
		
	title : None or str
		title to be used on top of each plot
		The (Default) None triggers  self.filename to be used. Setting a specific title as string will
		be used in all plots. To remove the title all together set an empty string with this command title="" 
		
	intensity_range : None, float or list [of two floats]
		intensity_range is a general switch that governs what intensity range the plots show. 
		For the 1d plots this is the y-axis for the 2d-plots this is the colour scale. 
		This parameter recognizes three settings. If set to "None" (Default) this uses the minimum and
		maximum of the data. A single value like in the example below and the intended use is the symmetric
		scale while a list with two entries an assymmetric scale e.g. 
		intensity_range=3e-3 is converted into intensity_range=[-3e-3,3e-3]
		
	data_type : str
		this is the datatype and effectively the unit put on the intensity axis 
		(Default)'differential Absorption in $\mathregular{\Delta OD}$
	
	baseunit : str
		baseunit is a neat way to change the unit on the time axis of the plots. (Default) "ps", but they 
		can be frames or something similarly. This is changing only the label of the axis. 
		During the import there is the option to divide the numbers by a factor. 
		I have also used frames or fs as units. Important is that all time units will be labeled with
		this unit.
		
	timelimits : None or list (of 2 floats), optional
		cut times at the low and high time limit. (Default) None uses the limits of measurement
		Important: If either the background or the chirp is to be fit this must include the 
		time before zero! Useful: It is useful to work on different regions, starting with
		the longest (then use the ta.Backgound function prior to fit) and expand from there
		
	scattercut : None or iterable (of floats or other iterable, always pairs!), optional
		intented to "cut" one or multiple scatter regions. (if (Default) None nothing
		happens) If it is set the spectral region between the limits is set to zero. 
		Usage single region: [lower region limit,upper region limit], 
		use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]
		
	bordercut : None or iterable (with two floats), optional
		cut spectra at the low and high wavelength limit. (Default) None 
		uses the limits of measurement 
		
	wave_nm_bin : None or float, optional
		rebins the original data into even intervals. If set to None the original data will be used. 
		If set to a width (e.g. 2nm), the wavelength axis will be divided into steps of this size
		and the mean of all measurements in the interval is taken. The re-binning stops as soon as
		the measured stepsize is wider than given here, then the original bins are used. 
		This function is particularly useful for spectrometer with non-linear dispersion, 
		like a prism in the infrared.
	
	equal_energy_bin : None or float(optional)
		if this is set the wave_nm_bin is ignored and the data is rebinned into equal energy bins (based upon that the data is in nm. 
		If dual axis  is on then the lower axis is energy and the upper is wavelength
	
	width : float, optional
		the width used in kinetics, see below (Default) 10nm
		
	rel_wave : float or list (of floats), optional 
		"rel_wave" and "width" (in the object called "wavelength_bin" work together for the creation 
		of kinetic plots. When plotting kinetic spectra one line will be plotted for each entrance
		in the list/vector rel_wave. During object generation the vector np.arange(300,1000,100) 
		is set as standard. Another typical using style would be to define a list of interesting 
		wavelength at which a kinetic development is to be plotted. At each selected wavelength 
		the data between wavelength+ta.wavelength_bin and wavelength-ta.wavelength_bin is averaged 
		for each timepoint returned 
		
	rel_time : float or list/vector (of floats), optional
		For each entry in rel_time a spectrum is plotted. If time_width_percent=0 (Default) the 
		nearest measured timepoint is chosen. For other values see "time_width_percent"
		
	time_width_percent : float
		"rel_time" and "time_width_percent" work together for creating spectral plots at 
		specific timepoints. For each entry in rel_time a spectrum is plotted. 
		If however e.g. time_width_percent=10 the region between the timepoint closest 
		to the  1.1 x timepoint and 0.9 x timepoint is averaged and shown 
		(and the legend adjusted accordingly). This is particularly useful for the densly
		sampled region close to t=0. Typically for a logarithmic recorded kinetics, the 
		timepoints at later times will be further appart than 10 percent of the value, 
		but this allows to elegantly combine values around time=0 for better statistics. 
		This averaging is only applied for the plotting function and not for the fits.
		
	ignore_time_region : None or list (of two floats or of lists), optional
		cut set a time range with a low and high limit from the fits. (Default) None nothing happens
		The region will be removed during the fitting process (and will be missing in the fit-result
		plots) 
		Usage single region: [lower region limit,upper region limit],
		use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]
		
	time_bin : None or int, optional
		is dividing the points on the time-axis in even bins and averages the found values in between. 
		This is a hard approach that also affects the fits. I do recommend to use this carefully, 
		it is most useful for modulated data. A better choice for transient absorption that only 
		affects the kinetics is "time_width_percent"
		
	cmap : None or matplotlib color map, optional
		is a powerfull variable that chooses the colour map applied for all plots. If set to 
		None (Default) then the self.cmap is used.
		As standard I use the color map "jet" from matplotlib. There are a variety of colormaps 
		available that are very usefull. Beside "jet", "viridis" is a good choice as it is well 
		visible under red-green blindness. Other useful maps are "prism" for high fluctuations 
		or diverging color maps like "seismic". 
		See https://matplotlib.org/3.1.0/tutorials/colors/colormaps.html for a comprehensive 
		selection. In the code the colormaps are imported so if plot_func is imported as pf then 
		self.cmap=pf.cm.viridis sets viridis as the map to use. Internally the colors are chosen 
		with the "colm" function. The 2d plots require a continuous color map so if something 
		else is give 2d plots are shown automatically with "jet". For all of the 1d plots however 
		I first select a number of colors before each plot. If cmap is a continous map then these
		are sampled evenly over the colourmap. Manual iterables of colours 
		cmap=[(1,0,0),(0,1,0),(0,0,1),...] are also accepted, as are vectors or dataframes that 
		contain as rows the colors. There must be of course sufficient colors present for 
		the numbers of lines that will be plotted. So I recommend to provide at least 10 colours 
		(e.g. your university colors). colours are always given as a, list or tuple with RGA or RGBA
		(with the last A beeing the Alpha=transparency. All numbers are between 0 and 1. 
		If a list/vector/DataFrame is given for the colours they will be used in the order provided.
		 
	color_offset : int, optional
		At the (Default) 0 the colours are chose from the beginning, for a larger value Color_offset 
		colors are skipped. Usually only used if multiple plots are created, and the data/or fit is
		only shown for some of them.
		 
	log_scale : bool, optional
		If True (Default), The 2D plots (Matrix) is plotted with a pseudo logarithmic intensity scale. 
		This usually does not give good results unless the intensity scale is symmetric 
		
	Scale_type : None or str
		is a general setting that can influences what time axis will be used for the plots. 
		"symlog" (linear around zero and logarithmic otherwise) "lin" and "log" are valid options.
		
	lintresh : float
		The pseudo logratihmic range "symlog" is used for most time axis. Symlog plots a range around
		time zero linear and beyond this linear treshold 'lintresh' on a logarithmic scale. (Default) 0.3 
		
	times : None or int
		are the number of components to be used in the SVD (Default) is None (which is seen as 6)
		
	save_figures_to_folder : bool, optional
		(Default) is False, if True the Figures are automatically saved 
	
	return_figures_handles : bool, optional
		(Default) is False, if True the Figure handles are returned as a dictionary.
		
	savetype : str or iterable (of str), optional 
		matplotlib allows the saving of figures in various formats. (Default) "png", 
		typical and recommendable options are "svg" and "pdf".  
		
	path : None, str or path object, optional
		This defines where the files are saved if the safe_figures_to_folder parameter is True, 
		quite useful if a lot of data sets are to be printed fast. 
		If a path is given, this is used. If a string like the (Default) "result_figures" is given, 
		then a subfolder of this name will be used (an generated if necessary) 
		relative to self.path. Use and empty string to use the self.path
		If set to None, the location of the plot_func will be used and
		a subfolder with title "result_figures" be generated here
		
	filename : str, optional
		offers to replace the base-name used for all plots (to e.g.~specify what sample was used). 
		if (Default) None is used, the self.filename is used as a base name. The filename plays only a 
		role during saving, as does the path and savetype
	
	print_click_position : bool, optional
		if True then the click position is printed for the spectral plots 
		
	return_plots : bool, optional 
		(Default) False, return is ignoriert. For True a dictionary with the handles to the figures is returned  
		
	'''
	global halfsize
	if ds is None:raise ValueError('We need something to plot!!!')
	if baseunit != 'ps':
		if baseunit == 'ns':
			ds.index.name = 'Time in ns'
		else:
			ds.index.name=baseunit
															
	if path is None:path=check_folder(path='result_figures',current_path=os.path.dirname(os.path.realpath(__file__)))
	if filename is None:filename='standard.sia'
	if not hasattr(plotting,'__iter__'):plotting=[plotting]
	debug=False
	plt.ion()
	if 0 in plotting:#MAtrix
		fig1 = plot2d(ds = ds, cmap = cmap, ignore_time_region = ignore_time_region, plot_type = plot_type, 
					baseunit = baseunit, intensity_range = intensity_range, scattercut = scattercut, 
					bordercut = bordercut, wave_nm_bin = wave_nm_bin, levels = 200, lintresh = lintresh, 
					timelimits = timelimits, time_bin = time_bin, title = title, log_scale = log_scale, 
					data_type = data_type, equal_energy_bin = equal_energy_bin)
		fig1.tight_layout()
		if debug:print('plotted Matrix')
	if 1 in plotting:#kinetics
		if halfsize:
			fig2,ax2=plt.subplots(figsize=(5,3),dpi=100)
		else:
			fig2,ax2=plt.subplots(figsize=(10,6),dpi=100)
		
		_ = plot1d(ds = ds, ax = ax2, subplot = True, cmap = cmap, width = width, wavelength = rel_wave,
					title = title, lines_are = 'data' ,	  plot_type = plot_type, lintresh = lintresh, 
					timelimits = timelimits, intensity_range = intensity_range, scattercut = scattercut, 
					ignore_time_region = ignore_time_region, baseunit = baseunit, data_type = data_type, 
					units = units)
		_ = plot1d(ds = ds, ax = ax2, subplot = False, cmap = cmap, width = width, wavelength = rel_wave,
					title = title, lines_are = 'smoothed', plot_type = plot_type, lintresh = lintresh, 
					timelimits = timelimits, intensity_range = intensity_range, scattercut = scattercut, 
					ignore_time_region = ignore_time_region, baseunit = baseunit, data_type = data_type, 
					units = units )
		if debug:print('plotted kinetics')
	if 2 in plotting:#Spectra
		if halfsize:
			fig3,ax3 = plt.subplots(figsize = (5,3),dpi = 100)
		else:
			fig3,ax3 = plt.subplots(figsize = (10,6),dpi = 100)

		_ = plot_time(ds, subplot = True, ax = ax3, plot_second_as_energy = False, cmap = cmap, 
						rel_time = rel_time, time_width_percent = time_width_percent, title = title, 
						baseunit = baseunit, lines_are = 'data'	  , scattercut = scattercut, 
						wave_nm_bin = wave_nm_bin, bordercut = bordercut, intensity_range = intensity_range, 
						ignore_time_region = ignore_time_region, data_type = data_type, units = units, 
						equal_energy_bin = equal_energy_bin)
		if plot_second_as_energy:
			_ = plot_time(ds ,subplot = False, ax = ax3, plot_second_as_energy = True, cmap = cmap, 
					rel_time = rel_time, time_width_percent = time_width_percent, title = title, 
					baseunit = baseunit, lines_are = 'smoothed', scattercut = scattercut, 
					wave_nm_bin = wave_nm_bin, bordercut = bordercut, intensity_range = intensity_range, 
					ignore_time_region = ignore_time_region, data_type = data_type, units = units, 
					equal_energy_bin = equal_energy_bin)
		else:
			_ = plot_time(ds ,subplot = False, ax = ax3, plot_second_as_energy = False, cmap = cmap, 
					rel_time = rel_time, time_width_percent = time_width_percent, title = title, 
					baseunit = baseunit, lines_are = 'smoothed', scattercut = scattercut, 
					wave_nm_bin = wave_nm_bin, bordercut = bordercut, intensity_range = intensity_range, 
					ignore_time_region = ignore_time_region, data_type = data_type, units = units, 
					equal_energy_bin = equal_energy_bin)

		fig3.tight_layout()
		if debug:print('plotted Spectra')
	if 3 in plotting:		#---plot at set time----------
		try:
			fig4 = SVD(ds ,  times = times ,  timelimits = timelimits ,  scattercut = scattercut ,  
					bordercut = bordercut ,  wave_nm_bin = wave_nm_bin, ignore_time_region = ignore_time_region,
					cmap = cmap)
											  
		except:
			print("SVD failed with:",sys.exc_info()[0])
		if debug:print('plotted SVD')
	if print_click_position:
		plt.connect('button_press_event', mouse_move)
	plt.show()
	if save_figures_to_folder:
		fi=filename.split('.')[0]
		name_extension=['RAW_MAT','RAW_SEL','RAW_SPEK','RAW_SVD']
		for a in range(4):
			try:
				eval('fig%i'%(a+1)).savefig(check_folder(path=path,filename='%s_%s.%s'%(fi,name_extension[a],savetype)),bbox_inches='tight',dpi=300)
			except Exception as e:
				print(e)
	if return_figures_handles:
		return_dicten={}
		name_extension=['RAW_MAT','RAW_SEL','RAW_SPEK','RAW_SVD']
		for i in range(4):
			try:
				return_dicten[name_extension[i]]=eval('fig%i'%(i+1))
			except:
				pass
		return return_dicten

def plot_time(ds, ax = None, rel_time = None, time_width_percent = 10, ignore_time_region = None, 
			wave_nm_bin = None, title = None, text_in_legend = None, baseunit = 'ps', 
			lines_are = 'smoothed', scattercut = None, bordercut = None, subplot = False, linewidth = 1, 
			color_offset = 0, intensity_range = None, plot_second_as_energy = True, cmap = standard_map, 
			data_type = 'differential Absorption in $\mathregular{\Delta OD}$', units = 'nm', equal_energy_bin = None, from_fit = None):
	'''Function to create plots at a certain time. In general you give under rel_time a 
		list of times at which yu do want to plot the time width percentage means that 
		this function integrates ewverything plus minus 10% at this time. lines_are is a 
		switch that regulates what is plotted. data plots the data only, 
		smoothed plots the data and a smoothed version of the data, fitted plots only the fit. 
		the subplot switch is for using this to plot e.g. multiple different datasets.
		
	Parameters
	---------------
	
	ds : DataFrame
		This dataframe contains the data to be plotted. It is copied and sliced into the
		regions defined. The dataframe expects the time to be in Index and the wavelength/energy 
		to be in the columns. The spectra is plotted with a second (energy) axis
	
	ax : None or matplotlib axis object, optional
		if None (Default), a figure and axis will be generated for the plot, if axis is given the plot will 
		placed in there.
	
	data_type : str
		this is the datatype and effectively the unit put on the intensity axis 
		(Default)'differential Absorption in $\mathregular{\Delta OD}$
	
	rel_time : float or list/vector (of floats), optional
		For each entry in rel_time a spectrum is plotted. If time_width_percent=0 (Default) the 
		nearest measured timepoint is chosen. For other values see 'time_width_percent'
	
	time_width_percent : float
		"rel_time" and "time_width_percent" work together for creating spectral plots at 
		specific timepoints. For each entry in rel_time a spectrum is plotted. 
		If however e.g. time_width_percent=10 the region between the timepoint closest 
		to the  1.1 x timepoint and 0.9 x timepoint is averaged and shown 
		(and the legend adjusted accordingly). This is particularly useful for the densly
		sampled region close to t=0. Typically for a logarithmic recorded kinetics, the 
		timepoints at later times will be further appart than 10 percent of the value, 
		but this allows to elegantly combine values around time=0 for better statistics. 
		This averaging is only applied for the plotting function and not for the fits.

	ignore_time_region : None or list (of two floats or of lists), optional
		cut set a time range with a low and high limit from the fits. (Default) None nothing happens
		The region will be removed during the fitting process (and will be missing in the fit-result
		plots) 
		Usage single region: [lower region limit,upper region limit],
		use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]
		
	wave_nm_bin : None or float, optional
		rebins the original data into even intervals. If set to None the original data will be used. 
		If set to a width (e.g. 2nm), the wavelength axis will be divided into steps of this size
		and the mean of all measurements in the interval is taken. The re-binning stops as soon as
		the measured stepsize is wider than given here, then the original bins are used. 
		This function is particularly useful for spectrometer with non-linear dispersion, 
		like a prism in the infrared.
		
	title : None or str, optional
		title to be used on top of each plot
		The (Default) None triggers  self.filename to be used. Setting a specific title as string will
		be used in all plots. To remove the title all together set an empty string with this command title="" 
		
	linewidth : float, optional
		linewidth to be used for plotting
				
	text_in_legend : str, optional
		text to be used in legend before the actually lines and colours (set as heasder)
	
	baseunit : str
		baseunit is a neat way to change the unit on the time axis of the plots. (Default) 'ps', but they 
		can be frames or something similarly. This is changing only the label of the axis. 
		During the import there is the option to divide the numbers by a factor. 
		I have also used frames or fs as units. Important is that all time units will be labeled with
		this unit.
				
	scattercut : None or iterable (of floats or other iterable, always pairs!), optional
		intented to "cut" one or multiple scatter regions. (if (Default) None nothing
		happens) If it is set the spectral region between the limits is set to zero. 
		Usage single region: [lower region limit,upper region limit], 
		use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]
		
	bordercut : None or iterable (with two floats), optional
		cut spectra at the low and high wavelength limit. (Default) None 
		uses the limits of measurement 

	subplot ; bool, optional
		False (Default) means this is a main plot in this axis! if True then this is the second plot in the axis 
		and things like axis ticks should not be reset
		this also avoids adding the object to the legend

	color_offset : int, optional
		At the (Default) 0 the colours are chose from the beginning, for a larger value Color_offset 
		colors are skipped. Usually only used if multiple plots are created, and the data/or fit is
		only shown for some of them.
		 		
	intensity_range : None, float or list [of two floats]
		intensity_range is a general switch that governs what intensity range the plots show. 
		For the 1d plots this is the y-axis for the 2d-plots this is the colour scale. 
		This parameter recognizes three settings. If set to "None" (Default) this uses the minimum and
		maximum of the data. A single value like in the example below and the intended use is the symmetric
		scale while a list with two entries an assymmetric scale e.g. 
		intensity_range=3e-3 is converted into intensity_range=[-3e-3,3e-3]
				
	plot_second_as_energy : bool, optional
		For (Default) True a second x-axis is plotted  with "eV" as unit
	
	cmap : None or matplotlib color map, optional
		is a powerfull variable that chooses the colour map applied for all plots. If set to 
		None (Default) then the self.cmap is used.
		As standard I use the color map "jet" from matplotlib. There are a variety of colormaps 
		available that are very usefull. Beside "jet", "viridis" is a good choice as it is well 
		visible under red-green blindness. Other useful maps are "prism" for high fluctuations 
		or diverging color maps like "seismic". 
		See https://matplotlib.org/3.1.0/tutorials/colors/colormaps.html for a comprehensive 
		selection. In the code the colormaps are imported so if plot_func is imported as pf then 
		self.cmap=pf.cm.viridis sets viridis as the map to use. Internally the colors are chosen 
		with the "colm" function. The 2d plots require a continuous color map so if something 
		else is give 2d plots are shown automatically with "jet". For all of the 1d plots however 
		I first select a number of colors before each plot. If cmap is a continous map then these
		are sampled evenly over the colourmap. Manual iterables of colours 
		cmap=[(1,0,0),(0,1,0),(0,0,1),...] are also accepted, as are vectors or dataframes that 
		contain as rows the colors. There must be of course sufficient colors present for 
		the numbers of lines that will be plotted. So I recommend to provide at least 10 colours 
		(e.g.~your university colors). colours are always given as a, list or tuple with RGA or RGBA
		(with the last A beeing the Alpha=transparency. All numbers are between 0 and 1. 
		If a list/vector/DataFrame is given for the colours they will be used in the order provided.
		 

		'''
	global halfsize
	if not hasattr(rel_time,'__iter__'):rel_time=[rel_time]
	rel_time=[a for a in rel_time if a<ds.index.values.astype('float').max()]
	
	if isinstance(cmap,list):
		colors=cmap[color_offset:]
	else:
		colors=colm(np.arange(color_offset,len(rel_time)+color_offset),cmap=cmap)
	if ax is None:
		if halfsize:
			fig,ax1=plt.subplots(figsize=(5,3),dpi=100)
		else:
			fig,ax1=plt.subplots(figsize=(10,6),dpi=100)
	else:
		ax1=ax
	ds = sub_ds(ds = ds, times = rel_time, time_width_percent = time_width_percent, 
				scattercut = scattercut, drop_scatter=True, bordercut = bordercut, baseunit=baseunit,
				ignore_time_region = ignore_time_region, wave_nm_bin = wave_nm_bin, equal_energy_bin = equal_energy_bin, from_fit = from_fit)
	if 'smoothed' in lines_are:
		if scattercut is None:
			smoothed=Frame_golay(ds, window = 5, order = 3,transpose=False)
			smoothed.plot(ax = ax1, style = '-', color = colors, legend = False, lw = linewidth)
		elif isinstance(scattercut[0], numbers.Number):#handling single scattercut
			if equal_energy_bin is not None:
				scattercut=[scipy.constants.h*scipy.constants.c/(a*1e-9*scipy.constants.electron_volt) for a in scattercut]
				scattercut=scattercut[::-1]
			smoothed=Frame_golay(ds.loc[:scattercut[0],:], window = 5, order = 3,transpose=False)
			smoothed.plot(ax = ax1, style = '-', color = colors, legend = False, lw = linewidth)
			smoothed=Frame_golay(ds.loc[scattercut[1]:,:], window = 5, order = 3,transpose=False)
			smoothed.plot(ax = ax1, style = '-', color = colors, legend = False, lw = linewidth, label='_nolegend_')
		else:#handling multiple scattercuts
			try:
				scattercut=flatten(scattercut)
				if equal_energy_bin is not None:
					scattercut=[scipy.constants.h*scipy.constants.c/(a*1e-9*scipy.constants.electron_volt) for a in scattercut]
					scattercut=scattercut[::-1]
				for i in range(len(scattercut)):
					if i==0:
						smoothed=Frame_golay(ds.loc[:scattercut[0],:], window = 5, order = 3,transpose=False)
						smoothed.plot(ax = ax1, style = '-', color = colors, legend = False, lw = linewidth)
					elif i<(len(scattercut)/2):
						smoothed=Frame_golay(ds.loc[scattercut[2*i-1]:scattercut[2*i],:], window = 5, order = 3,transpose=False)
						smoothed.plot(ax = ax1, style = '-', color = colors, legend = False, lw = linewidth, label='_nolegend_')
					else:
						smoothed=Frame_golay(ds.loc[scattercut[-1]:,:], window = 5, order = 3,transpose=False)
						smoothed.plot(ax = ax1, style = '-', color = colors, legend = False, lw = linewidth, label='_nolegend_')
			except:
				print('printing the smoothed scatter interpolation created an error, using default')
				smoothed=Frame_golay(window = 5, order = 3,transpose=False)
				smoothed.plot(ax = ax1, style = '-', color = colors, legend = False, lw = linewidth)
		if not subplot:
			leg = ax1.legend(ds,title = 'lines = smoothed', loc='best', labelspacing = 0, ncol = 2, columnspacing = 1, handlelength = 1, frameon = False)
	elif 'data' in lines_are:
		if subplot:
			ax1 = ds.plot(ax = ax1, legend = False, style = '*', color = colors, zorder = 0)
		else:
			ax1 = ds.plot(ax = ax1, legend = False, style = '*', color = colors)
			leg = ax1.legend(ds,labelspacing = 0, ncol = 2, columnspacing = 1, handlelength = 1, loc = 'best', frameon = False)
	elif 'fitted' in lines_are:
		if scattercut is None:
			ds.plot(ax = ax1, style = '-', color = colors, legend = False, lw = linewidth, alpha = 0.7)
		elif isinstance(scattercut[0], numbers.Number):
			if equal_energy_bin is not None:
				scattercut=[scipy.constants.h*scipy.constants.c/(a*1e-9*scipy.constants.electron_volt) for a in scattercut]
			ds.loc[:scattercut[0],:].plot(ax = ax1, legend = False, style = '-', color = colors, alpha = 0.7, lw = linewidth)
			ds.loc[scattercut[1]:,:].plot(ax = ax1, legend = False, style = '-', color = colors, alpha = 0.7, lw = linewidth, label='_nolegend_')
		else:
			try:
				scattercut=flatten(scattercut)
				if equal_energy_bin is not None:
					scattercut=[scipy.constants.h*scipy.constants.c/(a*1e-9*scipy.constants.electron_volt) for a in scattercut]
				for i in range(len(scattercut)):
					if i==0:
						ds.loc[:scattercut[0],:].plot(ax = ax1, legend = False, style = '-', color = colors, alpha = 0.7, lw = linewidth)
					elif i<(len(scattercut)/2):
						ds.loc[scattercut[2*i-1]:scattercut[2*i],:].plot(ax = ax1, legend = False, style = '-', color = colors, alpha = 0.7, lw = linewidth, label='_nolegend_')
					else:
						ds.loc[scattercut[-1]:,:].plot(ax = ax1, legend = False, style = '-', color = colors, alpha = 0.7, lw = linewidth, label='_nolegend_')
			except:
				ds.plot(ax = ax1, legend = False, style = '-', color = colors, alpha = 0.7, lw = linewidth)
		if not subplot:leg = ax1.legend(ds,title = 'lines = fit', loc = 'best', labelspacing = 0, ncol = 2, columnspacing = 1, handlelength = 1, frameon = False)
	if not subplot:
		if text_in_legend is not None:
			stringen=leg.get_title().get_text()
			texten=text_in_legend
			leg.set_title(texten + '\n' +stringen)

	else:#for multiple plotting
		return ax1
	if bordercut is None:
		ax1.autoscale(axis='x',tight=True)
	else:
		if equal_energy_bin is not None:
			bordercut=[scipy.constants.h*scipy.constants.c/(a*1e-9*scipy.constants.electron_volt) for a in bordercut]
			#bordercut=bordercut[::-1]
		ax1.set_xlim(bordercut)	
	if (not subplot) and plot_second_as_energy:
		ax2=ax1.twiny()
		ax2.set_xlim(ax1.get_xlim())
		ax2.set_xticks(ax1.get_xticks())
		if equal_energy_bin is not None:
			labels=['%.1f'%(scipy.constants.h*scipy.constants.c/(a*1e-9*scipy.constants.electron_volt)) for a in ax2.get_xticks()]
		else:
			labels=['%.2f'%(scipy.constants.h*scipy.constants.c/(a*1e-9*scipy.constants.electron_volt)) for a in ax2.get_xticks()]
		_=ax2.set_xticklabels(labels)
		if equal_energy_bin is not None:
			_=ax2.set_xlabel('Wavelength in nm')
		else:
			_=ax2.set_xlabel('Energy in eV')
		ax1.set_zorder(ax2.get_zorder()+1)

	if not subplot:	
		if not len(title)==0:		
			try:
				ax2.set_title(title,pad=10)
			except:
				ax1.set_title(title,pad=10)
		ax1.set_ylabel(data_type)
		ax1.yaxis.set_major_formatter(FuncFormatter(lambda x, pos: '%.2g'%(x)))
		ax1.set_xlabel(ds.index.name)
		ax1.minorticks_on()
		#ax1.xaxis.set_minor_locator(AutoMinorLocator(6))
		ax1.plot(ax1.get_xlim(),[0,0],color='black',lw=0.5,zorder=0, label='_nolegend_')
		if intensity_range is None:
			ax1.autoscale(axis='y',tight=True)
		else:
			if not hasattr(intensity_range,'__iter__'):#lets have an lazy option
				intensity_range=np.array([-intensity_range,intensity_range])
			ax1.set_ylim(intensity_range)
	if ax is None:
		return fig
	else:
		return ax1


def plot1d(ds = None, wavelength = None, width = None, ax = None, subplot = False, title = None, intensity_range = None, 
			baseunit = 'ps', timelimits = None, scattercut = None, bordercut = None, cmap = standard_map, plot_type = 'symlog',
			lintresh = 0.1, text_in_legend = None, lines_are = 'smoothed', color_offset = 0, ignore_time_region = None,
			linewidth = 1, data_type = 'differential Absorption in $\mathregular{\Delta OD}$', units = 'nm', from_fit = False):																			  
	'''Plots the single line kinetic for specific wavelength given with the parameter wavelength.

	Parameters
	---------------
	
	ds : DataFrame
		This dataframe contains the data to be plotted. It is copied and sliced into the
		regions defined. The dataframe expects the time to be in Index and the wavelength/energy 
		to be in the columns. The spectra is plotted with a second (energy) axis
	
	wavelength : float or list (of floats)
		wavelength is in the object called "rel_wave" and works with "width" 
		(in the object called "wavelength_bin") together for the creation 
		of kinetic plots. When plotting kinetic spectra one line will be plotted for each entrance
		in the list/vector rel_wave. During object generation the vector np.arange(300,1000,100) 
		is set as standard. Another typical using style would be to define a list of interesting 
		wavelength at which a kinetic development is to be plotted. At each selected wavelength 
		the data between wavelength+ta.wavelength_bin and wavelength-ta.wavelength_bin is averaged 
		for each timepoint returned 
		
	data_type : str
		this is the datatype and effectively the unit put on the intensity axis 
		(Default)'differential Absorption in $\mathregular{\Delta OD}$
		
	width : float, optional
		the width used in kinetics, see below (Default) 10nm		
		
	ax : None, matplotlib axis object optional
		If None (Default) a new plot is is created and a new axis, otherwise ax needs to be Matplotlib Axis
	
	subplot : bool, optional
		If False (Default) axis labels and such are set. If True, we plot into the same axis and 
		do not set labels
	
	text_in_legend : None, str, optional
		extra text to be put into the legend (above the lines)
	
	lines_are : str, optional
		Depending on this parameter the plots contain:
		'smoothed' = data lines of golay smoothed data (Default) 
		'data' = dots are data, 
		'fitted' = not data, just lines shown
		
	title : None or str
		title to be used on top of each plot
		The (Default) None triggers  self.filename to be used. Setting a specific title as string will
		be used in all plots. To remove the title all together set an empty string
		
	linewidth : float, optional
		linewidht to be used for plotting
		
	intensity_range : None, float or list [of two floats]
		intensity_range is a general switch that governs what intensity range the plots show. 
		For the 1d plots this is the y-axis for the 2d-plots this is the colour scale. 
		This parameter recognizes three settings. If set to "None" (Default) this uses the minimum and
		maximum of the data. A single value like in the example below and the intended use is the symmetric
		scale while a list with two entries an assymmetric scale e.g. 
		intensity_range=3e-3 is converted into intensity_range=[-3e-3,3e-3]
		
	baseunit : str
		baseunit is a neat way to change the unit on the time axis of the plots. (Default) 'ps', but they 
		can be frames or something similarly. This is changing only the label of the axis. 
		During the import there is the option to divide the numbers by a factor. 
		I have also used frames or fs as units. Important is that all time units will be labeled with
		this unit.
		
	timelimits : None or list (of 2 floats), optional
		cut times at the low and high time limit. (Default) None uses the limits of measurement
		Important: If either the background or the chirp is to be fit this must include the 
		time before zero! Useful: It is useful to work on different regions, starting with
		the longest (then use the ta.Backgound function prior to fit) and expand from there
		
	scattercut : None or iterable (of floats or other iterable, always pairs!), optional
		intented to "cut" one or multiple scatter regions. (if (Default) None nothing
		happens) If it is set the spectral region between the limits is set to zero. 
		Usage single region: [lower region limit,upper region limit], 
		use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]
		
	bordercut : None or iterable (with two floats), optional
		cut spectra at the low and high wavelength limit. (Default) None 
		uses the limits of measurement 
		
	ignore_time_region : None or list (of two floats or of lists), optional
		cut set a time range with a low and high limit from the fits. (Default) None nothing happens
		The region will be removed during the fitting process (and will be missing in the fit-result
		plots) 
		Usage single region: [lower region limit,upper region limit],
		use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]
		
	cmap : None or matplotlib color map, optional
		is a powerfull variable that chooses the colour map applied for all plots. If set to 
		None (Default) then the self.cmap is used.
		As standard I use the color map "jet" from matplotlib. There are a variety of colormaps 
		available that are very usefull. Beside "jet", "viridis" is a good choice as it is well 
		visible under red-green blindness. Other useful maps are "prism" for high fluctuations 
		or diverging color maps like "seismic". 
		See https://matplotlib.org/3.1.0/tutorials/colors/colormaps.html for a comprehensive 
		selection. In the code the colormaps are imported so if plot_func is imported as pf then 
		self.cmap=pf.cm.viridis sets viridis as the map to use. Internally the colors are chosen 
		with the "colm" function. The 2d plots require a continuous color map so if something 
		else is give 2d plots are shown automatically with "jet". For all of the 1d plots however 
		I first select a number of colors before each plot. If cmap is a continous map then these
		are sampled evenly over the colourmap. Manual iterables of colours 
		cmap=[(1,0,0),(0,1,0),(0,0,1),...] are also accepted, as are vectors or dataframes that 
		contain as rows the colors. There must be of course sufficient colors present for 
		the numbers of lines that will be plotted. So I recommend to provide at least 10 colours 
		(e.g.your university colors). colours are always given as a, list or tuple with RGA or RGBA
		(with the last A beeing the Alpha=transparency. All numbers are between 0 and 1. 
		If a list/vector/DataFrame is given for the colours they will be used in the order provided.
		 
	color_offset : int, optional
		At the (Default) 0 the colours are chose from the beginning, for a larger value Color_offset 
		colors are skipped. Usually only used if multiple plots are created, and the data/or fit is
		only shown for some of them.
		
	plot_type : None or str
		is a general setting that can influences what time axis will be used for the plots. 
		"symlog" (linear around zero and logarithmic otherwise) "lin" and "log" are valid options.
		
	lintresh : float
		The pseudo logratihmic range "symlog" is used for most time axis. Symlog plots a range around
		time zero linear and beyond this linear treshold 'lintresh' on a logarithmic scale. (Default) 0.3 
	
	from_fit : bool optional
		i needed this swtich to avoid re-slicing of data in spectal axis for equal energy bins

	'''
	global halfsize
	if not isinstance(ds,pandas.DataFrame):
		print("input format wrong")
	if ax is None:
		if halfsize:
			fig,ax1=plt.subplots(figsize=(5,3),dpi=100)
		else:
			fig,ax1=plt.subplots(figsize=(10,6),dpi=100)
	else:
		ax1=ax
	if width is None:width=1
	if not hasattr(wavelength, '__iter__'):wavelength = [wavelength]
	if isinstance(cmap,list):
		colors=cmap[color_offset:]
	else:
		colors = colm(np.arange(color_offset, len(wavelength)+color_offset), cmap = cmap)
	ds = sub_ds(ds = ds, wavelength = wavelength, wavelength_bin = width, scattercut = scattercut, 
				ignore_time_region = ignore_time_region, drop_ignore = True, from_fit = from_fit)
	if 'smoothed' in lines_are:
		if ignore_time_region is None:
			smoothed=Frame_golay(ds, window = 5, order = 3)
			smoothed.plot(ax = ax1, style = '-', color = colors, legend = False, lw = linewidth)
		elif isinstance(ignore_time_region[0], numbers.Number):
			smoothed=Frame_golay(ds.loc[:ignore_time_region[0],:], window = 5, order = 3)
			smoothed.plot(ax = ax1, style = '-', color = colors, legend = False, lw = linewidth)
			smoothed=Frame_golay(ds.loc[ignore_time_region[1]:,:], window = 5, order = 3)
			smoothed.plot(ax = ax1, style = '-', color = colors, legend = False, lw = linewidth, label='_nolegend_')
		else:
			try:
				ignore_time_region=flatten(ignore_time_region)
				for i in range(len(ignore_time_region)/2+1):
					if i==0:
						smoothed=Frame_golay(ds.loc[:ignore_time_region[0],:], window = 5, order = 3,transpose=True)
						smoothed.plot(ax = ax1, style = '-', color = colors, legend = False, lw = linewidth)
					elif i<(len(ignore_time_region)/2):
						smoothed=Frame_golay(ds.loc[ignore_time_region[2*i-1]:ignore_time_region[2*i],:], window = 5, order = 3,transpose=True)
						smoothed.plot(ax = ax1, style = '-', color = colors, legend = False, lw = linewidth, label='_nolegend_')
					else:
						smoothed=Frame_golay(ds.loc[ignore_time_region[-1]:,:], window = 5, order = 3,transpose=True)
						smoothed.plot(ax = ax1, style = '-', color = colors, legend = False, lw = linewidth, label='_nolegend_')
			except:
				smoothed=Frame_golay(ds, window = 5, order = 3)
				smoothed.plot(ax = ax1, style = '-', color = colors, legend = False, lw = linewidth)
		
	elif 'data' in lines_are:
		if subplot:ds.plot(ax = ax1, style = '*', color = colors, legend = False, zorder = 0, label='_nolegend_')
		else:	   ds.plot(ax = ax1, style = '*', color = colors, legend = False)
	elif 'fitted' in lines_are:
		if ignore_time_region is None:
			ds.plot(ax = ax1, style='-', color = colors, legend = False, lw = linewidth)
		elif isinstance(ignore_time_region[0], numbers.Number):
			ds.loc[:ignore_time_region[0],:].plot(ax = ax1, style='-', color = colors, legend = False, lw = linewidth)
			ds.loc[ignore_time_region[1]:,:].plot(ax = ax1, style='-', color = colors, legend = False, lw = linewidth, label='_nolegend_')
		else:
			try:
				ignore_time_region=flatten(ignore_time_region)
				for i in range(len(ignore_time_region)/2+1):
					if i==0:
						ds.loc[:ignore_time_region[0],:].plot(ax = ax1, style='-', color = colors, legend = False, lw = linewidth)
					elif i<(len(ignore_time_region)/2):
						ds.loc[ignore_time_region[2*i-1]:ignore_time_region[2*i],:].plot(ax = ax1, style='-', color = colors, legend = False, lw = linewidth, label='_nolegend_')
					else:
						ds.loc[ignore_time_region[-1]:,:].plot(ax = ax1, style='-', color = colors, legend = False, lw = linewidth, label='_nolegend_')
			except:
				ds.plot(ax = ax1, style='-', color = colors, legend = False, lw = linewidth)
	#Legend
	if not subplot:
		handles, labels = ax1.get_legend_handles_labels()
		handles=handles[:len(wavelength)]
		labels=labels[:len(wavelength)]
		for i,entry in enumerate(labels):
			labels[i]=entry + ' %s'%units
		if 'smoothed' in lines_are:leg=ax1.legend(labels,title='%g %s width, lines=smoothed'%(width,units),labelspacing=0,ncol=2,columnspacing=1,handlelength=1,frameon=False)
		elif 'data' in lines_are:  leg=ax1.legend(labels,title='%g %s width'%(width,units)				  ,labelspacing=0,ncol=2,columnspacing=1,handlelength=1,frameon=False)
		elif 'fitted' in lines_are:leg=ax1.legend(labels,title='%g %s width, lines=fit'%(width,units)	  ,labelspacing=0,ncol=2,columnspacing=1,handlelength=1,frameon=False)
		if text_in_legend is not None:
			stringen=leg.get_title().get_text()
			texten=text_in_legend
			leg.set_title(texten + '\n' +stringen)
		
	x=ds.index.values.astype('float')
	#limits and ticks
	if timelimits is None:timelimits=[min(x),max(x)]
	if "symlog" in plot_type:
		lintresh=lintresh
		ax1.set_xscale('symlog', linthresh=lintresh,subs=range(2,9),linscale=lintresh/4.)
		try:
			if lintresh>0.5:
				ticks=np.concatenate((np.arange(-100,0,10,),[-5,-3,-2,-1,-0.5,0,0.5],np.logspace(0,10,11)))
			elif lintresh>=0.3:
				ticks=np.concatenate((np.arange(-100,0,10,),[-5,-3,-2,-1,-0.3,0,0.3],np.logspace(0,10,11)))
			elif lintresh>=0.1:
				ticks=np.concatenate((np.arange(-100,0,10,),[-5,-3,-2,-1,-0.1,0,0.1],np.logspace(0,10,11)))
			else:
				ticks=np.concatenate((np.arange(-100,0,10,),[-5,-3,-2,-1,0],np.logspace(0,10,11)))
			ticks=ticks[ticks>timelimits[0]]
			ticks=ticks[ticks<timelimits[1]]
			ax1.set_xticks(ticks)
		except:
			pass
		ax1.set_xlim(timelimits[0],timelimits[1])
	elif "log" in plot_type:
		lower_time=max(1e-6,timelimits[0])
		ax1.set_xlim(lower_time,timelimits[1])
		ax1.set_xscale('log')
	elif "lin" in plot_type:
		ax1.set_xlim(timelimits[0],timelimits[1])
	if intensity_range is None:
		ax1.autoscale(axis='y',tight=True)
	else:
		if not hasattr(intensity_range,'__iter__'):#lets have an lazy option
			intensity_range=np.array([-intensity_range,intensity_range])
		ax1.set_ylim(intensity_range)			
	if not subplot:
		ax1.plot(ax1.get_xlim(),[0,0],'black',lw=1,zorder=10, label='_nolegend_')
		if title is not None:
			if title:		
				try:
					ax1.set_title(title,pad=10)
				except:
					print('title could not be set')
		if "symlog" in plot_type:
			ax1.plot([lintresh,lintresh],ax1.get_ylim(),color='black',linestyle='dashed',alpha=0.5)
			ax1.plot([-lintresh,-lintresh],ax1.get_ylim(),color='black',linestyle='dashed',alpha=0.5)
				
	ax1.set_xlabel(ds.index.name)
	ax1.set_ylabel(data_type)		
	#ax1.set_xlabel('time in %s'%baseunit)
	#ax1.set_ylabel('differential Absorption in $\mathregular{\Delta OD}$')
	ax1.yaxis.set_major_formatter(FuncFormatter(lambda x, pos: '%.2g'%(x)))
	if ax is None:
		return fig
	else:
		return ax1


def SVD(ds, times = None, scattercut = None, bordercut = None, timelimits = [5e-1, 150], wave_nm_bin = 10, 
		time_bin = None, wavelength_bin = None, plotting = True, baseunit = 'ps', title = None, ignore_time_region = None, 
		cmap = None, equal_energy_bin = None, data_type = 'differential Absorption in $\mathregular{\Delta OD}$'):
	'''This function calculates the SVD and plots an overview.
	
	Parameters
	------------
	
	ds : DataFrame
		This dataframe contains the data to be plotted. It is copied and sliced into the
		regions defined. The dataframe expects the time to be in Index and the wavelength/energy 
		to be in the columns. The spectra is plotted with a second (energy) axis
	
	times : None or int
		are the number of components to be used in the SVD (Default) is None (which is seen as 6)
	
	plotting : bool
		if True (Default) the functions plots the SVD, if False it returns the vectors
	
	title : None or str
		title to be used on top of each plot
		The (Default) None triggers  self.filename to be used. Setting a specific title as string will
		be used in all plots. To remove the title all together set an empty string with this command title="" 
		
	baseunit : str
		baseunit is a neat way to change the unit on the time axis of the plots. (Default) 'ps', but they 
		can be frames or something similarly. This is changing only the label of the axis. 
		During the import there is the option to divide the numbers by a factor. 
		I have also used frames or fs as units. Important is that all time units will be labeled with
		this unit.
		
	timelimits : None or list (of 2 floats), optional
		cut times at the low and high time limit. (Default) [5e-1 , 150] uses the limits of measurement
		Important: If either the background or the chirp is to be fit this must include the 
		time before zero! Useful: It is useful to work on different regions, starting with
		the longest (then use the ta.Backgound function prior to fit) and expand from there
		
	scattercut : None or iterable (of floats or other iterable, always pairs!), optional
		intented to "cut" one or multiple scatter regions. (if (Default) None nothing
		happens) If it is set the spectral region between the limits is set to zero. 
		Usage single region: [lower region limit,upper region limit], 
		use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]
		
	bordercut : None or iterable (with two floats), optional
		cut spectra at the low and high wavelength limit. (Default) None 
		uses the limits of measurement 
		
	wave_nm_bin : None or float, optional
		rebins the original data into even intervals. If set to None the original data will be used. 
		If set to a width (e.g. 2nm), the wavelength axis will be divided into steps of this size
		and the mean of all measurements in the interval is taken. The re-binning stops as soon as
		the measured stepsize is wider than given here, then the original bins are used. 
		This function is particularly useful for spectrometer with non-linear dispersion, 
		like a prism in the infrared. (Default = 10)
	
	wavelength_bin : float, optional
		the width used in kinetics, see below (Default) 10nm
		
	ignore_time_region : None or list (of two floats or of lists), optional
		cut set a time range with a low and high limit from the fits. (Default) None nothing happens
		The region will be removed during the fitting process (and will be missing in the fit-result
		plots) 
		Usage single region: [lower region limit,upper region limit],
		use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]
	
	time_bin : None or int, optional
		is dividing the points on the time-axis in even bins and averages the found values in between. 
		This is a hard approach that also affects the fits. I do recommend to use this carefully, 
		it is most useful for modulated data. A better choice for transient absorption that only 
		affects the kinetics is 'time_width_percent'	
		
	cmap : None or matplotlib color map, optional
		is a powerfull variable that chooses the colour map applied for all plots. If set to 
		None (Default) then the self.cmap is used.
		As standard I use the color map "jet" from matplotlib. There are a variety of colormaps 
		available that are very usefull. Beside "jet", "viridis" is a good choice as it is well 
		visible under red-green blindness. Other useful maps are "prism" for high fluctuations 
		or diverging color maps like "seismic". 
		See https://matplotlib.org/3.1.0/tutorials/colors/colormaps.html for a comprehensive 
		selection. In the code the colormaps are imported so if plot_func is imported as pf then 
		self.cmap=pf.cm.viridis sets viridis as the map to use. Internally the colors are chosen 
		with the "colm" function. The 2d plots require a continuous color map so if something 
		else is give 2d plots are shown automatically with "jet". For all of the 1d plots however 
		I first select a number of colors before each plot. If cmap is a continous map then these
		are sampled evenly over the colourmap. Manual iterables of colours 
		cmap=[(1,0,0),(0,1,0),(0,0,1),...] are also accepted, as are vectors or dataframes that 
		contain as rows the colors. There must be of course sufficient colors present for 
		the numbers of lines that will be plotted. So I recommend to provide at least 10 colours 
		(e.g.~your university colors). colours are always given as a, list or tuple with RGA or RGBA
		(with the last A beeing the Alpha=transparency. All numbers are between 0 and 1. 
		If a list/vector/DataFrame is given for the colours they will be used in the order provided.
	
	'''
	global halfsize
	if times is None:
		max_order=6
	else:
		max_order=times
	if cmap is None:cmap=standard_map
	colors=colm(np.arange(0,max_order,1),cmap=cmap)
	ds = sub_ds(ds, scattercut = scattercut, bordercut = bordercut, timelimits = timelimits, wave_nm_bin = wave_nm_bin, 
				wavelength_bin = wavelength_bin, time_bin = time_bin, ignore_time_region = ignore_time_region)
	U, s, V = np.linalg.svd(ds.values)
	if plotting:
		if halfsize:
			fig=plt.figure(figsize=(4,4),dpi=100)
		else:
			fig=plt.figure(figsize=(8,8),dpi=100)
		G = GridSpec(2, 6)
		ax1=fig.add_subplot(G[0,:2])
		ax2=fig.add_subplot(G[1,:])
		ax3=fig.add_subplot(G[0,2:])
		
		if title is not None:
			ax2.set_title(title,fontsize=plt.rcParams['figure.titlesize']-4)
		else:
			ax1.set_title("Component\nstrength",fontsize=plt.rcParams['figure.titlesize']-4)
			ax2.set_title("Spectral component",fontsize=plt.rcParams['figure.titlesize']-4)
			ax3.set_title("Temporal development\nof spectral component",fontsize=plt.rcParams['figure.titlesize']-4)
		s/=s.max()
		ax1.scatter(np.arange(max_order)+1,s[:max_order],c=colors,s=100)
		ax1.set_xlabel('SVD order',fontsize=plt.rcParams['axes.labelsize']-2)
		ax1.set_ylabel('Singular values norm.',fontsize=plt.rcParams['axes.labelsize']-2)
		ax1.set_xlim(0.5,max_order+0.5)
		if max_order == 6:
			ax1.set_xticks([round(a) for a in np.linspace(1,max_order,6)])
		else:
			ax1.set_xticks([round(a) for a in np.linspace(1,max_order,5)])
		V2=pandas.DataFrame(V.T,index=ds.columns.values.astype('float'))
		U2=pandas.DataFrame(U,index=ds.index.values.astype('float'))
		U2=U2.iloc[:,:len(s)].multiply(-s)
		V2=V2.iloc[:,:len(s)].multiply(-s)
		names=['SVD vector %i'%(a+1) for a in range(max_order)]
		U2=U2.iloc[:,:max_order]
		U2.columns=names
		V2=V2.iloc[:,:max_order]	
		V2.columns=names
		V2/=V2.abs().max(axis=1).max()
		V2.plot(ax=ax2,color=colors)
		ax2.set_ylabel('Intensity norm.',fontsize=plt.rcParams['axes.labelsize']-2)
		U2/=U2.abs().max(axis=1).max()
		U2.plot(ax=ax3,color=colors)
		ax3.set_ylabel('Intensity norm.',fontsize=plt.rcParams['axes.labelsize']-2)
		ax2.set_xlabel(ds.columns.name,fontsize=plt.rcParams['axes.labelsize']-2)
		lims=V2.index.values.astype(float)
		ax2.set_xlim(lims.min(),lims.max())
		#ax2.set_xticks(np.linspace(round(lims.min(),-2),round(lims.max()-2),5))
		ax2.xaxis.set_major_formatter(FuncFormatter(lambda x, pos: '%.4g'%(x)))

		ax3.set_xlabel(ds.index.name,fontsize=plt.rcParams['axes.labelsize']-2)
		tims=U2.index.values.astype(float)
		ax3.set_xlim(max([0.01,tims.min()]),tims.max())
		ax3.set_xscale('log')
		#ax3.set_xticks(np.logspace(-1,2,round(np.log())))
		ax3.xaxis.set_major_formatter(FuncFormatter(lambda x, pos: '%1g'%(x)))
		ax2.legend(frameon=False,labelspacing=0,borderpad=0,numpoints=1,handlelength=1)
		ax3.legend(frameon=False,fontsize=plt.rcParams['legend.fontsize'],labelspacing=0,borderpad=0,numpoints=1,handlelength=1)
		
		fig.tight_layout()
		return fig
	else:
		return U, s, V2,ds


def Species_Spectra(ta=None,conc=None,das=None):
	'''useful help function that returns a dictionary that has DataFrame as entries and the names of the 
	components as keys
	
	Parameters
	-----------
	
	ta : plot_func.TA object, optional 
		This object should contain a successful fit. The function will cycle through the fitted species 
		and return the matrix that is formed from the dynamics and the species associated spectrum
		If this given, then "conc" and "das" are ignored. We cycle through the columns of the concentration 
		and take the same column from the das Frame. 
	
	conc : DataFrame, optional
		Is read only if ta_object is None. This should contain the concentration matrix with the species as
		as columns
		
	das : DataFrame, optional
		This should contain the spectra of the species with one column per spectrum. The position of the columns 
		must match the columns in the conc (at least this is what is assumed)
		
	Examples
	---------
	dicten=Species_Spectra(ta)
	
	'''
	if ta is not None:
		try:
			time=ta.re['c'].index.values
			WL=ta.re['DAC'].index.values
			conc=ta.re['c']
			das=ta.re['DAC']
		except:
			print('the TA object must contain a successful fit')
			print(ta.re)
			return False
	else:
		if (conc is None) or (das is None):
			print('If the ta object is None, then we need both the conc and the das')
			return False			
		else:
			time=conc.index.values
			WL=das.index.values
	results={}
	for i in range(len(conc.columns)):
		A,B=np.meshgrid(conc.iloc[:,i].values,das.iloc[:,i].values)
		C=pandas.DataFrame((A*B).T,index=time,columns=WL)
		results[conc.columns[i]]=C
	try:
		for key in results.keys():
			results[key].index.name=ta.re['c'].index.name
			results[key].columns.name=ta.re['DAC'].index.name
	except Exception as e:
		print(e) 
	return results


def Fix_Chirp(ds, save_file = None, scattercut = None, intensity_range = 5e-3, wave_nm_bin = 10, bordercut=None,
			  shown_window = [-1.5, 1.5], filename = None, path = None, fitcoeff = None, max_points = 40, 
			  cmap = cm.prism, just_shift=False):										
	'''Manual selecting polynom for chirp. 	This function is opening 
		a plot and allows the user to select a number of points, which are then 
		approximated with a 4th order polynomial and finally to select a point 
		that is declared as time zero. The observed window as well as the intensities 
		and the colour map can be chosen to enable a good correction. Here a fast 
		iterating colour scheme such as "prism" is often a good choice. In all of the 
		selections a left click selects, a right click removes the last point and 
		a middle click (sometime appreviated by clicking left and right together) 
		finishes the selection. If no middle click exists, the process
		automatically ends after max_points (40 preset).
		Many of the parameters are from the raw plotting part
		
	Parameters
	-----------
	
	ds : DataFrame
		This dataframe contains the data to be plotted. It is copied and sliced into the
		regions defined. The dataframe expects the time to be in Index and the wavelength/energy 
		to be in the columns. The spectra is plotted with a second (energy) axis
	
	save_file : None or str, optional
		   If a raw file was read(e.g. "data.SIA") and the chirp correction was
		   completed, a file with the attached word "chirp" is created and
		   stored in the same location. ("data_chirp.dat") This file contains
		   the 5 values of the chirp correction. By selecting such a file
		   (e.g. from another raw data) a specific chirp is applied. If a
		   specific name is given with **chirp_file** (and optional **path**)
		   then this file is used.\n
		   GUI\n
		   The word *'gui'* can be used instead of a filename to open a gui that
		   allows the selection of a chrip file
	
	scattercut : None or iterable (of floats or other iterable, always pairs!), optional
		intented to "cut" one or multiple scatter regions. (if (Default) None nothing
		happens) If it is set the spectral region between the limits is set to zero. 
		Usage single region: [lower region limit,upper region limit], 
		use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]
	
	intensity_range : None, float or list [of two floats]
		intensity_range is a general switch that governs what intensity range the plots show. 
		For the 1d plots this is the y-axis for the 2d-plots this is the colour scale. 
		This parameter recognizes three settings. If set to "None" (Default) this uses the minimum and
		maximum of the data. A single value like in the example below and the intended use is the symmetric
		scale while a list with two entries an assymmetric scale e.g. 
		intensity_range=3e-3 is converted into intensity_range=[-3e-3,3e-3]	
	
	wave_nm_bin : None or float, optional
		rebins the original data into even intervals. If set to None the original data will be used. 
		If set to a width (e.g. 2nm), the wavelength axis will be divided into steps of this size
		and the mean of all measurements in the interval is taken. The re-binning stops as soon as
		the measured stepsize is wider than given here, then the original bins are used. 
		This function is particularly useful for spectrometer with non-linear dispersion, 
		like a prism in the infrared.	
		
	shown_window : list (with two floats), optional
		Defines the window that is shown during chirp correction. If the t=0 is not visible, adjust this parameter
		to suit the experiment. If problems arise, I recomment to use Plot_Raw to check where t=0 is located
		
	filename : str, optional
		name of the original file, that will be used to save the results later (with attached "_chirp")	
		
	path : str or path object (optional)
		if path is a string without the operation system dependent separator, it is treated as a relative path, 
		e.g. data will look from the working directory in the sub director data. Otherwise this has to be a 
		full path in either strong or path object form.
			
	fitcoeff : list or vector (5 floats), optional
		One can give a vector/list with 5 numbers representing the parameter
		of a 4th order polynomial (in the order
		:math:`(a4*x^4 + a3*x^3+a2*x^2+a1*x1+a0)`. The chirp parameter are
		stored in ta.fitcoeff and can thus be used in other TA objects. This
		vector is also stored with the file and automatically applied during
		re-loading of a hdf5-object

	max_points : int, optional
		Default = 40 max numbers of points to use in Gui selection. Useful option in case no middle mouse button
		is available. (e.g. touchpad)
		
	cmap : matplotlib colourmap, optional
		Colourmap to be used for the chirp correction. While there is a large selection here I recommend to choose
		a different map than is used for the normal 2d plotting.\n
		cm.prism (Default) has proofen to be very usefull
	
	just_shift: bool, optional
		This will turn of the polynomial part of the fitting and lets you just select the new time zero
	
		
		''' 
	ds=ds.fillna(0)
	if fitcoeff is not None:#we loaded a previous project this is a dublication, but I'm currently to lazy to make this tighter
		if isinstance(fitcoeff,str):fitcoeff=np.array(fitcoeff.split(','),dtype='float')
		if len(fitcoeff)==6:#old style parameter
			fitcoeff[-2]-=fitcoeff[-1]
			fitcoeff=fitcoeff[:5]
		wl=ds.columns.values.astype('float')#extract the wavelength
		time=ds.index.values.astype('float')#extract the time
		for i in range(0, len(wl), 1):
			correcttimeval = np.polyval(fitcoeff, wl[i])
			f = scipy.interpolate.interp1d((time-correcttimeval), ds.values[:,i], bounds_error=False, fill_value=0)
			fixed_wave = f(time)
			ds.values[:, i] = fixed_wave
		return ds
	else:
		if save_file is None:
			#lets start by choosing a good intensity
			if hasattr(intensity_range,'__iter__'):
				maxim=max(abs(np.array(intensity_range)))
				intensity_range=maxim
			elif intensity_range is None:
				intensity_range=ds.abs().max().max()
			intensities=(2**np.arange(-6.,4.,1.))*intensity_range
			window_difference=np.abs(shown_window[1]-shown_window[0])*0.1
			cutoff_window=np.abs(shown_window[1]-shown_window[0])
			timelimits=shown_window+np.asarray([-window_difference,window_difference])
			for repeat in range(30):
				if halfsize:
					fig,ax=plt.subplots(figsize=(6,6))
				else:
					fig,ax=plt.subplots(figsize=(12,12))
				ax = plot2d(ax = ax, cmap = cmap, ds = ds, wave_nm_bin = wave_nm_bin, scattercut = scattercut, bordercut = bordercut,
							timelimits = timelimits, intensity_range = intensity_range, 
							title = 'select intensity where we can work,\n  if happy,  choose confirm or abort', 
							use_colorbar = False, plot_type = "linear", log_scale = False)
				w=(ax.get_xlim()[1]-ax.get_xlim()[0])
				ax.add_patch(matplotlib.patches.Rectangle((ax.get_xlim()[0],shown_window[1]-0.1),w,0.1,facecolor='white'))
				for i,ent in enumerate(intensities):
					if halfsize:
						ax.text(ax.get_xlim()[0]+i*w/10.,shown_window[1]-0.1,'%.1g'%(2**np.arange(-6.,4,1.))[i],fontsize=10)
					else:
						ax.text(ax.get_xlim()[0]+i*w/10.,shown_window[1]-0.1,'%.1g'%(2**np.arange(-6.,4,1.))[i],fontsize=20)
				ax.add_patch( matplotlib.patches.Rectangle((ax.get_xlim()[0],shown_window[0]),w/4.,0.1,facecolor='white'))
				ax.text(ax.get_xlim()[0],shown_window[0],'Accept',fontsize=20)
				ax.add_patch(matplotlib.patches.Rectangle((ax.get_xlim()[0]+w*3./4.,shown_window[0]),w/4.,0.1,facecolor='white'))
				ax.text(ax.get_xlim()[0]+w*3./4.,shown_window[0],'Cancel all',fontsize=20)
				ax.set_ylim(shown_window)
				fig.tight_layout()
				choice =plt.ginput(1,timeout=15)
				factor=int((choice[0][0]-ax.get_xlim()[0])/(w/10.))
				
				if choice[0][1] > (shown_window[0]+window_difference):
					intensity_range=intensities[factor]
					print((2**np.arange(-6.,4,1.))[factor])
					intensity_range=[-intensity_range,intensity_range]
					plt.close(fig)
					continue		
				elif choice[0][1] < (shown_window[0]+window_difference):#we choice to finish the choices
					if choice[0][0] < ax.get_xlim()[0]+w/2:#
						print('accept')
						plt.close(fig)
						break
					else:
						plt.close(fig)
						return False
				else:
					print('click better please')
					plt.close(fig)
					continue
			if not just_shift:
				for repeat in range(10):
					if halfsize:
						fig,ax=plt.subplots(figsize=(6,6))
					else:
						fig,ax=plt.subplots(figsize=(12,12))
					ax = plot2d(ax = ax, cmap = cmap, ds = ds, wave_nm_bin = wave_nm_bin, scattercut = scattercut, bordercut = bordercut, 
								timelimits = shown_window, intensity_range = intensity_range, 
								title = 'select points,  rightclick  =  remove last,  \n middle click (or both at once finishes ', 
								use_colorbar = False, plot_type = "linear", log_scale = False)
					fig.tight_layout()
					polypts=np.asarray(plt.ginput(n=max_points,timeout=300, show_clicks=True,mouse_add=1, mouse_pop=3, mouse_stop=2))
					plt.close(fig)
					if halfsize:
						fig,ax=plt.subplots(figsize=(6,6))
					else:
						fig,ax=plt.subplots(figsize=(12,12))					
					ax = plot2d(ax = ax, ds = ds, cmap = cmap, wave_nm_bin = wave_nm_bin, scattercut = scattercut, bordercut = bordercut, 
								timelimits = shown_window, intensity_range = intensity_range, 
								title = 'like it? %i more attempts'%(9-repeat), use_colorbar = False, 
								plot_type = "linear", log_scale = False)
					#Fit a polynomial of the form p(x) = p[2] + p[1] + p[0]
					fitcoeff= np.polyfit(polypts[:, 0], polypts[:, 1], 4, full=False)
					
					correcttimeval = np.polyval(fitcoeff, ds.columns.values.astype('float'))
					ax.plot(ds.columns.values.astype('float'),correcttimeval)
					 
					ax.add_patch( matplotlib.patches.Rectangle((ax.get_xlim()[0],ax.get_ylim()[0]),w/4,0.2,facecolor='white'))
					ax.text(ax.get_xlim()[0],ax.get_ylim()[0]+0.05,'Save',fontsize=20)
					ax.add_patch(matplotlib.patches.Rectangle((ax.get_xlim()[0]+w*3/4,ax.get_ylim()[0]),w/4,0.2,facecolor='white'))
									   
					ax.text(ax.get_xlim()[0]+w*3/4,ax.get_ylim()[0]+0.05,'Redo',fontsize=20)
					fig.tight_layout()
					satisfied =plt.ginput(1)
					plt.close(fig)
					if satisfied[0][0] < ax.get_xlim()[0]+w/2:
						print('accepted')
						plt.close(fig)
						break
					elif repeat<8:
						plt.close(fig)
						continue
					else:
						plt.close(fig)
						return False
			else:
				fitcoeff=np.array([0,0,0,0,0])
		else:
			with open(save_file,'r') as f:
				fitcoeff=f.readline()
			fitcoeff=np.array(fitcoeff.split(','),dtype='float')
			if len(fitcoeff)==6:#old style params
				fitcoeff[-2]-=fitcoeff[-1]
				fitcoeff=fitcoeff[:5]
		time=ds.index.values.astype('float')#extract the time
		print(fitcoeff)
		ds_new=ds.apply(lambda x:np.interp(x=time+np.polyval(fitcoeff,float(x.name)),xp=time,fp=x),axis=0,raw=False)
		if save_file is None:
			#finding where zero time is
			for repeat in range(10):
				if halfsize:
					fig,ax=plt.subplots(figsize=(6,6))
				else:
					fig,ax=plt.subplots(figsize=(12,12))				
				if just_shift:
					ax = plot2d(ax = ax, cmap = cmap, ds = ds_new, wave_nm_bin = wave_nm_bin, scattercut = scattercut, bordercut = bordercut, 
							lintresh = np.max(timelimits), timelimits =timelimits, intensity_range = intensity_range, 
							title = 'select new time zero', plot_type = 'lin', use_colorbar = False, log_scale = False)
				else:
					ax = plot2d(ax = ax, cmap = cmap, ds = ds_new, wave_nm_bin = wave_nm_bin, scattercut = scattercut, bordercut = bordercut, 
							lintresh = np.max(timelimits), timelimits =np.array(timelimits)/5, intensity_range = intensity_range, 
							title = 'corrected select new zero', plot_type = 'lin', use_colorbar = False, log_scale = False)
				ax.plot(ax.get_xlim(),[0,0],'black',lw=0.5)
				#ax.set_ylim(-cutoff_window,cutoff_window)
				fig.tight_layout()
				fittingto = np.array(plt.ginput(1)[0])[1]
				print(fittingto)
				fitcoeff[-1]+=fittingto
				ds_new=ds.apply(lambda x:np.interp(x=time+np.polyval(fitcoeff,float(x.name)),xp=time,fp=x),axis=0,raw=False)
				plt.close(fig)
				if halfsize:
					fig,ax=plt.subplots(figsize=(6,6))
				else:
					fig,ax=plt.subplots(figsize=(12,12))
				if just_shift:
					ax = plot2d(ax = ax, ds = ds_new, cmap = cmap, wave_nm_bin = wave_nm_bin, scattercut = scattercut, bordercut = bordercut, 
							lintresh = np.max(timelimits), timelimits = np.array(timelimits)-fitcoeff[-1], intensity_range = intensity_range, 
							title = 'corrected,  please select', plot_type = 'lin', use_colorbar = False, log_scale = False)
				else:
					ax = plot2d(ax = ax, ds = ds_new, cmap = cmap, wave_nm_bin = wave_nm_bin, scattercut = scattercut, bordercut = bordercut, 
							lintresh = np.max(timelimits), timelimits = np.array(timelimits)/5, intensity_range = intensity_range, 
							title = 'corrected,  please select', plot_type = 'lin', use_colorbar = False, log_scale = False)
				ax.plot(ax.get_xlim(),[0,0],'black',lw=0.5)
				w=(ax.get_xlim()[1]-ax.get_xlim()[0])
				 
				ax.add_patch( matplotlib.patches.Rectangle((ax.get_xlim()[0],ax.get_ylim()[0]),w/4,0.2,facecolor='white'))
				ax.text(ax.get_xlim()[0],ax.get_ylim()[0]+0.05,'Save',fontsize=30)
				ax.add_patch(matplotlib.patches.Rectangle((ax.get_xlim()[0]+w*3/4,ax.get_ylim()[0]),w/4,0.2,facecolor='white'))
				ax.text(ax.get_xlim()[0]+w*3/4,ax.get_ylim()[0]+0.05,'Redo',fontsize=30)
				satisfied =plt.ginput(1)
				if satisfied[0][0] < ax.get_xlim()[0]+w/2:
					print('accepted')
					plt.close(fig)
					break
				elif repeat<8:
					plt.close(fig)
					continue
				else:
					plt.close(fig)
					return False
			print(fitcoeff)
			if filename is None:
				f='chirp.dat'
			else:
				f=filename.split('.')[0]
				f=f+'_chirp' + '.dat'
			if path is None:
				with open(f, 'w') as opened_file:
					opened_file.write(','.join(map(str,np.array(fitcoeff))))
			else:
				with open(check_folder(path=path,filename=f), 'w') as opened_file:
					opened_file.write(','.join(map(str,np.array(fitcoeff))))	
		return ds_new


def build_c(times, mod = 'paral', pardf = None, sub_steps = None):
	'''
	Build concentration matrix after model the parameters are:
	resolution is the width of the rise time (at sigma 50% intensity) 
	This function can also be used to create illustration dynamics.
	The parallel decays are created explicit, while the consecutive decays are 
	created by sampling the populations at the times given in the first vector and
	evaluate the progression at a number of substeps defined bu sub_samples (10 by default) 
	
	Parameters
	-----------
	
	times : np.array
		array with the times at which the dataframe should be generated. In general the 
		experimental times
	mod : str, optional
		this selects the model that is used to generate the concentrations.
		
			1. 'paral' (Default) or 'exponential' both are equivalent
			2. 'consecutive' or 'full_consecutive'
		
		In 2 the 'consecutive' and 'full_consecutive' are different in that for consecutive 
		the optimization is done using 'exponential' (as it shoudl give the same times)
		and then only in the last (final) iteration the 'full consecutive' differential
		equation is used. This has significant speed advantages, but can lead to errors particularly
		for the very fast times.
	sub_step : int, optional
		defines how many times the iterative loop (used in consecutive only) is sampling the concentrations 
		between the times given in "times"
	pardf : pd.DataFrame
		This dataframe must contain the parameter that are used for creating the dynamics
		the parameter must be named with the index.
		For the internal functions this must contain these keys:
		
			* 't0' = zero time, mandatory
			* 'resolution' = instrument response function, mandatory
			* 'background',optional = if this keyword is present a flat constant background is created (=1 over the whole time)
			* 'infinite',optional = if this keyword is present a new non decaying component is formed with the last decay time.
			* 'explicit_GS',optional = if this keyword is present the pulse function (= ground state) is added explicitly to the data
			* 'k0,k1,...' = with increasing integers are taken as decay times. te number of these components is used to determine how many shall be generated.
	sub_sample: bool or integer
		Default(None) does nothing
		This switch turns on a additional sampling of the kinetics, meaning that we add the number of steps between each measured steps for the model formation 
		usage: sub_sample=10
	
	Examples
	---------

	'''
	
	if 'sub_steps' in list(pardf.index.values):
		sub_steps=pardf['sub_steps']
	elif sub_steps is None:
		sub_steps=10 
	
	choices = {'paral':0,'parallel':0,'decays':0,'exponential':0,'consecutive':1,'sequential':1,'full_consecutive':1,'full_sequential':1}
	model=choices[mod]
	param=pardf.loc[pardf.is_rate,'value'].values.astype(float)
	t0=float(pardf.loc['t0','value'])
	resolution=float(pardf.loc['resolution','value'])
	if model == 0:#parallel
		c=np.exp(-1*np.tile(times-t0,(len(param),1)).T*param)
		c[(times-t0)<0]=1
		c*=np.tile(rise(x=times,sigma=resolution,begin=t0),(len(param),1)).T
		c=pandas.DataFrame(c,index=times)
		c.index.name='time'
		if 'background' in list(pardf.index.values):
			c['background']=1
		if 'infinite' in list(pardf.index.values):
			c['infinite']=rise(x=times,sigma=resolution,begin=t0)
		return c
	if model == 1:#consecutive decays
		n_decays=len(param)
		g=gauss(times,sigma=resolution/FWHM,mu=t0)
		if 'infinite' in list(pardf.index.values):
			infinite=True
			n_decays+=1
		else:
			infinite=False

		decays=param
		c=np.zeros((len(times),n_decays),dtype='float')
		
		if 'explicit_GS' in list(pardf.index.values):
			GS=True
			gs=np.zeros((len(times),1),dtype='float')
		else:
			GS=False
		for i in range(1,len(times)):
			dc=np.zeros(n_decays,dtype='float')
			dt=(times[i]-times[i-1])/(sub_steps)
			c_temp=c[i-1,:]
			for j in range(int(sub_steps)):
				for l in range(0,n_decays):
					if l>0:
						if infinite:
							if l<(n_decays-1):
								dc[l]=decays[l-1]*dt*c_temp[l-1]-decays[l]*dt*c_temp[l]
							else:
								dc[l]=decays[l-1]*dt*c_temp[l-1]
						else:
							dc[l]=decays[l-1]*dt*c_temp[l-1]-decays[l]*dt*c_temp[l]
					else:
						if infinite and n_decays==1:
							dc[l]=g[i]*dt
						else:
							dc[l]=g[i]*dt-decays[l]*dt*c_temp[l]
				c_temp=c_temp+dc
				c_temp[c_temp<0]=0
				#for b in range(c.shape[1]):
				#	c_temp[b] =np.nanmax([(c_temp[b]+float(dc[b])),0.])					
			c[i,:] =c_temp
			if GS and not infinite:
				gs[i]=-c[i,:].sum()
			elif GS:
				gs[i]=-c[i,:-1].sum()
		c=pandas.DataFrame(c,index=times)
		c.index.name='time'
		if infinite:
			labels=list(c.columns.values)
			labels[-1]='Non Decaying'
			c.columns=labels
			if 'background' in list(pardf.index.values):
				c['background']=1
		else:
			if 'background' in list(pardf.index.values):
				c['background']=1
		if GS:
			c['GS']=gs
		return c


def fill_int(ds,c,final=True,baseunit='ps',return_shapes=False):
	'''solving the intensity an equation_way, takes the target dataframe and the concentration frame 
	prepares the matrixes(c) the tries to solve this equation system using 
	eps=np.linalg.lstsq(AA,Af,rcond=-1)[0]
	if failes it returns a dictionary with 1000 as error (only entry) if successful
	it returns a dictionary that contains the 
	fit_error = (AE**2).sum() with AE beeing the difference of measured and calcuated matrix
	
	Parameters
	-----------
	
	ds : DataFrame
		DataFrame to be fitted
		
	c: DataFrame
		DataFrame oontaining the concentration matrix (the concentrations as with the times as index. 
		Each different species has a column with the species name as column name
		
	final : bool,optional
		if True (Default) the complete solutions will be attached otherwise only the error is attached
		
	baseunit : str,optional
		this string is used as unit for the time axis
	
	return_shapes : bool,optional
		Default = False, if True, then the concentrations and spectra are added to the re (even if not final)

	Returns
	------------------
	
	re : dict
		the dictionary "re" attached to the object containing all the matrixes and parameter. 
		
		if "final" is True:
		
			* "A" Shaped measured Matrix
			* "AC" Shaped calculated Matrix 
			* "AE" Difference between A and AC = linear error 
			* "DAC" DAS or SAS, labeled after the names given in the function (the columns of c) Care must be taken that this mesured intensity is C * DAS, the product. For exponential model the concentrations are normalized
			* "c" The Concentrations (meaning the evolution of the concentrations over time. Care must be taken that this mesured intensity is C * DAS, the product. For exponential model the concentrations are normalized
			* "error" is the S2, meaning AE**2.sum().sum()
			
		else:
		
			* "error" is the S2, meaning AE**2.sum()
			

	'''
	time=ds.index.values.astype('float')
	wl=ds.columns.values.astype('float')
	time_label=ds.index.name
	energy_label=ds.columns.name
	A=ds.values
	er=c.values
	ert = er.T
	AA = np.matmul(ert,er)
	Af = np.matmul(ert,A)
	try:
		eps=np.linalg.lstsq(AA,Af,rcond=-1)[0]
	except:
		re={'error':1000}
		return re
	eps[np.isnan(eps)]=0
	eps[np.isinf(eps)]=0
	AC = np.matmul(er,eps)
	AE = A-AC
	fit_error = (AE**2).sum()
	if final:
		A=pandas.DataFrame(A,index=time,columns=wl)
		AC=pandas.DataFrame(AC,index=time,columns=wl)
		AE=pandas.DataFrame(AE,index=time,columns=wl)
		DAC=pandas.DataFrame(eps.T,index=wl)
		A.index.name=time_label
		A.columns.name=energy_label
		AC.index.name=time_label
		AC.columns.name=energy_label
		AE.index.name=time_label
		AE.columns.name=energy_label
		DAC.index.name=energy_label
		try:
			DAC.columns=c.columns.values
		except:
			pass
		re={'A':A,'AC':AC,'AE':AE,'DAC':DAC,'error':fit_error,'c':c}
	elif return_shapes:
		re={'DAC':DAC,'error':fit_error,'c':c}
	else:
		re={'error':fit_error}
	return re


def err_func(paras, ds, mod = 'paral', final = False, log_fit = False, dump_paras = False, write_paras = False, 
			filename = None, ext_spectra = None, dump_shapes = False, sub_sample=None,pulse_sample=None):
	'''function that calculates and returns the error for the global fit. This function is intended for
	fitting a single dataset.
	
	Parameters
	--------------
	ds : DataFrame
		This dataframe contains the data to be fitted. This has to be shaped as it is intended to (so all shping parameters
		already applied. The dataframe expects the time to be in Index and the wavelength/energy 
		to be in the columns. The spectra is plotted with a second (energy) axis
	
	paras : lmfit parameter oject
		The parameter object that defines what is calculated 
		
	mod : str or function, optional
		The model selection is depending if it is an internal or external model. 
		The internal functions are triggered by calling their name 
		Two main are currently implemented 
		
			1. 'paral' (Default) or 'exponential'
			2. 'consecutive' or 'full_consecutive'
		
		In 2 the 'consecutive' and 'full_consecutive' are different in that for consecutive 
		the optimization is done using 'exponential' (as it shoudl give the same times)
		and then only in the last (final) iteration the 'full consecutive' differential
		equation is used. This has significant speed advantages, but can lead to errors particularly
		for the very fast times.
		
		As external model a function is handed to this parameter, this function 
		must accept the times and an paramater Dataframe and return a DataFrame 
		with the concentrations (similar to build_c) 
		
		for the internal functions:
		This datafram must contain the parameter that are used for creadting the dynamics
		the parameter must be named with the index.
		't0' = zero time, mandatory
		'resolution' = instrument response function, mandatory
		'background',optional = if this keyword is present a flat constant background is created (=1 over the whole time)
		'infinite',optional = if this keyword is present a new non decaying component is formed with the last decay time.
		'explicit_GS' = if this keyword is present thenthe ground state (including the bleach) will be added as a explicit component
		'k0,k1,...' = with increasing integers are taken as decay times. te number of these components is used to determine how many shall be generated.
	final : bool, optional
		this switch decides if just the squared error is returned (for False) (Default) or if the full
		matrixes are returned, including the r2 are returned.
		
	log_fit : bool, optional
		if False (Default) then the parameter are handed to the fitting function as they are, if true 
		then all times are first converted to log space. 
	
	dump_paras : bool, optional
		(Default) is False, If True creates two files in the working folder, one with the 
		currently used parameter created at the end of each optimisation step, and one with
		the set of parameter that up to now gave the lowest error. Intented to store
		the optimisation results if the fit needs to be interrupted 
		(if e.g. Ampgo simply needs to long to optimize.) useful option if things are slow
	
	filename : None or str, optional
		Only used in conjunction with 'dump_paras'. The program uses this filename to dump the 
		parameter to disk 
	
	sub_sample: False or iter, option
		Default is False: if a whole number value is given the time vector that comes from the measurements is sampled finer. The purpose is to in 
		the model catch faster dynamics than would be possible. Careful this is bloating up the size of the time vector and the simulation time is 
		to some extend proportional to this number
	
	pulse_sample: False or iter, option
		Default is False, but will be activated if the zero time is not included in the fit. If True then additional points in the 
		pump_region=np.linspace(t0-4*resolution,t0+4*resolution,20) are added.  This is better than sub sample, as it only adds 20 time points 
		to the data. This was necessary since otherwise the pulse (that is creating the intensity) is not sampled at all or only fractional. 
		This switch mainly has an influence on the absolute intensities of the species (the concentration matrix)    
	
	ext_spectra : DataFrame, optional
		(Default) is None, if given substract this spectra from the DataMatrix using the intensity 
		given in "C(t)" this function will only work for external models. The name of the spectral column 
		must be same as the name of the column used. If not the spectrum will be ignored. The spectrum will 
		be interpolated to the spectral points of the model ds before the substraction.
		a number of parameters can be defined to aid this process. These parameter are defined as normal parameters.
		"ext_spectra_scale" multiplies all spectra by this value (e.g. -1 to put the steady state absorption spectra in)
		"ext_spectra_shift" shifts all spectra by this value to compensate for calibration differences
		"ext_spectra_guide" (from version 7.1.0) This is a switch, if this keyword is present, then the spectra are used as guides 
		and not exclusively. This means the code will assume that these spectra are correct and substract them, then calulate the 
		difference and return as DAS the provided spectra plus the difference spectra
		
	write_paras : bool, optional 
		if True(Default) writes the currently varried values to screen

	'''
	global start_time
	time_label=ds.index.name
	energy_label=ds.columns.name
	pardf=par_to_pardf(paras)
	if log_fit:
		pardf.loc[pardf.is_rate,'value']=pardf.loc[pardf.is_rate,'value'].apply(lambda x: 10**x)
	if isinstance(mod,type('hello')):#did we use a build in model?
		times=ds.index.values.astype('float')
		times_ori=times.copy()	
		if pulse_sample is not None:
			t0=float(pardf.loc['t0','value'])
			resolution=float(pardf.loc['resolution','value'])
			if hasattr(pulse_sample,'__iter__'):pump_region=pulse_sample
			else:
				pump_region=np.linspace(t0-4*resolution,t0+4*resolution,20)
			if pump_region.max()<times_ori.min():
				connection_region=np.arange(pump_region[-1],times_ori.min(),resolution/10)
				times=np.unique(np.sort(np.hstack((pump_region,connection_region,times_ori))))
			else:
				times=np.unique(np.sort(np.hstack((pump_region,times_ori))))
		if sub_sample is not None:
			listen=[times]
			for i in range(1,sub_sample,1):
				listen.append(times_ori[:-1]+((times_ori[1:]-times_ori[:-1])*i/sub_sample))
			times=np.unique(np.hstack(listen))
			times.sort()
		if final:#for final we really want the model
			c=build_c(times=times,mod=mod,pardf=pardf)
		elif 'full_consecutive' in mod:# here we force the full consecutive modelling
			c=build_c(times=times,mod=mod,pardf=pardf)
		elif 'full_sequential' in mod:
			c=build_c(times=times,mod=mod,pardf=pardf)											 
		else:#here we "bypass" the full consecutive and optimize the rates with the decays
			c=build_c(times=times,mod='paral',pardf=pardf)
		c=c.loc[times_ori,:]
		c.index.name=time_label
		if ext_spectra is None:
			re=fill_int(ds=ds,c=c, return_shapes = dump_shapes)
		else:
			if 'ext_spectra_shift' in list(pardf.index.values):
				ext_spectra.index=ext_spectra.index.values+pardf.loc['ext_spectra_shift','value']
				ext_spectra=rebin(ext_spectra,ds.columns.values.astype(float))
			else:
				ext_spectra=rebin(ext_spectra,ds.columns.values.astype(float))
				
			if "ext_spectra_scale" in list(pardf.index.values):
				ext_spectra=ext_spectra*pardf.loc['ext_spectra_scale','value']
			c_temp=c.copy()
			for col in ext_spectra.columns.values:
				A,B=np.meshgrid(c.loc[:,col].values,ext_spectra.loc[:,col].values)
				C=pandas.DataFrame((A*B).T,index=c.index,columns=ext_spectra.index.values)
				ds=ds-C
				if "ext_spectra_guide" not in list(pardf.index.values):
					c_temp.drop(col,axis=1,inplace=True)
			re=fill_int(ds=ds,c=c_temp, return_shapes = dump_shapes)
		if final:
			labels=list(re['DAC'].columns.values)
			changed=True
			if 'background' in list(pardf.index.values):
				if 'infinite' in list(pardf.index.values):
					labels[-2]='Non Decaying'
				labels[-1]='background'
			else:
				if 'infinite' in list(pardf.index.values):
					labels[-1]='Non Decaying'
				else:
					changed=False
			if changed:
				re['DAC'].columns=labels
				re['c'].columns=labels
			if ext_spectra is not None:
				for col in ext_spectra.columns.values:
					if "ext_spectra_guide" in list(pardf.index.values):
						re['DAC'][col]=re['DAC'][col]+ext_spectra.loc[:,col].values
					else:
						re['DAC'][col]=ext_spectra.loc[:,col].values
						re['c'][col]=c.loc[:,col].values
					A,B=np.meshgrid(c.loc[:,col].values,ext_spectra.loc[:,col].values)
					C=pandas.DataFrame((A*B).T,index=c.index,columns=ext_spectra.index.values)
					re['A']=re['A']+C
					re['AC']=re['AC']+C
			re['r2']=1-re['error']/((re['A']-re['A'].mean().mean())**2).sum().sum()
			if dump_paras:
				try:
					pardf.loc['error','value']=re['error']
				except:
					pass
				try:
					pardf.loc['r2','value']=re['r2']
				except:
					pass
				try:
					if filename is None:
						store_name='minimal_dump_paras.par'
					else:
						store_name='minimal_dump_paras_%s.par'%filename
					min_df=pandas.read_csv(store_name,sep=',',header=None,skiprows=1)
					if float(min_df.iloc[-1,1])>float(re['error']):
						pardf.to_csv(store_name)
				except:
					pass
				if filename is None:
					store_name='dump_paras.par'
				else:
					store_name='dump_paras_%s.par'%filename
				pardf.to_csv(store_name)			
			return re
		else:
			if dump_paras:
				try:
					pardf.loc['error','value']=re['error']
				except:
					pass
				try:
					pardf.loc['r2','value']=re['r2']
				except:
					pass
				try:
					if filename is None:
						store_name='minimal_dump_paras.par'
					else:
						store_name='minimal_dump_paras_%s.par'%filename
					min_df=pandas.read_csv(store_name,sep=',',header=None,skiprows=1)
					if float(min_df.iloc[-1,1])>float(re['error']):
						pardf.to_csv(store_name)
				except:
					pass
				if filename is None:
					store_name='dump_paras.par'
				else:
					store_name='dump_paras_%s.par'%filename
				pardf.to_csv(store_name)
			if mod not in ['paral','exponential','consecutive']:
				if write_paras:
					print('----------------------------------')
					print(pardf)
				else:
					if np.abs(tm.time()-start_time)>10:
						start_time=tm.time()
						print(re['error'])
			if dump_shapes:
				if ext_spectra is not None:
					for col in ext_spectra.columns.values:
						if "ext_spectra_guide" in list(pardf.index.values):
							re['DAC'][col]=re['DAC'][col]+ext_spectra.loc[:,col].values
						else:
							re['DAC'][col]=ext_spectra.loc[:,col].values
							re['c'][col]=c.loc[:,col].values
				re['c'].to_csv(path_or_buf=filename + '_c')
				re['DAC'].to_csv(path_or_buf=filename + '_DAC')
			return re['error']
	else:# Nope we used an external model (sorry for the duplication)
		times=ds.index.values.astype('float')
		times_ori=times.copy()
		if pulse_sample is not None:
			t0=float(pardf.loc['t0','value'])
			resolution=float(pardf.loc['resolution','value'])
			if hasattr(pulse_sample,'__iter__'):pump_region=pulse_sample
			else:
				pump_region=np.linspace(t0-4*resolution,t0+4*resolution,20)
			if pump_region.max()<times_ori.min():
				connection_region=np.arange(pump_region[-1],times_ori.min(),resolution/10)
				times=np.unique(np.sort(np.hstack((pump_region,connection_region,times_ori))))
			else:
				times=np.unique(np.sort(np.hstack((pump_region,times_ori))))
		if sub_sample is not None:
			listen=[times]
			for i in range(1,sub_sample,1):
				listen.append(times_ori[:-1]+((times_ori[1:]-times_ori[:-1])*i/sub_sample))
			times=np.unique(np.hstack(listen))
			times.sort()
		try:
			c=mod(times=times,pardf=pardf.loc[:,'value'])
		except Exception as e:
			print(e)
		c=c.loc[times_ori,:]
		if ext_spectra is None:
			try:
				re=fill_int(ds=ds,c=c, return_shapes = dump_shapes)
			except Exception as e:
				print(e)
		else:
			ext_spectra.sort_index(inplace=True)
			if 'ext_spectra_shift' in list(pardf.index.values):
				ext_spectra.index=ext_spectra.index.values+pardf.loc['ext_spectra_shift','value']
				ext_spectra=rebin(ext_spectra,ds.columns.values.astype(float))
			else:
				ext_spectra=rebin(ext_spectra,ds.columns.values.astype(float))
			if "ext_spectra_scale" in list(pardf.index.values):
				ext_spectra=ext_spectra*pardf.loc['ext_spectra_scale','value']
			c_temp=c.copy()
			for col in ext_spectra.columns.values:
				A,B=np.meshgrid(c.loc[:,col].values,ext_spectra.loc[:,col].values)
				C=pandas.DataFrame((A*B).T,index=c.index,columns=ext_spectra.index.values)
				ds=ds-C
				if "ext_spectra_guide" not in list(pardf.index.values):
					c_temp.drop(col,axis=1,inplace=True)
			re=fill_int(ds=ds,c=c_temp, return_shapes = dump_shapes)
		if final:
			if len(re.keys())<3:#
				print('error in the calculation')
				return re
			if ext_spectra is None:
				re['DAC'].columns=c.columns.values
				re['c'].columns=c.columns.values
			else:
				re['DAC'].columns=c_temp.columns.values
				re['c'].columns=c_temp.columns.values
				for col in ext_spectra.columns.values:
					if "ext_spectra_guide" in list(pardf.index.values):
						re['DAC'][col]=re['DAC'][col]+ext_spectra.loc[:,col].values
					else:
						re['DAC'][col]=ext_spectra.loc[:,col].values
						re['c'][col]=c.loc[:,col].values
					A,B=np.meshgrid(c.loc[:,col].values,ext_spectra.loc[:,col].values)
					C=pandas.DataFrame((A*B).T,index=c.index,columns=ext_spectra.index.values)
					re['A']=re['A']+C
					re['AC']=re['AC']+C
			re['r2']=1-re['error']/((re['A']-re['A'].mean().mean())**2).sum().sum()
			if dump_paras:
				try:
					pardf.loc['error','value']=re['error']
				except:
					pass
				try:
					min_df=pandas.read_csv('minimal_dump_paras.par',sep=',',header=None,skiprows=1)
					if float(min_df.iloc[-1,1])>float(re['error']):
						pardf.to_csv('minimal_dump_paras.par')
				except:
					pass
				try:
					pardf.to_csv('dump_paras.par')
				except:
					print(pardf)
			return re
		else:
			if dump_paras:
				try:
					pardf.loc['error','value']=re['error']
				except:
					pass
				try:
					min_df=pandas.read_csv('minimal_dump_paras.par',sep=',',header=None,skiprows=1)
					if float(min_df.iloc[-1,1])>float(re['error']):
						pardf.to_csv('minimal_dump_paras.par')
				except Exception as e:
					print(e)
					pass
				try:
					pardf.to_csv('dump_paras.par')
				except Exception as e:
					print(e)
					print(pardf)
			if write_paras:
				print('----------------------------------')
				print(pardf)
			else:
				if np.abs(tm.time()-start_time)>30:
					start_time=tm.time()
					print(re['error'])
			if dump_shapes:
				if ext_spectra is not None:
					for col in ext_spectra.columns.values:
						if "ext_spectra_guide" in list(pardf.index.values):
							re['DAC'][col]=re['DAC'][col]+ext_spectra.loc[:,col].values
						else:
							re['DAC'][col]=ext_spectra.loc[:,col].values
							re['c'][col]=c.loc[:,col].values
				re['c'].to_csv(path_or_buf=filename + '_c')
				re['DAC'].to_csv(path_or_buf=filename + '_DAC')
			return re['error']


def err_func_multi(paras, mod = 'paral', final = False, log_fit = False, multi_project = None, 
					unique_parameter = None, weights = None, dump_paras = False, filename = None, 
					ext_spectra = None, dump_shapes = False, same_DAS = False,
					write_paras = True,sub_sample=None,pulse_sample=None,
					same_shape_params=True):
	'''function that calculates and returns the error for the global fit. This function is intended for
	fitting of multiple datasets
	
	Parameters
	--------------
	
	paras : lmfit parameter oject
		The parameter object that defines what is calculated 
		
	mod : str or function, optional
		The model selection is depending if it is an internal or external model. 
		The internal functions are triggered by calling their name 
		Two main are currently implemented 
		
			1. 'paral' (Default) or 'exponential'
			2. 'consecutive' or 'full_consecutive'
		
		In 2 the 'consecutive' and 'full_consecutive' are different in that for consecutive 
		the optimization is done using 'exponential' (as it shoudl give the same times)
		and then only in the last (final) iteration the 'full consecutive' differential
		equation is used. This has significant speed advantages, but can lead to errors particularly
		for the very fast times.
		
		for the internal functions:
		This datafram must contain the parameter that are used for creadting the dynamics
		the parameter must be named with the index.
		't0' = zero time, mandatory
		'resolution' = instrument response function, mandatory
		'background',optional = if this keyword is present a flat constant background is created (=1 over the whole time)
		'infinite',optional = if this keyword is present a new non decaying component is formed with the last decay time.
		'explicit_GS' = if this keyword is present thenthe ground state (including the bleach) will be added as a explicit component

		'k0,k1,...' = with increasing integers are taken as decay times. te number of these components is used to determine how many shall be generated.
		
		As external model a function is handed to this parameter, this function 
		must accept the times and an paramater Dataframe and return a DataFrame 
		with the concentrations (similar to build_c) 
		
	final : bool, optional
		this switch decides if just the squared error is returned (for False) (Default) or if the full
		matrixes are returned, including the r2 are returned.
		
	log_fit : bool, optional
		if False (Default) then the parameter are handed to the fitting function as they are, if true 
		then all times are first converted to log space. 
	
	dump_paras : bool, optional
		(Default) is False, If True creates two files in the working folder, one with the 
		currently used parameter created at the end of each optimisation step, and one with
		the set of parameter that up to now gave the lowest error. Intented to store
		the optimisation results if the fit needs to be interrupted 
		(if e.g. Ampgo simply needs to long to optimize.) useful option if things are slow
	
	filename : None or str, optional
		Only used in conjunction with 'dump_paras'. The program uses this filename to dump the 
		parameter to disk 
	
	multi_project : None or list (of TA projects), optional
		This switch is triggering the simultaneous optimisation of multiple datasets.  
		multi_project is as (Default) None. it expects an iterable (typically list) with other 
		TA projects (like ta) that are then optimised with the same parameter. 
		This means that all projects get the same parameter object for each iteration 
		of the fit and return their individual error, which is summed linearly. 
		The "weights" option allows to give each multi_project a specific weight (number) 
		that is multiplied to the error. If the weight object has the same number of items 
		as the multi_project it is assumed that the triggering object (the embedded project) 
		has the weight of 1, otherwise the first weight is for the embedded project. 
		The option 'unique_parameter' takes (a list) of parameter that are not 
		to be shared between the projects (and that are not optimized either) 
		The intended use of this is to give e.g. the pump power for multiple experiments to 
		study non linear behaviour. Returned will be only the parameter set for the optimium
		combination of all parameter. Internally, we iterate through the projects and calculate 
		for each project the error for each iteration. Important to note is that currently this 
		means that each DAS/SAS is calculated independently! For performing the same calculation
		with a single DAS, the Matrixes need to be concatenated before the run and an external
		function used to create a combined model. As this is very difficult to implement reliably
		For general use (think e.g. different pump wavelength) this has to be done manually. 
	
	unique_parameter : None or str or list (of strings), optional 
		only used in conjunction with 'multi_project', it takes (a list) of parameter that 
		are not to be shared between the projects (and that are not optimized either) 
		The intended use of this is to give e.g. the pump power for multiple experiments 
		to study non linear behaviour. (Default) None
		
	same_DAS : bool,optional
		changes the fit behavior and uses the same DAS for the optimization. This means that the ds are stacked before the fill int rounds
	
	weights : list of floats, optional
		only used in conjunction with 'multi_project'. The "weights" option allows to 
		give each multi\_project a specific weight (number) that is multiplied to the error. 
		If the weight object has the same number of items as the 'multi_project' it is assumed 
		that ta (the embedded project) has the weight of 1, otherwise the first weight is for the 
		embedded object
		
	ext_spectra : DataFrame, optional
		(Default) is None, if given substract this spectra from the DataMatrix using the intensity 
		given in "C(t)" this function will only work for external models. The name of the spectral column 
		must be same as the name of the column used. If not the spectrum will be ignored. The spectrum will 
		be interpolated to the spectral points of the model ds before the substraction.
		a number of parameters can be defined to aid this process. These parameter are defined as normal parameters.
		"ext_spectra_scale" multiplies all spectra by this value (e.g. -1 to put the steady state absorption spectra in)
		"ext_spectra_shift" shifts all spectra by this value to compensate for calibration differences
		"ext_spectra_guide" (from version 7.1.0) This is a switch, if this keyword is present, then the spectra are 
		used as guides and not exclusively. This means the code will assume that these spectra are correct and 
		substract them, then calulate the difference and return as DAS the provided spectra plus the difference spectra
		
	sub_sample: False or iter, option
		Default is False: if a whole number value is given the time vector that comes from the measurements is sampled finer. The purpose is to in 
		the model catch faster dynamics than would be possible. Careful this is bloating up the size of the time vector and the simulation time is 
		to some extend proportional to this number
	
	pulse_sample: False or iter, option
		Default is False, but will be activated if the zero time is not included in the fit. If True then additional points in the 
		pump_region=np.linspace(t0-4*resolution,t0+4*resolution,20) are added.  This is better than sub sample, as it only adds 20 time points 
		to the data. This was necessary since otherwise the pulse (that is creating the intensity) is not sampled at all or only fractional. 
		This switch mainly has an influence on the absolute intensities of the species (the concentration matrix)
	
	
	write_paras : bool, optional 
		if True(Default) writes the currently varried values to screen


	'''									   
	pardf_changing=par_to_pardf(paras)
	error_listen=[]
	r2_listen=[]
	if same_shape_params:
		slice_setting_object=multi_project[0].Copy()
		
	
	####### new same DAS, I'm lazy and will doublicate te loop. ###########
	if same_DAS:
		c_stack=[]
		ds_stack=[]
		par_stack=[]
		height_stack=[]
		for i,ta in enumerate(multi_project):
			if not same_shape_params:
				slice_setting_object=ta
			ds = sub_ds(ds = ta.ds, scattercut = slice_setting_object.scattercut, bordercut = slice_setting_object.bordercut, 
						timelimits = slice_setting_object.timelimits, wave_nm_bin = slice_setting_object.wave_nm_bin, 
						time_bin = slice_setting_object.time_bin, ignore_time_region = slice_setting_object.ignore_time_region)
			pardf=pardf_changing.copy()
			try:#let's see if the project has an parameter object
				pardf_ori=par_to_pardf(ta.par)
			except:
				pardf_ori=ta.pardf.copy()
			if unique_parameter is not None:
				for key in unique_parameter:
					pardf.loc[key,'value']=pardf_ori.loc[key,'value']
			par_stack.append(pardf)
			if log_fit:
				pardf.loc[pardf.is_rate,'value']=pardf.loc[pardf.is_rate,'value'].apply(lambda x: 10**x)

			times=ds.index.values.astype('float')
			times_ori=times.copy()
			if pulse_sample is not None:
				t0=float(pardf.loc['t0','value'])
				resolution=float(pardf.loc['resolution','value'])
				if hasattr(pulse_sample,'__iter__'):pump_region=pulse_sample
				else:
					pump_region=np.linspace(t0-4*resolution,t0+4*resolution,20)
				if pump_region.max()<times_ori.min():
					connection_region=np.arange(pump_region[-1],times_ori.min(),resolution/10)
					times=np.unique(np.sort(np.hstack((pump_region,connection_region,times_ori))))
				else:
					times=np.unique(np.sort(np.hstack((pump_region,times_ori))))
			if sub_sample is not None:
				listen=[times]
				for i in range(1,sub_sample,1):
					listen.append(times_ori[:-1]+((times_ori[1:]-times_ori[:-1])*i/sub_sample))
				times=np.unique(np.hstack(listen))
				times.sort()

			if isinstance(mod,type('hello')):#did we use a build in model?
				c=build_c(times=ds.index.values.astype('float'),mod=mod,pardf=pardf)
			else:
				c=mod(times=ds.index.values.astype('float'),pardf=pardf.loc[:,'value'])
			c=c.loc[times_ori,:]
			
			if ext_spectra is None:
				c_temp=c.copy()
			else:
				if 'ext_spectra_shift' in list(pardf.index.values):
					ext_spectra.index=ext_spectra.index.values+pardf.loc['ext_spectra_shift','value']
					ext_spectra=rebin(ext_spectra,ds.columns.values.astype(float))
				else:
					ext_spectra=rebin(ext_spectra,ds.columns.values.astype(float))
				if "ext_spectra_scale" in list(pardf.index.values):
					ext_spectra=ext_spectra*pardf.loc['ext_spectra_scale','value']
				c_temp=c.copy()
				for col in ext_spectra.columns.values:
					A,B=np.meshgrid(c.loc[:,col].values,ext_spectra.loc[:,col].values)
					C=pandas.DataFrame((A*B).T,index=c.index,columns=ext_spectra.index.values)
					ds=ds-C
					if "ext_spectra_guide" not in list(pardf.index.values):
						c_temp.drop(col,axis=1,inplace=True)
			if weights is not None:
				if len(weights)==len(multi_project)-1:
					weights=list(weights)
					weights.insert(0,1)
				elif len(weights)!=len(multi_project):
					Ex = ValueError()
					Ex.strerror='The number of entries i the list must either be the number of all elements (including \"TA\" or the number of elements in other. In this case the element ta gets the weight=1'
					raise Ex
				ds_stack.append(ds*weights[i])
			else:
				ds_stack.append(ds)
			c_stack.append(c_temp)
			height_stack.append(len(c_temp.index.values))
			
		A_con=pandas.concat(ds_stack)
		c_con=pandas.concat(c_stack)
		re=fill_int(ds=A_con,c=c_con, return_shapes = dump_shapes, final =final)
		if final:
			re['c']=c
		if dump_paras:
			try:
				pardf.loc['error','value']=re['error']
			except:
				pass
			if final:
				try:
					pardf.loc['r2','value']=1-re['error']/((re['A']-re['A'].mean().mean())**2).sum().sum()
				except:
					pass
			try:
				if filename is None:
					store_name='minimal_dump_paras.par'
				else:
					store_name='minimal_dump_paras_%s.par'%filename
				min_df=pandas.read_csv(store_name,sep=',',header=None,skiprows=1)
				if float(min_df.iloc[-1,1])>float(combined_error):
					pardf.to_csv(store_name)
			except:
				pass
			if filename is None:
				store_name='dump_paras.par'
			else:
				store_name='dump_paras_%s.par'%filename
			try:
				pardf.to_csv(store_name)
			except:
				print('Saving of %s failed'%store_name)
		
		if final:
			if isinstance(mod,type('hello')):#did we use a build in model?	
				labels=list(re['DAC'].columns.values)
				changed=True
				if 'background' in list(pardf.index.values):
					if 'infinite' in list(pardf.index.values):
						labels[-1]='Non Decaying'
						labels[-2]='background'
					else:
						labels[-1]='background'
				else:
					if 'infinite' in list(pardf.index.values):
						labels[-1]='Non Decaying'
					else:changed=False
				if changed:
					re['DAC'].columns=labels
					re['c'].columns=labels
				if ext_spectra is not None:
					for col in ext_spectra.columns.values:
						if "ext_spectra_guide" in list(pardf.index.values):
							re['DAC'][col]=re['DAC'][col]+ext_spectra.loc[:,col].values
						else:
							re['DAC'][col]=ext_spectra.loc[:,col].values
							re['c'][col]=c.loc[:,col].values
						A,B=np.meshgrid(c.loc[:,col].values,ext_spectra.loc[:,col].values)
						C=pandas.DataFrame((A*B).T,index=c.index,columns=ext_spectra.index.values)
						re['A']=re['A']+C
						re['AC']=re['AC']+C
			else:
				if ext_spectra is not None:
					for col in ext_spectra.columns.values:
						if "ext_spectra_guide" in list(pardf.index.values):
							re['DAC'][col]=re['DAC'][col]+ext_spectra.loc[:,col].values
						else:
							re['DAC'][col]=ext_spectra.loc[:,col].values
							re['c'][col]=c.loc[:,col].values
						A,B=np.meshgrid(c.loc[:,col].values,ext_spectra.loc[:,col].values)
						C=pandas.DataFrame((A*B).T,index=c.index,columns=ext_spectra.index.values)
						re['A']=re['A']+C
						re['AC']=re['AC']+C
				else:
					re['DAC'].columns=c.columns.values
					re['c'].columns=c.columns.values
			return_listen=[]
			for i,ta in enumerate(multi_project):
				re_local={}
				if i==0:
					lower=0
				else:
					lower=np.array(height_stack)[:i].sum()
				re_local['A']=re['A'].copy().iloc[lower:lower+height_stack[i],:]
				re_local['AC']=re['AC'].copy().iloc[lower:lower+height_stack[i],:]
				re_local['AE']=re['AE'].copy().iloc[lower:lower+height_stack[i],:]
				re_local['c']=re['c'].copy().iloc[lower:lower+height_stack[i],:]
				re_local['error_total']=re['error']
				re_local['error']=(re['AE']**2).sum().sum()
				re_local['DAC']=re['DAC'].copy()
				re_local['r2']=1-re_local['error']/((re_local['A']-re_local['A'].mean().mean())**2).sum().sum()
				re_local['r2_total']=1-re['error']/((re['A']-re['A'].mean().mean())**2).sum().sum()
				re_local['pardf']=par_stack[i]
				try:
					re_local['filename']=filename
				except:
					pass
				return_listen.append(re_local)
		if mod not in ['paral','exponential','consecutive']:
			if write_paras:
				print('----------------------------------')
				print(pardf)
			else:
				if np.abs(tm.time()-start_time)>10:
					start_time=tm.time()
					print(re['error'])
		if final:
			return return_listen
		else:
			return re['error']
	###################   not same DAS####################
	else:
		
		for i,ta in enumerate(multi_project):
			if not same_shape_params:
				slice_setting_object=ta
			ds = sub_ds(ds = ta.ds, scattercut = slice_setting_object.scattercut, bordercut = slice_setting_object.bordercut, 
						timelimits = slice_setting_object.timelimits, wave_nm_bin = slice_setting_object.wave_nm_bin, 
						time_bin = slice_setting_object.time_bin, ignore_time_region = slice_setting_object.ignore_time_region)	
			pardf=pardf_changing.copy()
			try:#let's see if the project has an parameter object
				pardf_ori=par_to_pardf(ta.par)
			except:
				pardf_ori=pardf
			if unique_parameter is not None:
				for key in unique_parameter:
					pardf.loc[key,'value']=pardf_ori.loc[key,'value']	
			if log_fit:
				pardf.loc[pardf.is_rate,'value']=pardf.loc[pardf.is_rate,'value'].apply(lambda x: 10**x)

			if isinstance(mod,type('hello')):#did we use a build in model?
				
				times=ds.index.values.astype('float')
				times_ori=times.copy()
				if pulse_sample is not None:
					t0=float(pardf.loc['t0','value'])
					resolution=float(pardf.loc['resolution','value'])
					if hasattr(pulse_sample,'__iter__'):pump_region=pulse_sample
					else:
						pump_region=np.linspace(t0-4*resolution,t0+4*resolution,20)
					if pump_region.max()<times_ori.min():
						connection_region=np.arange(pump_region[-1],times_ori.min(),resolution/10)
						times=np.unique(np.sort(np.hstack((pump_region,connection_region,times_ori))))
					else:
						times=np.unique(np.sort(np.hstack((pump_region,times_ori))))
				if sub_sample is not None:
					listen=[times]
					for i in range(1,sub_sample,1):
						listen.append(times_ori[:-1]+((times_ori[1:]-times_ori[:-1])*i/sub_sample))
					times=np.unique(np.hstack(listen))
					times.sort()
				
				c=build_c(times=ds.index.values.astype('float'),mod=mod,pardf=pardf)
				c=c.loc[times_ori,:]
				
				if ext_spectra is None:
					re=fill_int(ds=ds,c=c, return_shapes = dump_shapes)
				else:
					ext_spectra=rebin(ext_spectra,ds.columns.values.astype(float))
					c_temp=c.copy()
					for col in ext_spectra.columns.values:
						A,B=np.meshgrid(c.loc[:,col].values,ext_spectra.loc[:,col].values)
						C=pandas.DataFrame((A*B).T,index=c.index,columns=ext_spectra.index.values)
						ds=ds-C
						c_temp.drop(col,axis=1,inplace=True)
					re=fill_int(ds=ds,c=c_temp, return_shapes = dump_shapes)
				if final:
					if i==0:
						labels=list(re['DAC'].columns.values)
						changed=True
						if 'background' in list(pardf.index.values):
							if 'infinite' in list(pardf.index.values):
								labels[-1]='Non Decaying'
								labels[-2]='background'
							else:
								labels[-1]='background'
						else:
							if 'infinite' in list(pardf.index.values):
								labels[-1]='Non Decaying'
			   
							else:changed=False
						if changed:
							re['DAC'].columns=labels
							re['c'].columns=labels
						if ext_spectra is not None:
							for col in ext_spectra.columns.values:
								if "ext_spectra_guide" in list(pardf.index.values):
									re['DAC'][col]=re['DAC'][col]+ext_spectra.loc[:,col].values
								else:
									re['DAC'][col]=ext_spectra.loc[:,col].values
									re['c'][col]=c.loc[:,col].values
								A,B=np.meshgrid(c.loc[:,col].values,ext_spectra.loc[:,col].values)
								C=pandas.DataFrame((A*B).T,index=c.index,columns=ext_spectra.index.values)
								re['A']=re['A']+C
								re['AC']=re['AC']+C	
						re_final=re.copy()
					error_listen.append(re['error'])
					r2_listen.append(1-re['error']/((re['A']-re['A'].mean().mean())**2).sum().sum())
				else:
					if dump_shapes:
						if ext_spectra is not None:
							for col in ext_spectra.columns.values:
								if "ext_spectra_guide" in list(pardf.index.values):
									re['DAC'][col]=re['DAC'][col]+ext_spectra.loc[:,col].values
								else:
									re['DAC'][col]=ext_spectra.loc[:,col].values
									re['c'][col]=c.loc[:,col].values
						re['c'].to_csv(path_or_buf=ta.filename + '_c')
						re['DAC'].to_csv(path_or_buf=ta.filename + '_DAC')
					error_listen.append(re['error'])
			else:
				
				times=ds.index.values.astype('float')
				times_ori=times.copy()
				if pulse_sample is not None:
					t0=float(pardf.loc['t0','value'])
					resolution=float(pardf.loc['resolution','value'])
					if hasattr(pulse_sample,'__iter__'):pump_region=pulse_sample
					else:
						pump_region=np.linspace(t0-4*resolution,t0+4*resolution,20)
					if pump_region.max()<times_ori.min():
						connection_region=np.arange(pump_region[-1],times_ori.min(),resolution/10)
						times=np.unique(np.sort(np.hstack((pump_region,connection_region,times_ori))))
					else:
						times=np.unique(np.sort(np.hstack((pump_region,times_ori))))
				if sub_sample is not None:
					listen=[times]
					for i in range(1,sub_sample,1):
						listen.append(times_ori[:-1]+((times_ori[1:]-times_ori[:-1])*i/sub_sample))
					times=np.unique(np.hstack(listen))
					times.sort()
				
				c=mod(times=ds.index.values.astype('float'),pardf=pardf.loc[:,'value'])
				c=c.loc[times_ori,:]
				
				if ext_spectra is None:
					re=fill_int(ds=ds,c=c, return_shapes = dump_shapes)
				else:
					ext_spectra=rebin(ext_spectra,ds.columns.values.astype(float))
					c_temp=c.copy()
					for col in ext_spectra.columns.values:
						A,B=np.meshgrid(c.loc[:,col].values,ext_spectra.loc[:,col].values)
						C=pandas.DataFrame((A*B).T,index=c.index,columns=ext_spectra.index.values)
						ds=ds-C
						if "ext_spectra_guide" not in list(pardf.index.values):
							c_temp.drop(col,axis=1,inplace=True)
					re=fill_int(ds=ds,c=c_temp, return_shapes = dump_shapes)
				if final:
					if i==0:
						re['DAC'].columns=c.columns.values
						re['c'].columns=c.columns.values
						if ext_spectra is not None:
							for col in ext_spectra.columns.values:
								if "ext_spectra_guide" in list(pardf.index.values):
									re['DAC'][col]=re['DAC'][col]+ext_spectra.loc[:,col].values
								else:
									re['DAC'][col]=ext_spectra.loc[:,col].values
									re['c'][col]=c.loc[:,col].values
								A,B=np.meshgrid(c.loc[:,col].values,ext_spectra.loc[:,col].values)
								C=pandas.DataFrame((A*B).T,index=c.index,columns=ext_spectra.index.values)
								re['A']=re['A']+C
								re['AC']=re['AC']+C	
						re_final=re.copy()
					error_listen.append(re['error'])
					r2_listen.append(1-re['error']/((re['A']-re['A'].mean().mean())**2).sum().sum())
				else:
					if dump_shapes:
						if ext_spectra is not None:
							for col in ext_spectra.columns.values:
								if "ext_spectra_guide" in list(pardf.index.values):
									re['DAC'][col]=re['DAC'][col]+ext_spectra.loc[:,col].values
								else:
									re['DAC'][col]=ext_spectra.loc[:,col].values
									re['c'][col]=c.loc[:,col].values
						re['c'].to_csv(path_or_buf=filename + '_c')
						re['DAC'].to_csv(path_or_buf=filename + '_DAC')
					error_listen.append(re['error'])
		if weights is not None:
			if len(weights)==len(error_listen)-1:
				weights=list(weights)
				weights.insert(0,1)
			elif len(weights)!=len(error_listen):
				Ex = ValueError()
				Ex.strerror='The number of entries i the list must either be the number of all elements (including \"TA\" or the number of elements in other. In this case the element ta gets the weight=1'
				raise Ex
			combined_error=np.sqrt(((np.array(error_listen)*np.array(weights))**2).mean())
			if final:
				combined_r2=np.sqrt(((np.array(r2_listen)*np.array(weights))**2).mean())
		else:
			combined_error=np.sqrt((np.array(error_listen)**2).mean())
			if final:
				combined_r2=np.sqrt(((np.array(r2_listen))**2).mean())
		if final:
			re_final['error']=combined_error
			re_final['r2']=combined_r2
		if dump_paras:
			try:
				pardf.loc['error','value']=combined_error
			except:
				pass
			try:
				pardf.loc['r2','value']=combined_r2
			except:
				pass
			try:
				if filename is None:
					store_name='minimal_dump_paras.par'
				else:
					store_name='minimal_dump_paras_%s.par'%filename
				min_df=pandas.read_csv(store_name,sep=',',header=None,skiprows=1)
				if float(min_df.iloc[-1,1])>float(combined_error):
					pardf.to_csv(store_name)
			except:
				pass
			if filename is None:
				store_name='dump_paras.par'
			else:
				store_name='dump_paras_%s.par'%filename
			try:
				pardf.to_csv(store_name)
			except:
				print('Saving of %s failed'%store_name)
		if mod not in ['paral','exponential','consecutive']:
			print(combined_error)
		if final:
			return re_final
		else:
			return combined_error


def par_to_pardf(par):
	'''function to convert a parameters object into a pretty DataFrame, it expects par to be a lmfit parameters object and loops through the keys'''
	out_dicten={}
	for key in par.keys():
		out_dicten[key]={'value':par[key].value}
		if key[0] == 'k':#its a time parameter
			out_dicten[key]['is_rate']=True
		elif key[:2] == 'tk':#its a time parameter
			out_dicten[key]['is_rate']=True
		else:
			out_dicten[key]['is_rate']=False
		out_dicten[key]['min']=par[key].min
		out_dicten[key]['max']=par[key].max
		out_dicten[key]['vary']=par[key].vary
		out_dicten[key]['expr']=par[key].expr
	return pandas.DataFrame(out_dicten).T


def pardf_to_par(par_df):
	'''converts a dataframe to lmfit object
	set(value=None, vary=None, min=None, max=None, expr=None, brute_step=None)'''
	par=lmfit.Parameters()
	for key in par_df.index.values:
		par.add(key, value=par_df.loc[key,'value'], vary=par_df.loc[key,'vary'], min=par_df.loc[key,'min'], max=par_df.loc[key,'max'], expr=par_df.loc[key,'expr'])															 
	return par
	
	
def pardf_to_timedf(pardf):
	'''inverts all the rates to times in a dataframe'''
	timedf=pardf.copy()
	if 'upper_limit' in pardf.keys():
		for key in ['init_value','value','min','max','lower_limit','upper_limit']:
			for row in pardf.index.values:
				if timedf.loc[row,'is_rate']:
					if key == 'min':key_in='max'
					elif key == 'max':key_in='min'
					elif key == 'lower_limit':key_in='upper_limit'
					elif key == 'upper_limit':key_in='lower_limit'
					else:key_in=key
					try:
						if pardf.loc[row,key] !=0:
							timedf.loc[row,key_in]=1/pardf.loc[row,key] 
						else:
							timedf.loc[row,key_in]='inf'
					except:
						if key == 'init_value':pass#we don't save the init values, so we get an error when converting the saved file
						elif pardf.loc[row,key] is None:continue
						else:print('conversion of this key failed: %s %s'%(row,key))
	else:
		for key in ['init_value','value','min','max']:
			if key == 'min':key_in='max'		 
			elif key == 'max':key_in='min'
			else:key_in=key
			try:
				timedf.loc[pardf.is_rate,key_in]=pardf.loc[pardf.is_rate,key].apply(lambda x: 1/x if x!=0  else 'inf')
			except:
				if key == 'init_value':pass#we don't save the init values, so we get an error when converting the saved file
				else:print('conversion of this key failed:' + key)
	return timedf


class TA():	# object wrapper for the whole
	def __init__(self, filename=None, path = None, sep = "\t", decimal = '.', index_is_energy = False, transpose = False,
				sort_indexes = False, divide_times_by = None, shift_times_by = None, external_time = None, external_wave = None, 
				use_same_name = True, data_type = None , units = None, baseunit = None, ds = None, conversion_function = None):		 
		'''Function that opens and imports data into an TA object
		it is designed to open combined files that contain both the wavelength and the time. (e.g. SIA files as recorded by Pascher instruments software) or hdf5 projects saved by this software
		There are however a lot of additional options to open other ascii type files and adapt their format internally
		Attention times with Nan will be completely removed during the import
		
		Parameters
		----------
		
		filename : str
			
			* expects a filename in string form for opening a single file. 
			* alternatively 'gui' can be set as filename, then a TKinter gui is opened for select.
			* alternatively 'recent' can given as key word. in this case it tries to find a text file named "recent.dat" that should contain the path to the last file opened with the GUI. this file is then opened. if this file is not found the GUI is opened instead
		
		path : str or path object (optional)
			if path is a string without the operation system dependent separator, it is treated as a relative path, 
			e.g. data will look from the working directory in the sub director data. Otherwise this has to be a 
			full path in either strong or path object form.
		
		sep : str (optional)
			is the separator between different numbers, typical are tap (Backslash t) (Default) ,one or 
			multiple white spaces 'backslash s+' or comma ','.
		
		decimal : str (optional) 
			sets the ascii symbol that is used for the decimal sign. In most countries this is '.'(Default) 
			but it can be ',' in countries like Sweden or Germany
		
		index_is_energy : bool (optional)
			switches if the wavelength is given in nm (Default) or in eV (if True), currently everything 
			is handled as wavelength in nm internally
		
		data_type: str (optional)
			data_type is the string that represents the intensity measurements. Usually this contains if absolute 
			of differential data. This is used for the color intensity in the 2d plots and the y-axis for the 1d plots
			
		units: str (optional)
			this is used to identify the units on the energy axis and to label the slices, recognized is 'nm', 'eV' and 'keV' 
			but if another unit like 'cm^-1' is used it will state energy in 'cm^-1'. Pleas observe that if you use the index_is_energy
			switch the program tries to convert this energy into wavelength. 
			
		baseunit: str (optional)
			this is used to identify the units on the developing/time axis. This is name that is attached to the index of the dataframe. 
			setting this during import is equivalent to ta.baseunit
		
		transpose : bool (optional)
			if this switch is False (Default) the wavelength are the columns and the rows the times.
		
		sort_indexes : bool (optional)
			For False (Default) I assume that the times and energies are already in a rising order. 
			with this switch, both are sorted again. 
		
		divide_times_by : None or float (optional) 
			here a number can be given that scales the time by an arbitary factor. This is actually dividing
			the times by this value. Alternatively there is the variable self.baseunit. The latter only affects 
			what is written on the axis, while this value is actually used to scale the times. None (Default) 
			ignores this
		
		shift_times_by : None, float (optional)
			This a value by which the time axis is shifted during import. This is a useful option of e.g. 
			the recording software does not compensate for t0 and the data is always shifted. 
			None (Default) ignores this setting
		
		data_type : str, None
			this is the datatype and effectively the unit put on the intensity axis 
			(Default)'differential Absorption in $\mathregular{\Delta OD}$
		
		external_time : None or str (optional)
			Here a filename extension (string) can be given that contains the time vector. 
			The file is assumed to be at the same path as the data and to contain a single 
			type of separated data without header. 
			If use_same_name = True (default)
			It assumes that this is the ending for the file. The filename itself is taken from the filename. 
			e.g. if samp1.txt is the filename and external_time='.tid' the program searches 
			samp1.tid for the times. The transpose setting is applied and sets where the times are 
			to be inserted (row or column indexes)
			If use_same_name = False this should be the file containing the vector for the time (in the same format as the main file)
			
		external_wave : None or str (optional) 
			Here a filename extension (string) can be given that contains the wavelength vector. 
			If use_same_name = True (default)
			The file is assumed to be at the same path as the data and to contain a single type 
			of separated data without header. This is the ending for the file. The filename itself 
			is taken from the filename. e.g. if samp1.txt is the filename and external_wave='.wav' 
			then the program searches samp1.wav for the wavelength. The transpose setting is applied 
			and sets where the wavelength are to be inserted (columns or row indexes)
			If use_same_name = False
			this should be a full filename that contains the vector
			
		use_same_name : bool, optional
			this switches if the external filename included the loaded filename or is a separate file True(default)
			
		ds: pandas.DataFrame (optional)
			feed in an external dataframe instead of opening a file
			
		conversion_function: function(optional)
			function that receives should have the shape:
			return pandas Dataframe with time/frames  in rows and wavelength/energy in columns,
			The function is tested to accept (in that order) a 
			my_function(filename, external_time,external_wave), 
			my_function(filename, external_time), 
			my_function(filename,external_wave), 
			my_function(filename) and 
			return: the dataframe ds with the time_axis as rows and spectral axis as columns 
			if the ds.index.name ia not empty the "time axis" is in to that name the spectral axis is in ds.columns.name
			the return is investigated if it is one, two, or three things. 
			if two are returned then the second must be the name of what the intensity axis is. This value will then be set to data_type
			if three are returned the third is the baseunit (for the time axis) this allows to use the automatic naming in ps or nanosecond
			If the values units, data_type or baseunit are (manually) set in the import function the corresponding entries in
			datafram will be overwritten
			shift_times_by and divide_times_by will be applied if not None (useful to adjust for offset before chirp correction)
			 
			
			
		Returns
		-------
		
		A TA object with all parameter initialized
		
		Examples
		--------------
		
		Typical useage:
		
		>>> import plot_func as pf #import the module and give it a shorter name
		>>> ta=pf.TA('gui') #use a GUI to open
		>>> ta=pf.TA('sample_1.SIA') #use a filename in the same folder
		>>> ta=pf.TA('sample_1.hdf5',path='Data') #use a filename in the folder 'Data'
		
		Opening a list of files with external time vector (of the same name) so it looks for a data
		file "fite1.txt" and a file with the time information "file1.tid"
		
		>>>ta=pf.TA('file1.txt', external_time = 'tid')
		
		'''
		
		
		self.filename=filename
		self.path=check_folder(path=path,current_path=os.getcwd())
		if ds is not None:
			if filename is None:
				filename = 'external_ds'
				self.filename=filename
			else:
				if not isinstance(filename,str):
					filename = 'external_ds'
					self.filename=filename
			ext_data_switch=True
		else:
			ext_data_switch=False
		if filename == 'gui' and not ext_data_switch:
			root_window = tkinter.Tk()			
			root_window.withdraw()
			root_window.attributes('-topmost',True)
			root_window.after(1000, lambda: root_window.focus_force())
			complete_path = filedialog.askopenfilename(initialdir=os.getcwd())
			listen=os.path.split(complete_path)
			path=os.path.normpath(listen[0])
			self.path=path
			filename=listen[1]
			self.filename=filename
			with open('recent.dat','w') as f:
				f.write(complete_path)
		elif not ext_data_switch and filename == 'recent':
			try:
				with open('recent.dat','r') as f:
					complete_path = f.readline()
					listen=os.path.split(complete_path)
					path=os.path.normpath(listen[0])
					self.path=path
					filename=listen[1]
					self.filename=filename
			except:
				root_window = tkinter.Tk()
				root_window.withdraw()
				root_window.attributes('-topmost',True)
				root_window.after(1000, lambda: root_window.focus_force())
				complete_path = filedialog.askopenfilename(initialdir=os.getcwd())
				listen=os.path.split(complete_path)
				path=os.path.normpath(listen[0])
				self.path=path
				filename=listen[1]
				self.filename=filename
				with open('recent.dat','w') as f:
					f.write(complete_path)
		if filename == 'external' or ext_data_switch:#use a provided dataframe (ds) instead
			if data_type is not None:
				self.data_type = data_type
			if units is not None:
				self.units = units
				try:
					if len(ds.columns.name)==0:
						ds.columns.name= units
				except:
					pass
			else:
				try:
					if len(ds.columns.name)!=0:
						self.units = ds.columns.name 
				except:
					pass
			if baseunit is not None:
				self.baseunit = baseunit
				try:
					if len(ds.index.name)==0:
						if (baseunit == 'ps') or (baseunit == 'ns'):
							ds.index.name='Time in %s'%baseunit
						else:
							ds.index.name= baseunit
				except:
					pass
			else:
				try:
					if len(ds.index.name)!=0:
						self.baseunit = ds.index.name
				except:
						pass
			self.ds_ori=ds
			self.ds=ds
			self.__make_standard_parameter()
		elif ('hdf5' in filename) and (conversion_function is None):#we load a conversion function to deal with the file:#we read in data from previous run
			self.__read_project(saved_project=check_folder(path=self.path,filename=self.filename))
			self.__make_standard_parameter()
			self.Cor_Chirp(fitcoeff=self.fitcoeff)
		else:#we read in raw data from sia File
			if conversion_function is not None:
				filename=check_folder(path=self.path,filename=filename)										   
				try:
					ret=conversion_function(filename = filename, external_time = external_time, external_wave = external_wave)
				except:
					try:
						ret=conversion_function(filename = filename, external_time = external_time)
					except:
						try:
							ret=conversion_function(filename = filename, external_wave = external_wave)
						except:
							try:
								ret=conversion_function(filename = filename)
							except Exception as e:
								print(e)
								return False
				if isinstance(ret,pandas.DataFrame):
					##import is what we wanted
					ds=ret
				elif isinstance(ret,pandas.Series):
					ds=ret.as_frame()
				else:
					if len(ret) == 2:
						if data_type is None:
							ds,data_type=ret
						else:
							ds,_=ret
					elif len(ret) == 3:
						if data_type is None:
							ds,data_type,baseunit=ret
						else:
							ds,_,baseunit=ret
					else:
						print('sorry the return format of the conversion_function was not understood')
						print('return: the dataframe ds with the time_axis as rows and spectral axis as columns\n')
						print('if the ds.index.name ia not empty the "time axis" is in to that name the spectral axis is in ds.columns.name\n')
						print('the return is investigated if it is one, two, or three things.\n ')
						print('if two are returned then the second must be the name of what the intensity axis is. This value will then be set to data_type\n')
						print('if three are returned the third is the baseunit (for the time axis) this allows to use the automatic naming in ps or ns ' )
						return False
				## see if we have the name a data types in the data 
				if data_type is not None:
					self.data_type = data_type
				if units is not None:
					self.units = units
					try:
						if len(ds.columns.name)==0:
							ds.columns.name= units
					except:
						pass
				else:
					try:
						if len(ds.columns.name)!=0:
							self.units = ds.columns.name 
					except:
						pass
				if baseunit is not None:
					self.baseunit = baseunit
					try:
						if len(ds.index.name)==0:
							if (baseunit == 'ps') or (baseunit == 'ns'):
								ds.index.name='Time in %s'%baseunit
							else:
								ds.index.name= baseunit
					except:
						pass
				else:
					try:
						if len(ds.index.name)!=0:
							self.baseunit = ds.index.name
					except:
							pass
				if shift_times_by is not None:
					ds.index=ds.index.values+shift_times_by
				if divide_times_by is not None:
					ds.index=ds.index.values/divide_times_by
				self.ds_ori=ds
				self.ds=ds
			else:
				self.__read_ascii_data(sep = sep, decimal = decimal, index_is_energy = index_is_energy, 
									transpose = transpose, sort_indexes = sort_indexes, 
									divide_times_by = divide_times_by, shift_times_by = shift_times_by, 
									external_time = external_time, external_wave = external_wave, 
									use_same_name = use_same_name, data_type = data_type, units = units,  
									baseunit = baseunit)
			self.__make_standard_parameter()
			if len(self.ds.columns.values)<2:
				self.rel_wave=[self.ds.columns.values[0]]									
	def __call__(self):
		print('''
		Hello, this is a transient absorption project, 
		this project contains all the parameter and
		the main functions you can use to analyze data.
		-----------------------------------------------
		The main functions are:
		Plotting: ta.Plot_RAW(), ta.Plot_Interactive(), ta.Plot_fit_output()
		Comparing: ta.Compare_at_time(), ta.Compare_at_wave(), ta.Compare_DAC()
		Saving: ta.Save_Powerpoint(), ta.Save_project(), ta.Save_data(), ta.Save_Plots()
		Fitting: ta.Fit_Global()
		Shaping: ta.Background(), ta.Cor_Chirp(), ta.Man_Chirp()
		Parameter: call "ta.__dict__.keys()" to see all implemented shaping parameter
		-------------------------------------------------------------------------------
		you can get help and inspiration how to use each function by typing the function 
		name followed by a question mark like "ta.Fit_Global?"
		or by going to the documentation webpage:
		https://kimopack.readthedocs.io/en/latest
		--------------------------------------------------------------------------------
		In general I recommend to start by using one of the workflow notebooks that you can 
		download by calling pf.download_notebooks() or by going to 
		https://github.com/erdzeichen/KiMoPack/tree/main/Workflow_tools''')


	def __read_ascii_data(self, sep = "\t", decimal = '.', index_is_energy = False, transpose = False,
							sort_indexes = False, divide_times_by = None, shift_times_by = None, 
							external_time = None, external_wave = None, use_same_name = True, correct_ascii_errors = True,
							data_type = None, units = None,  baseunit = None):
		'''Fancy function that handles the import of pure ascii files.
		
		Parameters
		----------
		
		sep : str (optional)
			is the separator between different numbers, typical are tap (Backslash t) (Default) ,one or 
			multiple white spaces 'backslash s+' or comma ','.
		
		decimal : str (optional) 
			sets the ascii symbol that is used for the decimal sign. In most countries this is '.'(Default) 
			but it can be ',' in countries like Sweden or Germany
		
		index_is_energy : bool (optional)
			switches if the wavelength is given in nm (Default) or in eV (if True), currently everything 
			is handled as wavelength in nm internally
		
		data_type: str (optional)
			data_type is the string that represents the intensity measurements. Usually this contains if absolute 
			of differential data. This is used for the color intensity in the 2d plots and the y-axis for the 1d plots
		
		units: str (optional)
			this is used to identify the units on the energy axis and to label the slices, recognized is 'nm', 'eV' and 'keV' 
			but if another unit like 'cm^-1' is used it will state energy in 'cm^-1'. Pleas observe that if you use the index_is_energy
			switch the program tries to convert this energy into wavelength. 
		
		baseunit: str (optional)
			this is used to identify the units on the developing/time axis. This is name that is attached to the index of the dataframe. 
			setting this during import is equivalent to ta.baseunit
		
		transpose : bool (optional)
			if this switch is False (Default) the wavelength are the columns and the rows the times.
		
		sort_indexes : bool (optional)
			For False (Default) I assume that the times and energies are already in a rising order. 
			with this switch, both are sorted again. 
		
		divide_times_by : None or float (optional) 
			here a number can be given that scales the time by an arbitary factor. This is actually dividing
			the times by this value. Alternatively there is the variable self.baseunit. The latter only affects 
			what is written on the axis, while this value is actually used to scale the times. None (Default) 
			ignores this
		
		shift_times_by : None, float (optional)
			This a value by which the time axis is shifted during import. This is a useful option of e.g. 
			the recording software does not compensate for t0 and the data is always shifted. 
			None (Default) ignores this setting
		
		external_time : None or str (optional)
			Here a filename extension (string) can be given that contains the time vector. 
			The file is assumed to be at the same path as the data and to contain a single 
			type of separated data without header. 
			If use_same_name = True (default)
			It assumes that this is the ending for the file. The filename itself is taken from the filename. 
			e.g. if samp1.txt is the filename and external_time='.tid' the program searches 
			samp1.tid for the times. The transpose setting is applied and sets where the times are 
			to be inserted (row or column indexes)
			If use_same_name = False this should be the file containing the vector for the time (in the same format as the main file)
		
		external_wave : None or str (optional) 
			Here a filename extension (string) can be given that contains the wavelength vector. 
			If use_same_name = True (default)
			The file is assumed to be at the same path as the data and to contain a single type 
			of separated data without header. This is the ending for the file. The filename itself 
			is taken from the filename. e.g. if samp1.txt is the filename and external_wave='.wav' 
			then the program searches samp1.wav for the wavelength. The transpose setting is applied 
			and sets where the wavelength are to be inserted (columns or row indexes)
			If use_same_name = False
			this should be a full filename that contains the vector
		
		use_same_name : bool, optional
			this switches if the external filename included the loaded filename or is a separate file
		
		correct_ascii_errors :  bool (optional)
			If True (Default) then the code tries to catch some stuff like double minus signs and double dots
			
		'''
		
		self.ds_ori=pandas.read_csv(check_folder(path=self.path,filename=self.filename), sep=sep, index_col=0)
		if correct_ascii_errors:
			if (self.ds_ori.applymap(type) == float).all().all():
				pass#all columns were converted to float,nice
			else:
				print('some data bad, try filtering')
				try:# try forced conversion
					self.ds_ori=self.ds_ori.applymap(lambda x:  re.sub('--', '-',x) if type(x) is str else x)
					self.ds_ori=self.ds_ori.applymap(lambda x: re.sub(r'\.+', '.',x) if type(x) is str else x)
					self.ds_ori=self.ds_ori.astype(np.float64)
				except Exception as e:
					print('force cleaning went wrong and the file %s can not be read. Error message is:'%self.filename)
					print(e)
					return False
		if external_time is not None:
			if use_same_name:
				time_file=check_folder(path=self.path,filename=self.filename.split('.')[0]+'.'+external_time)
			else:
				time_file=check_folder(path=self.path,filename=external_time)
		if external_wave is not None:
			if use_same_name:
				wave_file=check_folder(path=self.path,filename=self.filename.split('.')[0]+'.'+external_wave)
			else:
				wave_file=check_folder(path=self.path,filename=external_wave)
		
		if external_time is not None:
			data_file_name=check_folder(path=self.path,filename=self.filename)
			times=pandas.read_csv(time_file,header=None,decimal=decimal).values.ravel()
			if transpose:
				if external_wave is not None:
					self.ds_ori=pandas.read_csv(check_folder(path=self.path,filename=self.filename), sep=sep , decimal=decimal, header=None)
					waves=pandas.read_csv(wave_file,header=None,decimal=decimal).values.ravel()
					self.ds_ori.index=waves
				else:
					self.ds_ori=pandas.read_csv(check_folder(path=self.path,filename=self.filename), sep=sep , decimal=decimal, index_col=0,header=None)
				self.ds_ori.columns=times
			else:
				if external_wave is not None:
					self.ds_ori=pandas.read_csv(check_folder(path=self.path,filename=self.filename), sep=sep, decimal=decimal,header=None)
					waves=pandas.read_csv(wave_file,header=None,decimal=decimal).values.ravel()
					self.ds_ori.columns=waves
				else:
					self.ds_ori=pandas.read_csv(check_folder(path=self.path,filename=self.filename), sep=sep, decimal=decimal)
				self.ds_ori.index=times
		elif external_wave is not None:
			data_file_name=check_folder(path=self.path,filename=self.filename)
			waves=pandas.read_csv(wave_file,header=None,decimal=decimal).values.ravel()
			if transpose:
				self.ds_ori=pandas.read_csv(check_folder(path=self.path,filename=self.filename), sep=sep, decimal=decimal)
				self.ds_ori.index=waves
			else:
				self.ds_ori=pandas.read_csv(check_folder(path=self.path,filename=self.filename), sep=sep, decimal=decimal,index_col=0,header=None)
				self.ds_ori.columns=waves
		self.ds_ori.columns=self.ds_ori.columns.astype('float')#Make columns indexes numbers
		self.ds_ori.index=self.ds_ori.index.astype('float')#Make row indexes numbers
		if index_is_energy:
			self.ds_ori.index=scipy.constants.h*scipy.constants.c/(self.ds_ori.index*1e-9*scipy.constants.electron_volt)
															  
		if transpose:
			self.ds_ori=self.ds_ori.T
		if sort_indexes:
			self.ds_ori.sort_index(axis=0,inplace=True)
			self.ds_ori.sort_index(axis=1,inplace=True)
		if shift_times_by is not None:
			self.ds_ori.index=self.ds_ori.index.values+shift_times_by
		if divide_times_by is not None:
			self.ds_ori.index=self.ds_ori.index.values/divide_times_by
		if data_type is not None:
			self.data_type = data_type
		if units is not None:
			self.units = units
		if baseunit is not None:
			self.baseunit = baseunit

		
	def __make_standard_parameter(self):
		'''function that sets the standard parameter. The function takes no input, but we use this docstring to explain the parameter.

		Parameters
		-------------

		log_scale : bool, optional
			If False (Default), The 2D plots (Matrix) is plotted with a pseudo logarithmic intensity scale. 
			This usually does not give good results unless the intensity scale is symmetric 
		self.cmap  : matplotlib.cm
			(Default)  standard_map - global parameter
			cmap is a powerfull variable that chooses the colour map applied for all plots. If set to 
			None (Default) then the self.cmap is used.
			As standard I use the color map "jet" from matplotlib. There are a variety of colormaps 
			available that are very usefull. Beside "jet", "viridis" is a good choice as it is well 
			visible under red-green blindness. Other useful maps are "prism" for high fluctuations 
			or diverging color maps like "seismic". 
			See https://matplotlib.org/3.1.0/tutorials/colors/colormaps.html for a comprehensive 
			selection. In the code the colormaps are imported so if plot_func is imported as pf then 
			self.cmap=pf.cm.viridis sets viridis as the map to use. Internally the colors are chosen 
			with the "colm" function. The 2d plots require a continuous color map so if something 
			else is give 2d plots are shown automatically with "jet". For all of the 1d plots however 
			I first select a number of colors before each plot. If cmap is a continous map then these
			are sampled evenly over the colourmap. Manual iterables of colours 
			cmap=[(1,0,0),(0,1,0),(0,0,1),...] are also accepted, as are vectors or dataframes that 
			contain as rows the colors. There must be of course sufficient colors present for 
			the numbers of lines that will be plotted. So I recommend to provide at least 10 colours 
			(e.g.~your university colors). colours are always given as a, list or tuple with RGA or RGBA
			(with the last A beeing the Alpha=transparency. All numbers are between 0 and 1. 
			If a list/vector/DataFrame is given for the colours they will be used in the order provided.
		self.lintresh : float
			The pseudo logratihmic range "symlog" is used for most time axis. Symlog plots a range around
			time zero linear and beyond this linear treshold 'lintresh' on a logarithmic scale. (Default) 0.3 
		self.log_fit  : 
			(Default)  False\n
			Transfer all the time-fitting parameters into log-space before the fit
		self.ignore_time_region  : None or list (of two floats or of lists)
			(Default)  None
			cut set a time range with a low and high limit from the fits. (Default) None nothing happens
			The region will be removed during the fitting process (and will be missing in the fit-result
			plots)\n
			Usage single region: [lower region limit,upper region limit]\n
			use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]
		self.error_matrix_amplification  : 
			(Default)  10
		self.rel_wave  : float or list (of floats) 
			(Default)  np.arange(300,1000,100)\n
			'rel_wave' and 'width' (in the object called 'wavelength_bin' work together for the creation 
			of kinetic plots. When plotting kinetic spectra one line will be plotted for each entrance
			in the list/vector rel_wave. During object generation the vector np.arange(300,1000,100) 
			is set as standard. Another typical using style would be to define a list of interesting 
			wavelength at which a kinetic development is to be plotted. At each selected wavelength 
			the data between wavelength+ta.wavelength_bin and wavelength-ta.wavelength_bin is averaged 
			for each timepoint returned 
		self.rel_time  : float or list/vector (of floats)
			(Default)  [0.2,0.3,0.5,1,3,10,30,100,300,1000,3000,9000]\n
			For each entry in rel_time a spectrum is plotted. If time_width_percent=0 (Default) the 
			nearest measured timepoint is chosen. For other values see 'time_width_percent'
		self.time_width_percent  : float
			(Default)  0 "rel_time" and "time_width_percent" work together for 
			creating spectral plots at specific timepoints. For each entry 
			in rel_time a spectrum is plotted. If however e.g. time_width_percent=10 
			the region between the timepoint closest to the  1.1 x timepoint 
			and 0.9 x timepoint is averaged and shown (and the legend adjusted accordingly). 
			This is particularly useful for the densly
			sampled region close to t=0. Typically for a logarithmic recorded kinetics, the 
			timepoints at later times will be further appart than 10 percent of the value, 
			but this allows to elegantly combine values around time=0 for better statistics. 
			This averaging is only applied for the plotting function and not for the fits.
		self.baseunit  : str
			(Default)  'ps'\n
			baseunit is a neat way to change the unit on the time axis of the plots. (Default) 'ps', but they 
			can be frames or something similarly. This is changing only the label of the axis. 
			During the import there is the option to divide the numbers by a factor. 
			I have also used frames or fs as units. Important is that all time units will be labeled with
			this unit.
		self.mod  : 
			(Default)  'exponential'\n
			This is the default fitting function, in general this is discussed in the fitting section
		self.scattercut  : None or iterable (of floats or other iterable, always pairs!)
			(Default)  None\n
			intented to "cut" one or multiple scatter regions. (if (Default) None nothing
			happens) If it is set the spectral region between the limits is set to zero. 
			Usage single region: [lower region limit,upper region limit], 
			use for multiple regions:[[lower limit 1,upper limit 1],[lower limit 2,upper limit 2],...]
		self.bordercut  : None or iterable (with two floats)
			(Default)  None\n
			cut spectra at the low and high wavelength limit. (Default) None 
			uses the limits of measurement 
		self.time_bin  : None or int
			(Default)  None
			is dividing the points on the time-axis in even bins and averages the found values in between. 
			This is a hard approach that also affects the fits. I do recommend to use this carefully, 
			it is most useful for modulated data. A better choice for transient absorption that only 
			affects the kinetics is 'time_width_percent'
		self.timelimits  : None or list (of 2 floats)
			(Default)  None\n
			cut times at the low and high time limit. (Default) None uses the limits of measurement
			Important: If either the background or the chirp is to be fit this must include the 
			time before zero! Useful: It is useful to work on different regions, starting with
			the longest (then use the ta.Backgound function prior to fit) and expand from there
		data_type : str
			this is the datatype and effectively the unit put on the intensity axis 
			(Default)'differential Absorption in $\mathregular{\Delta OD}$
		self.wave_nm_bin  : None or float 
			(Default)  None\n
			rebins the original data into even intervals. If set to None the original data will be used. 
			If set to a width (e.g. 2nm), the wavelength axis will be divided into steps of this size
			and the mean of all measurements in the interval is taken. The re-binning stops as soon as
			the measured stepsize is wider than given here, then the original bins are used. 
			This function is particularly useful for spectrometer with non-linear dispersion, 
			like a prism in the infrared.
		self.wavelength_bin  :  float, optional
			(Default)  10nm the width used in kinetics, see below
		self.intensity_range  : None, float or list [of two floats]
			(Default)  None - intensity_range is a general switch that governs what intensity range the plots show. 
			For the 1d plots this is the y-axis for the 2d-plots this is the colour scale. 
			This parameter recognizes three settings. If set to "None" (Default) this uses the minimum and
			maximum of the data. A single value like in the example below and the intended use is the symmetric
			scale while a list with two entries an assymmetric scale e.g. 
			intensity_range=3e-3 is converted into intensity_range=[-3e-3,3e-3]
		self.ds_ori.columns.name : str, optional 
			(Default)  'Wavelength in nm'\n
			This is the general energy axis. here we define it with the unit. Change this to energy for use in e.g x-ray science
		self.ds_ori.index.name : str, optional
			Standard 'Time in %s' % self.baseunit 
		self.data_type: str (optional)
			self.data_type='diff. Absorption in $\mathregular{\Delta OD}$'
		self.fitcoeff : list (5 floats)
			chirp correction polynom
		self.chirp_file : str
			if there is a file withthe right name write it here, otherwise None
		self.figure_path : str
			Path for saving figures, if set
		self.save_figures_to_folder : bool
			if True all figures are automatically saved when any plotfunction is called

		Examples
		-----------

		>>> ta.bordercut=[350,1200]  #remove all data outside this limit
		>>> ta.scattercut=[522,605]  #set data inside this limit to zero
		>>> ta.timelimits=[0.2,5000]  #remove all data outside this limit
		>>> ta.wave_nm_bin=5  #rebin the data to this width
		>>> ta.intensity_range=3e-3  #equivalent to [-3e-3,3e-3]
		>>> ta.intensity_range=[-1e-3,3e-3]  #intensity that is plotted in 2d plot and y-axis in 1d plots
		>>> ta.cmap=matplotlib.cm.prism  #choose different colour map
		>>> ta.ignore_time_region=[-0.1,0.1] #ignore -0.1ps to 0.1ps


		'''
		self.log_scale = False if not hasattr(self, 'log_scale') else self.log_scale
		self.cmap = standard_map if not hasattr(self, 'cmap') else self.cmap
		self.lintresh = 0.3 if not hasattr(self, 'lintresh') else self.lintresh
		self.log_fit = False if not hasattr(self, 'log_fit') else self.log_fit
		self.ignore_time_region = None if not hasattr(self, 'ignore_time_region') else self.ignore_time_region
		self.error_matrix_amplification = 10 if not hasattr(self, 'error_matrix_amplificatio') else self.error_matrix_amplification
		try:
			self.rel_wave = self.rel_wave
		except:
			self.rel_wave = np.arange(300,1000,100) if not hasattr(self, 'rel_wave') else self.rel_wave
		self.rel_time = [0.2,0.3,0.5,1,3,10,30,100,300,1000,3000,9000] if not hasattr(self, 'rel_time') else self.rel_time
		self.time_width_percent = 0 if not hasattr(self, 'time_width_percent') else self.time_width_percent
		self.baseunit = 'ps' if not hasattr(self, 'baseunit') else self.baseunit
		self.mod = 'exponential' if not hasattr(self, 'mod') else self.mod
		self.scattercut = None if not hasattr(self, 'scattercut') else self.scattercut
		self.bordercut = None if not hasattr(self, 'bordercut') else self.bordercut
		self.time_bin = None if not hasattr(self, 'time_bin') else self.time_bin
		self.timelimits = None if not hasattr(self, 'timelimits') else self.timelimits
		self.wave_nm_bin = None if not hasattr(self, 'wave_nm_bin') else self.wave_nm_bin
		self.wavelength_bin = 10 if not hasattr(self, 'wavelength_bin') else self.wavelength_bin
		self.save_figures_to_folder = False if not hasattr(self, 'save_figures_to_folder') else self.save_figures_to_folder
		self.intensity_range = None if not hasattr(self, 'intensity_range') else self.intensity_range
		self.ds_ori.index.name = 'Time in %s' % self.baseunit if not hasattr(self, 'ds_ori.index.name') else self.ds_ori.index.name
		self.equal_energy_bin = None if not hasattr(self, 'equal_energy_bin') else self.equal_energy_bin
		self.units='nm' if not hasattr(self, 'units') else self.units
		if self.units == 'nm':
			self.ds_ori.columns.name = 'Wavelength in %s'%self.units if not hasattr(self, 'ds_ori.columns.name') else self.ds_ori.columns.name
		elif self.units == 'eV':
			self.ds_ori.columns.name = 'Energy in %s'%self.units if not hasattr(self, 'ds_ori.columns.name') else self.ds_ori.columns.name
		elif self.units == 'keV':
			self.ds_ori.columns.name = 'Energy in %s'%self.units if not hasattr(self, 'ds_ori.columns.name') else self.ds_ori.columns.name
		else:
			self.ds_ori.columns.name = 'Energy in %s'%self.units if not hasattr(self, 'ds_ori.columns.name') else self.ds_ori.columns.name
		self.data_type= 'diff. Absorption in $\mathregular{\Delta OD}$' if not hasattr(self, 'data_type') else self.data_type
		
		try:#self.fitcoeff
			self.fitcoeff
			if len(list(self.fitcoeff))<5:raise
		except:
			self.fitcoeff=[0,0,0,0,0] #: test comment here
		try:#self.chirp_file
			self.chirp_file
		except:
			if os.path.isfile(check_folder(path=self.path,filename=self.filename.split('.')[0] + '_chirp.dat')):
				self.chirp_file=self.filename.split('.')[0] + '_chirp.dat'
			else:
				self.chirp_file=None
		try:#self.figure_path
			self.figure_path
		except:
			if self.save_figures_to_folder:
				self.figure_path=check_folder(path="result_figures",current_path=self.path)
			else:
				self.figure_path=None

		self.ds=self.ds_ori.copy()

	
	def Filter_data(self, ds=None, cut_bad_times = False, replace_bad_values = 0, value = 20, uppervalue = None, lowervalue = None, upper_matrix = None, lower_matrix = None):
		'''Filteres the data by applying hard replacements. if both replace_bad_values and 
		cut_bad_times are false or None, the times above "value" are replaced by zero
		
		Parameters
		------------
		
		ds : pandas Dataframe, optional
			if this is None (default) then the self.ds and self.ds_ori wil be filtered
		
		value : float, optional
			all values above this (absolute) value are considered to be corrupted. (Default 20) as classically the setup 
			reports optical DEnsity, an OD of 20 would be far above the typically expected
			OD 1e-3. Pascher instrument software uses a value of 21 to indicate an error.
		
		uppervalue : float, optional
			all values above this number are considered to be corrupted. (Default 20) as classically the setup 
			reports optical DEnsity, an OD of 20 would be far above the typically expected
			OD 1e-3. Pascher instrument software uses a value of 21 to indicate an error.
		
		lowervalue : float, optional
			all values below this number are considered to be corrupted. (Default -20) as classically the setup 
			reports optical DEnsity, an OD of -20 would be far above the typically expected
			OD 1e-3. Pascher instrument software uses a value of 21 to indicate an error.
		
		replace_bad_values : None of float, optional
			values above the treshold are replaced with this value. Ignored of None (Default)
			
		cut_bad_times = bool, optional
			True (Default=False) removes the whole time where this is true
		
		upper_matrix : Pandas DataFrame, optional
			all values above this treshold will be put N/A or replace by the value in replace_bad_values
		
		lower_matrix Pandas DataFrame, optional
			all values below this treshold will be put N/A or replace by the value in replace_bad_values
		
		the value is the upper bound. everything 
		above will be filtered. Standard is to drop the rows(=times) where something went wrong
		
		Examples
		---------
		
		typical usage
		
		>>> import plotfunc as pf
		>>> ta=pf.TA('testfile.SIA')
		>>> ta.Filter_data()
		>>> ta.Filter_data(value=1) #to filter times with at least one point with OD 1
		
		'''
		if uppervalue is None: uppervalue = np.abs(value)
		if lowervalue is None: lowervalue = -np.abs(value)
		
		if replace_bad_values is not None:
			cut_bad_times=False
		
		if ds is None:
			filtering=[self.ds,self.ds_ori]
		else:
			filtering=[ds]
		
		for dataset in filtering:
			if any([self.ignore_time_region is not None, self.scattercut is not None, self.bordercut is not None, self.timelimits is not None]):
				dataset=sub_ds(dataset, ignore_time_region = self.ignore_time_region, scattercut = self.scattercut, bordercut = self.bordercut, timelimits = self.timelimits)
			if cut_bad_times: #timepoint filter, delete the timepoints where value is stupid
				matrix_size=len(dataset.index.values)
				if upper_matrix is None:
					damaged_times=dataset[np.any(dataset.values>uppervalue,axis=1)].index
				else:
					damaged_times=dataset[np.any(dataset.values>upper_matrix,axis=1)].index
				dataset.drop(damaged_times,inplace = True)
				if lower_matrix is None:
					damaged_times=dataset[np.any(dataset.values<lowervalue,axis=1)].index
				else:
					damaged_times=dataset[np.any(dataset.values<lower_matrix,axis=1)].index
				dataset.drop(damaged_times,inplace = True)
				if len(dataset.index.values)<matrix_size*0.8:
					print('attention, more than 20% of the data was removed by this filter.') 
					print('Please check with if the spectal borders contain regions without light (and high noise)')
					print('Setting a bordercut and scattercut before the filtering might be useful')
			else: 
				if replace_bad_values is None: #individual data filter
					replace_bad_values=np.nan
				if upper_matrix is None:
					dataset.values[dataset.values>uppervalue]=replace_bad_values
				else:
					dataset.values[dataset.values>upper_matrix]=replace_bad_values
				if lower_matrix is None:
					dataset.values[dataset.values<lowervalue]=replace_bad_values
				else:
					dataset.values[dataset.values<lower_matrix]=replace_bad_values
				if replace_bad_values == np.nan: 
					if dataset.isna().sum().sum()>0.2 * dataset.notna().sum().sum():
						print('attention, more than 20% of the data was removed by this filter.') 
						print('Please check with if the spectal borders contain regions without light (and high noise)')
						print('Setting a bordercut and scattercut before the filtering might be useful')
				else:
					if dataset[dataset==replace_bad_values].notna().sum().sum()> 0.2* dataset[dataset!=replace_bad_values].notna().sum().sum():
						print('attention, more than 20% of the data was removed by this filter.') 
						print('Please check with if the spectal borders contain regions without light (and high noise)')
						print('Setting a bordercut and scattercut before the filtering might be useful')
		
		if ds is not None:return filtering[0]

	
	
	def Background(self, lowlimit=None,uplimit=-1, use_median=False, ds=None, correction=None):
		'''This is the background correction. In general it for each measured 
		wavelength averages the values from 'lowlimit' to 'uplimit' and 
		subtracts it from the data. It rund on the object (global) or if 
		given a specific ds local. 
		The low and uplimit can be set anywhere to substract any background.
		It is important to note that many problems during measurements might
		be visible in the data before time zero. So I recommend to first
		plot without background correction and only after this inspection 
		apply the background correction. 
		The fit function has its own way to calculcate and apply a background 
		That could be used instead (but making the fit less stable) 
		
		Parameters
		------------
		
		lowlimit : None or float, optional
			this is the lower limit from which the average (or median) is taken
			(Default) is None, in which case the lower limit of the data is used.
		
		uplimit : None or float, optional 
			this is the upper limit until which the average (or median) is taken
			(Default) is -1 (usually ps), in which case the lower limit of the data is used.
		
		use_median : bool, optional
			the Median is a more outlier resistant metric in comparision to 
			the Mean (Average). However the values are not quite as close 
			to the distribution center in case of very few values. False 
			(Default) means the Mean is used
			
		ds : None or DataFrame, optional
			if None (Default) the internal Dataframe self.ds is used, 
			otherwise the pandas DataFrame ds is corrected and returned
			
		correction : None or DataFrame, optional
			this is the correction applied. It must be a DataFrame with 
			the same numbers of columns (spectral points) as the used ds
		
		Examples
		--------
		if the object self has the name "ta"
		
		typical useage:  		
				
		>>> ta.Background()
		
		specify inegrated are to - inf (Default) up to -0.5ps and use the Median for computation	
		
		>>> ta.Background(uplimit = -0.5, use_median = True)
		
		'''
		if ds is None:
			run_global=True
			ds=self.ds
		else:
			run_global=False
			if correction is None:raise ValueError('We must have correction given, to slow otherhwise')
		if (lowlimit is None) and (correction is None):
			if use_median:
				correction=ds.loc[:uplimit].median(axis=0)
			else:
				correction=ds.loc[:uplimit].mean(axis=0)
		elif (lowlimit is not None) and (correction is None):
			if use_median:
				correction=ds.loc[lowlimit:uplimit].median(axis=0)
			else:
				correction=ds.loc[lowlimit:uplimit].mean(axis=0)
		if run_global:
			self.ds=ds-correction
			self.background_par=[lowlimit,uplimit,use_median,correction]
		else:
			return ds-correction


	def Man_Chirp(self,shown_window=[-1,1],path=None,max_points=40,cmap=cm.prism,ds=None,just_shift=False):
		'''Triggering of Manuel Fix_Chirp. usually used when Cor_Chirp has run already. 
		Alternatively delete the chirp file. This Function opens a plot in which the user manually selects a number of points
		These points will then be interpolated with a 4th order polynomial
		The user can then select a new t=0 point. 
		The first option allows to fine select an intensity setting for this chirp correction.
		However sometimes spikes are making this things difficult. In this case set a guessed intensity with self.intensity_range=1e-3
		
		Parameters
		-------------
		   
		path : str or path object (optional)
			if path is a string without the operation system dependent separator, it is treated as a relative path, 
			e.g. data will look from the working directory in the sub director data. Otherwise this has to be a 
			full path in either strong or path object form.
		
		shown_window : list (with two floats), optional
			Defines the window that is shown during chirp correction. If the t=0 is not visible, adjust this parameter
			to suit the experiment. If problems arise, I recomment to use Plot_Raw to check where t=0 is located

		max_points : int, optional
			Default = 40 max numbers of points to use in Gui selection. Useful option in case no middle mouse button
			is available. (e.g. touchpad)
		
		cmap : matplotlib colourmap, optional
			Colourmap to be used for the chirp correction. While there is a large selection here I recommend to choose
			a different map than is used for the normal 2d plotting.\n
			cm.prism (Default) has proofen to be very usefull
		
		ds: pandas dataframe,optional
			this allows to hand in an external ds, if this is done then the on disk saved fitcoeff are the new ones only and the 
			function returns the new fitcoeff and the combined fitcoeff, self also has a new variable called self.combined_fitcoeff
			the original file on dis and self.fitcoeff are NOT overwritten (are the old ones)
			the self.ds is the NEW one (with the correction applied)
			to reverse simply run Cor_Chirp()
			to permanently apply change self.fitcoeff with self.combined_fitcoeff and rename the file with 'filename_second_chirp' to filename_chirp
		'''
		if ds is None:
			ds=self.ds_ori
			original=True
		else:
			original=False
		if original:
			temp_ds = Fix_Chirp(ds, cmap = cmap, save_file = None, intensity_range = self.intensity_range, 
								wave_nm_bin = 10, shown_window = shown_window, filename = self.filename, 
								scattercut = self.scattercut, bordercut = self.bordercut, 
								path = check_folder(path = path, current_path = self.path), max_points = max_points, just_shift=just_shift)
		else:
			temp_ds = Fix_Chirp(ds, cmap = cmap, save_file = None, intensity_range = self.intensity_range, 
								wave_nm_bin = 10, shown_window = shown_window, filename = self.filename+'_second_chirp', 
								scattercut = self.scattercut, bordercut = self.bordercut, 
								path = check_folder(path = path, current_path = self.path), max_points = max_points, just_shift=just_shift)
		if isinstance(temp_ds,pandas.DataFrame):
			self.ds=temp_ds
			self.chirp_file=self.filename.split('.')[0] + '_chirp.dat'
			if original:#we have run from scratch
				self.Cor_Chirp(shown_window=shown_window,path=path,max_points=max_points,cmap=cmap,just_shift=just_shift)
			else:
				print('you provided a separate ds file. returned are the new fitcoeff and the combined fitcoeff, ta also has a new variable called ta.combined_fitcoeff')
				save_file=check_folder(path=path,current_path = self.path, filename=self.filename+'_second_chirp')
				with open(save_file,'r') as f:
					new_fitcoeff=f.readline()
				new_fitcoeff=np.array(new_fitcoeff.split(','),dtype='float')
				self.combined_fitcoeff=self.fitcoeff+new_fitcoeff
				return new_fitcoeff,self.combined_fitcoeff
		else:
			raise Warning('Man Chirp interrupted')
	def Check_Chirp(self, chirp_file = None, fitcoeff = None, cmap = cm.prism, path=None, ds=None, 
					shown_window = [-2, 2]):
		'''*Check_Chirp* is a function to check a provided chirp correction. It is is intended as 
		an option when Cor_Chirp fails due to lacking GUI. 
		A 4th order polynomial is plotted over the data, printed and returned.
		The intended use is that the components of this polynomial can be adjusted and then handed to 
		Cor_chirp.
		
		If neither a chirp_file not fitcoeff are provided a general purpose chrip correction is suggested
		as starting value. The function does not alter anything.
		use the internal intensity range to adjust for an optimum representation. ta.intensity_range=1e-3
		
		Parameters
		-------------
		
		chirp-file : None or str, optional
		   If a raw file was read(e.g. "data.SIA") and the chirp correction was
		   completed, a file with the attached word "chirp" is created and
		   stored in the same location. ("data_chirp.dat") This file contains
		   the 5 values of the chirp correction. By selecting such a file
		   (e.g. from another raw data) a specific chirp is applied. If a
		   specific name is given with **chirp_file** (and optional **path**)
		   then this file is used.\n
		   GUI\n
		   The word *'gui'* can be used instead of a filename to open a gui that
		   allows the selection of a chrip file
		   
		path : str or path object (optional)
			if path is a string without the operation system dependent separator, it is treated as a relative path, 
			e.g. data will look from the working directory in the sub director data. Otherwise this has to be a 
			full path in either strong or path object form.
		
		shown_window : list (with two floats), optional
			Defines the window that is shown during chirp correction. If the t=0 is not visible, adjust this parameter
			to suit the experiment. If problems arise, I recomment to use Plot_Raw to check where t=0 is located
			
		fitcoeff : list or vector (5 floats), optional
			One can give a vector/list with 5 numbers representing the parameter
			of a 4th order polynomial (in the order
			:math:`(a4*x^4 + a3*x^3+a2*x^2+a1*x1+a0)`. 

		cmap : matplotlib colourmap, optional
			Colourmap to be used for the chirp correction. While there is a large selection here I recommend to choose
			a different map than is used for the normal 2d plotting.\n
			cm.prism (Default) has proofen to be very usefull

		Examples
		----------
		
		In most cases:
		
		>>> import plot_func as pf
		>>> ta = pf.TA('test1.SIA') #open the original project, 
		>>> fitcoeff=ta.Check_Chirp() # if no specific correction is found
		>>> fitcoeff[1]+=0.1e-8
		>>> ta.Check_Chirp(fitcoeff=fitcoeff)
		
		'''
		
		if ds is None:ds=self.ds_ori
		if (chirp_file is None) and (fitcoeff is None):
			fitcoeff=[-3.384e-12,1.456e-08,-2.366e-05,0.0172,-4.306]
		elif chirp_file is not None:
			if 'gui' in chirp_file:
				root_window = tkinter.Tk()
				root_window.withdraw()
				root_window.attributes('-topmost',True)
				root_window.after(1000, lambda: root_window.focus_force())
				complete_path = filedialog.askopenfilename(initialdir=os.getcwd())
				listen=os.path.split(complete_path)
				path=os.path.normpath(listen[0])
				chirp_file=listen[1]
			path=check_folder(path,self.path)
			with open(check_folder(path=path,filename=chirp_file),'r') as f:
				fitcoeff=[float(a) for a in f.readline().split(',')]
		fig,ax=plt.subplots()
		ax = plot2d(ax = ax, ds = ds, cmap = cmap, wave_nm_bin = self.wave_nm_bin, scattercut = self.scattercut, bordercut = self.bordercut, 
					timelimits = shown_window, intensity_range = self.intensity_range, 
					title = 'This plot shows the set chirp', use_colorbar = False, 
					plot_type = "linear", log_scale = False)
		correcttimeval = np.polyval(fitcoeff, ds.columns.values.astype('float'))
		ax.plot(ds.columns.values.astype('float'),correcttimeval)	
		print('used fitcoeff:')
		print(fitcoeff)
		return fitcoeff
			
			

	def Cor_Chirp(self, chirp_file = None, path = None, shown_window = [-1, 1], fitcoeff = None, max_points = 40, cmap = cm.prism, just_shift=False):
		'''*Cor_Chirp* is a powerful Function to correct for a different arrival times of 
		different wavelength (sometimes call chirp). 
		In general if a file is opened for the first time this function is opening 
		a plot and allows the user to select a number of points, which are then 
		approximated with a 4th order polynomial and finally to select a point 
		that is declared as time zero. The observed window as well as the intensities 
		and the colour map can be chosen to enable a good correction. Here a fast 
		iterating colour scheme such as "prism" is often a good choice. In all of the 
		selections a left click selects, a right click removes the last point and 
		a middle click (sometime appreviated by clicking left and right together) 
		finishes the selection. If no middle click exists, the process
		automatically ends after max_points (40 preset).
		
		The first option allows to fine select an intensity setting for this chirp correction.
		However sometimes spikes are making this things difficult. 
		In this case set a guessed intensity with self.intensity_range=1e-3\n 
		Note that scattercut, bordercut and intensity_range can be used

		After the first run the polynom is stored in self.fitcoeff, a new matrix 
		calculated from self.ds_ori that is stored as self.ds and a file stored in the 
		same location as the original data. The second time the function *Cor_Chirp* is 
		run the function will find the file and apply the chirp correction automatically.

		If one does want to re-run the chirp correction the function *Man_Chirp* does
		not look for this file, but creates after finishing a new file.

		Alternatively the polynom or a filename can be given that load a chirp correction
		(e.g. from a different run with the same sample).
		The function *Cor_Chirp* selects in the order: 

		# "fitcoeff"
		# "other files"
		# "stored_file"
		# call Man_Chirp (clicking by hand)
		

		
		Parameters
		-------------
		
		chirp-file : None or str, optional
		   If a raw file was read(e.g. "data.SIA") and the chirp correction was
		   completed, a file with the attached word "chirp" is created and
		   stored in the same location. ("data_chirp.dat") This file contains
		   the 5 values of the chirp correction. By selecting such a file
		   (e.g. from another raw data) a specific chirp is applied. If a
		   specific name is given with **chirp_file** (and optional **path**)
		   then this file is used.\n
		   GUI\n
		   The word *'gui'* can be used instead of a filename to open a gui that
		   allows the selection of a chrip file
		   
		path : str or path object (optional)
			if path is a string without the operation system dependent separator, it is treated as a relative path, 
			e.g. data will look from the working directory in the sub director data. Otherwise this has to be a 
			full path in either strong or path object form.
		
		shown_window : list (with two floats), optional
			Defines the window that is shown during chirp correction. If the t=0 is not visible, adjust this parameter
			to suit the experiment. If problems arise, I recomment to use Plot_Raw to check where t=0 is located
			
		fitcoeff : list or vector (5 floats), optional
			One can give a vector/list with 5 numbers representing the parameter
			of a 4th order polynomial (in the order
			:math:`(a4*x^4 + a3*x^3+a2*x^2+a1*x1+a0)`. The chirp parameter are
			stored in ta.fitcoeff and can thus be used in other TA objects. This
			vector is also stored with the file and automatically applied during
			re-loading of a hdf5-object

		max_points : int, optional
			Default = 40 max numbers of points to use in Gui selection. Useful option in case no middle mouse button
			is available. (e.g. touchpad)
		
		cmap : matplotlib colourmap, optional
			Colourmap to be used for the chirp correction. While there is a large selection here I recommend to choose
			a different map than is used for the normal 2d plotting.\n
			cm.prism (Default) has proofen to be very usefull
			
		just_shift: bool, optional
			This switch turns of the polynomial selection and only permits a shift.

		Examples
		----------
		
		In most cases:
		
		>>> import plot_func as pf
		>>> ta = pf.TA('test1.SIA') #open the original project, 
		>>> ta.Cor_Chirp()
		
		Selecting a specific correction
		
		>>> ta.Cor_Chirp(‘gui’)
		>>> ta.Cor_Chirp(chirp_file = 'older_data_chirp.dat')
		>>> #use the coefficients from a different project
		>>> ta.Cor_Chirp(fitcoeff = ta_old.fitcoeff) #use the coefficients from a different project
		
		'''
		
		if chirp_file is None:
			chirp_file=self.chirp_file
		elif 'gui' in chirp_file:
			root_window = tkinter.Tk()
			root_window.withdraw()
			root_window.attributes('-topmost',True)
			root_window.after(1000, lambda: root_window.focus_force())
			complete_path = filedialog.askopenfilename(initialdir=os.getcwd())
			listen=os.path.split(complete_path)
			path=os.path.normpath(listen[0])
			chirp_file=listen[1]
		path=check_folder(path,self.path)
		if fitcoeff is not None:#we use a stored project
			try:
				if len(fitcoeff)==5 or len(fitcoeff)==6:#we provide a valid list/vector
					if all(elem == 0 for elem in fitcoeff):
						self.ds=self.ds_ori
						print('all chirp coefficients are zero so no chirp correction applied')
					else:	
						self.ds=Fix_Chirp(self.ds_ori,fitcoeff=fitcoeff)
					self.fitcoeff=fitcoeff #we came to here so fitcoeff must be right
				else:
					raise
			except:
				self.ds=self.ds_ori
		  
				print('something went wrong with the provided fitcoeff. This should be either a list/array with 5-6 parameter or the object should contain the parameter')
				print('fitcoeff is currently:' + fitcoeff)
		else:
			try:
				if chirp_file is None:
					save_file=None
				else:
					save_file=check_folder(path = path, filename = chirp_file)
				self.ds = Fix_Chirp(self.ds_ori, cmap = cmap, save_file = save_file,
									filename = self.filename, path = self.path,
									scattercut = self.scattercut, bordercut = self.bordercut, 
									intensity_range = self.intensity_range, wave_nm_bin = 10, shown_window = shown_window, 
									fitcoeff = fitcoeff, max_points = max_points, just_shift=just_shift)
				if save_file is None:
					if self.filename is None:
						chirp_file='chirp.dat'
					else:
						f=self.filename.split('.')[0]
						chirp_file=f+'_chirp' + '.dat'
				self.chirp_file=chirp_file
				with open(check_folder(path=path,filename=chirp_file),'r') as f:
					self.fitcoeff=[float(a) for a in f.readline().split(',')]
			except:
				print(check_folder(path=self.path,filename=self.filename.split('.')[0] + '_chirp.dat'))
				if os.path.isfile(check_folder(path=self.path,filename=self.filename.split('.')[0] + '_chirp.dat')):
					print('somehting is wrong, try deleting old chirp file')
					raise
				else:
					print('No old chirp file')
					self.Man_Chirp(path=path,cmap=cmap,shown_window=shown_window,max_points=max_points)
					chirp_file=self.chirp_file
					with open(check_folder(path=path,filename=chirp_file),'r') as f:
						self.fitcoeff=[float(a) for a in f.readline().split(',')]		
		self.ds.columns.name=self.ds_ori.columns.name
		self.ds.index.name=self.ds_ori.index.name

	def Plot_Interactive(self, fitted = False, ds = None, cmap = None, plot_on_move = False):
		'''interactive plotting function. it plots the matrix in the middle and two slices that are selected by the mouse (click) 
		
		Parameters
		---------------
		fitted : bool, optional
			this switch decides if the fitted or the RAW data is plotted with this widget to 
			inspect the data data. If fitted is False (Default) then the raw data and an interpolation 
			is used to plot. 
		
		cmap : None or matplotlib color map, optional
			is a powerfull variable that chooses the colour map applied for all plots. If set to 
			None (Default) then the self.cmap is used.
			As standard I use the color map "jet" from matplotlib. There are a variety of colormaps 
			available that are very usefull. Beside "jet", "viridis" is a good choice as it is well 
			visible under red-green blindness. Other useful maps are "prism" for high fluctuations 
			or diverging color maps like "seismic". 
			See https://matplotlib.org/3.1.0/tutorials/colors/colormaps.html for a comprehensive 
			selection. In the code the colormaps are imported so if plot_func is imported as pf then 
			self.cmap=pf.cm.viridis sets viridis as the map to use. Internally the colors are chosen 
			with the "colm" function. The 2d plots require a continuous color map so if something 
			else is give 2d plots are shown automatically with "jet". For all of the 1d plots however 
			I first select a number of colors before each plot. If cmap is a continous map then these
			are sampled evenly over the colourmap. Manual iterables of colours 
			cmap=[(1,0,0),(0,1,0),(0,0,1),...] are also accepted, as are vectors or dataframes that 
			contain as rows the colors. There must be of course sufficient colors present for 
			the numbers of lines that will be plotted. So I recommend to provide at least 10 colours 
			(e.g.~your university colors). colours are always given as a, list or tuple with RGA or RGBA
			(with the last A beeing the Alpha=transparency. All numbers are between 0 and 1. 
			If a list/vector/DataFrame is given for the colours they will be used in the order provided.		
		
		ds : DataFrame, optional
			if None (Default), the program first tests self.ds and if this is not there then self.ds_ori.
			This option was introduced to allow plotting of other matrixes with the same parameter
			
		plot_on_move : bool, optional
			Default: False plots the slices after click, on True the plot constantly reslices and on click 
			The current position is written down.
			
		 '''
		from matplotlib.widgets import Cursor
		if cmap is None:cmap=self.cmap
		if ds is None:
			if not fitted:
				if self.ds is None:
					ds=self.ds_ori.copy()
				else:
					ds=self.ds.copy()
			else:
				ds=self.re['A']
				modelled=self.re['AC']
		intensity_range=self.intensity_range
		if intensity_range is None:
			try:
				maxim=max([abs(ds.values.min()),abs(ds.values.max())])
				intensity_range=[-maxim,maxim]
			except:
				intensity_range=[-1e-2,1e-2]
		else:
			if not hasattr(intensity_range,'__iter__'):#lets have an lazy option
				intensity_range=[-intensity_range,intensity_range]
		class MouseMove:
			# initialization
			def __init__(self, ds, cmap, intensity_range, log_scale, baseunit, 
							timelimits, scattercut, bordercut, wave_nm_bin, equal_energy_bin, ignore_time_region,
							time_bin, lintresh, data_type, width, time_width_percent):
				fig = plt.figure(tight_layout=True,figsize=(14,8))
				gs = GridSpec(5, 4)
				self.ax= fig.add_subplot(gs[1:, :3])
				self.ds=ds
				self.cmap=cmap
				self.intensity_range=intensity_range
				self.log_scale=log_scale
				self.baseunit=baseunit
				self.timelimits=timelimits
				self.scattercut=scattercut
				self.bordercut=bordercut
				self.wave_nm_bin=wave_nm_bin
				self.equal_energy_bin=equal_energy_bin
				self.ignore_time_region=ignore_time_region
				self.time_bin=time_bin
				self.data_type=data_type
				self.width=width
				self.lintresh=lintresh
				self.time_width_percent=time_width_percent
				
				
				self.ax = plot2d(ds=ds, ax=self.ax, cmap=cmap, intensity_range=self.intensity_range, 
					log_scale=self.log_scale, baseunit=self.baseunit, timelimits=self.timelimits, 
					scattercut=self.scattercut, bordercut=self.bordercut, wave_nm_bin=self.wave_nm_bin, equal_energy_bin=self.equal_energy_bin,
					ignore_time_region=self.ignore_time_region, time_bin=self.time_bin, 
					lintresh=self.lintresh, data_type = self.data_type, use_colorbar = False)
				self.ax_time= fig.add_subplot(gs[0, :3],sharex=self.ax)
				self.ax_kinetic= fig.add_subplot(gs[1:, -1],sharey=self.ax)
				plt.subplots_adjust(wspace=0,hspace=0)
				if plot_on_move:
					fig.canvas.mpl_connect('motion_notify_event', self.move)
					fig.canvas.mpl_connect('button_press_event', self.click)
				else:
					fig.canvas.mpl_connect('button_press_event', self.move)
			
			def click(self, event):
				
				x, y = event.xdata, event.ydata
				if self.equal_energy_bin is not None:
					x=scipy.constants.h*scipy.constants.c/(x*1e-9*scipy.constants.electron_volt) 
				print('x=%g, y=%g\n'%(x,y))
			
			def move(self, event):
				x, y = event.xdata, event.ydata
				if self.equal_energy_bin is not None:
					x=scipy.constants.h*scipy.constants.c/(x*1e-9*scipy.constants.electron_volt) 
				try:
					self.ax_time.cla()
				except:
					pass
					
				if not fitted:	
					ds_temp1 = sub_ds(ds = Frame_golay(ds,5,3), times = y, time_width_percent = self.time_width_percent, 
									scattercut = self.scattercut, drop_scatter=True, bordercut = self.bordercut, 
									ignore_time_region = self.ignore_time_region, wave_nm_bin = self.wave_nm_bin, equal_energy_bin=self.equal_energy_bin,
									wavelength_bin = self.width)
					ds_temp1.plot(ax=self.ax_time,style='-',color='red')
					
				else:
					ds_temp1 = sub_ds(ds = modelled, times = y, time_width_percent = self.time_width_percent, 
									scattercut = self.scattercut, drop_scatter=True, bordercut = self.bordercut, 
									ignore_time_region = self.ignore_time_region, wave_nm_bin = self.wave_nm_bin,  equal_energy_bin=self.equal_energy_bin,
									wavelength_bin = self.width)
					ds_temp1.plot(ax=self.ax_time,style='-',color='red')
					
				ds_temp = sub_ds(ds = ds, times = y, time_width_percent = self.time_width_percent, 
									scattercut = self.scattercut, drop_scatter=True, bordercut = self.bordercut, 
									ignore_time_region = self.ignore_time_region, wave_nm_bin = self.wave_nm_bin,  equal_energy_bin=self.equal_energy_bin,
									wavelength_bin = self.width)
				ds_temp.plot(ax=self.ax_time,style='*',color='black')
				self.ax_time.plot(self.ax_time.get_xlim(),[0,0],'gray')
				if not fitted:
					self.ax_time.legend(['%.3g %s smoothed'%(y,self.baseunit)])
				else:
					self.ax_time.legend(['%.3g %s fitted'%(y,self.baseunit)])
				self.ax_time.set_yticks(self.ax_time.get_ylim())
				self.ax_time.set_yticklabels(['%.1e'%f for f in self.ax_time.get_ylim()])
				
				for i in range(3):
					try:
						self.ax_kinetic.lines.pop(0)
					except:
						pass
				if self.width is None:
					self.width = 10
				
				if not fitted:
					ds_temp1 = sub_ds(ds = Frame_golay(ds), wavelength = x, scattercut = self.scattercut, drop_scatter=True, 
								bordercut = self.bordercut, ignore_time_region = self.ignore_time_region, 
								wave_nm_bin = self.wave_nm_bin,wavelength_bin = self.width)
					self.ax_kinetic.plot(ds_temp1.values,ds_temp1.index.values,'-',label='%.0f smoothed'%x,color='red')
				else:
					ds_temp1 = sub_ds(ds = modelled, wavelength = x, scattercut = self.scattercut, drop_scatter=True, 
								bordercut = self.bordercut, ignore_time_region = self.ignore_time_region, 
								wave_nm_bin = self.wave_nm_bin, wavelength_bin = self.width)
					self.ax_kinetic.plot(ds_temp1.values,ds_temp1.index.values,'-',label='%.0f fitted'%x,color='red')
				
				ds_temp = sub_ds(ds = ds, wavelength = x, scattercut = self.scattercut, drop_scatter=True, 
								bordercut = self.bordercut, ignore_time_region = self.ignore_time_region, 
								wave_nm_bin = self.wave_nm_bin, wavelength_bin = self.width)
				
				self.ax_kinetic.set_xlim(min([0,min(ds_temp.values)]),max([max(ds_temp.values),0]))

				self.ax_kinetic.plot(ds_temp.values,ds_temp.index.values,'*',label='%.0f'%x, color='black')
				self.ax_kinetic.plot([0,0],self.ax_kinetic.get_ylim(),'gray')
				self.ax_kinetic.legend(['%.0f'%x])				
				self.ax_kinetic.set_xticks(self.ax_kinetic.get_xlim())
				self.ax_kinetic.set_xticklabels(['%.1e'%f for f in self.ax_kinetic.get_xlim()])
				self.ax_kinetic.set_yticklabels(self.ax.get_yticklabels())
				plt.subplots_adjust(wspace=0,hspace=0)
				
		eve=MouseMove(ds, cmap, self.intensity_range, self.log_scale, self.baseunit, self.timelimits, 
						self.scattercut, self.bordercut, self.wave_nm_bin, self.equal_energy_bin, self.ignore_time_region,
						self.time_bin, self.lintresh, self.data_type, self.wavelength_bin, self.time_width_percent)
		cursor = Cursor(eve.ax, useblit=True, color='red', linewidth=2)
		return eve,cursor
		


	def Plot_RAW(self, plotting = range(4), title = None, scale_type = 'symlog', times = None,
				cmap = None, filename = None, path = "result_figures", savetype = 'png' , print_click_position = False,
				plot_second_as_energy = True, ds = None, return_figures_handles=False):
		'''This is a wrapper function that triggers the plotting of various RAW (non fitted) plots. 
		The shaping parameter are taken from the object and should be defined before.
		The parameter in this plot call are to control the general look and features of the plot.
		Which plots are printed is defined byt the first command (plotting)
		The plots are generated on the fly using self.ds and all the shaping parameter
		In all plots the RAW data is plotted as dots and interpolated with lines 
		(using Savitzky-Golay window=5, order=3 interpolation). As defined by the internal parameters 
		at selected time-points and the kinetics for selected wavelength are shaped by the 
		object parameter. The SVD is performed using the same shaping parameter and is commonly 
		used as an orientation for the number of components in the data.
		Everything is handed over to 'plot_raw' function that can be used for extended RAW plotting.
		
		Parameters
		---------------
		
		plotting : int or iterable (of integers), optional
			This parameter determines which figures are plotted 
			the figures can be called separately with plotting = 1
			or with a list of plots (Default) e.g. plotting=range(4) calls plots 0,1,2,3.
			The plots have the following numbers: 
			
				0. Matrix
				1. Kinetics
				2. Spectra
				3. SVD 
			
			The plotting takes all parameter from the "ta" object.
		
		title : None or str
			title to be used on top of each plot
			The (Default) None triggers  self.filename to be used. Setting a specific title as string will.
			be used in all plots. To remove the title all together set an empty string with this command title="" .
		
		Scale_type : None or str
			is a general setting that can influences what time axis will be used for the plots. 
			"symlog" (linear around zero and logarithmic otherwise) "lin" and "log" are valid options.
		
		times : int
			are the number of components to be used in the SVD (Default) is 6.
		
		cmap : None or matplotlib color map, optional
			is a powerfull variable that chooses the colour map applied for all plots. If set to 
			None (Default) then the self.cmap is used.
			As standard I use the color map "jet" from matplotlib. There are a variety of colormaps 
			available that are very usefull. Beside "jet", "viridis" is a good choice as it is well 
			visible under red-green blindness. Other useful maps are "prism" for high fluctuations 
			or diverging color maps like "seismic". 
			See https://matplotlib.org/3.1.0/tutorials/colors/colormaps.html for a comprehensive 
			selection. In the code the colormaps are imported so if plot_func is imported as pf then 
			self.cmap=pf.cm.viridis sets viridis as the map to use. Internally the colors are chosen 
			with the "colm" function. The 2d plots require a continuous color map so if something 
			else is give 2d plots are shown automatically with "jet". For all of the 1d plots however 
			I first select a number of colors before each plot. If cmap is a continous map then these
			are sampled evenly over the colourmap. Manual iterables of colours 
			cmap=[(1,0,0),(0,1,0),(0,0,1),...] are also accepted, as are vectors or dataframes that 
			contain as rows the colors. There must be of course sufficient colors present for 
			the numbers of lines that will be plotted. So I recommend to provide at least 10 colours 
			(e.g.~your university colors). colours are always given as a, list or tuple with RGA or RGBA
			(with the last A beeing the Alpha=transparency. All numbers are between 0 and 1. 
			If a list/vector/DataFrame is given for the colours they will be used in the order provided.		
		
		filename : str, optional
			offers to replace the base-name used for all plots (to e.g.~specify what sample was used). 
			if (Default) None is used, the self.filename is used as a base name. The filename plays only a 
			role during saving, as does the path and savetype.
		
		path : None or str or path object, optional
			This defines where the files are saved if the safe_figures_to_folder parameter is True, 
			quite useful if a lot of data sets are to be printed fast. 
			If a path is given, this is used. If a string like the (Default) "result_figures" is given, 
			then a subfolder of this name will be used (an generated if necessary) 
			relative to self.path. Use and empty string to use the self.path
			If set to None, the location of the plot_func will be used and
			a subfolder with title "result_figures" be generated here.
		
		savetype : str or iterable (of str), optional 
			matplotlib allows the saving of figures in various formats. (Default) "png", 
			typical and recommendable options are "svg" and "pdf".
			
		return_figures_handles : bool, optional
			(Default) is False, if True the Figure handles are returned as a dictionary. 
			
		print_click_position : bool, optional
			if True then the click position is printed for the spectral plots 
			
		ds : DataFrame, optional
			if None (Default), the program first tests self.ds and if this is not there then self.ds_ori.
			This option was introduced to allow plotting of other matrixes with the same parameter
		
		Examples
		------------
		
		Typically one would call this function empty for an overview. We name the object "ta" so with
		
		>>> ta=pf.TA('testfile.SIA') 
		
		This would trigger the plotting of the 4 mayor plots for an overview.
		
		>>> ta.Plot_RAW()
		
		This would plot only the kinetics.
		
		>>> ta.Plot_RAW(1)
		>>> ta.Plot_RAW(plotting = 1)
 
		'''
		
		path=check_folder(path=path,current_path=self.path)
		if self.save_figures_to_folder:
			self.figure_path=path
		if cmap is None:cmap=self.cmap
		if ds is None:
			if self.ds is None:
				ds=self.ds_ori.copy()
			else:
				ds=self.ds.copy()
		if filename is None: filename=self.filename
		if not hasattr(plotting,"__iter__"):plotting=[plotting]
		if title is None:
			if filename is None:
				title=self.filename
			else:
				title=filename
		r=plot_raw(ds=ds, plotting=plotting, cmap=cmap, title=title, path=path, filename=filename, 
				intensity_range=self.intensity_range, log_scale=self.log_scale, baseunit=self.baseunit, 
				timelimits=self.timelimits, scattercut=self.scattercut, bordercut=self.bordercut, 
				wave_nm_bin=self.wave_nm_bin, rel_wave=self.rel_wave, width=self.wavelength_bin, 
				time_width_percent=self.time_width_percent, ignore_time_region=self.ignore_time_region, 
				time_bin=self.time_bin, rel_time=self.rel_time, save_figures_to_folder=self.save_figures_to_folder, 
				savetype=savetype,plot_type=scale_type,lintresh=self.lintresh, times=times, 
				print_click_position = print_click_position, data_type = self.data_type, 
				plot_second_as_energy = plot_second_as_energy, units=self.units, equal_energy_bin = self.equal_energy_bin,
				return_figures_handles=return_figures_handles)
		if return_figures_handles:
			return r


	def Save_Plots(self, path = 'result_figures', savetype = None, title = None, filename = None, scale_type = 'symlog', 
					patches = False, cmap = None):
		'''Convenience function that sets save_plots_to_folder temporarily to true and replots everything
		
		Parameters
		----------
		
		path : None, str or path, optional
			(Default) None, if left on None, then a folder "result_figures" is created in the folder 
			of the data (self.path)
		
		savetype : str or iterable (of str), optional 
			matplotlib allows the saving of figures in various formats. (Default) "png", 
			typical and recommendable options are "svg" and "pdf". 
			
		title : None or str, optional
			(Default) None, Use this title on all plots. if None, use self.filename
			
		filename : str, optional
			(Default) None, Base name for all plots. If None, then self.filename will be used
			
		scale_type : str, optional
			"symlog" (Default), "linear", "log" time axis
			
		patches : bool, optional
			For true use white patches to label things in the 2d matrixes, to safe
			space for publication
			
		cmap : None or matplotlib color map, optional
			is a powerfull variable that chooses the colour map applied for all plots. If set to 
			None (Default) then the self.cmap is used.
			As standard I use the color map "jet" from matplotlib. There are a variety of colormaps 
			available that are very usefull. Beside "jet", "viridis" is a good choice as it is well 
			visible under red-green blindness. Other useful maps are "prism" for high fluctuations 
			or diverging color maps like "seismic". 
			See https://matplotlib.org/3.1.0/tutorials/colors/colormaps.html for a comprehensive 
			selection. In the code the colormaps are imported so if plot_func is imported as pf then 
			self.cmap=pf.cm.viridis sets viridis as the map to use. Internally the colors are chosen 
			with the "colm" function. The 2d plots require a continuous color map so if something 
			else is give 2d plots are shown automatically with "jet". For all of the 1d plots however 
			I first select a number of colors before each plot. If cmap is a continous map then these
			are sampled evenly over the colourmap. Manual iterables of colours 
			cmap=[(1,0,0),(0,1,0),(0,0,1),...] are also accepted, as are vectors or dataframes that 
			contain as rows the colors. There must be of course sufficient colors present for 
			the numbers of lines that will be plotted. So I recommend to provide at least 10 colours 
			(e.g.~your university colors). colours are always given as a, list or tuple with RGA or RGBA
			(with the last A beeing the Alpha=transparency. All numbers are between 0 and 1. 
			If a list/vector/DataFrame is given for the colours they will be used in the order provided.
		
		Examples
		---------
		
		>>> ta.Save_Plots()
		>>> ta.Save_Plots(patches = True)
		
		'''
		if cmap is None:cmap=self.cmap
		if savetype is None:
			savetype=['png']
		elif savetype in ['png','pdf','svg']:savetype=[savetype]
		elif hasattr(savetype,"__iter__"):savetype=list(savetype)
		else:
			print('Please specify a single filetype from \'png\',\'pdf\',\'svg\' or a list of those. Nothing was saved')
			return False
		if cmap is None:cmap=standard_map
		origin=self.save_figures_to_folder
		self.save_figures_to_folder=True
		try:
			for t in savetype:
				plt.close('all')
				self.Plot_RAW(savetype = t, path = path, cmap = cmap, title = title, 
								scale_type = scale_type, filename = filename, units=self.units, 
								equal_energy_bin = self.equal_energy_bin)
				plt.close('all')
				print('saved RAW plots type %s to %s'%(t,check_folder(path=path,current_path=self.path)))
		except:
			print('Saving of Raw plots for filetype %s failed'%t)
		try:
			for t in savetype:
				plt.close('all')
				self.Plot_fit_output(savetype=t,path=path,cmap=cmap,title=title,scale_type=scale_type,patches=patches,filename=filename)
				plt.close('all')
				print('Saved Fit plots of type %s to %s'%(t,check_folder(path=path,current_path=self.path)))
		except:
			print('Saving of Fit plots for filetype %s failed'%t)
		self.save_figures_to_folder=origin


	def __Fit_Chirp_inner(	self, opt_coeff, initial_fit_coeff = None, params = None, scattercut = None, bordercut = None, 
							timelimits = None, wave_nm_bin = None, time_bin = None, mod = None, log_fit = None, ds_back_corr = None):
		''' Function to calculate a new chirp corrected matrix and return an error value,  The "cost function" for the chirp optimization 
		'''
		fitcoeff = np.array([opt_coeff['p4'].value, opt_coeff['p3'].value, opt_coeff['p2'].value, opt_coeff['p1'].value, opt_coeff['p0'].value])
		#fitcoeff = __shift_polynom(fitcoeff, -500) #this was an arbitrary shift of the data by 500 to make the chirp parameter 
		time = ds_back_corr.index.values.astype('float')#extract the time
		ds_new = ds_back_corr.apply(lambda x:np.interp(x = time+np.polyval(fitcoeff, x.name), xp = time, fp = x), axis = 0, raw = False)
		re = err_func(paras = params, ds = ds_new, mod = mod, final = False, log_fit = log_fit)
		return re
		
		
	def __Fit_Chirp_outer(self, pardf, results, fit_ds, fit_chirp_iterations, mod, deep_iteration = False):
		'''Broken out Chirp optimization,  takes the fitted parameters and performs 'fit_chirp_iterations' times the loop,  
		(optimise chirp + optimize global) after each global iteration the error is compared to the previous. It continues until no improvement is made or 
		until the 'fit_chirp_iterations' is reached. If the error is reduced by more than a factor of 100 in a single step,  it is assumed that something fishy is going on and we restart the fit,  but with a 10x smaller simplex stepsize and deep_iteraction FAlse
		
		Parameters
		-----------
		pardf
					# deep_iteration uses the previous kinetic optimized parameter as the input into the next fit. 
			#Can be great but can also run away, in general not needed and can be triggered by feeding the 
			#results back into the global fit'''
		if pardf.vary.any():
			initial_error = [results.residual[0]]
			par_into_chirpfit = results.params
		else:
			initial_error = [err_func(paras = self.par_fit, ds = fit_ds, mod = self.mod, final = False, log_fit = self.log_fit)]
			par_into_chirpfit = self.par
		par_into_chirpfit['t0'].vary = False
		initial_fit_coeff = self.fitcoeff
		if len(initial_fit_coeff) == 6:
			initial_fit_coeff[4] = self.fitcoeff[4]+self.fitcoeff[5]
			initial_fit_coeff = initial_fit_coeff[:5]
		chirp_par = lmfit.Parameters()
															   
		for i in range(5):
			chirp_par.add('p%i'%(4-i), value = initial_fit_coeff[i])
		chirp_par['p4'].set(min = chirp_par['p4']-0.5, max = chirp_par['p4']+0.5)
		try:#lets send in the background corrected matrix fails if no prior background was done
			correction = self.background_par[3]
			ds_back_corr = self.ds_ori-correction
		except:
			ds_back_corr = self.ds_ori
		ds_back_corr = sub_ds(ds = ds_back_corr, scattercut = self.scattercut, bordercut = self.bordercut, 
								timelimits = self.timelimits, wave_nm_bin = self.wave_nm_bin, time_bin = self.time_bin, 
								equal_energy_bin = self.equal_energy_bin)																	   
		print('Before chirpfit the error is:{:.6e}'.format(initial_error[-1]))
		
		#################################################################################################################
		#----Chirp fit loop---------------------------------------------------------------------------------
		#################################################################################################################
		
		for loop in range(fit_chirp_iterations):
			
			chirpmini  =  lmfit.Minimizer(self.__Fit_Chirp_inner, chirp_par, 
											fcn_kws = {'ds_back_corr':ds_back_corr.copy(), 'initial_fit_coeff':initial_fit_coeff, 
														'params':par_into_chirpfit, 'mod':self.mod, 'log_fit':self.log_fit, 
														'scattercut':self.scattercut, 'bordercut':self.bordercut, 
														'timelimits':self.timelimits, 'wave_nm_bin':self.wave_nm_bin, 
														'time_bin':self.time_bin}) 
			step_size = 5e-2
			try:
				start  =  tm.time()
				simp = np.array([chirp_par['p4'].value, chirp_par['p3'].value, chirp_par['p2'].value, chirp_par['p1'].value, chirp_par['p0'].value])
				simp = np.tile(simp.reshape(5, 1), 6).T
				for i in range(5):
					if simp[i+1, i] != 0:
						if i<4:
							simp[i+1, i] = simp[i+1, i]*(step_size)
						else:
							simp[i+1, i] = simp[i+1, i]+0.1
					else:
						simp[i+1, i] = 1e-4	
				#we start by optimizing the chirp with fixed Global fit
				chirp_results  =  chirpmini.minimize('nelder',  options = {'maxfev':1e4, 'fatol':initial_error[-1]*1e-6, 'initial_simplex':simp})
				end  =  tm.time()
				opt_coeff = chirp_results.params
				temp = np.array([opt_coeff['p4'].value, opt_coeff['p3'].value, opt_coeff['p2'].value, opt_coeff['p1'].value, opt_coeff['p0'].value])
				#Create the new chirp corrected data
				time = ds_back_corr.index.values.astype('float')#extract the time		 
				new_ds = ds_back_corr.copy().apply(lambda x:np.interp(x = time+np.polyval(temp, float(x.name)), xp = time, fp = x), axis = 0, raw = False)
				#New Global Fit
				fit_ds_loop = sub_ds(ds = new_ds, scattercut = self.scattercut, bordercut = self.bordercut, timelimits = self.timelimits, wave_nm_bin = self.wave_nm_bin, equal_energy_bin = self.equal_energy_bin, time_bin = self.time_bin)	
				if pardf.vary.any():							 
					mini  =  lmfit.Minimizer(err_func, par_into_chirpfit, fcn_kws = {'ds':fit_ds_loop, 'mod':mod, 'log_fit':self.log_fit, 'final':False})
					results_in_chirp  =  mini.minimize('nelder', options = {'maxiter':1e5})
					initial_error.append(results_in_chirp.residual[0])
				else:
					initial_error.append(err_func(paras = par_into_chirpfit, ds = fit_ds_loop, mod = mod, final = False, log_fit = self.log_fit))	
				if initial_error[-1]<initial_error[-2]:
					self.opt_chirp=temp
					if initial_error[-2]/initial_error[-1]>100:#something fishy going on. lets try again
						print('Chirp_loop {:02d} strange decrease step size'.format(loop+1))
						initial_error[-1] = initial_error[-2]
						step_size = step_size/100
						if len(initial_error)>4:
							if initial_error[-4] == initial_error[-1]:#we have run this trick now three times,  time to break
								raise StopIteration
						deep_iteration=False
					else:
						print('Chirp_loop {:02d} resulted in   :{:.8e}'.format(loop+1, initial_error[-1]))
						if deep_iteration:			#This results in a very deep iteration of the starting parameter
							if pardf.vary.any():
								par_into_chirpfit = results_in_chirp.params				
						chirp_par = chirp_results.params
				else:
					raise StopIteration
			except StopIteration:
				print('iteration is not smaller finished chirp looping')
				break
			except:
				print('failure in chirp optimisation in iteration %i'%(loop+1))
				import sys
				print("Unexpected error:",  sys.exc_info()[0])
				initial_error.append(initial_error[0])#to avoid that numbers are written
				break
		
		#################################################################################################################
		#-----------------------------------------end chrip fit loop-------------------------------------------------
		#################################################################################################################
		
		if initial_error[-1]<initial_error[0]:#lets check if we improved anything
			print('chirp fit improved error by %.2g percent'%(100*(1-initial_error[-1]/initial_error[0])))
			if isinstance(temp, list) or isinstance(temp, type(np.arange(1))):
				self.fitcoeff = temp
				
			else:
				raise
			time = ds_back_corr.index.values.astype('float')#extract the time
			self.ds = ds_back_corr.apply(lambda x:np.interp(x = time+np.polyval(temp, float(x.name)), xp = time, fp = x), axis = 0, raw = False)			 
			fit_ds = sub_ds(ds = self.ds, scattercut = self.scattercut, bordercut = self.bordercut, timelimits = self.timelimits, wave_nm_bin = self.wave_nm_bin, equal_energy_bin = self.equal_energy_bin, time_bin = self.time_bin)
			if pardf.vary.any():
				results.params = results_in_chirp.params
			
		return results, fit_ds
		
		
	def Fit_Global(self, par = None, mod = None, confidence_level = None, use_ampgo = False, other_optimizers=None, fit_chirp = False,
						fit_chirp_iterations = 10, multi_project = None, unique_parameter = None, weights = None, same_DAS = False,
						dump_paras = False, dump_shapes = False, filename = None, ext_spectra = None,
					    write_paras=False, tol = 1e-5, sub_sample=None,pulse_sample=None,same_shape_params=True):
		"""This function is performing a global fit of the data. As embedded object it uses 
		the parameter control options of the lmfit project as an essential tool. 
		(my thanks to Matthew Newville and colleagues for creating this phantastic tool) 
		[M. Newville, T. Stensitzki, D. B. Allen, A. Ingargiola, 2014. DOI: 10.5281/ZENODO.11813.]. 
		The what type of fitting is performed is controlled by setting of the parameter here.
		
		The general fitting follows this routine:
			1. 	create a copy of the Data-Matrix self.ds is created with the shaping parameters
			2. 	Then a Matrix is created that represents the fractional population of each species 
				(or processes in case of the paral model). 
				This Matrix contains one entry for each timepoint and represents the kinetic model 
				based upon the starting parameter. (see below for a description of the models). 
				This model formation can by done by using a build in or a user supplied function. 
				(handled in the function "pf.build_c")
				-> If an ext_spectra is provided this its intensity is substacted from the matrix (only for external models)
			3. 	Then the process/species associated spectra for each of the species is calculated 
				using the linalg.lstsq algorithm from numpy 
				(https://numpy.org/doc/stable/reference/generated/numpy.linalg.lstsq.html)
			4. 	From the convoluted calculated species concentrations and spectra a calculated matrix 
				is formed (handled in the function "pf.fill_int")
			5. 	The difference between calculated and measured spectra is calculated, point-wise squared 
				and summed together. (function "err_func" or "err_func_multi" if multiple datasets are fitted)
			6. 	This difference is minimized by iterating 2-4 with changing parameters using an 
				optimization algorithm (generally nelder-mead simplex)
			7. 	Finally in a last run of 2-5 the final spectra are calculated (using the "final" flag) 
				and the optimized parameter, the matrixes 
				("A"-measured, "AC" - calculated, "AE" - linear error), 
				spectra (always called "DAS") the concentrations (called "c") 
				are written in the dictionary "ta.re" together with a few result representations 
				and other fit outputs. The optimized parameter are also written into ta.par_fit 
				(as an parameter object) that can be re-used as input into further optimization steps.
			8.	Under Windows we load the keyboard library and the Fit can be interrupted by pressing 
				the "ctrl+shift+q" key simultaneously (hardcoded). Consider using the parameter write_paras or 
				dump_paras to observe details during the fit.
				
		All mandatory parameters are in general taken from the internal oject (self) The optional parameter control the behaviour of the fitting function  
		
		Parameters
		------------------
			
			par : lmfit parameter oject, optional
				Here another parameter object could be given,overwriting the (Default is self.par) 
				
			mod : str or function, optional
				Give a extra model selection (Default uses self.mod)
				internal modells: 'paral','exponential','consecutive','full_consecutive'
				see also :meth:`plot_func.build_c` and :meth:`plot_func.err_func`
				
			confidence_level: None or float (0.5-1), optional
				If this is changed from None (Default) to a value between 0.5 and 1 the code will
				try to calculate the error of the parameter for the fit. For each parameter that 
				can vary a separate optimization is	performed, that attempts to find the upper 
				and lower bound at which the total error of the re-optimized globally fitted results
				reaches the by F-statistics defined confidence bound. See :meth:`plot_func.s2_vs_smin2` for details 
				on how this level is determined. Careful, this option might run for very long time.
				Meaning that it typically takes 50 optimization per variable parameter (hard coded limit 200)
				The confidence level is to be understood that it defines the e.g. 0.65 * 100\% area that the
				parameter with this set of values is within this bounds.
				Normal behaviour for this is to re-optimise the parameter during the optimization. if the parameter 
				par['error_param_fix'] is present, this will be suppressed.
			
			use_ampgo : bool, optional
				(Default) is False
				Changes the optimizer from a pure Nelder mead to Ampgo with a local Nelder Mead. 
				For using this powerfull tool all parameter need to have a "min" and a "max" set. 
				Typically takes 10-40x longer than a standard optimization, but can due to its 
				tunneling algorithm more reliably find global minima. 
				see:https://lmfit.github.io/lmfit-py/fitting.html for further details
			
			other_optimizers : str, optional
				(Default) is None
				if this is changed from None to a string that exists in lmfit, then this optimizer will be used instead. 
				Useful choices are e.g. **least_squares** or similar words. This is particularly useful if the 
				problem does not lend to be solved with nelder-mead. This includes e.g. osciallations.
			
			fit_chirp : bool, optional
				(Default) is False
				a powerful optimization of the chirp parameter. For this to work the data 
				needs to include timepoints before and after t=0 and one should have reached 
				a decent fit of most features in the spectrum. We perform an Nelder-Mead optimisation 
				of the parameter followed by a Nelder-Mead optimization of the chirp parameter 
				as one iteration. After each consecutive optimization it is checked if the total error 
				improved. If not the fit is ended, if yes the maximum number of iterations 
				'fit_chirp_iterations' is performed. Warning, this goes well in many cases, 
				but can lead to very strange results in others, always carefully check the results.
				I recommend to make a copy of the object before runnning a chirp optimization.
			
			fit_chirp_iterations : int, optional
				maximum number of times the global - chirp loop is repeated. 
				Typically this iterations run 2-5 times, (Default) is 10
			
			dump_paras : bool, optional
				(Default) is False, If True creates two files in the working folder, one with the 
				currently used parameter created at the end of each optimisation step, and one with
				the set of parameter that up to now gave the lowest error. Intented to store
				the optimisation results if the fit needs to be interrupted 
				(if e.g. Ampgo simply needs to long to optimize.) useful option if things are slow
				this parameter also triggers the writing of fitout to a textfile on disc
				
			dump_shapes : bool, optional
				this option dumps the concentratoin matrix and the DAS onto disk for each round of optimization,
				mostly useful for multi-project fitting that wants to use the spectral or temporal intensity
			
			write_paras : bool, optional 
				if True(Default) writes the currently varried values to screen
			
			filename : None or str, optional
				Only used in conjunction with 'dump_paras'. The program uses this filename to dump the 
				parameter to disk 
			
			multi_project : None or list (of TA projects), optional
				This switch is triggering the simultaneous optimisation of multiple datasets.  
				multi_project is as (Default) None. it expects an iterable (typically list) with other 
				TA projects (like ta) that are then optimised with the same parameter. 
				This means that all projects get the same parameter object for each iteration 
				of the fit and return their individual error, which is summed linearly. 
				The "weights" option allows to give each multi_project a specific weight (number) 
				that is multiplied to the error. If the weight object has the same number of items 
				as the multi_project it is assumed that the triggering object (the embedded project) 
				has the weight of 1, otherwise the first weight is for the embedded project. 
				The option 'unique_parameter' takes (a list) of parameter that are not 
				to be shared between the projects (and that are not optimized either) 
				The intended use of this is to give e.g. the pump power for multiple experiments to 
				study non linear behaviour. Returned will be only the parameter set for the optimium
				combination of all parameter. Internally, we iterate through the projects and calculate 
				for each project the error for each iteration. Important to note is that currently this 
				means that each DAS/SAS is calculated independently! For performing the same calculation
				with a single DAS, the Matrixes need to be concatenated before the run and an external
				function used to create a combined model. As this is very difficult to implement reliably
				For general use (think e.g. different pump wavelength) this has to be done manually. 
			
			unique_parameter : None or str or list (of strings), optional 
				only used in conjunction with 'multi_project', it takes (a list) of parameter that 
				are not to be shared between the projects (and that are not optimized either) 
				The intended use of this is to give e.g. the pump power for multiple experiments 
				to study non linear behaviour. (Default) None
			
			same_DAS : bool,optional
				changes the fit behavior and uses the same DAS for the optimization. 
				This means that the ds are stacked before the fill_int rounds. This option is only used in multi-project fitting
			
			weights : list of floats, optional
				only used in conjunction with 'multi_project'. The "weights" option allows to 
				give each multi\_project a specific weight (number) that is multiplied to the error. 
				If the weight object has the same number of items as the 'multi_project' it is assumed 
				that ta (the embedded project) has the weight of 1, otherwise the first weight is for the 
				embedded object
				
			ext_spectra : DataFrame, optional
				(Default) is None, if given substract this spectra from the DataMatrix using the intensity 
				given in "C(t)" this function will only work for external models. The name of the spectral column 
				must be same as the name of the column used. If not the spectrum will be ignored. The spectrum will 
				be interpolated to the spectral points of the model ds before the substraction.
				a number of parameters can be defined to aid this process. These parameter are defined as normal parameters.
				"ext_spectra_scale" multiplies all spectra by this value (e.g. -1 to put the steady state absorption spectra in)
				"ext_spectra_shift" shifts all spectra by this value to compensate for calibration differences
				"ext_spectra_guide" (from version 7.1.0) This is a switch, if this keyword is present, then the spectra are 
				used as guides and not exclusively. This means the code will assume that these spectra are correct and substract 
				them, then calulate the difference and return as DAS the provided spectra plus the difference spectra
				
			tol : float, optional
				the tolerance value that is handed to the optimizer (absolute) for nelder-mead the moment this means:
				df < tol  (corresponds to fatol)
				number_of_function_evaluations < maxfev (default 200 * n variables)
				number_of_iterations < maxiter           (default 200 * n variables)
			
			sub_sample: False or iter, option
				Default is False: if a whole number value is given the time vector that comes from the measurements is sampled finer. The purpose is to in 
				the model catch faster dynamics than would be possible. Careful this is bloating up the size of the time vector and the simulation time is 
				to some extend proportional to this number
			
			pulse_sample: False or iter, option
				Default is False, but will be activated if the zero time is not included in the fit. If True then additional points in the 
				pump_region=np.linspace(t0-4*resolution,t0+4*resolution,20) are added.  This is better than sub sample, as it only adds 20 time points 
				to the data. This was necessary since otherwise the pulse (that is creating the intensity) is not sampled at all or only fractional. 
				This switch mainly has an influence on the absolute intensities of the species (the concentration matrix)

			same_shape_params: bool, optional
				This parameter decides if the same shaping parameter (scattercut,rebinning are used in the multiprojects
				This should be set to False for multimodal fitting.

		Returns
		------------------
		
		re : dict
			the dictionary "re" attached to the object containing all the matrixes and parameter. 
			The usual keys are:
			"A" Shaped measured Matrix
			"AC" Shaped calculated Matrix 
			"AE" Difference between A and AC = linear error 
			"DAS" DAS or SAS, labeled after the names given in the function (the columns of c) Care must be taken that this mesured intensity is C * DAS, the product. For exponential model the concentrations are normalized
			"c" The Concentrations (meaning the evolution of the concentrations over time. Care must be taken that this mesured intensity is C * DAS, the product. For exponential model the concentrations are normalized
			"fit_results_rates" DataFrame with the fitted rates (and the confidence intervals if calculated)
			"fit_results_times" DataFrame with the fitted decay times (and the confidence intervals if calculated)
			"fit_output" The Fit object as returned from lmfit. (This is not saved with the project!)
			"error" is the S2, meaning AE**2.sum().sum()
			"r2"=1-"error"/(('A'-'A'.mean())**2).sum(), so the residuals scaled with the signal size
		
		par_fit : lmfit parameter object
			is written into the object as a lmfit parameter object with the optimized results (that can be use further)
		
		fitcoeff : list, if chirpfit is done
			The chirp parameter are updated
		
		ds : DataFrame, if chirpfit is done
			A new ds is calculated form ds_ori if ChripFit is done
			
		The rest is mainly printed on screen.
		
		
		Examples
		--------------------
		
		Non optional:
		
		>>> ta=pf.TA('testfile.SIA') #load data
 		>>> ta.mod='exponential'    #define model
		>>> ta.par=lmfit.Parameters()  #create empty parameter object
		>>> ta.par.add('k0',value=1/0.1,vary=True) #add at least one parameter to optimize
		  
		Trigger simple fit:
		
		>>> ta.Fit_Global()
		
		Trigger fit with Chrip Fit:
		
		>>> ta.Fit_Global(fit_chirp=True)
		
		Trigger iterative Chirp fitting with fresh refinement of the Global kinetic parametersfor i in range(5):
		
		>>> for i in range(5):
		>>> 	start_error=ta.re['error']
		>>> 	ta.par=ta.par_fit
		>>> 	ta.Fit_Global(fit_chirp=True)
		>>> 	if not ta.re['error'] < start_error:break
		
		Trigger fit fit error calculations
		
		>>> ta.Fit_Global(confidence_level=0.66)
		
		Trigger fit of multiple projects
		#use the GUI_open function to open a list of objects (leave empty for using the GUI)
		
		>>> other_projects=pf.GUI_open(['sample_1.hdf5','sample_2.hdf5'],path='Data')
		>>> ta.Fit_Global(multi_project=other_projects)
		
		
		For more examples please see the complete documentation under :ref:`Fitting, Parameter optimization and Error estimation`
		or :ref:`Fitting multiple measured files at once`
		"""
		if par is None:par=self.par
		if mod is None:mod=self.mod
		try:
			t0=par['t0']
		except:
			try:
				par.add('t0',value=0,min=-0.5,max=0.5,vary=False)
			except:
				print("Unexpected error:", sys.exc_info()[0])
		try:
			resolution=par['resolution']
		except:
			try:
				par.add('resolution',value=0.086,min=0.04,max=0.5,vary=False)
			except:
				print("Unexpected error:", sys.exc_info()[0])
		try:
			par['infinite'].value=1
			par['infinite'].vary=False
		except:
			pass
		try:
			par['background'].value=1
			par['background'].vary=False
		except:
			pass
		try:
			par['explicit_GS'].value=1
			par['explicit_GS'].vary=False
		except:
			pass
			
		try: # this is either freezing or enabling the re-optimization of all other parameter during confidence interval calculation
			par['error_param_fix'].value=1
			par['error_param_fix'].vary=False
			vary_error_parameter=False
		except:
			vary_error_parameter=True
		
		pardf=par_to_pardf(par)
		pardf.loc[np.logical_and(pardf.loc[:,'min'].values<0,pardf.is_rate),'min']=0
		pardf.loc[np.logical_and(pardf.loc[:,'max'].values<0,pardf.is_rate),'max']=0

		pardf['init_value']=pardf['value']
		if dump_paras:
			pardf_temp=pardf.copy()
			pardf_temp.loc['error','value']=1000
			pardf_temp.to_csv('minimal_dump_paras.par')
	
		if self.log_fit:
			for key in ['value','min','max']:
				pardf.loc[pardf.is_rate,key]=pardf.loc[pardf.is_rate,key].apply(lambda x: np.log10(x))	
		
		#create-shape the data to be fitted	
		fit_ds = sub_ds(ds = self.ds.copy(), scattercut = self.scattercut, bordercut = self.bordercut, 
						timelimits = self.timelimits, wave_nm_bin = self.wave_nm_bin, equal_energy_bin = self.equal_energy_bin, 
						time_bin = self.time_bin, ignore_time_region = self.ignore_time_region, drop_scatter = True, drop_ignore = True)
		time_label=fit_ds.index.name
		energy_label=fit_ds.columns.name
		if pulse_sample is None:
			if self.ignore_time_region is not None:
				pulse_sample=True
			if self.timelimits is not None:
				if min(self.timelimits)>0:
					pulse_sample=True
		
		############################################################################
		#----Global optimisation------------------------------------------------------
		############################################################################
		global start_time
		start_time=start_time-30
		try:
			keyboard.__package__
			def iter_cb(params, iterative, resid, ds=None,mod=None,log_fit=None,final=None,dump_paras=None,
						filename=None,ext_spectra=None,dump_shapes=None, 
						write_paras=None,multi_project=None,unique_parameter=None,
						weights=None,same_DAS=None,sub_sample=None,pulse_sample=None,same_shape_params=None):
				if keyboard.is_pressed("ctrl+shift+q"):
					print('---------------------------------------------')
					print('---------  Interupted by user          ------')
					print('---------------------------------------------')
					print('-----------   Last fitted parameter    ------')
					print(par_to_pardf(params))
					return True
				else: 
					return None
		except:
			def iter_cb(params, iterative, resid, ds=None,mod=None,log_fit=None,final=None,dump_paras=None,filename=None,ext_spectra=None,dump_shapes=None, 
												write_paras=None,multi_project=None,unique_parameter=None,
												weights=None,same_DAS=None,sub_sample=None,pulse_sample=None,same_shape_params=None):
				return None
		print('Optimizing, after the starting error the new error values will be displayed every 10s \n to interrupt try to press ctrl+shift+q')
		if multi_project is None:
			#check if there is any concentration to optimise
			if (filename is None) and dump_shapes: filename = self.filename
			if pardf.vary.any():#ok we have something to optimize
				mini = lmfit.Minimizer(err_func,pardf_to_par(pardf),iter_cb=iter_cb,
										fcn_kws={'ds':fit_ds,'mod':mod,'log_fit':self.log_fit,'final':False,
												'dump_paras':dump_paras,'filename':filename,'ext_spectra':ext_spectra,
												'dump_shapes':dump_shapes, 'write_paras':write_paras,'sub_sample':sub_sample,
												'pulse_sample':pulse_sample})
				if other_optimizers is None:
					if not use_ampgo:
						if len(pardf[pardf.vary].index)>3:
							print('we use adaptive mode for nelder')
							#results = mini.minimize('nelder',options={'adaptive':True,'fatol':tol})
							results = mini.minimize('nelder',tol=tol,options={'adaptive':True})												 
						else:
							results = mini.minimize('nelder',tol=tol)
							#results = mini.minimize('nelder',options={'fatol':tol})
					else:
						results = mini.minimize('ampgo',**{'local':'Nelder-Mead'})
				else:
					results = mini.minimize(other_optimizers)
					
		############################################################################
		#----Multi project	Global optimisation----------------------------------------
		##########################################################################
		
		else:
			fit_chirp=False #chirp fitting currently only works for single problems
			if pardf.vary.any():#ok we have something to optimize lets return the spectra
				multi_project.insert(0,self)
				mini = lmfit.Minimizer(err_func_multi,pardf_to_par(pardf),iter_cb=iter_cb,
										fcn_kws={'multi_project':multi_project,'unique_parameter':unique_parameter,
										'weights':weights,'mod':mod,'log_fit':self.log_fit,'final':False,
										'dump_paras':dump_paras,'filename':filename,'ext_spectra':ext_spectra,
										'dump_shapes':dump_shapes,'write_paras':write_paras,'same_DAS':same_DAS,'sub_sample':sub_sample,
										'pulse_sample':pulse_sample,'same_shape_params':same_shape_params})
				if other_optimizers is None:
					if len(pardf[pardf.vary].index)>3:
						print('we use adaptive mode for nelder')
						results = mini.minimize('nelder',options={'adaptive':True,'fatol':tol})
					else:
						results = mini.minimize('nelder',options={'fatol':tol})
				else:
					results = mini.minimize(other_optimizers)

		#######################################################################
		#----Fit chirp----------------------------------------------------------------------------------
		####################################################################
		
		if self.ignore_time_region is not None: 
			if fit_chirp:
				print('sorry but currently you can not both ignore a time region and fit the chirp (assuming that you ignore the time-zero region)')
			fit_chirp=False
		if fit_chirp:
			print('Done initial fitting now chirpfit')
			results,fit_ds=self.__Fit_Chirp_outer(pardf,results,fit_ds,fit_chirp_iterations,mod)
			
			
		####################################################################
		#------Write results to parameter------------------------
		############################################################
		
		if pardf.vary.any():#we actually have optimised something
			pardf['value']=par_to_pardf(results.params)['value']
			if self.log_fit:
				for key in ['value','min','max']:
					pardf.loc[pardf.is_rate,key]=pardf.loc[pardf.is_rate,key].apply(lambda x: 10**x)
			self.par_fit=pardf_to_par(pardf)
		else:
			print('ATTENTION: we have not optimized anything but just returned the parameters')
			self.par_fit=self.par
		if multi_project is None:
			re=err_func(paras=self.par_fit,ds=fit_ds,mod=self.mod,final=True,log_fit=self.log_fit,ext_spectra=ext_spectra,sub_sample=sub_sample,pulse_sample=pulse_sample)
		else:
			if same_DAS:
				re_listen = err_func_multi(paras = self.par_fit, mod = mod, final = True, log_fit = self.log_fit, 
									multi_project = multi_project, unique_parameter = unique_parameter, same_DAS = same_DAS, weights = weights, 
									ext_spectra = ext_spectra,sub_sample=sub_sample,pulse_sample=pulse_sample,same_shape_params=same_shape_params)
				re=re_listen[0]
			else:
				re = err_func_multi(paras = self.par_fit, mod = mod, final = True, log_fit = self.log_fit, 
									multi_project = multi_project, unique_parameter = unique_parameter, same_DAS = same_DAS, weights = weights, 
									ext_spectra = ext_spectra,sub_sample=sub_sample,pulse_sample=pulse_sample,same_shape_params=same_shape_params)
		
		############################################################################
		#----Estimate errors---------------------------------------------------------------------
		############################################################################
		
		if confidence_level is not None:#ok we calculate errors to the level of the confidence_level
			if self.log_fit:
				for key in ['value','min','max']:
					pardf.loc[pardf.is_rate,key]=pardf.loc[pardf.is_rate,key].apply(lambda x: np.log10(x))
			if pardf.vary.any():#we actually have optimised something
				if (0.6 < confidence_level < 1) or (1 < confidence_level < 0.6):
					if multi_project is None:
						target=s2_vs_smin2(Spectral_points=len(re['A'].columns),Time_points=len(re['A'].index),number_of_species=len(re['DAC'].columns),fitted_kinetic_pars=len(pardf[pardf.vary].index),target_quality=confidence_level)
					else:
						multi_project.insert(0,self)
						# we assume that we have the same number of spectal points but are stacking the times
						total_time_points=np.array([len(t.re['A'].index) for t in multi_project]).sum()					
						target=s2_vs_smin2(Spectral_points=len(re['A'].columns),Time_points=total_time_points,number_of_species=len(re['DAC'].columns),fitted_kinetic_pars=len(pardf[pardf.vary].index),target_quality=confidence_level)
					target_s2=re['error']*target
					list_of_variable_parameter=pardf[pardf.vary].index.values
					conf_limits={}
					iterative_calls=0
					for fixed_par in list_of_variable_parameter:
						conf_limits[fixed_par]={'upper':None,'lower':None}
						for i in ['lower','upper']:
							print('Trying to find %s, %s confidence limit'%(fixed_par,i))
							pardf_local=self.par_fit.copy()
							pardf_local[fixed_par].vary=False
							par_local=lmfit.Parameters()
							if 'lower' in i:#go below min
								if par_to_pardf(pardf_local).loc[fixed_par,'is_rate']:
									par_local.add(fixed_par,value=pardf_local[fixed_par].value*0.95,min=0,max=pardf_local[fixed_par].value,vary=vary_error_parameter)
								else:
									par_local.add(fixed_par,value=pardf_local[fixed_par].value*0.95,max=pardf_local[fixed_par].value,vary=vary_error_parameter)
							else: #go above min
								par_local.add(fixed_par,value=pardf_local[fixed_par].value*1.05,min=pardf_local[fixed_par].value,vary=vary_error_parameter)
							
							def sub_problem(par_local,varied_par,pardf_local,fit_ds=None,mod=None,log_fit=None,multi_project=None,unique_parameter=None,weights=None,target_s2=None,ext_spectra=None,same_DAS=False,sub_sample=None,pulse_sample=None):
								pardf_local[varied_par].value=par_local[varied_par].value
								if par_to_pardf(pardf_local).vary.any():
									if multi_project is None:
										mini_sub = lmfit.Minimizer(err_func,pardf_local,fcn_kws={'ds':fit_ds,'mod':mod,'log_fit':log_fit,'ext_spectra':ext_spectra,'sub_sample':sub_sample,'pulse_sample':pulse_sample, 'dump_paras':False})
									else:
										mini_sub = lmfit.Minimizer(err_func_multi,pardf_local,fcn_kws={'multi_project':multi_project,'unique_parameter':unique_parameter,'weights':weights,'same_DAS':same_DAS,'mod':mod,'log_fit':log_fit,'ext_spectra':ext_spectra,'sub_sample':sub_sample,'pulse_sample':pulse_sample,'same_shape_params':same_shape_params})
									if other_optimizers is None:
										if len(pardf[pardf.vary].index)>3:
											results_sub = mini_sub.minimize('Nelder',options={'xatol':0.01,'adaptive':True})
										else:
											results_sub = mini_sub.minimize('Nelder',options={'xatol':0.01})
									else:
										results_sub = mini_sub.minimize(other_optimizers)
									local_error=(results_sub.residual[0]-target_s2)**2
									return local_error
								else:
									if multi_project is None:
										return err_func(pardf_local,ds=fit_ds,mod=mod,log_fit=log_fit,ext_spectra=ext_spectra)
									else:
										return err_func_multi(pardf_local,multi_project=multi_project,unique_parameter=unique_parameter,weights=weights,mod=mod,log_fit=log_fit,ext_spectra=ext_spectra)
							try:
								
								mini_local = lmfit.Minimizer(sub_problem,par_local,fcn_kws={'varied_par':fixed_par,'pardf_local':pardf_local,'fit_ds':fit_ds,'multi_project':multi_project, 'unique_parameter':unique_parameter,'same_DAS':same_DAS,'weights':weights,								'mod':mod,'log_fit':self.log_fit,'target_s2':target_s2,'ext_spectra':ext_spectra,'sub_sample':sub_sample,'pulse_sample':pulse_sample})							
								one_percent_precission=(target-1)*0.01*re['error']
								#results_local = mini_local.minimize('least_squares',ftol=one_percent_precission)
								if other_optimizers is None:
									results_local = mini_local.minimize(method='nelder',options={'maxiter':100,'fatol':one_percent_precission})
								else:
									results_local = mini_local.minimize(method=other_optimizers,options={'maxiter':100,'fatol':one_percent_precission})
								iterative_calls+=results_local.nfev
								if results_local.success:
									conf_limits[fixed_par][i]=results_local.params[fixed_par].value
								else:
									print("tried to optimise %i times achieved residual %g with targeted %g"%(results_local.nfev,(np.sqrt(results_local.residual[0])+target_s2),target_s2))
							except Exception:
								#print("Unexpected error:", sys.exc_info()[0])
								#print('##############################################')
								#print('The error was:')
								#print(e)
								print("error in %s at %s limit"%(fixed_par,i))
								continue
				else:
					print("please use a confidence level between 0.6 and 1")
					return False
				print("it took %i optimisations to get the confidence"%iterative_calls)
		
		############################################################################		
		#-----prepare frames for storage without confidence and store them------------------------
		############################################################################
		
		if pardf.vary.any():
			re['fit_output']=results#let's store the fit results in the re_object for now.
			if confidence_level is not None:
				re['confidence']=conf_limits
				pardf.insert(len(pardf.columns),'lower_limit',None)
				pardf.insert(len(pardf.columns),'upper_limit',None)
				for key in conf_limits.keys():
					pardf.loc[key,'lower_limit']=conf_limits[key]['lower']
					pardf.loc[key,'upper_limit']=conf_limits[key]['upper']
				if self.log_fit:
					for key in ['value','min','max','lower_limit','upper_limit']:
						for row in pardf[pardf.is_rate].index.values:
							try:
								pardf.loc[row,key]=10**pardf.loc[row,key]
							except:
								if pardf.loc[row,key] is None:
									continue
								elif pardf.loc[row,key].isnan():
									continue
								else:
									print('%s,%s has could not be converted and has value'%(row,key))
									print(pardf.loc[row,key])
									continue
				re['confidence']['target-level']='%.1f\n'%((confidence_level)*100)
		re['fit_results_rates']=pardf
		timedf=pardf_to_timedf(pardf)
		re['fit_results_times']=timedf			
		if same_DAS:
			for i,re_local in enumerate(re_listen):
				for name in ['fit_output','fit_results_rates','fit_results_times']:
					try:
						re_listen[i][name]=re[name]
					except:
						print(name + 'not found')
		
		###############################################
		##convert energy back to wavelength#############
		################################################
		if 1:
			if self.equal_energy_bin is not None:
				if same_DAS:
					for i,re_local in enumerate(re_listen):
						for name in ['A','AC','AE']:
							re_local[name].columns=(scipy.constants.h*scipy.constants.c/(re_local[name].columns.values*1e-9*scipy.constants.electron_volt))
							re_local[name].columns.name='wavelength in nm'
							re_local[name].sort_index(inplace=True,axis=1,ascending=True)
						re_local['DAC'].index=(scipy.constants.h*scipy.constants.c/(re_local['DAC'].index.values*1e-9*scipy.constants.electron_volt))
						re_local['DAC'].index.name='wavelength in nm'
						re_local['DAC'].sort_index(inplace=True,axis=0,ascending=True)
						re_listen[i]=re_local
				else:
					for name in ['A','AC','AE']:
						re[name].columns=(scipy.constants.h*scipy.constants.c/(re[name].columns.values*1e-9*scipy.constants.electron_volt))
						re[name].columns.name='wavelength in nm'
						re[name].sort_index(inplace=True,axis=1,ascending=True)
					re['DAC'].index=(scipy.constants.h*scipy.constants.c/(re['DAC'].index.values*1e-9*scipy.constants.electron_volt))
					re['DAC'].index.name='wavelength in nm'
					re['DAC'].sort_index(inplace=True,axis=0,ascending=True)
		
		
		############################################################################
		#---print the output---------------------------------------------------
		############################################################################
		self.re=re
		if same_DAS:
			re_listen[0]=re
			self.multi_projects=re_listen

		Result_string='\nFit Results:\n'
		
		if isinstance(mod,type('hello')):
			Result_string+='Model Used: %s\n\n'%mod
		else:
			Result_string+='Model Used: External function\n\n'
		
		
		if self.ignore_time_region is not None:
			try:
				Result_string+='the time between %.3f %s and %.3f %s was excluded from the optimization\n\n'%(self.ignore_time_region[0],self.baseunit,self.ignore_time_region[1],self.baseunit)
			except:#we got a list
				for entry in self.ignore_time_region:
					Result_string+='the time between %.3f %s and %.3f %s was excluded from the optimization\n\n'%(entry[0],self.baseunit,entry[1],self.baseunit)
		Result_string+='The minimum error is:{:.8e}\n'.format(re['error'])
		try:
			Result_string+='The minimum R2-value is:{:.8e}\n'.format(re['r2'])
		except:
			pass
		if same_DAS:
			Result_string+='The minimum global error is:{:.8e}\n'.format(re['error_total'])
			Result_string+='The minimum global R2-value is:{:.8e}\n'.format(re['r2_total'])
		if confidence_level is not None:
			Result_string+='\nIn Rates with confidence interval to level of %.1f\n\n'%((confidence_level)*100)
			Result_string+=pardf.loc[:,['value','lower_limit','upper_limit','init_value','vary','min','max','expr']].to_markdown(tablefmt="grid")
			Result_string+='\n\nThe rates converted to times with unit %s with confidence interval to level of %.1f\n\n'%(self.baseunit,(confidence_level)*100)
			Result_string+=timedf.loc[:,['value','lower_limit','upper_limit','init_value','vary','min','max','expr']].to_markdown(tablefmt="grid")
		else:
			Result_string+='\nIn Rates\n\n'
			Result_string+=pardf.loc[:,['value','init_value','vary','min','max','expr']].to_markdown(tablefmt="grid")
			Result_string+='\n\nThe rates converted to times with unit %s\n\n'%self.baseunit
			Result_string+=timedf.loc[:,['value','init_value','vary','min','max','expr']].to_markdown(tablefmt="grid")
		if same_DAS:
			Result_string+='\n\nthe other objects were layed into self.multi_projects as list with the local re on position 0.\n By replacing assuming that self = ta write: \n ta.re = ta.multi_projects[1] and then ta.Plot_fit_output to look on the other fits\n '
		try:
			if not results.aborted:
				print(Result_string)
		except:
			print(Result_string)
		if same_DAS:
			for i,re_local in enumerate(re_listen):
				re_listen[i]['Result_string']=Result_string
		else:
			re['Result_string']=Result_string
			
		if dump_paras:
			with open("Fit_results_print.par", "w") as text_file:
				text_file.write(Result_string)
			
			
	def Plot_fit_output(self, plotting = range(7), path = 'result_figures', savetype = 'png', 
						evaluation_style = False, title = None, scale_type = 'symlog', 
						patches = False, filename = None, cmap = None , print_click_position = False,
						plot_second_as_energy = True, return_figures_handles=False):
																	 
		'''plots all the fit output figures. The figures can be called separately 
		or with a list of plots. e.g. range(6) call plots 0-5 Manual plotting of certain type:
		
		This is a wrapper function that triggers the plotting of all the fitted plots.
		The parameter in this plot call are to control the general look and features of the plot.
		Which plots are printed is defined by the first command (plotting)
		The plots are generated from the fitted Matrixes and as such only will work after a fit was actually
		completed (and the "re" dictionary attached to the object.)
		In all plots the RAW data is plotted as dots and the fit with lines 
		
		Contents of the plots

			0. DAC contains the assigned spectra for each component of the fit. For
			   a modelling with independent exponential decays this corresponds to
			   the "Decay Associated Spectra" (DAS). For all other models this
			   contains the "Species Associated Spectra" (SAS). According to the
			   model the separate spectra are labeled by time (process) or name, if
			   a name is associated in the fitting model. The spectra are shown in
			   the extracted strength in the right pane and normalized in the left.
			   Extracted strength means that the measured spectral strength is the
			   intensity (concentration matrix) times this spectral strength. As the
			   concentration maxima for all DAS are 1 this corresponds to the
			   spectral strength for the DAS. (please see the documentation for the
			   fitting algorithm for further details).

			1. summed intensity. All wavelength of the spectral axis are summed for
			   data and fit. The data is plotted in a number of ways vs linear and
			   logarithmic axis. This plot is not ment for publication but very
			   useful to evaluate the quality of a fit.

			2. plot kinetics for selected wavelength (see corresponding RAW plot).

			3. plot spectra at selected times (see corresponding RAW plot).

			4. plots matrix (measured, modelled and error Matrix). The parameter are
			   the same as used for the corresponding RAW plot with the addition of
			   "error_matrix_amplification" which is a scaling factor multiplied
			   onto the error matrix. I recommend to play with different "cmap",
			   "log_scale" and "intensity_scale" to create a pleasing plot.

			5. concentrations. In the progress of the modelling/fitting a matrix is
			   generated that contains the relative concentrations of the species
			   modelled. This plot is showing the temporal development of these
			   species. Further details on how this matrix is generated can be found
			   in the documentation of the fitting function. The modeled spectra are
			   the convolution of these vectors (giving the time-development) and
			   the DAS/SAS (giving the spectral development).
		
		Parameters
		---------------
		
		plotting : int or iterable (of integers), optional
			This parameter determines which figures are plotted 
			the figures can be called separately with plotting = 1
			or with a list of plots (Default) e.g. plotting=range(6) calls plots 0,1,2,3,4,5
			The plots have the following numbers:
			
				0. DAS or SAS
				1. summed intensity
				2. Kinetics
				3. Spectra
				4. Matrixes
				5. Concentrations (the c-object)
				6. Residuals (the difference between the measured and the fitted data)
			
			The plotting takes all parameter from the "ta" object unless otherwise specified
		
		path : None, str or path object, optional
			This defines where the files are saved if the safe_figures_to_folder parameter is True, 
			quite useful if a lot of data sets are to be printed fast. 
			If a path is given, this is used. If a string like the (Default) "result_figures" is given, 
			then a subfolder of this name will be used (an generated if necessary) 
			relative to self.path. Use and empty string to use the self.path
			If set to None, the location of the plot_func will be used and
			a subfolder with title "result_figures" be generated here
		
		return_figures_handles : bool, optional
			(Default) is False, if True the Figure handles are returned as a dictionary.
		
		savetype : str or iterable (of str), optional 
			matplotlib allows the saving of figures in various formats. (Default) "png", 
			typical and recommendable options are "svg" and "pdf".  
			
		evaluation_style : bool, optional
			True (Default = False) adds a lot of extra information in the plot
		
		title : None or str, optional
		   "title=None" is in general the filename that was loaded. Setting a
		   specific title will be used in all plots. To remove the title all
		   together set an empty string with title=""
		   
		scale_type : str, optional
		   refers to the time-axis and takes, "symlog" (Default)(linear around zero and logarithmic otherwise)
		   and "lin" for linear and  "log" for logarithmic, switching all the time axis to this type
		
		patches : bool, optional
			If False (Default) the names "measured" "fitted" "difference" will be placed above the images.
			If True, then they will be included into the image (denser)

		filename : str, optional
			offers to replace the base-name used for all plots (to e.g.specify what sample was used). 
			if (Default) None is used, the self.filename is used as a base name. The filename plays only a 
			role during saving, as does the path and savetype	
			
		cmap : None or matplotlib color map, optional
			is a powerfull variable that chooses the colour map applied for all plots. If set to 
			None (Default) then the self.cmap is used.
			As standard I use the color map "jet" from matplotlib. There are a variety of colormaps 
			available that are very usefull. Beside "jet", "viridis" is a good choice as it is well 
			visible under red-green blindness. Other useful maps are "prism" for high fluctuations 
			or diverging color maps like "seismic". 
			See https://matplotlib.org/3.1.0/tutorials/colors/colormaps.html for a comprehensive 
			selection. In the code the colormaps are imported so if plot_func is imported as pf then 
			self.cmap=pf.cm.viridis sets viridis as the map to use. Internally the colors are chosen 
			with the "colm" function. The 2d plots require a continuous color map so if something 
			else is give 2d plots are shown automatically with "jet". For all of the 1d plots however 
			I first select a number of colors before each plot. If cmap is a continous map then these
			are sampled evenly over the colourmap. Manual iterables of colours 
			cmap=[(1,0,0),(0,1,0),(0,0,1),...] are also accepted, as are vectors or dataframes that 
			contain as rows the colors. There must be of course sufficient colors present for 
			the numbers of lines that will be plotted. So I recommend to provide at least 10 colours 
			(e.g.your university colors). colours are always given as a, list or tuple with RGA or RGBA
			(with the last A beeing the Alpha=transparency. All numbers are between 0 and 1. 
			If a list/vector/DataFrame is given for the colours they will be used in the order provided.		
		
		print_click_position : bool, optional
			if True then the click position is printed for the spectral plots 
		
		Examples
		------------
		
		Typically one would call this function empty for an overview:
		After the minimum fit
		
		>>> ta=pf.TA('testfile.SIA')
		>>> ta.par=lmfit.Parameters()
		>>> ta.par.add('k0',value=1/0.1,vary=True)
		>>> ta.Fit_Global()
		
		One usually plots the an overview
		
		>>> ta.Plot_fit_output()
		>>> ta.Plot_fit_output(plotting=range(6)) #is the same as before
		>>> ta.Plot_fit_output(2) #would plot only the kinetics
		>>> ta.Plot_fit_output(plotting = 2) #would plot only the kinetics
				

		'''
		try:
			re=self.re
		except:
			print('We need to have fitted something so that we can plot')
			return False
		path=check_folder(path=path,current_path=self.path)
		if self.save_figures_to_folder:
			self.figure_path=path
		if cmap is None:cmap=self.cmap
		if filename is None:filename=self.filename
		if title is None:
			if filename is None:
				title=self.filename
			else:
				title=filename
		if not hasattr(plotting,"__iter__"):plotting=[plotting]
		returning_dict=plot_fit_output(self.re, self.ds, cmap = cmap, plotting = plotting, title = title, 
						path = path, f = filename, intensity_range = self.intensity_range, 
						log_scale = self.log_scale, baseunit = self.baseunit, timelimits = self.timelimits, 
						scattercut = self.scattercut, bordercut = self.bordercut, 
						error_matrix_amplification = self.error_matrix_amplification, 
						wave_nm_bin = self.wave_nm_bin, rel_wave = self.rel_wave, width = self.wavelength_bin, 
						rel_time = self.rel_time, save_figures_to_folder = self.save_figures_to_folder, 
						log_fit = self.log_fit,mod = self.mod, savetype = savetype, 
						time_width_percent = self.time_width_percent, evaluation_style = evaluation_style, 
						filename = self.filename, scale_type = scale_type, patches = patches, lintresh = self.lintresh,
						print_click_position = print_click_position, ignore_time_region = self.ignore_time_region,
						data_type = self.data_type, plot_second_as_energy = plot_second_as_energy, units= self.units, 
						equal_energy_bin = self.equal_energy_bin, return_figures_handles=return_figures_handles)
		if return_figures_handles:
			return returning_dict



	def Save_data(self, save_RAW = True, save_Fit = True, save_slices = True, save_binned = False, 
					filename = None, save_fit_results = True, path = 'Data_export', sep = str('\t')):
		'''handy function to save the data on disk as dat files.
		The RAW labeled files contain the chirp corrected values (self.ds)
		
		the save_slices switch turns on the dump of the separate sliced figures (time and spectral) 

		
		Parameters
		----------
		
		save_binned : bool, optional
			is also the re-binned matrix to be saved.
			
		save_slices : bool, optional
			save the kinetics and spectra from the fitted data (with the fits)
			
		sep : str, optional
			what symbol is used to separate different number. (typical either 'tab' or comma
		
		save_RAW : bool, optional
			(Default) True then the first slide with the RAW data is created 
		
		save_Fit : bool, optional
			(Default) True then the second slide with the Fitted data is created 
		
		path : None, str or path, optional
			(Default) None, if left on None, then a folder "result_figures" is created in the folder 
			of the data (self.path)
			
		save_fit_results : bool, optional
			if True (Default)  a neatly formated file with the fit results is created and stored with the data
			
		filename : str, optional
			(Default) None, Base name for all plots. If None, then self.filename will be used
			
		Examples
		---------
		>>> ta.Save_Data
		
		'''
		
		if filename is None:filename = self.filename.split('.')[0]
		if save_RAW:
			self.ds.to_csv(check_folder(path = path, current_path = self.path, 
							filename = filename+'_chirp_corrected_raw_matrix.dat'), sep = sep)
			if save_binned:
				sub = sub_ds(self.ds, scattercut = self.scattercut, bordercut = self.bordercut, 
							timelimits = self.timelimits, wave_nm_bin = self.wave_nm_bin, 
							time_bin = self.time_bin)
				sub.to_csv(check_folder(path = path, current_path = self.path, 
										filename = filename+'_chirp_corrected_rebinned_matrix.dat'), sep = sep)
			if save_slices:
				sub = sub_ds(ds = self.ds.copy(), wavelength_bin = self.wavelength_bin, wavelength = self.rel_wave)
				#sub.columns.name = 'wavelength [nm] in %.0f bins'%self.wavelength_bin
				sub.to_csv(check_folder(path = path, current_path = self.path, 
							filename = filename+'_chirp_corrected_RAW_kinetics.dat'), sep = sep)
				sub = sub_ds(ds = self.ds.copy(), times = self.rel_time, time_width_percent = self.time_width_percent, 
								scattercut = self.scattercut, bordercut = self.bordercut, wave_nm_bin = self.wave_nm_bin)
				sub.to_csv(check_folder(path = path, current_path = self.path, 
							filename = filename+'_chirp_corrected_RAW_Spectra.dat'), sep = sep)
		if save_Fit:
			try:
				self.re.keys()
			except:
				print('no fit in data')
				save_Fit = False
		if save_Fit:
			self.re['A'].to_csv(check_folder(path = path, current_path = self.path, 
								filename = filename+'_matrix used as fit input.dat'), sep = sep)
			self.re['AC'].to_csv(check_folder(path = path, current_path = self.path, 
								filename = filename+'_matrix calculated during fit.dat'), sep = sep)
			self.re['AE'].to_csv(check_folder(path = path, current_path = self.path, 
								filename = filename+'_error_matrix calculated during fit.dat'), sep = sep)
			if save_slices:	
				sub = sub_ds(ds = self.re['AC'].copy(), wavelength_bin = self.wavelength_bin, wavelength = self.rel_wave)
				#sub.columns.name = 'wavelenth [nm] in %.0f bins'%self.wavelength_bin
				sub.to_csv(check_folder(path = path, current_path = self.path, 
							filename = filename+'_fitted_kinetics.dat'), sep = sep)
				sub = sub_ds(ds = self.re['A'].copy(), wavelength_bin = self.wavelength_bin, wavelength = self.rel_wave)
				#sub.columns.name = 'wavelenth [nm] in %.0f bins'%self.wavelength_bin
				sub.to_csv(check_folder(path = path, current_path = self.path, 
							filename = filename+'_measured_kinetics.dat'), sep = sep)
				sub = sub_ds(ds = self.re['AC'].copy(), times = self.rel_time, 
							time_width_percent = self.time_width_percent, scattercut = self.scattercut, 
							bordercut = self.bordercut, wave_nm_bin = self.wave_nm_bin)
				sub.to_csv(check_folder(path = path, current_path = self.path, 
							filename = filename+'_fitted_spectra.dat'), sep = sep)
				sub = sub_ds(ds = self.re['A'].copy(), times = self.rel_time, 
							time_width_percent = self.time_width_percent, scattercut = self.scattercut, 
							bordercut = self.bordercut, wave_nm_bin = self.wave_nm_bin)
				sub.to_csv(check_folder(path = path, current_path = self.path, 
							filename = filename+'_measured_spectra.dat'), sep = sep)
				self.re['DAC'].to_csv(check_folder(path = path, current_path = self.path, 
							filename = filename+'_DAS-SAS.dat'), sep = sep)
		
			if save_fit_results:
				
				Result_string='\nFit Results:\n'
				
				if isinstance(self.mod,type('hello')):
					Result_string+='Model Used: %s\n\n'%self.mod
				else:
					Result_string+='Model Used: External function\n\n'
				
				if self.ignore_time_region is not None:
					Result_string+='the time between %.3f %s and %.3f %s was excluded from the optimization\n'%(self.ignore_time_region[0],self.baseunit,self.ignore_time_region[1],self.baseunit)	
				Result_string+='The minimum error is:{:.8e}\n'.format(self.re['error'])
				Result_string+='The minimum R2-value is:{:.8e}\n'.format(self.re['r2'])
				
				if 'confidence' in self.re:
					Result_string+='\nIn Rates with confidence interval to level of %s\n'%self.re['confidence']['target-level']
					Result_string+=self.re['fit_results_rates'].to_string(columns=['value','lower_limit','upper_limit','init_value','vary','min','max','expr'])
					Result_string+='\n\nThe rates converted to times with unit %s with confidence interval to level of %s\n'%(self.baseunit,self.re['confidence']['target-level'])
					Result_string+=self.re['fit_results_times'].to_string(columns=['value','lower_limit','upper_limit','init_value','vary','min','max','expr'])
				else:
					Result_string+='\nIn Rates\n'
					Result_string+=self.re['fit_results_rates'].to_string(columns=['value','init_value','vary','min','max','expr'])
					Result_string+='\n\nThe rates converted to times with unit %s\n'%self.baseunit
					Result_string+=self.re['fit_results_times'].to_string(columns=['value','init_value','vary','min','max','expr'])

				with open(check_folder(path = path, current_path = self.path, filename = filename+'_fit_results_parameter.par'), "w") as text_file:
					text_file.write(Result_string)

	def Save_Powerpoint(self, save_RAW = True, save_Fit = True, filename = None, 
						path = 'result_figures', scale_type = 'symlog', title = None, patches = False, cmap=None , savetype = 'pptx'):
		'''This function creates two power point slides. On the first it summarizes the RAW plots and on 
		the second (if existent) it summarizes the fitted results
		
		Parameters
		----------
		
		save_RAW : bool, optional
			(Default) True then the first slide with the RAW data is created 
		
		save_Fit : bool, optional
			(Default) True then the second slide with the Fitted data is created 
		
		path : None, str or path, optional
			(Default) None, if left on None, then a folder "result_figures" is created in the folder 
			of the data (self.path)
		
		savetype : str or iterable (of str), optional 
			triggers the additional creation of a composite file in this format.
			matplotlib allows the saving of figures in various formats. (Default) "png", 
			typical and recommendable options are "svg" and "pdf".  
			
		title : None or str, optional
			(Default) None, Use this title on all plots. if None, use self.filename
			
		filename : str, optional
			(Default) None, Base name for all plots. If None, then self.filename will be used
			
		scale_type : str, optional
			'symlog' (Default), 'linear', 'log' time axis
			
		patches : bool, optional
			For true use white patches to label things in the 2d matrixes, to safe
			space for publication
			
		cmap : None or matplotlib color map, optional
			is a powerfull variable that chooses the colour map applied for all plots. If set to 
			None (Default) then the self.cmap is used.
			As standard I use the color map "jet" from matplotlib. There are a variety of colormaps 
			available that are very usefull. Beside "jet", "viridis" is a good choice as it is well 
			visible under red-green blindness. Other useful maps are "prism" for high fluctuations 
			or diverging color maps like "seismic". 
			See https://matplotlib.org/3.1.0/tutorials/colors/colormaps.html for a comprehensive 
			selection. In the code the colormaps are imported so if plot_func is imported as pf then 
			self.cmap=pf.cm.viridis sets viridis as the map to use. Internally the colors are chosen 
			with the "colm" function. The 2d plots require a continuous color map so if something 
			else is give 2d plots are shown automatically with "jet". For all of the 1d plots however 
			I first select a number of colors before each plot. If cmap is a continous map then these
			are sampled evenly over the colourmap. Manual iterables of colours 
			cmap=[(1,0,0),(0,1,0),(0,0,1),...] are also accepted, as are vectors or dataframes that 
			contain as rows the colors. There must be of course sufficient colors present for 
			the numbers of lines that will be plotted. So I recommend to provide at least 10 colours 
			(e.g.~your university colors). colours are always given as a, list or tuple with RGA or RGBA
			(with the last A beeing the Alpha=transparency. All numbers are between 0 and 1. 
			If a list/vector/DataFrame is given for the colours they will be used in the order provided.
		
		Examples
		---------
		
		>>> ta.Save_Powerpoint()
		>>> ta.Save_Powerpoint(patches = True)
		
		'''						 
		if isinstance(savetype,type('hello')):savetype=[savetype]
		if not hasattr(savetype,"__iter__"):savetype=[savetype]
		raw_names=["MAT","SEL","SPEK","SVD"]
		raw_names=[check_folder(current_path=self.path, path=path, filename=self.filename.split('.')[0] + "_RAW_"+str(a) +".png") for a in raw_names]
		fit_names=["FIG_MAT","SPEC","SEL","SUM","DAC"]
		fit_names=[check_folder(current_path=self.path, path=path, filename=self.filename.split('.')[0] + "_" +str(a) +".png") for a in fit_names]
		plt.close('all')
		origin=self.save_figures_to_folder
		if filename is None:
			filename=self.filename
			filename=filename.split('.')[0]
		if save_RAW:
			self.save_figures_to_folder=True
			self.Plot_RAW(savetype = 'png', scale_type = scale_type, title = title, cmap = cmap, path = path)	
			plt.close('all')
		if save_Fit:
			try:
				self.save_figures_to_folder=True
				self.Plot_fit_output(savetype = 'png', scale_type = scale_type, title = title, patches = patches, cmap = cmap , path = path)
				plt.close('all')
			except:
				save_Fit = False
				print('run into problems with adding the fit results. Have you fitted something?')
			try:
				Result_string=self.re['Result_string']
				Result_string=Result_string.replace('lower_limit','low_lim')
				Result_string=Result_string.replace('upper_limit','up_lim')
				Result_string.replace('===','=')
			except Exception as e:
				print(e)
		
		if ('pptx' in savetype) or ('ppt' in savetype):
			try:
				left=Inches(0.2)
				top=Inches(0.2)
				prs = Presentation()
				blank_slide_layout = prs.slide_layouts[6]
				if save_RAW:
					slide = prs.slides.add_slide(blank_slide_layout)
					left = top = Inches(0.5)
					pic = slide.shapes.add_picture(str(raw_names[0].resolve()), left=left+Inches(4.5), top=top, width=Inches(4.5))
					pic = slide.shapes.add_picture(str(raw_names[1].resolve()), left=left, top=top, width=Inches(4.5))
					pic = slide.shapes.add_picture(str(raw_names[2].resolve()), left=left, top=top+Inches(3), width=Inches(4.5))
					try:
						pic = slide.shapes.add_picture(str(raw_names[3].resolve()), left=left+Inches(4.5), top=top+Inches(3), height=Inches(3.4))
					except:
						pass
				if save_Fit:
					try:
						slide2 = prs.slides.add_slide(blank_slide_layout)
						left = top = Inches(0.1)
						pic = slide2.shapes.add_picture(str(fit_names[0].resolve()), left=left+Inches(7.0), top=top, height=Inches(3.4))#Matrix
						pic = slide2.shapes.add_picture(str(fit_names[1].resolve()), left=left, top=top, height=Inches(2))
						pic = slide2.shapes.add_picture(str(fit_names[2].resolve()), left=left, top=top+Inches(2), height=Inches(2))
						pic = slide2.shapes.add_picture(str(fit_names[3].resolve()), left=left, top=top+Inches(3.9), height=Inches(1.4))
						pic = slide2.shapes.add_picture(str(fit_names[4].resolve()), left=left, top=top+Inches(5.4), height=Inches(2))
						text1 = slide2.shapes.add_textbox(left=left+Inches(5.2), top=top+Inches(2.5), width=Inches(4.5), height=Inches(4.5))
						text1.text = '{}'.format(Result_string.replace('===','='))
						try:
							text1.text_frame.fit_text(font_family='Garamond', max_size=6, bold=True, italic=False)
							#text1.text_frame.fit_text(font_family='Haettenschweiler', max_size=6, bold=False, italic=False)
						except:
							text1.text_frame.fit_text(font_family='Arial', max_size=5.0, bold=False, italic=False)

					except Exception as e:
						print('exited when saving the fit plots')
						print(e)
				plt.close('all')
				self.save_figures_to_folder=origin
				prs.save(check_folder(path=path,current_path=self.path,filename=self.filename.split('.')[0] + '.pptx'))
				print('The images and a powerpoint was saved to %s'%check_folder(path=path,current_path=self.path))
			except Exception as e:
				print('Error in powerpoint generation. Most likely a module is missing.')
				print('We need python-pptx to create a powerpoint file.  Either use "pip install python-pptx" or "conda install -c conda-forge python-pptx" ')
				print('We will save the results as pdf format for now. Check th error if somehting else went wrong')
				print(e)
				savetype.append('pdf')
			
			
		if ('pdf' in savetype) or ('png' in savetype) or ('svg' in savetype):
			if save_RAW:
				fig,ax=plt.subplots(nrows=2,ncols=2,figsize=(10,7.5))
				ax[0,1].imshow(mpimg.imread(str(raw_names[0])))
				ax[0,0].imshow(mpimg.imread(str(raw_names[1])))
				ax[1,0].imshow(mpimg.imread(str(raw_names[2])))
				ax[1,1].imshow(mpimg.imread(str(raw_names[3])))
				ax[0,0].axis('off');ax[1,0].axis('off');ax[0,1].axis('off');ax[1,1].axis('off')
				for entry in savetype:
					if entry == "pptx":continue
					try:
						fig.tight_layout()
						fig.savefig(check_folder(path=path,current_path=self.path,filename=self.filename.split('.')[0] + '_RAW-summary.%s'%entry),dpi=600)
					except:
						print("saving in" + entry +"failed")
			if save_Fit:
				G = GridSpec(4, 8)
				fig1=plt.figure(figsize=(8,10))
				ax1=fig1.add_subplot(G[0,:5])
				ax2=fig1.add_subplot(G[1,:5])
				ax3=fig1.add_subplot(G[2,:6])
				ax4=fig1.add_subplot(G[3,:6])
				ax5=fig1.add_subplot(G[0:2,5:])
				ax6=fig1.add_subplot(G[1:,6:])
				ax1.imshow(mpimg.imread(str(fit_names[1])))
				ax2.imshow(mpimg.imread(str(fit_names[2])))
				ax3.imshow(mpimg.imread(str(fit_names[3])))
				ax4.imshow(mpimg.imread(str(fit_names[4])))
				ax5.imshow(mpimg.imread(str(fit_names[0])))
				ax6.text(0,0,Result_string.replace('===','='),font='Garamond',fontsize=6,fontweight='bold')
				ax1.axis('off');ax2.axis('off');ax3.axis('off');ax4.axis('off');ax5.axis('off');ax6.axis('off')
				for entry in savetype:
					if entry == "pptx": continue
					try:
						fig1.tight_layout()
						fig1.savefig(check_folder(path=path,current_path=self.path,filename=self.filename.split('.')[0] + '_Fit-summary.%s'%entry),dpi=600)
					except:
						print("saving in" + entry +"failed")
		
	def Print_Results(self,to_file=False):
		if 're' in self.__dict__:
			try:
				print('{}'.format(self.re['Result_string'].decode('utf-8')))
			except:
				try:
					print('{}'.format(self.re['Result_string']))
				except:
					print('printing of results failed.')
		if to_file:
			with open(self.filename+'_fitting_results','w') as f:
				f.write('{}'.format(self.re['Result_string']))


	def Save_project(self, filename=None,path=None):
		'''function to dump all the parameter of an analysis into an hdf5 file. 
		This file contains the ds_ori and all the parameter, including fitting parameter
		and results.
		One limitation is the fitting model. If the model is build in, so the model is
		'exponential' or 'parallel' then the safing works. If an external model is used then the 
		dostring of the external function is stored, but not the function itself.
				
		Parameters
		----------
		
		path : None, str or path, optional
			(Default) None, if left on None, then a folder "Data" is created in the folder 
			of the project (self.path)
				
		filename : str, optional
			(Default) None, Base name for all plots. If None, then self.filename will be used
			
		Examples
		--------
		
		>>> ta.Save_project()
		
		'''
		
		if filename is None:
			filename = self.filename
		hdf5_name =check_folder(path = path, current_path = self.path, filename = filename.split('.')[0]+'.hdf5')
		if os.path.exists(hdf5_name):
			try:
				os.remove(hdf5_name)
			except:
				try:
					hdf5_name.close()
					os.remove(hdf5_name)
				except:
					print('File exists but can not be deleted')
		re_switch = False
		with h5py.File(hdf5_name, 'w') as f:
			for key in self.__dict__.keys():
				if key == 'mod':
					if self.__dict__[key] in ['paral','exponential','consecutive','full_consecutive']:
						f.create_dataset(name=key, data=self.__dict__[key])
					else:
						try:
							docstring=self.__dict__[key].__doc__
							if isinstance(docstring,type('hello')):
								f.create_dataset(name=key, data=docstring)
						except:
							f.create_dataset(name=key, data='external_function_without_docstring')
				elif key in ['rel_wave','rel_time']:#need extra, as it is bypassed by the re-switch
					f.create_dataset(name=key, data=self.__dict__[key])
				elif key[:2] == 're' :
					re_switch = True
					for key2 in self.__dict__['re']:
						if key2 == 'fit_output':continue
						
						if key2 == 'error':
							data = self.__dict__['re'][key2]
							try:
								f.create_dataset(name='re_error', data=data)
							except:
								print('saving of ' + key2 + ' failed' )
						elif key2 == 'confidence':
							for key3 in self.__dict__['re']['confidence'].keys():
								try:
									f.create_dataset(name='re_confidence_%s_upper'%key3, data=self.__dict__['re'][key2][key3]['upper'])
									f.create_dataset(name='re_confidence_%s_lower'%key3, data=self.__dict__['re'][key2][key3]['lower'])
								except:
									try:
										f.create_dataset(name='re_confidence_%s'%key3, data=self.__dict__['re'][key2][key3])
									except:
										print('saving of' + key3 + 'in confidence failed')
						elif isinstance(self.__dict__['re'][key2], pandas.DataFrame):
							pass
						else:
							try:
								f.create_dataset(name='re_' + key2, data=self.__dict__['re'][key2])
							except:
								
								print('saving of ' + key2 + ' failed' )
				elif key == 'cmap':
					pass
				elif key == 'intensity_range':
					data=self.__dict__['intensity_range']
					if isinstance(data, type(1e-3)):
						data=[-data,data]
					if data is None:
						f.create_dataset(name='intensity_range', data='None')
					else:
						f.create_dataset(name='intensity_range', data=data)	
				elif key == 'background_par':
					f.create_dataset(name='back', data=self.__dict__['background_par'][3])
				
				elif key in ['par','par_fit']:
									 
					df=par_to_pardf(self.__dict__[key])#pandas has a bug and problems handling mixed type columns when saving. So we clean up.
					for sub_key in ['min','max','value']:
						try:
							df[sub_key]=df[sub_key].astype(float)
						except:
							pass
					df['is_rate']=df['is_rate'].astype(bool)
					df['vary']=df['vary'].astype(bool)
					df['expr']=df['expr'].apply(lambda x:'%s'%x)
					df.to_hdf(str(hdf5_name.resolve()), key=key, append=True, mode='r+', format='t')
				else:
					data = self.__dict__[key]
					if data is None:
						f.create_dataset(name=key, data='None')
					else:
						if isinstance(data, pandas.DataFrame):
							pass
						else:
							try:
								f.create_dataset(name=key, data=data)
							except:
								if key == 'path':
									pass
								elif key == 'figure_path':
									pass
								else:
									print('the saving of  %s  failed'%key)
			if 'fitcoeff' not in f:
				try:
					f.create_dataset(name='fitcoeff', data=self.fitcoeff)
				except:
					try:
						with open(self.chirp_file,'r') as f2:
							fitcoeff=f2.readline()
							f.create_dataset(name='fitcoeff', data=fitcoeff)
					except:
						pass
		self.ds_ori.to_hdf(str(hdf5_name.resolve()), key='ds_ori', append=True, mode='r+', format='t')#save_raw_data
		if re_switch:
			#print('re-switched')
			for key in ['A', 'AC', 'AE', 'DAC', 'c']:
				self.re[key].to_hdf(str(hdf5_name.resolve()), key='re_' + key, append=True, mode='r+', format='t')
		try:
			f.close()
		except:
			pass
		
		print('The project was saved to %s'%check_folder(path = path, current_path = self.path))


	def __read_project(self, saved_project=None,current_path=None):
		'''function to re-read all the parameter of a previous analysis 
		into an hdf5 file, current path is the path that the file should 
		assume is its "home directory" after successful loading. If not 
		set we take the filepath at which the file is currently stored as path'''
		if saved_project is None:
			raise ImportError('We do need a project to import')
		if current_path is None:current_path=os.path.dirname(os.path.abspath(saved_project))
		try:
			import h5py
		except:
		 
			print('could not import hdf5, current version requires that this is installed. IF running Anaconda open Conda promt and type: conda install h5py')
		data_frame_list=[]

		# we hav to handle the old and new type of saving
		with h5py.File(saved_project, 'r') as f:
			if 're_final_setup_par' in f.keys():old_switch=True
					 
		
			else:old_switch=False
					   
		if old_switch:print('we read an old style data_file and directly update it into the new file_type after loading')
		with h5py.File(saved_project, 'r') as f:
			for key in f.keys():
				try:
					if "re_" in key[:3]:
						if 're' not in self.__dict__.keys():
							self.__dict__['re']={}
						if 're_confidence_' in key:#_upper, _lower
							if 'confidence' not in self.__dict__['re']:
								self.__dict__['re']['confidence']={}
							if '_upper' in key[-6:]:
								if key[14:-6] not in self.__dict__['re']['confidence']:
									self.__dict__['re']['confidence'][key[14:-6]]={}
								self.__dict__['re']['confidence'][key[14:-6]]['upper']=f[key][()]
							elif '_lower' in key[-6:]:
								if key[14:-6] not in self.__dict__['re']['confidence']:
									self.__dict__['re']['confidence'][key[14:-6]]={}
								self.__dict__['re']['confidence'][key[14:-6]]['lower']=f[key][()]
							else:
								self.__dict__['re']['confidence'][key[14:]]=f[key][()]
						else:
							self.__dict__['re'][key[3:]]=f[key][()]
					elif "back" in key[:4]:
						rea=f[key][()]
						self.__dict__['background_par']=[None,-1,False]
						self.__dict__['background_par'].append(rea)
					elif "par" in key[:3]:
						if old_switch:#old type of saved data
							try:
								os.remove('temp_file.json')
							except:
								pass
							with open('temp_file.json','w') as g:
								g.write(f[key][()])
							with open('temp_file.json','r') as g:
								self.par=lmfit.Parameters()
								self.par.load(g)
							try:
								os.remove('temp_file.json')
							except:
								pass
						else:#new type of data
							raise		
					else:
						read=f[key][()]

						if isinstance(read,bytes):
							read=f[key].asstr()[()] 
						if isinstance(read,str):
							if (read=='None') or (read=='none'):
								read=None
						elif key in ['bordercut','timelimits','fitcoeff','scattercut']:
							read=[float(a) for a in read]
						elif key =='intensity_range':
							read=[float(a) for a in read]
						elif key in ['rel_time','rel_wave']:
							read=np.array(read,dtype=np.float64)
						elif key in ['scattercut']:
							try:
								read=[float(a) for a in read]
							except:#maybe we have  a list of scattercuts
								try:
									out_listen=[]
									for listen in read:
										out_listen.append([float(a) for a in listen])
								except:#no idea lets see what happens
									pass
						self.__dict__[key]=read
				except:#we'll get an exception every time there is an dataframe
					#print('Frame:'+key)
					data_frame_list.append(key)

		try:
			f.close()
		except:
			pass
		for key in data_frame_list:
			try:
				if "re_" in key[:3]:
					#print('re in list')
					self.__dict__['re'][key[3:]]=pandas.read_hdf(saved_project,key=key,mode='r',data_columns=True)
				elif key in ['ds_ori','par_fit','par']:
					self.__dict__[key]=pandas.read_hdf(saved_project,key=key,mode='r',data_columns=True)
				else:
					print("missing key:" + key)
			except Exception as e:
				if key == 'par' and old_switch:pass # we have read it before already and the error is ok
		 
				else:
					print("error in key:" + key)
					print(e)
		
		try:
			self.__dict__['re']['fit_results_rates']=self.__dict__['par_fit']
			self.__dict__['re']['fit_results_times']=pardf_to_timedf(self.__dict__['re']['fit_results_rates'])
		except:
			pass
		
		#the par conversion function failed, quickfix
		for over_key in ['par_fit','par']:
			try:
				par_df=self.__dict__[over_key].loc[:,['value','min','max','vary','expr']]
				par=lmfit.Parameters()
				for key in par_df.index.values:
					par.add(key, value=par_df.loc[key,'value'], vary=par_df.loc[key,'vary'], min=par_df.loc[key,'min'], max=par_df.loc[key,'max'])
															  
				self.__dict__[over_key]=par
			except:
				pass
		if old_switch:
			try:
				self.__dict__['re']['fit_results_rates']=par_to_pardf(self.par)
				self.__dict__['re']['fit_results_times']=pardf_to_timedf(par_to_pardf(self.par))
				self.__dict__['par_fit']=self.par
			except:
				pass
		self.save_figures_to_folder=False
		self.path=current_path
		if old_switch:#convert project into new type
			#clean old files that were read wrong
			for key in ['re_final_int_par','re_final_setup_par','re_final_time_par','re_int_error']:
				try:
					del self.__dict__[key]
				except KeyError:
					print(f'Key {key} is not in the dictionary')
			for key in ['final_int_par','final_setup_par','final_time_par','int_error']:
				try:
					del self.__dict__['re'][key]
				except KeyError:
					print(f'Key {key} is not in the dictionary')
			try:
				self.save_project()
				print("project converted into new data type and saved again")
			except:
				print("project converted ibut could not be saved")
		self.path=current_path
		for key in ['time_bin','rel_wave','rel_time','scattercut','bordercut','timelimits','intensity_range','wave_nm_bin','wavelength_bin','ignore_time_region']:
			try:
				if isinstance(self.__dict__[key],bytes):
					#print(key + ' set to None')
					self.__dict__[key]=None
				elif isinstance(self.__dict__[key],str):
					if self.__dict__[key]=='None':
						self.__dict__[key]=None
			except:
				continue
		try:
			self.figure_path=str(self.figure_path)
			if 'None' in self.figure_path:
				self.figure_path=None
		except:
			pass
		try:
			self.re['Result_string']=self.re['Result_string'].decode('utf-8')
		except:
			pass
		

	def Copy(self):
		'''returns a deep copy of the object.
		
		Examples
		--------
		>>>ta=plot_func.TA('testfile.hdf5') #open a project
		>>>ta1=ta.Copy() #make a copy for some tests or a differnet fit
		
		''' 
		import copy
		return copy.deepcopy(self)


	def Compare_at_time(self, rel_time = None, other = None, fitted = False, norm_window = None,
						time_width_percent = None, spectra = None, data_and_fit = False, cmap = None , 
						print_click_position = False, linewidth = 1, title='', plot_second_as_energy = True):
		'''This function plots multiple spectra into the same figure at a given rel_time (timepoints) and 
		allows for normalization. Very useful to compare the spectra for different solvents or quenchers, or
		e.g. different fits. The ta.time_width_percent parameter defines if this is a
		single time (if time_width_percent = 0) or an integrated window.
		Only "rel_time" is a mandatory, the rest can be taken from the original project (ta).
		
		The normalization is realized by giving a norm_window 
		at which the intensity in the triggering object is integrated (in ta.Compare_at_time(other..) 
		"ta" is the triggering object. The in each of the other curves the same window is
		integrated and the curve scaled by this value. Important to note is that this window
		does not need to be in the plot. e.g. the normalization can be done at a different time.
		
		Very often one would like to compare the measured spectra at a certain
		time to an external spectrum (e.g. spectro-electro-chemistry or steady
		state absorption). This can be done by loading a specific spectrum into
		a DataFrame and handing this data Frame to the comparision function. The
		function can also be used to plot e.g. the measured spectra vs. an
		external spectrum without giving any "other" Projects. (very useful for
		comparisions).
		 
		Parameters
		-------------
		
		rel_time : float or list/vector (of floats)
			Specify the times where to plot, single value or list/vector of values. 
			For each entry in rel_time a spectrum is plotted. 
			If time_width_percent=0 (Default) the nearest measured 
			timepoint is chosen. For other values see parameter "time_width_percent".
			
		other : TA object or list of those, optional 
			should be ta.plot_func objects (loaded or copied) and is what 
			is plotted against the data use a list [ta1,ta2,... ] or generate this
			list using the Gui function. See section :ref:`Opening multiple files` in 
			the documentation
		
		fitted : bool, optional
			True/False (Default) - use fitted data instead of raw data. 
			If True, the fitted datapoints (without interpolation) are used. 
			This is intended for comparing e.g. different fits
		
		norm_window : None or list/vector (with 4 floats), optional
			norm_window Give a list/tupel/vector with 4 entries in the order 
			[Start - time, End - time, Start - wavelength, End - Wavelength], 
			see section  :ref:`Normalization and Scaling`  in the documentation.
			If None (Default) no normalization is done.
		
		linewidth : float, optional
			linewidth to be used for plotting		
		
		time_width_percent : None or float, optional
			"rel_time" and "time_width_percent" work together for creating spectral plots at 
			specific timepoints. For each entry in rel_time a spectrum is plotted. 
			If however e.g. time_width_percent=10 the region between the timepoint closest 
			to :math:`timepoint+0.1xtimepoint´ and :math:`timepoint-0.1xtimepoint` is averaged and shown 
			(and the legend adjusted accordingly). If None (Default) is given, the value is 
			taken from the triggering object (self.time_width_percent) This is particularly useful for the densly
			sampled region close to t=0. Typically for a logarithmic recorded kinetics, the 
			timepoints at later times will be further appart than 10 percent of the value, 
			but this allows to elegantly combine values around time=0 for better statistics. 
			This averaging is only applied for the plotting function and not for the fits.
		
		spectra : None or DataFrame, optional
			If an DataFrame with the wavelength as index is provided, Then the spectra of each column
			is plotted into the differential spectra 1-1 and the column names are used in the legend
			Prior scaling is highly suggested. These spectra are not (in general) scaled with the 
			norm window. (see examples).
		
		data_and_fit : bool, optional
			True or False (Default), choose if for the Fitted plot the raw data of the 
			other projects is to be plotting in addition to the fitted line. For False (Default)
			Only the fit is plotted. 
		
		cmap : None or matplotlib color map, optional
			is a powerfull variable that chooses the colour map applied for all plots. If set to 
			None (Default) then the self.cmap is used.
			As standard I use the color map "jet" from matplotlib. There are a variety of colormaps 
			available that are very usefull. Beside "jet", "viridis" is a good choice as it is well 
			visible under red-green blindness. Other useful maps are "prism" for high fluctuations 
			or diverging color maps like "seismic". 
			See https://matplotlib.org/3.1.0/tutorials/colors/colormaps.html for a comprehensive 
			selection. In the code the colormaps are imported so if plot_func is imported as pf then 
			self.cmap=pf.cm.viridis sets viridis as the map to use. Internally the colors are chosen 
			with the "colm" function. The 2d plots require a continuous color map so if something 
			else is give 2d plots are shown automatically with "jet". For all of the 1d plots however 
			I first select a number of colors before each plot. If cmap is a continous map then these
			are sampled evenly over the colourmap. Manual iterables of colours 
			cmap=[(1,0,0),(0,1,0),(0,0,1),...] are also accepted, as are vectors or dataframes that 
			contain as rows the colors. There must be of course sufficient colors present for 
			the numbers of lines that will be plotted. So I recommend to provide at least 10 colours 
			(e.g.~your university colors). colours are always given as a, list or tuple with RGA or RGBA
			(with the last A beeing the Alpha=transparency. All numbers are between 0 and 1. 
			If a list/vector/DataFrame is given for the colours they will be used in the order provided.
		plot_second_as_energy : bool, optional
			For (Default) True a second x-axis is plotted  with "eV" as unit

		print_click_position : bool, optional
			if True then the click position is printed for the spectral plots 

		Examples
		----------
		
		>>> import plot_func as pf
		>>> ta = pf.TA("test1.hdf5") #open the original project
		
		Now open a bunch of other porjects to comare against
		
		>>> other_projects = pf.GUI_open(project_list = ["file1.SIA", "file2.SIA"])
		
		Typical use is compare the raw data without normalization at 1ps and 6ps.
		
		>>> ta.Compare_at_time(rel_time = [1,6], others = other_project)
		
		Compare the fit withput normalization at 1ps and 6ps.
		
		>>> ta.Compare_at_time(rel_time = [1,6], others = other_project, fitted = True)
		
		Compare with normalization window between 1ps and 2ps and 400nm and 450nm.
		
		>>> norm_window=[1,2,400,450]
		>>> ta.Compare_at_time(rel_time = [1,6], others = other_project, norm_window = norm_window)
		
		Compare the spectrum at 1ps and 6ps with an external spectrum.
		
		>>> ext_spec = pd.read_csv("Ascii_spectrum.dat", sep = ",")
		>>> ta.Compare_at_time(rel_time = [1,6], spectra = ext_spec)

		Use example -  Often there are a lot of different measurements to 
		compare at multiple time. The normlization is performed at the ground state bleach
		460 nm and early in time. Then it is better to make a new plot for each 
		timepoint. The normalization window stays fixed.
		
		>>> plt.close("all") #make some space
		>>> norm_window=[0.3,0.5,450,470] #define window in ground state bleach
		>>> for t in [0.3,0.5,1,3,10,30]: #iterate over the wavelength 
		>>> 	ta.Compare_at_time(rel_time = t, others = other_project, norm_window = norm_window)
		
		'''
		global halfsize
		if self.save_figures_to_folder:self.figure_path=check_folder(path='result_figures',current_path=self.path)
		if time_width_percent is None:time_width_percent=self.time_width_percent
		if rel_time is None:rel_time=self.rel_time
		if other is not None:
			if not hasattr(other,'__iter__'):other=[other]
		if rel_time is not None:
			if not hasattr(rel_time,'__iter__'):rel_time=[rel_time]		
		else:
			rel_time=[1]
		if cmap is None:cmap=self.cmap
		if fitted:
			try:
				re=self.re
			except:
				print("No fitted results present")
				return False
			if norm_window is not None:
				ref_scale=re['A'].loc[norm_window[0]:norm_window[1],norm_window[2]:norm_window[3]].mean().mean()
			
			if halfsize:
				fig,ax=plt.subplots(figsize=(5,3),dpi=100)
			else:
				fig,ax=plt.subplots(figsize=(10,6),dpi=100)
			objects=len(rel_time)*(1+len(other))
			colors=colm(cmap=cmap,k=range(objects))
			_=plot_time(re['A'], ax = ax, rel_time = rel_time, time_width_percent = time_width_percent, 
						baseunit = self.baseunit, lines_are = 'data', cmap = colors[:len(rel_time)], 
						title = '', linewidth = linewidth, subplot= True, scattercut = self.scattercut,
						plot_second_as_energy = plot_second_as_energy)
			_=plot_time(re['AC'], ax = ax, rel_time = rel_time, time_width_percent = time_width_percent, 
						baseunit = self.baseunit, lines_are = 'fitted', cmap = colors[:len(rel_time)], 
						title = '', subplot = False, linewidth = linewidth, scattercut = self.scattercut,
						plot_second_as_energy = plot_second_as_energy)
			handles, labels=ax.get_legend_handles_labels()
			lab=['%g %s'%(ent,self.baseunit) + '_' + str(self.filename) for ent in rel_time]
			han=handles[:len(rel_time)*2]
			for ent in rel_time:
				lab.append('%g %s fit'%(ent,self.baseunit) + '_' + str(self.filename))
			if other is not None:
				for i,o in enumerate(other):
					try:
						re=o.re
					except:
						print('%s has no fitted results'%o.filename)
						continue
					if norm_window is not None:
						rel_scale=re['A'].loc[norm_window[0]:norm_window[1],norm_window[2]:norm_window[3]].mean().mean()
						try:
							scaling=(rel_scale/ref_scale)
							ax=plot_time(re['AC']/scaling, cmap = colors, ax = ax, rel_time = rel_time, 
										time_width_percent = time_width_percent, title = '', 
										lines_are = 'fitted', subplot = True, 
										color_offset = len(rel_time)*(i+1), linewidth = linewidth, 
										scattercut = o.scattercut,plot_second_as_energy = plot_second_as_energy)
							handles, labels=ax.get_legend_handles_labels()
							for ent in rel_time:
								lab.append('%g %s fit'%(ent,o.baseunit) + '_' + str(o.filename))
							for a in handles[-len(rel_time):]:
								han.append(a)
							
							if data_and_fit:
								ax=plot_time(re['A']/scaling, cmap = self.cmap, ax = ax, rel_time = rel_time, 
											time_width_percent = time_width_percent, title = o.filename, 
											baseunit = self.baseunit, lines_are = 'data', subplot = True, 
											color_offset = len(rel_time)*(i+1), linewidth = linewidth, 
											scattercut = o.scattercut,plot_second_as_energy = plot_second_as_energy)
								handles, labels=ax.get_legend_handles_labels()
								for ent in rel_time:
									lab.append('%g %s'%(ent,o.baseunit) + '_' + str(o.filename))
								for a in handles[-len(rel_time):]:
									han.append(a)
							norm_failed = False
						except:
							print('scaling Failed!')
							norm_failed = True
					else: norm_failed=True	
					if norm_failed:
						ax=plot_time(re['AC'], cmap = colors, ax = ax, rel_time = rel_time, 
										time_width_percent = time_width_percent, title = '', 
										lines_are = 'fitted', subplot = True, 
										color_offset = len(rel_time)*(i+1), linewidth = linewidth, 
										scattercut = o.scattercut, plot_second_as_energy = plot_second_as_energy)
						handles, labels=ax.get_legend_handles_labels()
						for ent in rel_time:
							lab.append('%g %s fit'%(ent,o.baseunit) + '_' + str(o.filename))
						for a in handles[-len(rel_time):]:
							han.append(a)
						
						if data_and_fit:
							ax=plot_time(re['A'], cmap = self.cmap, ax = ax, rel_time = rel_time, 
										time_width_percent = time_width_percent, title = o.filename, 
										baseunit = self.baseunit, lines_are = 'data', subplot = True, 
										color_offset = len(rel_time)*(i+1), 
										linewidth = linewidth, scattercut = o.scattercut, 
										plot_second_as_energy = plot_second_as_energy)
							handles, labels=ax.get_legend_handles_labels()
							for ent in rel_time:
								lab.append('%g %s'%(ent,o.baseunit) + '_' + str(o.filename))
							for a in handles[-len(rel_time):]:
								han.append(a)
			if not norm_failed:
				ax.set_title('compare measured and fitted data at given times\n scaled to t=%g ps : %g ps , wl= %g nm: %g nm'%(norm_window[0],norm_window[1],norm_window[2],norm_window[3]))
			else:
				ax.set_title('compare measured and fitted data at given times')
			ax.set_xlim(re['A'].columns.values[0],re['A'].columns.values[-1])
			ax.legend(han, lab ,labelspacing = 0, ncol = 2, columnspacing = 1, handlelength = 1, frameon = False)
		else:
			if norm_window is not None:
				ref_scale=self.ds.loc[norm_window[0]:norm_window[1],norm_window[2]:norm_window[3]].mean().mean()
			objects=len(rel_time)*(1+len(other))
			colors=colm(cmap=cmap,k=range(objects))
			if halfsize:
				fig,ax=plt.subplots(figsize=(5,3),dpi=100)
			else:
				fig,ax=plt.subplots(figsize=(10,6),dpi=100)
			_=plot_time(self.ds, ax = ax, rel_time = rel_time, time_width_percent = time_width_percent, 
						title = title, lines_are = 'data', scattercut = self.scattercut, 
						bordercut = self.bordercut, wave_nm_bin = self.wave_nm_bin, cmap = colors, 
						subplot = True, linewidth = linewidth, baseunit=self.baseunit, 
						plot_second_as_energy = plot_second_as_energy)
			if 1:
				_=plot_time(self.ds, ax = ax, rel_time = rel_time, time_width_percent = time_width_percent, 
							title = title, lines_are = 'smoothed', scattercut = self.scattercut, 
							bordercut = self.bordercut,wave_nm_bin = self.wave_nm_bin, cmap = colors, 
							subplot = False, linewidth = linewidth, baseunit = self.baseunit,
							plot_second_as_energy = plot_second_as_energy)
			handles, labels=ax.get_legend_handles_labels()
			lab=['%g %s'%(ent,self.baseunit) + '_' + str(self.filename) for ent in rel_time]
			han=handles[:len(rel_time)]
			if other is not None:
				for i,o in enumerate(other):
					if norm_window is not None:
						rel_scale=o.ds.loc[norm_window[0]:norm_window[1],norm_window[2]:norm_window[3]].mean().mean()
						try:
							scaling = (rel_scale/ref_scale)
							ax=plot_time(o.ds/scaling, cmap = colors, ax = ax, rel_time = rel_time, 
										time_width_percent = time_width_percent, title = title, 
										lines_are = 'data', scattercut = o.scattercut, 
										bordercut = o.bordercut, linewidth = linewidth, 
										wave_nm_bin = o.wave_nm_bin, subplot = True, 
										color_offset = len(rel_time)*(i+1),
										plot_second_as_energy = plot_second_as_energy)
							handles, labels=ax.get_legend_handles_labels()
							for ent in rel_time:
								lab.append('%g %s'%(ent,o.baseunit) + '_' + str(o.filename))
							for a in handles[-len(rel_time):]:
								han.append(a)

							if data_and_fit:
								ax=plot_time(o.ds/scaling, cmap = colors, ax = ax, rel_time = rel_time, 
										time_width_percent = time_width_percent, title = title, 
										lines_are = 'smoothed', scattercut = o.scattercut, 
										bordercut = o.bordercut, linewidth = linewidth, 
										wave_nm_bin = o.wave_nm_bin, subplot = True, 
										color_offset = len(rel_time)*(i+1),
										plot_second_as_energy = plot_second_as_energy)
							scaling_failed=False
						except:
							print('scaling Failed!')
							scaling_failed=True
					else:
						scaling_failed=True
					if scaling_failed:
						ax=plot_time(o.ds, cmap = colors, ax = ax, rel_time = rel_time, 
									time_width_percent = time_width_percent, title = title, 
									lines_are = 'data', scattercut = o.scattercut, 
									bordercut = o.bordercut, linewidth = linewidth, 
									wave_nm_bin = o.wave_nm_bin, subplot = True, 
									color_offset = len(rel_time)*(i+1),
									plot_second_as_energy = plot_second_as_energy)
						handles, labels=ax.get_legend_handles_labels()
						for ent in rel_time:
							lab.append('%g %s'%(ent,o.baseunit) + '_' + str(o.filename))
						for a in handles[-len(rel_time):]:
							han.append(a)
						if data_and_fit:
							ax=plot_time(o.ds, cmap = colors, ax = ax, rel_time = rel_time, 
									time_width_percent = time_width_percent, title = title, 
									lines_are = 'smoothed', scattercut = o.scattercut, 
									bordercut = o.bordercut, linewidth = linewidth, 
									wave_nm_bin = o.wave_nm_bin, subplot = True, 
									color_offset = len(rel_time)*(i+1),
									plot_second_as_energy = plot_second_as_energy)
									
			if not scaling_failed:
				ax.set_title('compare measured and smoothed data at given times\n scaled to t=%g ps : %g ps , wl= %g nm: %g nm'%(norm_window[0],norm_window[1],norm_window[2],norm_window[3]))
			else:
				ax.set_title('compare measured and smoothed data at given times')
			ax.legend(han, lab, labelspacing = 0, ncol = 2, columnspacing = 1, handlelength = 1, frameon = False)
			if self.bordercut is None:
				ax.set_xlim(self.ds.columns.values[0],self.ds.columns.values[-1])
			else:
				ax.set_xlim(self.bordercut)
		if spectra is not None:
			spectra.plot(ax=ax,legend=False)
			handles, labels=ax.get_legend_handles_labels()
			han.append(handles[-1])
			lab.append(labels[-1])
			ax.legend(han, lab ,labelspacing = 0, ncol = 2, columnspacing = 1, handlelength = 1, frameon = False)
		if plot_second_as_energy:
			ax2=ax.twiny()
			ax2.set_xlim(ax.get_xlim())
			ax2.set_xticks(ax.get_xticks())
			labels=['%.2f'%(scipy.constants.h*scipy.constants.c/(a*1e-9*scipy.constants.electron_volt)) for a in ax2.get_xticks()]
			_=ax2.set_xticklabels(labels)
			_=ax2.set_xlabel('Energy in eV')
			ax.set_zorder(ax2.get_zorder()+1)
		fig=plt.gcf()
		fig.tight_layout()
		if self.save_figures_to_folder:
			fig.savefig(check_folder(path=self.figure_path,filename='compare_at_time_%s.png'%'_'.join(['%g'%a for a in rel_time])),bbox_inches='tight')
						  

	def Compare_at_wave(self, rel_wave = None, other = None, fitted = False, norm_window = None,
						width = None, cmap = None, data_and_fit = False, scale_type = 'symlog', linewidth = 1):
		'''This function plots multiple kinetics into the same figure at one or
		multiple given wavelength (rel_wave) and  allows for 
		:ref:`Normalization and Scaling` Very useful to compare the 
		kinetics for different quencher concentrations or pump powers, 
		or e.g. different fits. The parameter width or the general self.wavelength_bin
		which is used if width is None (Default) defines the width of
		the spectral window that is integrated and shown.

		A normalization window can be given at which all the plotted curves are normalized to. 
		This window does not have to be in the plotted region. See :ref:`Normalization and Scaling`
		
		Parameters
		--------------
		
		rel_wave : float or list/vector (of floats)
			Specify the wavelength where to plot the kinetics, single value or
			list/vector of values (only mandatory entry) For each entry in 
			rel_wave a kinetic is plotted. 'rel_wave' and 'width' 
			(in the object called 'wavelength_bin' work together for the creation 
			of kinetic plots. At each selected wavelength the data between 
			wavelength+width/2 and wavelength-width/2 is averaged 
			for each timepoint
			
		other : TA object or list of those, optional 
			should be ta.plot_func objects (loaded or copied) and is what 
			is plotted against the data use a list [ta1,ta2,... ] or generate this
			list using the Gui function. See section :ref:`Opening multiple files` in 
			the documentation

		fitted : bool, optional
			True/False (Default) - use fitted data instead of raw data. 
			If True, the fitted datapoints (without interpolation) are used. 
			This is intended for comparing e.g. different fits

		norm_window : None or list/vector (with 4 floats), optional
			norm_window Give a list/tupel/vector with 4 entries in the order 
			[Start - time, End - time, Start - wavelength, End - Wavelength], 
			see section :ref:`Normalization and Scaling` in the documentation.
			If None (Default) no normalization is done.

		width
		   Specify the width above and below the given wavelength that is
		   integrated as window. If left to (Default) "None" the value from ta is
		   used.

		data_and_fit : bool, optional
			True or False (Default), choose if for the Fitted plot the raw data of the 
			other projects is to be plotting in addition to the fitted line. For False (Default)
			Only the fit is plotted.
			
		linewidth : float, optional
			linewidth to be used for plotting
				
		cmap : None or matplotlib color map, optional
			is a powerfull variable that chooses the colour map applied for all plots. If set to 
			None (Default) then the self.cmap is used.
			As standard I use the color map "jet" from matplotlib. There are a variety of colormaps 
			available that are very usefull. Beside "jet", "viridis" is a good choice as it is well 
			visible under red-green blindness. Other useful maps are "prism" for high fluctuations 
			or diverging color maps like "seismic". 
			See https://matplotlib.org/3.1.0/tutorials/colors/colormaps.html for a comprehensive 
			selection. In the code the colormaps are imported so if plot_func is imported as pf then 
			self.cmap=pf.cm.viridis sets viridis as the map to use. Internally the colors are chosen 
			with the "colm" function. The 2d plots require a continuous color map so if something 
			else is give 2d plots are shown automatically with "jet". For all of the 1d plots however 
			I first select a number of colors before each plot. If cmap is a continous map then these
			are sampled evenly over the colourmap. Manual iterables of colours 
			cmap=[(1,0,0),(0,1,0),(0,0,1),...] are also accepted, as are vectors or dataframes that 
			contain as rows the colors. There must be of course sufficient colors present for 
			the numbers of lines that will be plotted. So I recommend to provide at least 10 colours 
			(e.g.~your university colors). colours are always given as a, list or tuple with RGA or RGBA
			(with the last A beeing the Alpha=transparency. All numbers are between 0 and 1. 
			If a list/vector/DataFrame is given for the colours they will be used in the order provided.
		
		Scale_type : None or str
			is a general setting that can influences what time axis will be used for the plots. 
			"symlog" (linear around zero and logarithmic otherwise) "lin" and "log" are valid options.	
					
		Examples
		--------
		
		>>> import plot_func as pf
		>>> ta = pf.TA('test1.hdf5') #open the original project
		
		Now open a bunch of other projects to compare against
		
		>>> other_projects = pf.GUI_open(project_list = ['file1.SIA', 'file2.SIA'])
		
		Typical use:
		Compare the raw data without normalization at 400 nm and 500 nm
		
		>>> ta.Compare_at_wave(rel_wave = [400, 500], others = other_project)
		
		Compare the quality of the fit data without normalization at 400 nm and 500 nm
		
		>>> ta.Compare_at_wave(rel_wave = [400, 500], others = other_project, fitted = True)
		
		Compare with normalization window between 1ps and 2ps and 400nm and 450nm
		
		>>> norm_window=[1,2,400,450]
		>>> ta.Compare_at_wave(rel_wave = [400, 500], others = other_project, norm_window = norm_window)
		
		Use example: Often there are a lot of different measurements to 
		compare at multiple wavelength. The normlization is performed at the ground state bleach
		460 nm and early in time. Then it is better to make a new plot for each 
		wavelength. The normalization window stays fixed.
		
		>>> plt.close('all') #make some space
		>>> norm_window=[0.3,0.5,450,470] #define window in ground state bleach
		>>> for wave in [300,400,500,600,700]: #iterate over the wavelength 
		>>> 	ta.Compare_at_wave(rel_wave = wave, others = other_project, norm_window = norm_window)
		
		'''	
		global halfsize
		if self.save_figures_to_folder:self.figure_path=check_folder(path='result_figures',current_path=self.path)
		if width is None:width=self.wavelength_bin
		if rel_wave is None:
			rel_wave=self.rel_wave
		if other is not None:
			if not hasattr(other,'__iter__'):other=[other]
		if not hasattr(rel_wave,'__iter__'):
			rel_wave=[rel_wave]
		if cmap is None:cmap=self.cmap
		if fitted:
			try:
				re=self.re
			except:
				print("No fitted results present")
				return False
			if norm_window is not None:
				ref_scale = re['A'].loc[norm_window[0]:norm_window[1], norm_window[2]:norm_window[3]].mean().mean()
			if halfsize:
				fig, ax = plt.subplots(figsize = (5,3), dpi = 100)
			else:
				fig, ax = plt.subplots(figsize = (10, 6), dpi = 100)
			colors = colm(cmap = cmap, k = range(len(rel_wave)*(2+len(other))))
			ax = plot1d(re['A'], ax = ax, wavelength = rel_wave, width = width, lines_are = 'data', 
						cmap = colors, title = '', plot_type = scale_type, linewidth = linewidth)
			ax = plot1d(re['AC'], ax = ax, wavelength = rel_wave, width = width, lines_are = 'fitted', 
						cmap = colors, title = '', subplot = True, plot_type = scale_type, linewidth = linewidth)
			#ax = plot1d(re['AC'], ax = ax, wavelength = rel_wave, width = width, lines_are = 'fitted', 
			#			cmap = colors, color_offset = len(rel_wave), title = '', subplot = True, plot_type = scale_type)
			hand,  labels = ax.get_legend_handles_labels()
			lab=['%g nm'%a + '_' + str(self.filename) for a in rel_wave]
			for ent in rel_wave:
				lab.append('%g nm'%ent + '_' + str(self.filename))
			
			if other is not None:
				for i,o in enumerate(other):
					i+=1
					color_offset=(i+1)*len(rel_wave)
					try:
						re=o.re
					except:
						print('%s has no fitted results'%o.filename)
						continue
					if norm_window is not None:
						rel_scale=re['A'].loc[norm_window[0]:norm_window[1],norm_window[2]:norm_window[3]].mean().mean()
													 
						try:
							scaling=(rel_scale / ref_scale)
							if data_and_fit:
								ax = plot1d(re['A']/scaling, cmap = colors, ax = ax, wavelength = rel_wave, width = width,
											title = '', lines_are = 'data', subplot = True, color_offset = color_offset,
											plot_type = scale_type, linewidth = linewidth)
							ax = plot1d(re['AC']/scaling, cmap = colors, ax = ax, wavelength = rel_wave, width = width,
										title = '', lines_are = 'fitted', subplot = True, color_offset = color_offset, 
										plot_type = scale_type, linewidth = linewidth)
																					 
						except:
							print('scaling Failed!')
							if data_and_fit:
								ax = plot1d(re['A'], cmap = colors, ax = ax, wavelength = rel_wave, width = width, 
												title = '', lines_are = 'data', subplot = True, color_offset = color_offset, 
												plot_type = scale_type, linewidth = linewidth)
							ax = plot1d(re['AC'], cmap = colors, ax = ax, wavelength = rel_wave, width = width, 
											title = '', lines_are = 'fitted', subplot = True, color_offset = color_offset, 
											plot_type = scale_type, linewidth = linewidth) 
					else:
						if data_and_fit:
							ax = plot1d(re['A'], cmap = colors, ax = ax, wavelength = rel_wave, width = width, 
											title = '', lines_are = 'data', subplot = True, color_offset = color_offset, 
											plot_type = scale_type, linewidth = linewidth)
						ax = plot1d(re['AC'], cmap = colors, ax = ax, wavelength = rel_wave, width = width, 
										title = '', lines_are = 'fitted', subplot = True, color_offset = color_offset, 
										plot_type = scale_type, linewidth = linewidth)
					for ent in rel_wave:
						if data_and_fit:
							lab.append('%g nm'%ent + '_' + str(o.filename))
						lab.append('%g nm'%ent + '_' + str(o.filename))
					handles, labels=ax.get_legend_handles_labels()
					if data_and_fit:
						for a in handles[-2*len(rel_wave):]:
							hand.append(a)
					else:
						for a in handles[-len(rel_wave):]:
							hand.append(a)
			if norm_window is not None:
				 
				ax.set_title('compare measured and fitted data at given wavelength \n scaled to t=%g ps : %g ps , wl= %g nm: %g nm'%(norm_window[0],norm_window[1],norm_window[2],norm_window[3]))
			else:
				ax.set_title('compare measured and fitted data at given wavelength')
			ax.set_xlim(re['A'].index.values[0],re['A'].index.values[-1])
			ax.legend(hand,lab)
		else:
			if halfsize:
				fig, ax = plt.subplots(figsize = (5,3), dpi = 100)
			else:
				fig, ax = plt.subplots(figsize = (10, 6), dpi = 100)
			colors = colm(cmap = cmap, k = range(len(rel_wave)*(2+len(other))))
			if norm_window is not None:
				ref_scale = self.ds.loc[norm_window[0]:norm_window[1], norm_window[2]:norm_window[3]].mean().mean()
			ax = plot1d(self.ds, cmap = colors, ax = ax, wavelength = rel_wave, width = width, title = self.filename,
						baseunit = self.baseunit, lines_are = 'data', scattercut = self.scattercut, 
						bordercut = self.bordercut, subplot = False, color_offset = 0, timelimits = self.timelimits, 
						intensity_range = self.intensity_range, plot_type = scale_type, linewidth = linewidth)
			ax = plot1d(self.ds, cmap = colors, ax = ax, wavelength = rel_wave, width = width, title = self.filename, 
						baseunit = self.baseunit, lines_are = 'smoothed', scattercut = self.scattercut, 
						bordercut = self.bordercut, subplot = False, color_offset = 0, timelimits = self.timelimits, 
						intensity_range = self.intensity_range, plot_type = scale_type, linewidth = linewidth)
			if 1:
				handles, labels=ax.get_legend_handles_labels()
				lab=['%g nm'%a + '_' + str(self.filename) for a in rel_wave]
				hand=handles[len(rel_wave):]
				if other is not None:
					for i,o in enumerate(other):
						i+=1
						color_offset=(i+1)*len(rel_wave)
						if norm_window is not None:
							rel_scale=o.ds.loc[norm_window[0]:norm_window[1],norm_window[2]:norm_window[3]].mean().mean()
							try:
								scaling=(rel_scale/ref_scale)
								ax = plot1d(o.ds/scaling, cmap = colors, ax = ax, wavelength = rel_wave, width = width,
											title = o.filename, baseunit = self.baseunit, timelimits = self.timelimits,
											lines_are = 'data', scattercut = self.scattercut, bordercut = self.bordercut,
											subplot = True, color_offset = color_offset, 
											intensity_range = self.intensity_range, plot_type = scale_type, linewidth = linewidth)
								ax = plot1d(o.ds/scaling, cmap = colors, ax = ax, wavelength = rel_wave, width = width,
											title = o.filename, baseunit = self.baseunit, timelimits = self.timelimits,
											lines_are = 'smoothed', scattercut = self.scattercut, bordercut = self.bordercut,
											subplot = True, color_offset = color_offset, 
											intensity_range = self.intensity_range, plot_type = scale_type, linewidth = linewidth)
							except:
								print('scaling failed')
								ax = plot1d(o.ds, cmap = colors, ax = ax, wavelength = rel_wave, width = width,
											title = o.filename, baseunit = self.baseunit, timelimits = self.timelimits,
											lines_are = 'data', scattercut = self.scattercut, bordercut = self.bordercut,
											subplot = True, color_offset = color_offset, 
											intensity_range = self.intensity_range, plot_type = scale_type, linewidth = linewidth)
								ax = plot1d(o.ds, cmap = colors, ax = ax, wavelength = rel_wave, width = width,
											title = o.filename, baseunit = self.baseunit, timelimits = self.timelimits,
											lines_are = 'smoothed', scattercut = self.scattercut, bordercut = self.bordercut,
											subplot = True, color_offset = color_offset, 
											intensity_range = self.intensity_range, plot_type = scale_type, linewidth = linewidth)
						else:
							ax = plot1d(o.ds, cmap = colors, ax = ax, wavelength = rel_wave, width = width, title = o.filename,
										baseunit = self.baseunit, timelimits = self.timelimits, lines_are = 'data',
										scattercut = self.scattercut, bordercut = self.bordercut, subplot = True,
										color_offset = color_offset, intensity_range = self.intensity_range, 
										plot_type = scale_type, linewidth = linewidth)
							ax = plot1d(o.ds, cmap = colors, ax = ax, wavelength = rel_wave, width = width, title = o.filename,
										baseunit = self.baseunit, timelimits = self.timelimits, lines_are = 'smoothed',
										scattercut = self.scattercut, bordercut = self.bordercut, subplot = True,
										color_offset = color_offset, intensity_range = self.intensity_range, 
										plot_type = scale_type, linewidth = linewidth)
						for ent in ['%g nm'%a + '_' + str(o.filename) for a in rel_wave]:
							lab.append(ent)
						handles, labels=ax.get_legend_handles_labels()
						for ha in handles[-len(rel_wave):]:
							hand.append(ha)
				ax.set_title('compare measured and smoothed data at given wavelength')
				if norm_window is not None:
					ax.set_title('compare measured and smoothed data at given wavelength \n scaled to t=%g ps : %g ps , wl= %g nm: %g nm'%(norm_window[0],norm_window[1],norm_window[2],norm_window[3]))
				ax.legend(hand,lab)
		if self.save_figures_to_folder:  
			fig.savefig(check_folder(path=self.figure_path,filename='compare_at_wave_%s.png'%'_'.join(['%g'%a for a in rel_wave])),
			bbox_inches='tight')
			return ax


	def Compare_DAC(self, other = None, spectra = None, separate_plots = False, cmap = None):
		'''This is a convenience function to plot multiple extracted spectra 
		(DAS or species associated) into the same figure or into a separate figure
		each. Other should be ta.plot_func objects (loaded or copied). By
		standard it plots all into the same window. If all project have the same
		number of components one can activate "separate_plots" and have each
		separated (in the order created in the projects).

		The "Spectra" parameter allows as before the inclusion of an external
		spectrum. Others is optional and I use this function often to compare
		species associated spectra with one or multiple steady state spectra.
		
		Parameters
		--------------
			
		other : TA object or list of those, optional 
			should be ta.plot_func objects (loaded or copied) and is what 
			is plotted against the data use a list [ta1,ta2,... ] or generate this
			list using the Gui function. See section :ref:`Opening multiple files` in 
			the documentation

		spectra : None or DataFrame, optional
			If an DataFrame with the wavelength as index is provided, Then the spectra of each column
			is plotted into the differential spectra 1:1 and the column names are used in the legend
			Prior scaling is highly suggested. These spectra are not (in general) scaled with the 
			norm window. (see examples)

		separate_plots : bool, optional
			True or False (Default), separate plots is the switch that decides if a axis or 
			multiple axis are used. This option will result in a crash unless all objects have the
			same number of DAS/SAS components
				
		cmap : None or matplotlib color map, optional
			is a powerfull variable that chooses the colour map applied for all plots. If set to 
			None (Default) then the self.cmap is used.
			As standard I use the color map "jet" from matplotlib. There are a variety of colormaps 
			available that are very usefull. Beside "jet", "viridis" is a good choice as it is well 
			visible under red-green blindness. Other useful maps are "prism" for high fluctuations 
			or diverging color maps like "seismic". 
			See https://matplotlib.org/3.1.0/tutorials/colors/colormaps.html for a comprehensive 
			selection. In the code the colormaps are imported so if plot_func is imported as pf then 
			self.cmap=pf.cm.viridis sets viridis as the map to use. Internally the colors are chosen 
			with the "colm" function. The 2d plots require a continuous color map so if something 
			else is give 2d plots are shown automatically with "jet". For all of the 1d plots however 
			I first select a number of colors before each plot. If cmap is a continous map then these
			are sampled evenly over the colourmap. Manual iterables of colours 
			cmap=[(1,0,0),(0,1,0),(0,0,1),...] are also accepted, as are vectors or dataframes that 
			contain as rows the colors. There must be of course sufficient colors present for 
			the numbers of lines that will be plotted. So I recommend to provide at least 10 colours 
			(e.g.~your university colors). colours are always given as a, list or tuple with RGA or RGBA
			(with the last A beeing the Alpha=transparency. All numbers are between 0 and 1. 
			If a list/vector/DataFrame is given for the colours they will be used in the order provided.
		
		Examples
		--------
		
		>>> import plot_func as pf
		>>> ta = pf.TA('test1.hdf5') #open the original project, 
		>>> this MUST contain a fit, otherwise this will raise an error
		
		Now open a bunch of other projects to compare against,
		
		>>> #compare in a single window
		>>> other_projects = pf.GUI_open(project_list = ['file1.hdf5', 'file2.hdf5'])
		>>> ta.Compare_DAC(others = other_project)
		>>> #comprare in separate windows, 
		>>> #the other projects must have the same number of components
		>>> ta.Compare_DAC(others = other_project, separate_plots = True)
		
		Compare the DAC to an external spectrum
		
		>>> ext_spec = pd.read_csv('Ascii_spectrum.dat', sep = ',')
		>>> ta.Compare_DAC(spectra = ext_spec) #compare just the current solution
		>>> ta.Compare_DAC(spectra = ext_spec, others = other_project) #compare multiple
		
		'''
		global halfsize
		if self.save_figures_to_folder:self.figure_path = check_folder(path = 'result_figures', current_path = self.path)
		if other is not None:
			if not hasattr(other, '__iter__'):other = [other]
		try:
			re = self.re.copy()
		except:
			print("No fitted results present")
			return False
		if cmap is None:cmap = self.cmap
		species=re['DAC'].columns.values
		if other is None:
			col = range(len(re['DAC'].columns.values))
			colors = colm(cmap = cmap, k = col)
			
		else:
			re['DAC'].columns = [self.filename + '\n' + '%s'%a for a in re['DAC'].columns]
			if separate_plots:
				colors = colm(cmap = cmap, k = np.arange(len(other)+1))
			else:
				colors = colm(cmap = cmap, k = np.arange((len(other)+1)*len(species)))
		DAC = re['DAC']
		hand=[]
		if separate_plots:
			n_cols = int(np.ceil(len(re['DAC'].columns)/2))
			col = [colors[0] for a in range(len(re['DAC'].columns))]
			if self.scattercut is None:
				if halfsize:
					ax = DAC.plot(subplots = separate_plots, figsize = (6, 5), layout = (n_cols, 2), 
									legend = False, color = col, sharex = False)
				else:
					ax = DAC.plot(subplots = separate_plots, figsize = (12, 10), layout = (n_cols, 2), 
									legend = False, color = col, sharex = False)

				a=ax.ravel()
				handles,labels=a[0].get_legend_handles_labels()
				hand.append(handles[-1])
			elif isinstance(self.scattercut[0],  numbers.Number):
				if halfsize:
					ax = DAC.loc[:self.scattercut[0], :].plot(subplots = separate_plots, figsize = (6, 5), layout = (n_cols, 2), 
															legend = False, color = col, sharex = False)
				else:
					ax = DAC.loc[:self.scattercut[0], :].plot(subplots = separate_plots, figsize = (12, 10), layout = (n_cols, 2), 
															legend = False, color = col, sharex = False)								
				a=ax.ravel()
				handles,labels=a[0].get_legend_handles_labels()
				hand.append(handles[-1])	
				DAC_cut=DAC.loc[self.scattercut[1]:, :]
				for i,am in enumerate(DAC_cut.columns):
					DAC_cut.iloc[:,i].plot(ax = a[i], legend = False, color = col)
			else:
				scattercut = flatten(self.scattercut)
				for i in range(len(scattercut)/2+1):
					if i == 0:
						if halfsize:
							ax = DAC.loc[:self.scattercut[0], :].plot(subplots = separate_plots, figsize = (6, 5), layout = (n_cols, 2), 
																		legend = False, color = col, sharex = False)						
						else:
							ax = DAC.loc[:self.scattercut[0], :].plot(subplots = separate_plots, figsize = (12, 10), layout = (n_cols, 2), 
																		legend = False, color = col, sharex = False)
		
						a=ax.ravel()
						handles,labels=a[0].get_legend_handles_labels()
						hand.append(handles[-1])	
					elif i<(len(scattercut)/2):
						for j,am in enumerate(ax):
							DAC.loc[scattercut[2*i-1]:scattercut[2*i], :].plot(ax = a[j], legend = False, color = col, label = '_nolegend_')
					else:
						for j,am in enumerate(ax):
							DAC.loc[scattercut[-1]:, :].plot(ax = a[j], legend = False, color = col, label = '_nolegend_')
		else:
			if self.scattercut is None:
				if halfsize:
					ax  =  DAC.plot(subplots = separate_plots, figsize = (8,4), legend = False, color = colors[:len(species)], label = '_nolegend_')
				else:
					ax  =  DAC.plot(subplots = separate_plots, figsize = (16, 8), legend = False, color = colors[:len(species)], label = '_nolegend_')
			elif isinstance(self.scattercut[0], numbers.Number):
				if halfsize:
					ax  =  DAC.loc[:self.scattercut[0], :].plot(subplots = separate_plots, figsize = (8, 4), legend = False, color = colors[:len(species)], label = '_nolegend_')
					ax  =  DAC.loc[self.scattercut[1]:,  :].plot(ax=ax, subplots = separate_plots, figsize = (8, 4), legend = False, color = colors[:len(species)], label = '_nolegend_')
				else:
					ax  =  DAC.loc[:self.scattercut[0], :].plot(subplots = separate_plots, figsize = (16, 8), legend = False, color = colors[:len(species)], label = '_nolegend_')
					ax  =  DAC.loc[self.scattercut[1]:,  :].plot(ax=ax, subplots = separate_plots, figsize = (16, 8), legend = False, color = colors[:len(species)], label = '_nolegend_')
			else:
				scattercut  =  flatten(self.scattercut)
				for i in range(math.ceil(len(scattercut)/2+1)):
					if i  ==  0:
						if halfsize:
							ax  =  DAC.loc[:scattercut[0],  :].plot(subplots = separate_plots, figsize = (8, 4), legend = False, color = colors[:len(species)], label = '_nolegend_')
						else:
							ax  =  DAC.loc[:scattercut[0],  :].plot(subplots = separate_plots, figsize = (16, 8), legend = False, color = colors[:len(species)], label = '_nolegend_')
					elif i<(len(scattercut)/2):
						if halfsize:
							ax  =  DAC.loc[scattercut[2*i-1]:scattercut[2*i],  :].plot(ax=ax, subplots = separate_plots, figsize = (8, 4), legend = False, color = colors[:len(species)], label = '_nolegend_')
						else:
							ax  =  DAC.loc[scattercut[2*i-1]:scattercut[2*i],  :].plot(ax=ax, subplots = separate_plots, figsize = (16, 8), legend = False, color = colors[:len(species)], label = '_nolegend_')
					else:
						if halfsize:
							ax  =  DAC.loc[scattercut[-1]:,  :].plot(ax=ax, subplots = separate_plots, figsize = (8,4), legend = False, color = colors[:len(species)], label = '_nolegend_')
						else:
							ax  =  DAC.loc[scattercut[-1]:,  :].plot(ax=ax, subplots = separate_plots, figsize = (16, 8), legend = False, color = colors[:len(species)], label = '_nolegend_')
		if other is not None:
			for i,o in enumerate(other):
				try:
					re=o.re.copy()
				except:
					print('%s has no fitted results'%o.filename)
					continue
				re['DAC'].columns=[o.filename + '\n' + '%s'%a for a in re['DAC'].columns]
				if separate_plots:				
					col=[colors[i+1] for a in range(len(re['DAC'].columns))]
					for j,am in enumerate(re['DAC'].columns):
						if o.scattercut is None:	
							re['DAC'].iloc[:,j].plot(subplots=False,ax=a[j],legend=False,color=col[i])
							if j==0:
								handles,labels=a[0].get_legend_handles_labels()
								hand.append(handles[-1])	
						elif isinstance(o.scattercut[0],  numbers.Number):
							re['DAC'].iloc[:,j].loc[:o.scattercut[0]].plot(subplots=False,ax=a[j],legend=False,color=col[i])
							if j==0:
								handles,labels=a[0].get_legend_handles_labels()
								hand.append(handles[-1])	
							re['DAC'].iloc[:,j].loc[o.scattercut[1]:].plot(subplots=False,ax=a[j],legend=False,color=col[i],label = '_nolegend_')
						else:
							scattercut = flatten(o.scattercut)
							for m in range(math.ceil(len(scattercut)/2+1)):
								if m == 0:
									re['DAC'].iloc[:,j].loc[:scattercut[0]].plot(subplots=False,ax=a[j],legend=False,color=col[i])
									if j==0:
										handles,labels=a[j].get_legend_handles_labels()
										hand.append(handles[-1])	
								elif m<(len(scattercut)/2):
									re['DAC'].iloc[:,j].loc[scattercut[2*m-1]:scattercut[2*m]].plot(subplots=False,ax=a[j],legend=False,color=col[i],label = '_nolegend_')
								else:
									re['DAC'].iloc[:,j].loc[scattercut[-1]:].plot(subplots=False,ax=a[j],legend=False,color=col[i],label = '_nolegend_')
						a[j].set_xlabel('Wavelength in nm')				
						a[j].set_ylabel('Spectral strength in arb. units')
						a[j].legend(fontsize=8,frameon=False)
			
				else:
					dacs=len(re['DAC'].columns)
					col=colors[(i+1)*dacs:(i+2)*dacs]	
					DAC=re['DAC']
					if o.scattercut is None:	
						ax = DAC.plot(subplots=separate_plots,ax=ax,legend=False,color=colors[(i+1)*len(species):(i+2)*len(species)])
					elif isinstance(o.scattercut[0],  numbers.Number):
						ax = DAC.loc[:o.scattercut[0], :].plot(subplots=separate_plots,ax=ax,legend=False,color=colors[(i+1)*len(species):(i+2)*len(species)])
						DAC.loc[o.scattercut[1]:, :].plot(subplots=separate_plots,ax=ax,legend=False,color=colors[(i+1)*len(species):(i+2)*len(species)])
					else:
						scattercut = flatten(o.scattercut)
						for i in range(math.ceil(len(scattercut)/2+1)):
							if i == 0:
								ax = DAC.loc[:scattercut[0], :].plot(subplots=separate_plots,ax=ax,legend=False,color=colors[(i+1)*len(species):(i+2)*len(species)])
							elif i<(len(scattercut)/2):
								ax = DAC.loc[scattercut[2*i-1]:scattercut[2*i], :].plot(subplots=separate_plots,ax=ax,legend=False,color=colors[(i+1)*len(species):(i+2)*len(species)])
							else:
								ax = DAC.loc[scattercut[-1]:, :].plot(subplots=separate_plots,ax=ax,legend=False,color=colors[(i+1)*len(species):(i+2)*len(species)])
					ax.set_xlabel('Wavelength in nm')				
					ax.set_ylabel('Spectral strength in arb. units')
					ax.legend(fontsize=8,frameon=False)
		if not hasattr(ax,'__iter__'):ax=np.array([ax])
		if spectra is not None:
			for a in ax:
				spectra.plot(ax=a,subplots=separate_plots)				
		fig=(ax.ravel())[0].figure
		if separate_plots:
			if halfsize:
				fig.set_size_inches(6,5)
			else:
				fig.set_size_inches(12,10)
			axes_number=fig.get_axes()
			names=[self.filename]
			if other is not None:
				for o in other:
					names.append(o.filename)
			for i,ax in enumerate(axes_number):
				try:
					nametemp=['%s'%species[i] + ' - ' + a for a in names]
					ax.legend(hand,nametemp)
				except:
					pass
		else:
			ax=fig.get_axes()[0]
			names=[self.filename]
			if other is not None:
				for o in other:
					names.append(o.filename)
			handles,labels=ax.get_legend_handles_labels()
			nametemp=[]
			try:
				for a in names:
					for b in species:
						nametemp.append('%s'%b + ' - ' + a)
				ax.legend(handles,nametemp)
			except:
				pass
			if halfsize:
				fig.set_size_inches(8,4)
			else:
				fig.set_size_inches(16,8)
		fig.tight_layout()
		if self.save_figures_to_folder:
			fig.savefig(check_folder(path=self.figure_path,filename='compare_DAC.png'),bbox_inches='tight')



